from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
rows = []
for path in ROOT.rglob("*.pptx"):
    try:
        prs = Presentation(path)
        pictures = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
        rows.append((path.stat().st_mtime, path, len(prs.slides), pictures))
    except Exception:
        continue

for _, path, slides, pictures in sorted(rows, reverse=True)[:40]:
    rel = path.relative_to(ROOT)
    marker = "NO_IMAGES" if pictures == 0 else ""
    print(f"{pictures:02d} images | {slides:02d} slides | {rel} {marker}")
