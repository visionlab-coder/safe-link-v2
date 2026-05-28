from __future__ import annotations

import math
import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
ASSETS = OUT / "assets"
PREVIEWS = OUT / "previews"
PPTX = OUT / "SAFE-LINK_관리자_교육용_사용설명서.pptx"

W, H = 16, 9
SLIDE_PX = (1600, 900)

INK = "242124"
MUTED = "6B6560"
PAPER = "FAF7F1"
WARM = "EFE7DA"
CORAL = "D96B4B"
TEAL = "188A8A"
SAFETY = "E6B84A"
CHARCOAL = "1E2329"
GREEN = "28A36A"
RED = "C94E4E"
BLUE = "3A72D8"


def ensure_dirs() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    PREVIEWS.mkdir(parents=True, exist_ok=True)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def rgb(value: str) -> RGBColor:
    r, g, b = hex_to_rgb(value)
    return RGBColor(r, g, b)


def add_textbox(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_chip(slide, text, x, y, w, color=TEAL, fill="FFFFFF", size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def draw_wrapped(draw, text, xy, max_width, fnt, fill, line_gap=8):
    x, y = xy
    lines = []
    for raw in text.split("\n"):
        words = raw.split(" ")
        line = ""
        for word in words:
            trial = word if not line else f"{line} {word}"
            if draw.textbbox((0, 0), trial, font=fnt)[2] <= max_width:
                line = trial
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += fnt.size + line_gap
    return y


def pill(draw, xy, text, fill, outline=None, txt="FFFFFF", w=None):
    x, y = xy
    f = font(22, True)
    bbox = draw.textbbox((0, 0), text, font=f)
    width = w or bbox[2] + 42
    draw.rounded_rectangle((x, y, x + width, y + 48), radius=22, fill=fill, outline=outline, width=2)
    draw.text((x + width / 2, y + 24), text, font=f, fill=txt, anchor="mm")


def screen_frame(draw, title, subtitle=None):
    draw.rounded_rectangle((70, 52, 1130, 728), radius=38, fill="#FDFBF7", outline="#D8CEC2", width=3)
    draw.rounded_rectangle((70, 52, 1130, 122), radius=38, fill="#282625")
    draw.rectangle((70, 90, 1130, 122), fill="#282625")
    for i, c in enumerate(["#F06A52", "#E7B84A", "#42B883"]):
        draw.ellipse((98 + i * 38, 78, 120 + i * 38, 100), fill=c)
    draw.text((210, 89), title, font=font(24, True), fill="#FFFFFF", anchor="lm")
    if subtitle:
        draw.text((1080, 89), subtitle, font=font(18, True), fill="#E9E1D7", anchor="rm")


def make_mockup(name: str, kind: str) -> Path:
    img = Image.new("RGB", (1200, 780), "#F4EEE4")
    d = ImageDraw.Draw(img)
    screen_frame(d, "SAFE-LINK", kind)

    if kind == "인증":
        d.text((120, 178), "1. 언어 선택", font=font(30, True), fill="#" + INK)
        for i, label in enumerate(["한국어", "English", "Tiếng Việt", "中文", "ไทย"]):
            x = 120 + (i % 3) * 225
            y = 235 + (i // 3) * 78
            d.rounded_rectangle((x, y, x + 190, y + 52), radius=16, fill="#FFFFFF", outline="#D4C8BA", width=2)
            d.text((x + 95, y + 26), label, font=font(22, True), fill="#" + INK, anchor="mm")
        d.text((120, 420), "2. 역할 선택", font=font(30, True), fill="#" + INK)
        for i, (label, color) in enumerate([("관리자", "#" + TEAL), ("작업자", "#" + CORAL)]):
            x = 120 + i * 270
            d.rounded_rectangle((x, 475, x + 230, 560), radius=22, fill=color)
            d.text((x + 115, 518), label, font=font(28, True), fill="#FFFFFF", anchor="mm")
        d.rounded_rectangle((760, 210, 1045, 575), radius=30, fill="#282625")
        d.text((902, 275), "로그인 / 회원가입", font=font(28, True), fill="#FFFFFF", anchor="mm")
        for i, label in enumerate(["이메일", "비밀번호"]):
            d.rounded_rectangle((805, 330 + i * 72, 1000, 380 + i * 72), radius=12, fill="#FFFFFF")
            d.text((825, 355 + i * 72), label, font=font(18), fill="#8A8178", anchor="lm")
        d.rounded_rectangle((805, 490, 1000, 545), radius=16, fill="#" + CORAL)
        d.text((902, 518), "시작하기", font=font(22, True), fill="#FFFFFF", anchor="mm")

    elif kind == "대시보드":
        d.text((120, 172), "관리자 대시보드", font=font(34, True), fill="#" + INK)
        d.text((120, 216), "오늘 확인할 일을 한 화면에서 봅니다.", font=font(21), fill="#" + MUTED)
        cards = [("TBM 발송", "새 교육 만들기", TEAL), ("서명 현황", "미확인자 추적", SAFETY), ("1:1 채팅", "번역 상담", CORAL), ("작업자", "명단 관리", BLUE)]
        for i, (a, b, c) in enumerate(cards):
            x = 120 + (i % 2) * 360
            y = 295 + (i // 2) * 150
            d.rounded_rectangle((x, y, x + 310, y + 110), radius=24, fill="#FFFFFF", outline="#D8CEC2", width=2)
            d.ellipse((x + 24, y + 24, x + 76, y + 76), fill="#" + c)
            d.text((x + 100, y + 38), a, font=font(25, True), fill="#" + INK)
            d.text((x + 100, y + 72), b, font=font(18), fill="#" + MUTED)
        d.rounded_rectangle((850, 200, 1050, 600), radius=26, fill="#282625")
        for i, pct in enumerate([82, 64, 38]):
            y = 270 + i * 100
            d.text((880, y), ["TBM 확인", "서명 완료", "퀴즈 참여"][i], font=font(19, True), fill="#FFFFFF")
            d.rounded_rectangle((880, y + 38, 1015, y + 52), radius=8, fill="#4B4745")
            d.rounded_rectangle((880, y + 38, 880 + int(135 * pct / 100), y + 52), radius=8, fill="#" + [GREEN, TEAL, SAFETY][i])
            d.text((1030, y + 45), f"{pct}%", font=font(18, True), fill="#FFFFFF", anchor="lm")

    elif kind == "TBM 작성":
        d.text((118, 174), "TBM 작성", font=font(34, True), fill="#" + INK)
        d.rounded_rectangle((118, 232, 745, 500), radius=24, fill="#FFFFFF", outline="#D8CEC2", width=2)
        d.text((150, 275), "오늘 작업 내용과 위험요인을 입력", font=font(25, True), fill="#" + INK)
        draw_wrapped(d, "예: 3층 철근 양중 작업. 이동 동선 통제, 안전고리 체결, 낙하물 주의.", (150, 335), 520, font(22), "#" + MUTED)
        d.rounded_rectangle((150, 430, 350, 475), radius=14, fill="#" + TEAL)
        d.text((250, 453), "AI 초안 생성", font=font(19, True), fill="#FFFFFF", anchor="mm")
        d.rounded_rectangle((810, 230, 1045, 500), radius=24, fill="#FFF8EA", outline="#E8D6A3", width=2)
        d.text((838, 278), "발송 전 확인", font=font(25, True), fill="#" + INK)
        for i, line in enumerate(["현장명", "대상 작업자", "번역 언어", "발송 시간"]):
            d.text((845, 335 + i * 40), f"✓ {line}", font=font(20, True), fill="#" + MUTED)
        d.rounded_rectangle((118, 555, 1045, 618), radius=18, fill="#" + CORAL)
        d.text((582, 587), "작업자에게 TBM 발송", font=font(25, True), fill="#FFFFFF", anchor="mm")

    elif kind == "서명 현황":
        d.text((120, 174), "TBM 서명 현황", font=font(34, True), fill="#" + INK)
        names = [("김민수", "완료", GREEN), ("Nguyen An", "완료", GREEN), ("Somchai", "미확인", SAFETY), ("Rustam", "미확인", SAFETY), ("Li Wei", "완료", GREEN)]
        d.rounded_rectangle((120, 235, 720, 590), radius=24, fill="#FFFFFF", outline="#D8CEC2", width=2)
        for i, (n, s, c) in enumerate(names):
            y = 275 + i * 58
            d.line((150, y + 42, 690, y + 42), fill="#EFE7DA", width=2)
            d.text((155, y + 18), n, font=font(22, True), fill="#" + INK, anchor="lm")
            d.rounded_rectangle((560, y, 670, y + 34), radius=16, fill="#" + c)
            d.text((615, y + 17), s, font=font(17, True), fill="#FFFFFF", anchor="mm")
        d.rounded_rectangle((790, 255, 1030, 540), radius=28, fill="#282625")
        d.text((910, 320), "64%", font=font(72, True), fill="#FFFFFF", anchor="mm")
        d.text((910, 382), "서명 완료", font=font(24, True), fill="#E9E1D7", anchor="mm")
        d.rounded_rectangle((835, 445, 985, 500), radius=18, fill="#" + SAFETY)
        d.text((910, 473), "미확인자 알림", font=font(20, True), fill="#FFFFFF", anchor="mm")

    elif kind == "채팅":
        d.text((120, 174), "1:1 AI 번역 채팅", font=font(34, True), fill="#" + INK)
        d.rounded_rectangle((120, 230, 1030, 585), radius=28, fill="#FFFFFF", outline="#D8CEC2", width=2)
        rows = [("관리자", "오늘 작업 위치 확인했나요?", TEAL, 160), ("작업자", "네, 3층으로 이동했습니다.", CORAL, 430), ("AI", "번역과 음성 안내가 자동 생성됩니다.", SAFETY, 250)]
        for i, (who, msg, c, x) in enumerate(rows):
            y = 285 + i * 80
            d.rounded_rectangle((x, y, x + 460, y + 54), radius=20, fill="#" + c)
            d.text((x + 22, y + 27), f"{who}: {msg}", font=font(20, True), fill="#FFFFFF", anchor="lm")
        d.rounded_rectangle((170, 520, 850, 560), radius=16, fill="#F4EEE4")
        d.text((190, 540), "메시지를 입력하고 Enter", font=font(18), fill="#8A8178", anchor="lm")
        d.rounded_rectangle((870, 515, 990, 565), radius=16, fill="#" + TEAL)
        d.text((930, 540), "전송", font=font(20, True), fill="#FFFFFF", anchor="mm")

    elif kind == "세부 메뉴":
        d.text((120, 174), "관리자 세부 섹션", font=font(34, True), fill="#" + INK)
        items = [("작업자 등록", "명단·언어·소속 관리", TEAL), ("NFC / QR", "현장 출입과 스티커 발급", CORAL), ("퀴즈", "교육 이해도 확인", SAFETY), ("ESG", "교육·서명 리포트", GREEN), ("용어집", "현장 은어 표준화", BLUE), ("라이브", "실시간 통역 지원", RED)]
        for i, (a, b, c) in enumerate(items):
            x = 125 + (i % 3) * 310
            y = 250 + (i // 3) * 150
            d.rounded_rectangle((x, y, x + 260, y + 105), radius=24, fill="#FFFFFF", outline="#D8CEC2", width=2)
            d.rectangle((x, y, x + 14, y + 105), fill="#" + c)
            d.text((x + 38, y + 38), a, font=font(23, True), fill="#" + INK)
            d.text((x + 38, y + 70), b, font=font(17), fill="#" + MUTED)

    path = ASSETS / f"{name}.png"
    img.save(path, quality=95)
    return path


def make_visual_assets() -> dict[str, Path]:
    return {
        "auth": make_mockup("auth_flow", "인증"),
        "dash": make_mockup("dashboard", "대시보드"),
        "tbm": make_mockup("tbm_create", "TBM 작성"),
        "status": make_mockup("status", "서명 현황"),
        "chat": make_mockup("chat", "채팅"),
        "sections": make_mockup("sections", "세부 메뉴"),
    }


def setup_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=PAPER):
    add_round_rect(slide, 0, 0, W, H, color, radius=False)


def footer(slide, idx: int):
    add_textbox(slide, f"SAFE-LINK 관리자 교육용 사용 설명서  |  {idx:02d}", 0.65, 8.36, 5, 0.26, 8, MUTED)


def title(slide, main, sub=None, idx=None):
    add_textbox(slide, main, 0.75, 0.62, 8.9, 0.72, 29, INK, True)
    if sub:
        add_textbox(slide, sub, 0.78, 1.28, 9.8, 0.42, 13, MUTED)
    if idx:
        footer(slide, idx)


def image_slide(prs, idx, main, sub, asset, bullets):
    slide = blank(prs)
    bg(slide)
    title(slide, main, sub, idx)
    slide.shapes.add_picture(str(asset), Inches(0.78), Inches(1.95), width=Inches(9.0))
    add_round_rect(slide, 10.25, 2.0, 4.75, 5.55, "FFFFFF", WARM)
    for i, (head, body) in enumerate(bullets):
        y = 2.35 + i * 1.28
        add_chip(slide, f"{i+1}", 10.58, y, 0.42, [TEAL, CORAL, SAFETY, GREEN][i % 4], "FFFFFF", 11)
        add_textbox(slide, head, 11.12, y - 0.03, 3.3, 0.28, 14, INK, True)
        add_textbox(slide, body, 11.12, y + 0.34, 3.35, 0.55, 10.5, MUTED)
    return slide


def build_deck() -> Presentation:
    ensure_dirs()
    assets = make_visual_assets()
    prs = setup_presentation()

    s = blank(prs)
    bg(s, CHARCOAL)
    add_round_rect(s, 0, 0, 16, 9, CHARCOAL, radius=False)
    add_textbox(s, "SAFE-LINK", 0.85, 0.72, 3.4, 0.48, 20, "FFFFFF", True)
    add_textbox(s, "관리자 교육용\n사용 설명서", 0.9, 2.0, 7.2, 1.95, 43, "FFFFFF", True)
    add_textbox(s, "회원가입·로그인부터 TBM 발송, 서명 현황, 채팅, 세부 관리 메뉴까지 초보자 기준으로 따라 하는 운영 가이드", 0.96, 4.22, 7.2, 0.8, 16, "E9E1D7")
    s.shapes.add_picture(str(assets["dash"]), Inches(8.3), Inches(1.0), width=Inches(6.9))
    add_chip(s, "Admin Onboarding", 0.96, 6.65, 2.25, CORAL, CHARCOAL, 12)
    add_textbox(s, "2026.05", 0.96, 7.65, 1.2, 0.28, 10, "D8CEC2")

    s = blank(prs)
    bg(s)
    title(s, "처음 관리자가 기억할 전체 흐름", "이 순서만 이해하면 대부분의 운영 업무를 시작할 수 있습니다.", 2)
    steps = [("1", "접속", "/auth에서 언어와 역할 선택"), ("2", "가입/로그인", "이메일과 비밀번호로 관리자 진입"), ("3", "대시보드", "오늘 할 일과 알림 확인"), ("4", "TBM 발송", "교육 내용을 작성하고 대상자에게 전송"), ("5", "확인/조치", "서명 현황, 채팅, 리포트 관리")]
    for i, (n, h, b) in enumerate(steps):
        x = 0.9 + i * 3.02
        add_round_rect(s, x, 3.15, 2.42, 2.34, "FFFFFF", WARM)
        add_textbox(s, n, x + 0.2, 3.34, 0.56, 0.42, 24, [TEAL, CORAL, SAFETY, GREEN, BLUE][i], True, PP_ALIGN.CENTER)
        add_textbox(s, h, x + 0.32, 4.0, 1.7, 0.38, 18, INK, True, PP_ALIGN.CENTER)
        add_textbox(s, b, x + 0.22, 4.58, 1.95, 0.52, 11, MUTED, False, PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "→", x + 2.47, 4.07, 0.35, 0.3, 22, CORAL, True, PP_ALIGN.CENTER)
    add_textbox(s, "TIP  처음에는 모든 메뉴를 다 보려 하지 말고, ‘TBM 발송 → 서명 확인 → 미확인자 연락’만 반복해서 익히면 됩니다.", 1.05, 6.55, 12.8, 0.46, 15, INK, True)

    image_slide(prs, 3, "회원가입과 로그인", "관리자는 먼저 언어를 선택하고 관리자 역할로 들어갑니다.", assets["auth"], [
        ("언어 선택", "한국어를 기본으로 선택합니다. 외국인 작업자 화면은 작업자 언어를 따릅니다."),
        ("관리자 선택", "역할 선택 화면에서 ‘관리자’를 눌러 관리자 로그인으로 이동합니다."),
        ("로그인 또는 가입", "기존 계정은 로그인, 신규 관리자는 회원가입을 진행합니다."),
        ("비밀번호 확인", "초기 비밀번호는 현장 책임자가 안전하게 공유한 값만 사용합니다."),
    ])

    s = image_slide(prs, 4, "최초 설정", "로그인 후 프로필과 현장 정보를 맞춰야 권한과 메뉴가 제대로 보입니다.", assets["auth"], [
        ("이름", "보고서와 채팅에 표시될 관리자 이름을 입력합니다."),
        ("역할", "안전관리자, 현장관리자, 본사 관리자 등 권한을 확인합니다."),
        ("현장", "담당 현장 코드가 틀리면 다른 현장 데이터가 보이지 않을 수 있습니다."),
        ("언어", "관리자 화면 언어와 알림 언어를 지정합니다."),
    ])
    add_textbox(s, "설정이 끝나면 자동으로 관리자 대시보드로 이동합니다.", 10.55, 7.1, 3.7, 0.35, 12, CORAL, True)

    image_slide(prs, 5, "관리자 대시보드 읽는 법", "대시보드는 ‘오늘 해야 할 일’을 빠르게 찾는 첫 화면입니다.", assets["dash"], [
        ("TBM 발송", "새 안전교육을 만들고 작업자에게 보냅니다."),
        ("서명 현황", "누가 읽고 서명했는지 확인합니다."),
        ("1:1 채팅", "작업자와 번역 채팅으로 문의를 처리합니다."),
        ("세부 메뉴", "작업자, NFC, QR, 퀴즈, ESG, 용어집을 관리합니다."),
    ])

    image_slide(prs, 6, "TBM 작성과 발송", "현장 상황을 짧게 입력하고, AI 초안과 번역 결과를 확인한 뒤 발송합니다.", assets["tbm"], [
        ("작업 내용 입력", "오늘 작업, 장소, 위험요인, 보호구를 구체적으로 적습니다."),
        ("AI 초안 생성", "현장 은어를 표준 문장으로 정리하고 교육 문장으로 다듬습니다."),
        ("대상자 확인", "현장·조·언어가 맞는지 발송 전 한 번 더 확인합니다."),
        ("발송", "발송 후에는 서명 현황 화면에서 확인률을 추적합니다."),
    ])

    image_slide(prs, 7, "서명 현황 확인", "발송 후 관리자의 핵심 업무는 미확인자를 빠르게 찾는 것입니다.", assets["status"], [
        ("완료자", "읽음과 서명 시간이 기록된 작업자입니다."),
        ("미확인자", "교육을 열지 않았거나 서명하지 않은 작업자입니다."),
        ("알림", "미확인자에게 재알림 또는 현장 구두 안내를 진행합니다."),
        ("기록", "교육·서명 기록은 보고서 근거로 활용됩니다."),
    ])

    image_slide(prs, 8, "1:1 AI 번역 채팅", "언어가 달라도 짧고 명확하게 질문하면 작업자에게 자연스럽게 전달됩니다.", assets["chat"], [
        ("작업자 선택", "연락할 작업자를 선택하고 대화를 시작합니다."),
        ("짧게 입력", "한 문장에 하나의 요청만 넣으면 번역 품질이 좋아집니다."),
        ("번역 확인", "위험 작업 지시는 번역 결과를 한 번 더 확인합니다."),
        ("기록 관리", "중요 대화는 현장 조치 기록으로 남길 수 있습니다."),
    ])

    image_slide(prs, 9, "작업자 관리", "작업자 등록은 언어와 현장 배정이 가장 중요합니다.", assets["sections"], [
        ("신규 등록", "이름, 연락처, 주 사용 언어, 소속 팀을 입력합니다."),
        ("언어 확인", "TBM 번역과 음성 안내가 이 언어 기준으로 제공됩니다."),
        ("현장 배정", "현장 코드가 맞아야 해당 TBM과 알림을 받습니다."),
        ("수정/비활성", "퇴근, 전출, 계약 종료자는 목록에서 정리합니다."),
    ])

    image_slide(prs, 10, "NFC와 QR", "현장 출입, 스티커 발급, 작업자 확인을 빠르게 처리하는 메뉴입니다.", assets["sections"], [
        ("NFC 스티커", "작업자별 스티커를 발급하고 태그 기록을 확인합니다."),
        ("QR 코드", "현장 입장 또는 TBM 접근용 QR을 생성합니다."),
        ("일일 로그", "태그 시간과 출입 이력을 확인합니다."),
        ("오류 처리", "스티커 분실·교체 시 기존 토큰을 비활성화합니다."),
    ])

    image_slide(prs, 11, "퀴즈와 ESG", "교육 이해도와 안전 이행 기록을 관리 리포트로 연결합니다.", assets["sections"], [
        ("퀴즈 생성", "TBM 내용 기반으로 짧은 확인 문제를 만듭니다."),
        ("참여 확인", "작업자별 응답 여부와 정답률을 확인합니다."),
        ("ESG 리포트", "교육, 서명, 참여 데이터를 보고용으로 정리합니다."),
        ("인센티브", "퀴즈 참여나 안전 활동에 보상을 부여할 수 있습니다."),
    ])

    image_slide(prs, 12, "용어집과 라이브 통역", "현장 은어를 표준화하고, 실시간 소통이 필요한 상황을 지원합니다.", assets["sections"], [
        ("용어 등록", "현장 은어, 표준어, 설명을 함께 저장합니다."),
        ("번역 품질", "용어집이 많을수록 TBM 번역이 현장에 맞게 개선됩니다."),
        ("라이브 통역", "짧은 현장 지시나 질의응답을 즉시 전달합니다."),
        ("검수", "위험 작업 지시는 관리자가 최종 의미를 확인합니다."),
    ])

    s = blank(prs)
    bg(s)
    title(s, "초보 관리자를 위한 운영 체크리스트", "매일 이 순서대로 확인하면 누락을 줄일 수 있습니다.", 13)
    checks = [
        ("출근 직후", "담당 현장과 오늘 작업자 명단 확인"),
        ("작업 전", "TBM 작성, 번역 확인, 대상자 발송"),
        ("작업 시작 전", "서명 미완료자 확인 후 재알림"),
        ("작업 중", "1:1 채팅·라이브 통역으로 문의 처리"),
        ("퇴근 전", "서명·퀴즈·출입 기록을 보고서로 정리"),
    ]
    for i, (h, b) in enumerate(checks):
        y = 2.0 + i * 1.05
        add_round_rect(s, 1.1, y, 13.4, 0.72, "FFFFFF", WARM)
        add_chip(s, f"{i+1}", 1.38, y + 0.18, 0.42, [TEAL, CORAL, SAFETY, GREEN, BLUE][i], "FFFFFF", 10)
        add_textbox(s, h, 2.0, y + 0.18, 2.0, 0.25, 14, INK, True)
        add_textbox(s, b, 4.1, y + 0.18, 8.6, 0.25, 13, MUTED)

    s = blank(prs)
    bg(s, CHARCOAL)
    add_textbox(s, "마지막으로", 0.92, 0.85, 2.2, 0.36, 18, "E9E1D7", True)
    add_textbox(s, "SAFE-LINK 운영의 핵심은\n‘보내고, 확인하고, 기록하는 것’입니다.", 0.92, 2.05, 10.2, 1.55, 39, "FFFFFF", True)
    add_textbox(s, "처음에는 메뉴를 모두 외우지 않아도 됩니다. TBM 발송, 서명 현황, 미확인자 조치, 작업자 문의 대응 흐름만 반복하면 관리자 업무의 대부분을 안정적으로 시작할 수 있습니다.", 0.98, 4.2, 9.7, 0.9, 17, "E9E1D7")
    add_chip(s, "교육 종료", 0.98, 6.35, 1.45, CORAL, CHARCOAL, 12)
    s.shapes.add_picture(str(assets["status"]), Inches(10.5), Inches(1.35), width=Inches(4.7))
    add_textbox(s, "SAFE-LINK 관리자 교육용 사용 설명서", 0.98, 7.78, 4.2, 0.26, 10, "D8CEC2")

    return prs


def preview_slide(i: int, title_text: str, subtitle: str = "", asset: Path | None = None, dark: bool = False) -> None:
    img = Image.new("RGB", SLIDE_PX, "#" + (CHARCOAL if dark else PAPER))
    d = ImageDraw.Draw(img)
    title_color = "#FFFFFF" if dark else "#" + INK
    muted_color = "#E9E1D7" if dark else "#" + MUTED
    d.text((72, 64), title_text, font=font(42, True), fill=title_color)
    if subtitle:
        d.text((76, 124), subtitle, font=font(22), fill=muted_color)
    if asset and asset.exists():
        mock = Image.open(asset).convert("RGB")
        mock.thumbnail((900, 585))
        x = 72
        y = 210
        d.rounded_rectangle((x - 8, y - 8, x + mock.width + 8, y + mock.height + 8), radius=26, fill="#FFFFFF", outline="#DDD0C4", width=3)
        img.paste(mock, (x, y))
        panel_x = x + mock.width + 64
        d.rounded_rectangle((panel_x, 230, 1510, 705), radius=30, fill="#FFFFFF" if not dark else "#302D2B", outline="#DDD0C4", width=2)
        notes = ["따라 하기", "확인하기", "기록하기"]
        for n, note in enumerate(notes):
            cy = 320 + n * 105
            d.rounded_rectangle((panel_x + 45, cy - 25, panel_x + 90, cy + 20), radius=16, fill="#" + [TEAL, CORAL, SAFETY][n])
            d.text((panel_x + 68, cy - 3), str(n + 1), font=font(20, True), fill="#FFFFFF", anchor="mm")
            d.text((panel_x + 118, cy - 4), note, font=font(26, True), fill=title_color if dark else "#" + INK, anchor="lm")
    else:
        d.rounded_rectangle((85, 220, 1510, 760), radius=34, fill="#FFFFFF" if not dark else "#302D2B", outline="#DDD0C4", width=3)
        steps = ["접속", "가입/로그인", "대시보드", "TBM 발송", "확인/조치"]
        for n, step in enumerate(steps):
            x = 165 + n * 265
            d.ellipse((x, 410, x + 92, 502), fill="#" + [TEAL, CORAL, SAFETY, GREEN, BLUE][n])
            d.text((x + 46, 456), str(n + 1), font=font(34, True), fill="#FFFFFF", anchor="mm")
            d.text((x + 46, 555), step, font=font(24, True), fill=title_color if dark else "#" + INK, anchor="mm")
            if n < 4:
                d.text((x + 147, 456), "→", font=font(34, True), fill="#" + CORAL, anchor="mm")
    d.text((72, 840), f"SAFE-LINK 관리자 교육용 사용 설명서 | {i:02d}", font=font(16), fill=muted_color)
    img.save(PREVIEWS / f"slide_{i:02d}.png")


def make_preview_contact_sheet() -> None:
    assets = {
        "auth": ASSETS / "auth_flow.png",
        "dash": ASSETS / "dashboard.png",
        "tbm": ASSETS / "tbm_create.png",
        "status": ASSETS / "status.png",
        "chat": ASSETS / "chat.png",
        "sections": ASSETS / "sections.png",
    }
    titles = [
        ("SAFE-LINK 관리자 교육용 사용 설명서", assets["dash"], True),
        ("처음 관리자가 기억할 전체 흐름", None, False),
        ("회원가입과 로그인", assets["auth"], False),
        ("최초 설정", assets["auth"], False),
        ("관리자 대시보드 읽는 법", assets["dash"], False),
        ("TBM 작성과 발송", assets["tbm"], False),
        ("서명 현황 확인", assets["status"], False),
        ("1:1 AI 번역 채팅", assets["chat"], False),
        ("작업자 관리", assets["sections"], False),
        ("NFC와 QR", assets["sections"], False),
        ("퀴즈와 ESG", assets["sections"], False),
        ("용어집과 라이브 통역", assets["sections"], False),
        ("운영 체크리스트", None, False),
        ("마지막으로", assets["status"], True),
    ]
    for idx, (t, asset, dark) in enumerate(titles, 1):
        preview_slide(idx, t, asset=asset, dark=dark)

    thumbs = []
    for idx in range(1, len(titles) + 1):
        im = Image.open(PREVIEWS / f"slide_{idx:02d}.png").resize((320, 180))
        thumbs.append(im)
    sheet = Image.new("RGB", (4 * 320, math.ceil(len(thumbs) / 4) * 180), "#FFFFFF")
    for i, im in enumerate(thumbs):
        sheet.paste(im, ((i % 4) * 320, (i // 4) * 180))
    sheet.save(PREVIEWS / "contact_sheet.png")


def main() -> None:
    ensure_dirs()
    prs = build_deck()
    prs.save(PPTX)
    make_preview_contact_sheet()
    print(f"pptx={PPTX}")
    print(f"previews={PREVIEWS}")


if __name__ == "__main__":
    main()
