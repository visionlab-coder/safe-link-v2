from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "generated" / "daewoo-hyper-safety-business-20260609" / "대우건설_HyperSafety_AI_SQ-LINK_사업제안서_PPT_v1_20260609.pptx"
REPORT = ROOT / "docs" / "generated" / "대우건설_HyperSafety_AI_사업보고서_v1_전무님발언반영_20260609.md"
PREVIEW = ROOT / "docs" / "generated" / "daewoo-hyper-safety-business-20260609" / "previews" / "contact_sheet.png"

prs = Presentation(PPTX)
all_text = " ".join(
    shape.text
    for slide in prs.slides
    for shape in slide.shapes
    if getattr(shape, "has_text_frame", False) and shape.text
)
pictures = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)

print(f"pptx_exists={PPTX.exists()}")
print(f"report_v1_exists={REPORT.exists()}")
print(f"preview_exists={PREVIEW.exists()}")
print(f"slide_count={len(prs.slides)}")
print(f"picture_count={pictures}")
print(f"has_sq_link={'SQ-LINK' in all_text}")
print(f"has_3d_coordinate={'3D 좌표' in all_text}")
print(f"has_poc={'8주 PoC' in all_text or 'PoC' in all_text}")
print(f"has_no_robot_direct={'로봇을 만드는 회사가 아니라' in all_text or '직접 개발' in all_text}")
print(f"has_quality={'품질' in all_text}")
