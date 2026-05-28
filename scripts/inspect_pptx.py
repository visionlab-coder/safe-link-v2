from __future__ import annotations

import sys
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation


def main() -> None:
    path = Path(sys.argv[1])
    prs = Presentation(path)
    pictures = 0
    text_shapes = 0
    empty_text = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pictures += 1
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text_shapes += 1
                if not shape.text.strip():
                    empty_text += 1
    with ZipFile(path) as z:
        xml_text = "\n".join(
            z.read(name).decode("utf-8", errors="ignore")
            for name in z.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
    placeholders = [token for token in ["Slide Number", "sldNum", "Click to add"] if token in xml_text]
    print(f"slides={len(prs.slides)}")
    print(f"pictures={pictures}")
    print(f"text_shapes={text_shapes}")
    print(f"empty_text_shapes={empty_text}")
    print(f"placeholders={','.join(placeholders) if placeholders else 'none'}")


if __name__ == "__main__":
    main()
