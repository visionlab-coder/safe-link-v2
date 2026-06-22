from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606"
ASSET = BASE / "assets_micro_drone_nest"
OUT = BASE / "SAFE-LINK_Micro_Drone_Nest_대우건설_Advanced_제출용_20260607.pptx"
PREVIEW = BASE / "previews-micro-drone-advanced"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8", "ink": "101828", "muted": "667085", "line": "D0D5DD",
    "blue": "175CD3", "green": "027A48", "orange": "B54708", "red": "B42318",
    "white": "FFFFFF", "soft": "F2F4F7", "navy": "0B1220", "note": "FFF7D6",
}


def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, transparency=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.05
        except Exception:
            pass
    return shp


def txt(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()


def footer(slide, n):
    txt(slide, "SAFE-LINK Micro Drone Nest / 서원토건", 0.55, 7.04, 4.8, 0.24, 10, "muted")
    txt(slide, f"{n:02d}", 12.25, 7.04, 0.45, 0.24, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    txt(slide, label, 0.62, 0.34, 5.6, 0.3, 12, "blue", True)
    txt(slide, title, 0.62, 0.74, 11.7, 0.72, 28, "ink", True)
    if sub:
        txt(slide, sub, 0.64, 1.5, 11.2, 0.38, 15, "muted")
    rule(slide, 0.62, 2.02, 12.05)
    footer(slide, n)


def img(slide, filename, x, y, w, h):
    p = ASSET / filename
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        txt(slide, filename, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, 13, "muted", True, "center")


def note(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True)
    rect(slide, x, y, 0.08, h, color, color, False)
    txt(slide, head, x + 0.22, y + 0.15, w - 0.42, 0.32, 15, color, True)
    txt(slide, body, x + 0.22, y + 0.55, w - 0.42, h - 0.66, 14, "ink")


def table(slide, x, y, rows, widths, row_h=0.5, head="navy"):
    for r, row in enumerate(rows):
        cx = x
        for c, value in enumerate(row):
            fill = head if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            rect(slide, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            txt(slide, value, cx + 0.08, y + r * row_h + 0.09, widths[c] - 0.16, row_h - 0.12,
                11 if r == 0 else 12.5, color, r == 0)
            cx += widths[c]


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, W, H, "paper", "paper", False)
        return s

    # 1 cover
    s = slide()
    img(s, "01-micro-drone-nest-hero.png", 6.52, 0, 6.82, 7.5)
    txt(s, "DAEWOO HYPER SAFETY & AI", 0.72, 0.68, 4.8, 0.3, 12, "blue", True)
    txt(s, "SAFE-LINK\nMicro Drone Nest", 0.72, 1.28, 5.25, 1.38, 35, "ink", True)
    txt(s, "초소형 드론 네스트 기반\n건설현장 사각공간 자동정찰 플랫폼", 0.74, 3.08, 5.2, 0.78, 21, "ink", True)
    rule(s, 0.74, 4.12, 4.4, "blue")
    txt(s, "손바닥보다 작은 드론 · 분산 네스트 · 자동복귀/충전 · 실시간 관제 · SAFE-LINK 조치보고", 0.74, 4.45, 5.2, 0.72, 15, "muted")
    txt(s, "서원토건 / 2026.06.07", 0.74, 6.78, 3.4, 0.25, 11, "muted")

    # 2 correction
    s = slide()
    header(s, "01 DIRECTION", "큰 드론이 아니라 작은 안전 감각기관입니다", "현장 곳곳의 네스트에서 초소형 드론이 자동 출발·복귀·충전합니다.", 2)
    rows = [
        ["구분", "기존 드론 순찰", "Micro Drone Nest"],
        ["크기", "중형·대형 드론", "손바닥보다 작은 초소형"],
        ["공간", "상공·외곽·넓은 구역", "샤프트·PIT·점검구·틈새"],
        ["운용", "사람이 조종", "정기순찰 + 관제센터 명령"],
        ["가치", "항공 촬영", "사각공간 반복 확인"],
        ["연동", "영상 업로드", "위험후보 → 조치 → 보고"],
    ]
    table(s, 0.9, 2.45, rows, [1.55, 4.4, 5.2], 0.58)
    txt(s, "이 제안은 건설현장 사각공간에 작은 안전 감각기관을 심는 것입니다.", 1.0, 6.08, 11.2, 0.36, 20, "blue", True, "center")

    # 3 blind spots
    s = slide()
    header(s, "02 TARGET", "진짜 타깃은 넓은 상공이 아니라 좁은 사각공간", "사람이 들어가기 어렵고 CCTV도 보지 못하는 공간을 반복 확인합니다.", 3)
    img(s, "02-hard-to-reach-space.png", 0.75, 2.35, 5.75, 3.6)
    rows = [
        ["사각공간", "주요 위험"],
        ["엘리베이터 샤프트", "추락, 낙하물, 난간 훼손"],
        ["지하 PIT·집수정", "침수, 조도 부족, 유해가스"],
        ["천장·배관 주변", "누수, 임시전기, 화재위험"],
        ["동바리·거푸집 틈", "지지재 이상, 발판 불량"],
        ["개구부 하부", "낙하물, 안전망 훼손"],
    ]
    table(s, 6.85, 2.45, rows, [2.4, 3.35], 0.53)

    # 4 network
    s = slide()
    header(s, "03 NEST NETWORK", "네스트는 보관함이 아니라 현장 센서 거점입니다", "각 네스트가 배터리, 충전, 영상, 통신 상태를 SAFE-LINK로 보냅니다.", 4)
    img(s, "03-distributed-nests.png", 0.75, 2.3, 5.85, 3.55)
    note(s, "보관·충전", "초소형 드론을 보호하고 접점식 또는 무선충전으로 대기", 6.95, 2.35, 4.95, 0.95, "blue")
    note(s, "통신·상태", "배터리, 온습도, 도어, 드론 장착 여부를 관제센터에 전송", 6.95, 3.55, 4.95, 0.95, "green")
    note(s, "미션 대기", "정기순찰 또는 관제센터 긴급 호출을 기다리는 현장 거점", 6.95, 4.75, 4.95, 0.95, "orange")

    # 5 operating scenario
    s = slide()
    header(s, "04 OPERATION", "정기 자동순찰과 긴급 호출이 함께 작동합니다", "자동순찰은 반복성을 만들고, 관제센터 명령은 현장 대응성을 만듭니다.", 5)
    steps = [("스케줄", "정기 순찰"), ("네스트", "상태 확인"), ("드론", "근접 촬영"), ("AI", "위험후보"), ("관리자", "조치 요청"), ("복귀", "자동충전")]
    for i, (a, b) in enumerate(steps):
        x = 0.55 + i * 2.05
        rect(s, x, 2.9, 1.55, 1.15, "white", ["blue", "green", "orange", "red", "blue", "green"][i], True)
        txt(s, a, x + 0.1, 3.08, 1.35, 0.26, 13.5, ["blue", "green", "orange", "red", "blue", "green"][i], True, "center")
        txt(s, b, x + 0.1, 3.55, 1.35, 0.24, 12, "ink", True, "center")
        if i < 5:
            txt(s, "→", x + 1.58, 3.32, 0.28, 0.24, 18, "muted", True, "center")
    note(s, "관제센터 긴급 호출", "특정 네스트 선택 → 즉시 이륙 → 재촬영/정지/복귀 명령 → 위험 발견 시 조치요청 생성", 1.0, 5.0, 5.5, 1.0, "blue")
    note(s, "조치보고 흐름", "위험후보 발견 → 안전관리자 확인 → 담당자 배정 → 완료사진 등록 → 보고서 생성", 6.9, 5.0, 5.3, 1.0, "green")

    # 6 control center
    s = slide()
    header(s, "05 CONTROL CENTER", "관제센터는 작은 드론들의 실시간 눈을 갖게 됩니다", "네스트 위치, 영상 피드, 위험후보, 조치상태를 한 화면에서 관리합니다.", 6)
    img(s, "04-control-center-micro-feeds.png", 0.75, 2.35, 5.85, 3.55)
    rows = [
        ["기능", "설명"],
        ["네스트 상태", "드론 유무, 배터리, 충전, 통신"],
        ["실시간 영상", "사각공간별 최근 촬영/라이브 피드"],
        ["원격 명령", "재촬영, 이동, 긴급복귀, 미션중지"],
        ["조치관리", "담당자, 기한, 완료사진, 재확인"],
        ["보고 자동화", "일일·주간 리포트 생성"],
    ]
    table(s, 6.85, 2.45, rows, [2.1, 3.65], 0.53)

    # 7 tech & AI
    s = slide()
    header(s, "06 TECH & AI", "1차 PoC는 완전자동 판정이 아니라 위험후보 태깅으로 시작합니다", "안전관리자가 최종 확인하는 human-in-the-loop 구조가 현실적입니다.", 7)
    rows1 = [
        ["초소형 드론", "요구사항"],
        ["크기·안전", "손바닥 이하, 보호가드/덕트형, 저속"],
        ["시야", "근접 카메라, LED 라이트"],
        ["위치", "Optical Flow, AprilTag, UWB"],
        ["복귀", "네스트 마커 기반 자동복귀"],
    ]
    rows2 = [
        ["AI 태깅", "위험후보"],
        ["샤프트", "난간 훼손, 낙하물, 하부 접근"],
        ["PIT", "침수, 조도 부족, 이물질"],
        ["천장", "누수, 전선 노출, 화재위험"],
        ["개구부 하부", "안전망 훼손, 낙하물"],
    ]
    table(s, 0.85, 2.35, rows1, [2.0, 3.75], 0.55)
    table(s, 6.75, 2.35, rows2, [2.0, 3.75], 0.55)
    txt(s, "촬영 → 프레임 추출 → 위험후보 태깅 → 관리자 확인 → 조치요청 → 데이터셋 누적", 1.0, 6.0, 11.2, 0.34, 17, "blue", True, "center")

    # 8 PoC + budget
    s = slide()
    header(s, "07 POC & BUDGET", "3~5개 사각공간이면 8주 PoC가 가능합니다", "대우건설 PoC 지원금 기준으로 5천만 원형 1차 실증안을 제안합니다.", 8)
    rows = [
        ["구성", "내용"],
        ["장비", "초소형 드론 2~5대, 소형 네스트 3~5개"],
        ["대상", "샤프트, PIT, 점검구, 동바리 틈, 개구부 하부"],
        ["운영", "일 3회 자동순찰 + 관제센터 수동 호출"],
        ["검증", "이륙, 근접촬영, 복귀, 충전, 위험후보 태깅"],
        ["기간", "8주 PoC"],
    ]
    table(s, 0.85, 2.35, rows, [2.0, 4.9], 0.52)
    note(s, "1차 PoC 예산", "초소형 드론 500만\n네스트 700만\nSAFE-LINK 모듈 1,000만\nAI 태깅 800만\n운영·분석·예비비 포함\n합계 5,000만 원", 8.05, 2.42, 3.75, 2.8, "green")
    note(s, "내부 선행검증", "제출 전 100만~500만 원 규모로 기체·네스트 목업·간이 관제화면 검증 가능", 8.05, 5.5, 3.75, 0.8, "orange")

    # 9 KPI & risk
    s = slide()
    header(s, "08 KPI & RISK", "성과지표와 안전운용 기준을 같이 제시해야 신뢰가 생깁니다", "초소형 드론은 안전성, 개인정보, 충전, 통신 리스크를 처음부터 관리해야 합니다.", 9)
    rows = [
        ["정량 KPI", "목표"],
        ["사각공간/네스트", "3~5개"],
        ["드론 운용", "2~5대, 일 3회"],
        ["자동복귀/충전", "80% 이상"],
        ["영상 전송", "90% 이상"],
        ["위험태깅", "70% 이상"],
        ["조치요청", "위험 발견 후 1분 이내"],
    ]
    table(s, 0.85, 2.35, rows, [2.8, 2.4], 0.48)
    rows2 = [
        ["리스크", "대응"],
        ["충돌/접촉", "보호가드, 저속, 통제시간 운용"],
        ["추락", "경량 기체, 비상정지, 회수 프로토콜"],
        ["개인정보", "촬영 고지, 얼굴 비식별, 접근권한 제한"],
        ["통신장애", "자동복귀, 수동 회수 프로토콜"],
        ["배터리", "과충전 차단, 네스트 온도감시"],
    ]
    table(s, 6.45, 2.35, rows2, [2.15, 3.75], 0.48)

    # 10 value
    s = slide()
    header(s, "09 VALUE", "대우건설이 얻는 것은 드론 장비가 아니라 사각공간의 상시 감각입니다", "현장 하나에서 검증하고 대우건설형 Hyper Safety 패키지로 확장합니다.", 10)
    img(s, "05-safelink-action-report.png", 0.75, 2.25, 5.75, 3.55)
    note(s, "대우건설 요청", "PoC 현장 1개소\n사각공간 협의\n안전관리자 피드백\n비행·촬영 시간 협의", 6.95, 2.35, 4.95, 1.1, "blue")
    note(s, "제공 가치", "작업자 진입 전 위험 확인\n조치보고 자동화\n위치·시간·영상 이력 축적", 6.95, 3.75, 4.95, 1.1, "green")
    note(s, "확장성", "현장별 네스트 추가\n우선구매·공동개발\n정부과제·SaaS화 연계", 6.95, 5.15, 4.95, 1.1, "orange")

    BASE.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles = [
        "SAFE-LINK Micro Drone Nest", "큰 드론이 아니라 작은 안전 감각기관", "진짜 타깃은 좁은 사각공간",
        "네스트는 현장 센서 거점", "정기순찰과 긴급 호출", "관제센터는 작은 눈을 갖는다",
        "기술요구사항과 AI 태깅", "3~5개 공간 8주 PoC", "KPI와 리스크", "대우건설 가치",
    ]
    try:
        f1 = ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf", 38)
        f2 = ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf", 20)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    paths = []
    img_map = {1: "01-micro-drone-nest-hero.png", 3: "02-hard-to-reach-space.png", 4: "03-distributed-nests.png", 6: "04-control-center-micro-feeds.png", 10: "05-safelink-action-report.png"}
    for i, title in enumerate(titles, 1):
        canvas = Image.new("RGB", (1600, 900), "#" + C["paper"])
        d = ImageDraw.Draw(canvas)
        d.text((80, 60), "SAFE-LINK Micro Drone Nest", font=f2, fill="#" + C["blue"])
        d.text((80, 135), title, font=f1, fill="#" + C["ink"])
        d.line((80, 235, 1460, 235), fill="#" + C["line"], width=3)
        if i in img_map and (ASSET / img_map[i]).exists():
            im = Image.open(ASSET / img_map[i]).convert("RGB")
            im.thumbnail((760, 430))
            canvas.paste(im, (80, 285))
        for k, col in enumerate(["blue", "green", "orange", "red"]):
            d.rounded_rectangle((980, 310 + k * 105, 1390, 370 + k * 105), 16, outline="#" + C[col], width=3, fill="white")
        d.text((80, 830), f"{i:02d}", font=f2, fill="#" + C["muted"])
        p = PREVIEW / f"slide_{i:02d}.png"
        canvas.save(p)
        paths.append(p)
    sheet = Image.new("RGB", (1600, 840), "white")
    d = ImageDraw.Draw(sheet)
    for i, p in enumerate(paths):
        im = Image.open(p).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 265
        sheet.paste(im, (x, y))
        d.text((x, y + 210), f"{i+1:02d}", font=f2, fill="#" + C["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build()
    preview()
    print(OUT)
    print(PREVIEW / "contact_sheet.png")
