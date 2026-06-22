from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-magbrake-fall-arrest-20260607"
ASSET = BASE / "assets"
OUT = BASE / "SAFE-LINK_MagBrake_Fall_Arrest_대우현대국가과제_제안서_20260607.pptx"
PREVIEW = BASE / "previews-magbrake-ppt"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8", "ink": "101828", "muted": "667085", "line": "D0D5DD",
    "blue": "175CD3", "green": "027A48", "orange": "B54708", "red": "B42318",
    "white": "FFFFFF", "soft": "F2F4F7", "navy": "0B1220", "note": "FFF7D6",
    "copper": "B86B2B"
}

def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))

def rect(slide, x, y, w, h, fill="white", line="line", radius=True, transparency=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill); shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    if radius:
        try: shp.adjustments[0] = 0.05
        except Exception: pass
    return shp

def txt(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = body; r.font.name = "맑은 고딕"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box

def rule(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color); shp.line.fill.background()

def footer(slide, n):
    txt(slide, "SAFE-LINK MagBrake Fall Arrest / 서원토건·서원이앤에이", 0.55, 7.04, 5.6, 0.24, 9.5, "muted")
    txt(slide, f"{n:02d}", 12.25, 7.04, 0.45, 0.24, 10, "muted", align="right")

def header(slide, label, title, sub, n):
    txt(slide, label, 0.62, 0.34, 6.0, 0.3, 12, "blue", True)
    txt(slide, title, 0.62, 0.74, 11.8, 0.72, 27, "ink", True)
    if sub: txt(slide, sub, 0.64, 1.5, 11.2, 0.4, 14.5, "muted")
    rule(slide, 0.62, 2.02, 12.05)
    footer(slide, n)

def img(slide, name, x, y, w, h):
    p = ASSET / name
    if p.exists(): slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True); txt(slide, name, x, y+h/2-.15, w, .3, 13, "muted", True, "center")

def note(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True); rect(slide, x, y, .08, h, color, color, False)
    txt(slide, head, x+.22, y+.15, w-.42, .32, 14.5, color, True)
    txt(slide, body, x+.22, y+.55, w-.42, h-.66, 13.5, "ink")

def table(slide, x, y, rows, widths, row_h=.48, head="navy"):
    for r,row in enumerate(rows):
        cx=x
        for c,value in enumerate(row):
            fill=head if r==0 else ("white" if r%2 else "soft")
            color="white" if r==0 else "ink"
            rect(slide,cx,y+r*row_h,widths[c],row_h,fill,"line",False)
            txt(slide,value,cx+.08,y+r*row_h+.085,widths[c]-.16,row_h-.12,10.8 if r==0 else 12,color,r==0)
            cx += widths[c]

def build():
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H); blank=prs.slide_layouts[6]
    def slide():
        s=prs.slides.add_slide(blank); rect(s,0,0,W,H,"paper","paper",False); return s

    # 1 Cover
    s=slide(); img(s,"01-site-fall-hazard-srl.png",6.45,0,6.88,7.5)
    txt(s,"DAEWOO · HYUNDAI E&C · NATIONAL R&D",.72,.64,5.0,.28,11.5,"blue",True)
    txt(s,"SAFE-LINK\nMagBrake\nFall Arrest",.72,1.16,5.4,1.82,34,"ink",True)
    txt(s,"와전류 자기제동 기반 스마트 추락방지 및 구조알림 시스템",.74,3.35,5.25,.62,20,"ink",True)
    rule(s,.74,4.18,4.3,"blue")
    txt(s,"기존 안전대·SRL을 대체하지 않고, 감속·감지·위치·구조알림·보고를 보강하는 현실적 안전장비 R&D",.74,4.52,5.25,.8,15,"muted")
    txt(s,"서원토건 / 서원이앤에이 / SAFE-LINK",.74,6.78,3.9,.25,11,"muted")

    # 2 Problem
    s=slide(); header(s,"01 PROBLEM","추락사고의 핵심은 충격과 구조 골든타임입니다","골조공사 고소작업은 개구부·단부·샤프트·발코니에서 반복 위험이 발생합니다.",2)
    img(s,"01-site-fall-hazard-srl.png",.75,2.35,5.75,3.55)
    note(s,"충격하중","추락 순간 작업자 신체와 앵커에 큰 하중이 전달될 수 있음",7.0,2.35,4.85,.88,"red")
    note(s,"구조 지연","매달림 상태에서 위치 파악과 구조 출동이 늦어질 수 있음",7.0,3.45,4.85,.88,"orange")
    note(s,"기록 부족","시간, 위치, 장비상태, 구조조치 이력이 수작업으로 남음",7.0,4.55,4.85,.88,"blue")
    note(s,"데이터 부족","어떤 구역·공정에서 반복 위험이 생기는지 축적이 약함",7.0,5.65,4.85,.88,"green")

    # 3 Correct positioning
    s=slide(); header(s,"02 POSITIONING","자기부상이 아니라 와전류 자기제동입니다","현실적 포지션은 기존 SRL/안전블록을 고도화하는 보조 감속·감지 시스템입니다.",3)
    rows=[["잘못된 표현","수정 표현"],["자기력으로 사람을 띄운다","와전류 자기제동으로 낙하 속도와 충격을 줄인다"],["기존 안전대를 대체한다","기존 안전대·SRL·충격흡수장치를 보조한다"],["바로 현장 사용 가능","시험·인증·더미 낙하검증이 필요한 R&D 시제품"],["자기부상 추락방지","자기제동형 스마트 추락방지"]]
    table(s,1.0,2.55,rows,[5.0,5.85],.64)
    txt(s,"제안의 신뢰도는 ‘무리한 약속’이 아니라 인증 가능한 개발 로드맵에서 나옵니다.",1.0,5.95,11.2,.38,19,"blue",True,"center")

    # 4 Principle
    s=slide(); header(s,"03 PRINCIPLE","와전류 제동은 비접촉 감속 원리를 활용합니다","자석과 도체의 상대운동이 커질수록 운동을 방해하는 방향의 자기장이 생기는 원리입니다.",4)
    img(s,"02-magbrake-device-cutaway.png",.75,2.25,5.75,3.55)
    steps=[("추락", "라인 급출"),("회전", "드럼/디스크"),("와전류", "도체 내부 유도"),("감속", "회전 저항 증가"),("로그", "센서 데이터")]
    for i,(a,b) in enumerate(steps):
        x=6.95+i*1.05
        rect(s,x,3.0,.85,1.15,"white",["red","orange","copper","blue","green"][i],True)
        txt(s,a,x+.06,3.18,.72,.25,11.5,["red","orange","copper","blue","green"][i],True,"center")
        txt(s,b,x+.06,3.62,.72,.25,9.5,"ink",True,"center")
        if i<4: txt(s,"→",x+.87,3.35,.25,.25,13,"muted",True,"center")
    note(s,"현실적 역할","완전 정지가 아니라 기존 SRL 잠김과 충격흡수 사이에서 감속을 보조",7.0,4.85,4.85,1.0,"blue")

    # 5 System concept
    s=slide(); header(s,"04 SYSTEM","추락방지 장비와 SAFE-LINK 구조알림을 하나로 묶습니다","감속 모듈, 하네스 센서, 위치태그, 앵커센서, 구조알림 게이트웨이로 구성합니다.",5)
    nodes=[("Harness","IMU 추락감지","blue"),("SRL","잠김 구조","green"),("MagBrake","와전류 감속","orange"),("Location","UWB/BLE/NFC","blue"),("SAFE-LINK","알림·조치보고","red")]
    for i,(a,b,c) in enumerate(nodes):
        x=.75+i*2.45
        rect(s,x,3.0,1.85,1.25,"white",c,True)
        txt(s,a,x+.12,3.18,1.6,.28,14,c,True,"center")
        txt(s,b,x+.12,3.64,1.6,.28,11.5,"ink",True,"center")
        if i<4: txt(s,"→",x+1.86,3.4,.35,.3,20,"muted",True,"center")
    txt(s,"추락 이벤트 → 위치/작업자 식별 → 긴급 알림 → 구조 상태 관리 → 초기 보고서 자동 생성",1.0,5.4,11.2,.45,20,"ink",True,"center")

    # 6 Modules
    s=slide(); header(s,"05 MODULES","제품이 아니라 모듈별 R&D 과제로 쪼개야 현실적입니다","기계·센서·위치·플랫폼을 분리해 단계별 실증합니다.",6)
    rows=[["모듈","기능","난이도"],["MagBrake 감속","회전 드럼/디스크 감속","높음"],["Smart Harness","IMU 자유낙하·급정지 감지","중간"],["SRL Interface","기존 안전블록 결합 실험","높음"],["Shock Absorber","웨빙/댐퍼 완충","중간"],["UWB/BLE Tag","작업자·구역 위치 식별","중간"],["SAFE-LINK Gateway","추락 이벤트 서버 전송","낮음~중간"],["Auto Report","초기·재발방지 보고서","낮음~중간"]]
    table(s,.9,2.35,rows,[3.05,5.7,2.0],.46)

    # 7 SAFE-LINK control
    s=slide(); header(s,"06 RESCUE FLOW","구조 골든타임은 자동 알림과 위치 식별에서 시작됩니다","추락감지 즉시 현장소장, 안전관리자, 협력업체 반장에게 구조 워크플로우를 전송합니다.",7)
    img(s,"04-safelink-rescue-control.png",.75,2.35,5.85,3.55)
    note(s,"즉시 알림","작업자, 위치, 앵커, 이벤트 시간을 자동 전송",6.95,2.4,4.95,.95,"red")
    note(s,"구조상태","접수, 출동, 구조중, 완료 상태를 SAFE-LINK에 기록",6.95,3.6,4.95,.95,"orange")
    note(s,"보고 자동화","사고 초기보고서와 재발방지 조치보고서를 자동 생성",6.95,4.8,4.95,.95,"blue")

    # 8 Lab/test
    s=slide(); header(s,"07 TEST STRATEGY","사람 대상 테스트가 아니라 더미 낙하시험부터 시작합니다","인증 전 현장 사용이 아니라 벤치·더미·교육장 테스트로 데이터를 확보합니다.",8)
    img(s,"03-dummy-drop-test-lab.png",.75,2.35,5.85,3.55)
    rows=[["시험","측정 데이터"],["벤치 회전시험","회전속도, 제동력, 발열"],["5/10/20kg 단계하중","감속거리, 온도, 반복성"],["40/60/80kg 더미낙하","최대하중, 정지시간, 충격량"],["기존 SRL 비교","SRL 단독 vs MagBrake 보조"],["반복시험","30회 이상 데이터 축적"]]
    table(s,6.85,2.42,rows,[2.45,3.35],.53)

    # 9 Roadmap
    s=slide(); header(s,"08 ROADMAP","18개월 개발, 3~6개월 내 가시 PoC", "건설사 제안은 단기 데모와 장기 인증 로드맵을 분리해야 합니다.",9)
    rows=[["단계","기간","목표","산출물"],["Phase 0","0~1개월","기술정의·IP·협력구조","개념서/NDA/역할표"],["Phase 1","1~3개월","센서형 하네스 PoC","IMU 감지·SAFE-LINK 알림"],["Phase 2","3~6개월","와전류 벤치 실험","감속모듈·데이터"],["Phase 3","6~9개월","더미 낙하 테스트","80kg 더미·하중데이터"],["Phase 4","9~12개월","SRL 통합 시제품","통합 구조·반복시험"],["Phase 5","12~18개월","현장 비작업 테스트","시험기관·인증 검토"]]
    table(s,.75,2.35,rows,[1.55,1.45,3.65,4.15],.49)

    # 10 Standards & safety realism
    s=slide(); header(s,"09 CERTIFICATION REALISM","안전장비는 기준·시험·인증을 우회할 수 없습니다","대외 제안에서는 ‘인증 전 현장 사용’이 아니라 ‘시험기관 검증을 전제로 한 R&D’로 표현해야 합니다.",10)
    note(s,"기준 검토","개인추락방지시스템, SRL, 충격흡수, 앵커, 하중 기준 검토",.9,2.45,3.55,1.35,"blue")
    note(s,"시험기관","KOSHA/KCL/KTR 등 시험·인증 가능성 사전 협의",4.9,2.45,3.55,1.35,"green")
    note(s,"위험관리","발열, 자석 이탈, 과하중, 반복사용, 센서 오탐을 시험 항목화",8.9,2.45,3.55,1.35,"orange")
    rows=[["표현 원칙","제안서 문구"],["대체품 아님","기존 안전대·SRL 보조 시스템"],["현장 즉시사용 아님","더미시험·시험기관 검증 후 단계 적용"],["성능 과장 금지","감속 보조·충격완화·구조알림 중심"],["SAFE-LINK 가치","위치·알림·조치·보고 자동화"]]
    table(s,1.45,4.35,rows,[3.4,6.8],.45)

    # 11 Construction company fit
    s=slide(); header(s,"10 BUILDER FIT","대우건설·현대엔지니어링에는 ‘현장 PoC + 공동개발’로 제안합니다","건설사는 안전장비 제조보다 현장성, PoC, 데이터, 브랜드 효과를 봅니다.",11)
    note(s,"대우건설","Hyper Safety & AI 공모전\n골조현장 추락위험 PoC\nSAFE-LINK 연동 보고 자동화",.85,2.5,3.65,1.65,"blue")
    note(s,"현대엔지니어링","스마트건설·안전기술 제안\n고소작업 안전장비 R&D\n현장 교육장/모형현장 실증",4.85,2.5,3.65,1.65,"green")
    note(s,"서원토건","철근콘크리트 골조 위험구역 정의\n현장 시나리오·데이터 제공\n전문건설업 관점 PoC",8.85,2.5,3.65,1.65,"orange")
    txt(s,"건설사에는 ‘제품 구매’보다 ‘공동개발 가능한 안전 R&D 과제’로 제안하는 것이 현실적입니다.",1.0,5.45,11.2,.45,19,"ink",True,"center")

    # 12 National R&D
    s=slide(); header(s,"11 NATIONAL R&D", "국가제안사업은 스마트 PPE·중대재해 예방·현장 데이터로 묶습니다","기계제동 단독보다 센서·플랫폼·구조알림을 결합해야 과제성이 커집니다.",12)
    img(s,"05-rnd-roadmap-collaboration.png",.75,2.3,5.85,3.55)
    rows=[["축","과제화 포인트"],["스마트 PPE","하네스 센서, 위치태그, 앵커 상태 감지"],["로봇/기계안전","와전류 제동 모듈, 더미 낙하 검증"],["AI/데이터","추락 이벤트 데이터셋, 위험구역 분석"],["중대재해 예방","구조 골든타임, 보고 자동화"],["산학협력","기계·전기·안전공학 연구실 공동개발"]]
    table(s,6.85,2.42,rows,[2.2,3.6],.53)

    # 13 IP
    s=slide(); header(s,"12 IP STRATEGY","특허는 ‘자기제동 + 센서 + 구조알림 + 보고’ 결합으로 잡습니다","단일 브레이크보다 건설현장 운용 시스템 청구항이 방어력이 큽니다.",13)
    rows=[["청구항 축","내용"],["기계 구조","SRL 드럼/디스크에 결합된 와전류 감속 모듈"],["센서 이벤트","IMU, 하중, 온도, 앵커 상태 기반 추락 판정"],["위치 식별","UWB/BLE/NFC로 작업자·구역·앵커 매칭"],["구조 알림","안전관리자·반장·본사 알림 및 상태 추적"],["보고 자동화","추락 이벤트 기반 초기보고서·재발방지 보고서 생성"],["데이터 검증","이벤트 로그, 해시, 시험데이터 저장"]]
    table(s,1.0,2.35,rows,[2.6,8.25],.54)

    # 14 Collaboration ask
    s=slide(); header(s,"13 COLLABORATION ASK","협업 요청은 역할을 명확히 나눠야 실행됩니다","대우·현대·국가과제 모두 현장·시험·제조·플랫폼 역할 분리가 필요합니다.",14)
    rows=[["주체","역할"],["건설사","PoC 현장/교육장 제공, 안전관리자 피드백, 현장 위험구역 선정"],["서원토건","골조공정 위험 시나리오, 전문건설업 현장성, 실증 운영"],["서원이앤에이/SAFE-LINK","알림·조치·보고 플랫폼, 데이터 구조, 대시보드"],["SRL/안전장비 협력사","기구 설계, 제품화, 인증 대응"],["대학/연구기관","와전류 해석, 센서 알고리즘, 시험 데이터 분석"],["시험기관","더미 낙하시험, 기준 검토, 인증 로드맵"]]
    table(s,.9,2.35,rows,[2.6,8.5],.53)

    # 15 Closing
    s=slide(); header(s,"14 CLOSING","MagBrake의 가치는 ‘추락을 없앤다’가 아니라 구조 가능한 시간과 데이터를 만든다는 데 있습니다","현실적이고 검증 가능한 스마트 추락방지 R&D로 포지셔닝합니다.",15)
    note(s,"1단계","센서형 하네스 + SAFE-LINK 긴급알림\n1~3개월 내 데모 가능",.9,2.55,3.55,1.55,"blue")
    note(s,"2단계","와전류 자기제동 벤치 실험\n3~6개월 내 데이터 확보",4.9,2.55,3.55,1.55,"orange")
    note(s,"3단계","더미 낙하·SRL 통합·시험기관 검토\n6~18개월 R&D",8.9,2.55,3.55,1.55,"green")
    txt(s,"SAFE-LINK MagBrake Fall Arrest System",1.0,5.15,11.2,.42,28,"ink",True,"center")
    txt(s,"기존 추락방지 장비를 더 똑똑하게 만들고, 추락 후 구조와 보고를 자동화하는 현실적 안전기술",1.3,5.78,10.6,.42,18,"muted",True,"center")

    BASE.mkdir(parents=True, exist_ok=True); prs.save(OUT)

def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles=["Cover","Problem","Positioning","Principle","System","Modules","Rescue Flow","Test Strategy","Roadmap","Certification","Builder Fit","National R&D","IP Strategy","Collaboration","Closing"]
    try:
        f1=ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf",34); f2=ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf",18)
    except Exception: f1=f2=ImageFont.load_default()
    image_map={1:"01-site-fall-hazard-srl.png",2:"01-site-fall-hazard-srl.png",4:"02-magbrake-device-cutaway.png",7:"04-safelink-rescue-control.png",8:"03-dummy-drop-test-lab.png",12:"05-rnd-roadmap-collaboration.png"}
    paths=[]
    for i,t in enumerate(titles,1):
        canvas=Image.new("RGB",(1600,900),"#"+C["paper"]); d=ImageDraw.Draw(canvas)
        d.text((80,60),"SAFE-LINK MagBrake Fall Arrest",font=f2,fill="#"+C["blue"]); d.text((80,130),t,font=f1,fill="#"+C["ink"])
        d.line((80,230,1460,230),fill="#"+C["line"],width=3)
        if i in image_map and (ASSET/image_map[i]).exists():
            im=Image.open(ASSET/image_map[i]).convert("RGB"); im.thumbnail((760,430)); canvas.paste(im,(80,285))
        for k,col in enumerate(["blue","green","orange","red"]):
            d.rounded_rectangle((980,310+k*105,1390,370+k*105),16,outline="#"+C[col],width=3,fill="white")
        d.text((80,830),f"{i:02d}",font=f2,fill="#"+C["muted"])
        p=PREVIEW/f"slide_{i:02d}.png"; canvas.save(p); paths.append(p)
    sheet=Image.new("RGB",(1600,1100),"white"); d=ImageDraw.Draw(sheet)
    for i,p in enumerate(paths):
        im=Image.open(p).resize((300,169)); x=25+(i%5)*315; y=25+(i//5)*335
        sheet.paste(im,(x,y)); d.text((x,y+176),f"{i+1:02d} {titles[i]}",font=f2,fill="#"+C["ink"])
    sheet.save(PREVIEW/"contact_sheet.png")

if __name__=="__main__":
    build(); preview(); print(OUT); print(PREVIEW/"contact_sheet.png")
