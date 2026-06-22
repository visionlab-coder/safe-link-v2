from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "generated" / "daewoo-hyper-safety-20260604"
PPTX_PATH = OUT_DIR / "SAFE-LINK_UNDERGROUND_대우건설_제안서_사람작업형_고도화_20260605.pptx"
PREVIEW_DIR = OUT_DIR / "previews-humanized-ppt"

W, H = 13.333, 7.5
PX_W, PX_H = 1600, 900

PALETTE = {
    "paper": "FAFAF8",
    "ink": "111827",
    "muted": "667085",
    "soft": "F1F3F5",
    "line": "D9DEE7",
    "blue": "1D4ED8",
    "green": "047857",
    "amber": "B45309",
    "red": "B91C1C",
    "violet": "6D28D9",
    "navy": "0F172A",
    "white": "FFFFFF",
    "note": "FFF7D6",
}


def rgb(name_or_hex):
    value = PALETTE.get(name_or_hex, name_or_hex).strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def text(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = body
    run.font.name = "맑은 고딕"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def shape(slide, x, y, w, h, fill="white", line="line", radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line)
    s.line.width = Pt(1)
    if radius:
        try:
            s.adjustments[0] = 0.06
        except Exception:
            pass
    return s


def line(slide, x, y, w, color="line"):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(color)
    s.line.fill.background()
    return s


def footer(slide, n):
    text(slide, "SAFE-LINK UNDERGROUND / 서원토건", 0.55, 7.05, 4.6, 0.22, 10, "muted")
    text(slide, f"{n:02d}", 12.25, 7.05, 0.45, 0.22, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    text(slide, label, 0.62, 0.36, 3.6, 0.26, 11, "blue", True)
    text(slide, title, 0.62, 0.72, 11.5, 0.68, 28, "ink", True)
    if sub:
        text(slide, sub, 0.64, 1.46, 10.8, 0.38, 14, "muted")
    line(slide, 0.62, 1.96, 12.05)
    footer(slide, n)


def note(slide, label, body, x, y, w, h, accent="blue"):
    shape(slide, x, y, w, h, "white", "line", True)
    shape(slide, x, y, 0.08, h, accent, accent, False)
    text(slide, label, x + 0.22, y + 0.18, w - 0.4, 0.28, 13, accent, True)
    text(slide, body, x + 0.22, y + 0.55, w - 0.38, h - 0.68, 15, "ink")


def tag(slide, label, x, y, w, color="blue"):
    shape(slide, x, y, w, 0.34, "white", color, True)
    text(slide, label, x + 0.08, y + 0.05, w - 0.16, 0.22, 10, color, True, "center")


def add_image(slide, path, x, y, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        return True
    shape(slide, x, y, w, h, "soft", "line", True)
    text(slide, "화면 캡처", x, y + h / 2 - 0.15, w, 0.3, 14, "muted", True, "center")
    return False


def find_capture(keyword):
    base = ROOT / "docs" / "generated"
    candidates = list(base.glob("특허사무실_제출용_SAFE-LINK_실제앱화면_청구항별_20260601/*.png"))
    for p in candidates:
        if keyword in p.name:
            return p
    return None


def table(slide, x, y, rows, widths, row_h=0.5, header_fill="navy"):
    for r, row in enumerate(rows):
        cx = x
        for c, value in enumerate(row):
            fill = header_fill if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            weight = r == 0
            shape(slide, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            text(slide, value, cx + 0.08, y + r * row_h + 0.1, widths[c] - 0.16, row_h - 0.14, 12 if r == 0 else 13, color, weight)
            cx += widths[c]


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        shape(s, 0, 0, W, H, "paper", "paper", False)
        return s

    dashboard = find_capture("관리자대시보드")
    tbm = find_capture("위험성평가")
    nfc = find_capture("NFC")
    live = find_capture("실시간통역")
    report = find_capture("보고서")

    # 1. Editorial cover
    s = slide()
    text(s, "대우건설 Hyper Safety & AI Open Innovation", 0.65, 0.58, 8.0, 0.32, 13, "blue", True)
    text(s, "지하 골조 안전 사각을\n작업면 데이터로 메우다", 0.65, 1.22, 8.2, 1.55, 36, "ink", True)
    text(s, "SAFE-LINK UNDERGROUND", 0.68, 3.12, 6.0, 0.34, 18, "navy", True)
    text(s, "AI-RAN · 로키드 스마트글라스 · AI 로봇개 · 다국어 TBM · 작업중지 증빙", 0.7, 3.55, 7.4, 0.56, 16, "muted")
    line(s, 0.72, 4.32, 5.6, "blue")
    shape(s, 8.7, 0.78, 3.45, 5.4, "white", "line", True)
    text(s, "제안의 한 점", 9.0, 1.08, 2.4, 0.3, 14, "blue", True)
    text(s, "대형 건설사의 스마트 안전은 강하지만,\n가장 위험한 지하 철근콘크리트 작업면은 통신·시야·언어·증빙이 동시에 약해집니다.", 9.0, 1.55, 2.75, 1.35, 18, "ink", True)
    text(s, "서원토건은 그 작업면을 아는 전문건설사입니다.", 9.0, 4.25, 2.75, 0.75, 17, "green", True)
    text(s, "2026.06.05 / 서원토건", 0.72, 6.83, 3.6, 0.25, 11, "muted")

    # 2. Field problem
    s = slide()
    header(s, "FIELD PAIN", "지하 철근콘크리트 현장은 안전관리의 마지막 사각입니다", "앱·CCTV·센서가 있어도, 지하 작업면에서 연결되지 않으면 작동하지 않습니다.", 2)
    note(s, "통신", "지하층, PIT, 코어 내부에서 모바일 안전앱과 신고 체계가 끊김", 0.8, 2.55, 3.55, 1.55, "blue")
    note(s, "시야", "CCTV는 고정형이라 동바리·거푸집·타설처럼 매일 바뀌는 작업면을 따라가지 못함", 4.9, 2.55, 3.55, 1.55, "amber")
    note(s, "언어", "외국인 근로자가 TBM을 ‘들었다’와 ‘이해했다’는 별개 문제", 9.0, 2.55, 3.3, 1.55, "green")
    text(s, "핵심 문제는 장비 부족이 아니라\n위험이 발생하는 작업면과 안전 데이터가 분리되어 있다는 점입니다.", 1.1, 5.0, 10.9, 0.8, 25, "ink", True, "center")

    # 3. Why Seowon
    s = slide()
    header(s, "WHY SEOWON", "철근콘크리트 전문건설업자가 가장 잘 아는 위험", "공정 단위 위험을 알아야 AI도, PoC도, 현장 수용성도 설득됩니다.", 3)
    rows = [
        ["공정", "현장 위험", "SAFE-LINK 대응"],
        ["철근 배근", "찔림·협착·철근 다발 전도", "작업조별 TBM + 구역 안전여권"],
        ["거푸집·동바리", "미고정·수평재 누락·해체순서 오류", "점검 게이트 + 글라스 촬영"],
        ["콘크리트 타설", "측압·호스 휘핑·동선 충돌", "Pour-Go / No-Go AI"],
        ["지하층·PIT", "통신단절·조도부족·산소결핍", "AI-RAN + 로봇개 선행 순찰"],
    ]
    table(s, 0.9, 2.45, rows, [2.2, 4.3, 5.0], 0.62)
    text(s, "원청 관제의 빈틈을 전문건설업 작업면 데이터로 보완", 1.0, 6.0, 11.2, 0.34, 20, "blue", True, "center")

    # 4. System architecture, not cards
    s = slide()
    header(s, "SYSTEM", "통신·시야·순찰·증빙을 하나의 안전 루프로 연결", "각 기술을 따로 파는 것이 아니라, 지하 작업허가 체계로 묶습니다.", 4)
    y = 3.0
    xs = [0.75, 3.85, 6.95, 10.05]
    labels = [("AI-RAN", "지하 통신 버블", "blue"), ("Rokid Glass", "작업조장 시야", "green"), ("Robot Dog", "선행 순찰", "violet"), ("SAFE-LINK", "TBM·신고·증빙", "amber")]
    for x, (a, b, c) in zip(xs, labels):
        shape(s, x, y, 2.25, 1.25, "white", c, True)
        text(s, a, x + 0.15, y + 0.18, 1.95, 0.3, 16, c, True, "center")
        text(s, b, x + 0.15, y + 0.62, 1.95, 0.3, 14, "ink", True, "center")
    for x in [3.1, 6.2, 9.3]:
        text(s, "→", x, y + 0.38, 0.45, 0.35, 24, "muted", True, "center")
    text(s, "작업 전 확인 → 작업 중 감지 → 위험신고/작업중지 → 조치 이력 → 보고서 자동화", 1.1, 5.25, 11.0, 0.46, 22, "ink", True, "center")

    # 5. Actual product proof
    s = slide()
    header(s, "CURRENT ASSET", "SAFE-LINK는 이미 현장 운영 화면을 갖고 있습니다", "제안은 백지 아이디어가 아니라 기존 플랫폼에 AI-RAN·글라스·로봇개를 붙이는 확장입니다.", 5)
    add_image(s, dashboard, 0.75, 2.35, 3.8, 2.15)
    add_image(s, tbm, 4.85, 2.35, 3.8, 2.15)
    add_image(s, nfc, 8.95, 2.35, 3.35, 2.15)
    text(s, "관리자 관제", 0.75, 4.72, 3.8, 0.28, 14, "ink", True, "center")
    text(s, "TBM·위험성평가", 4.85, 4.72, 3.8, 0.28, 14, "ink", True, "center")
    text(s, "NFC/QR 인증", 8.95, 4.72, 3.35, 0.28, 14, "ink", True, "center")
    text(s, "대우건설 PoC에서는 이 기존 기능을 지하 철근콘크리트 공정에 맞게 재구성합니다.", 1.05, 5.75, 11.0, 0.38, 18, "blue", True, "center")

    # 6. Three proposed operating routines
    s = slide()
    header(s, "OPERATING ROUTINES", "심사위원이 기억할 세 개의 운영 개념", "기술명보다 현장에서 어떻게 쓰이는지가 중요합니다.", 6)
    shape(s, 0.9, 2.55, 3.3, 2.9, "note", "line", True)
    text(s, "01", 1.15, 2.83, 0.55, 0.3, 15, "blue", True)
    text(s, "Rebar-Zone\nSafety Passport", 1.15, 3.22, 2.55, 0.8, 22, "ink", True)
    text(s, "작업구역별 TBM·퀴즈·서명·로봇순찰 완료 여부를 묶어 작업허가 상태를 표시", 1.15, 4.28, 2.6, 0.75, 15, "ink")
    shape(s, 5.0, 2.55, 3.3, 2.9, "white", "line", True)
    text(s, "02", 5.25, 2.83, 0.55, 0.3, 15, "green", True)
    text(s, "Pour-Go /\nNo-Go AI", 5.25, 3.22, 2.55, 0.8, 22, "ink", True)
    text(s, "타설 전 동바리·거푸집·개구부·TBM·통신·미처리 신고를 자동 확인", 5.25, 4.28, 2.6, 0.75, 15, "ink")
    shape(s, 9.1, 2.55, 3.3, 2.9, "white", "line", True)
    text(s, "03", 9.35, 2.83, 0.55, 0.3, 15, "violet", True)
    text(s, "Multilingual\nStop-Work Relay", 9.35, 3.22, 2.55, 0.8, 21, "ink", True)
    text(s, "외국어 위험신고를 관리자 조치와 모국어 회신, 해시 증빙까지 연결", 9.35, 4.28, 2.6, 0.75, 15, "ink")

    # 7. Rokid and robot dog
    s = slide()
    header(s, "FIELD DEVICES", "로키드 글라스와 로봇개는 보여주기 장비가 아닙니다", "지하 작업면의 시야와 선행 위험탐지를 SAFE-LINK 데이터로 연결합니다.", 7)
    shape(s, 0.95, 2.45, 5.3, 3.25, "white", "line", True)
    text(s, "로키드 스마트글라스", 1.25, 2.78, 4.5, 0.35, 21, "green", True)
    text(s, "안드로이드 기반 · 영상 촬영 · 열화상 카메라\n작업조장 시야를 원격 안전관리자와 공유\n위험 캡처를 SAFE-LINK 이벤트로 저장", 1.25, 3.45, 4.45, 1.1, 16, "ink")
    tag(s, "착용형 관제", 1.25, 4.95, 1.45, "green")
    tag(s, "열화상", 2.9, 4.95, 1.05, "green")
    tag(s, "영상증빙", 4.15, 4.95, 1.25, "green")
    shape(s, 7.05, 2.45, 5.3, 3.25, "white", "line", True)
    text(s, "AI 순찰 로봇개", 7.35, 2.78, 4.5, 0.35, 21, "violet", True)
    text(s, "모빌리오 협업 또는 상용 로봇 도입 검토\n지하 선행 순찰 · 통신 품질 확인 · 센서 확장\n순찰 결과를 TBM/작업중지/보고서로 연결", 7.35, 3.45, 4.45, 1.1, 16, "ink")
    tag(s, "선행 순찰", 7.35, 4.95, 1.35, "violet")
    tag(s, "이동 센서", 8.9, 4.95, 1.35, "violet")
    tag(s, "위험기록", 10.45, 4.95, 1.25, "violet")

    # 8. AI-RAN
    s = slide()
    header(s, "AI-RAN STRATEGY", "AI-RAN은 ‘통신망’이 아니라 안전 인프라로 제안합니다", "협력은 확정 표현이 아니라 추진 수위로 안전하게 기재합니다.", 8)
    text(s, "추진 표현", 0.95, 2.55, 2.0, 0.3, 16, "blue", True)
    shape(s, 0.95, 3.05, 11.4, 1.5, "white", "line", True)
    text(s, "경희대학교 홍인기 교수와 MOU 체결 추진\nAI-RAN/O-RAN 기반 지하 건설현장 안전통신 국책과제 또는 대우건설 공동제안 추진 예정", 1.25, 3.33, 10.8, 0.75, 20, "ink", True, "center")
    text(s, "주의: MOU 체결 전에는 ‘확정 협력기관’으로 단정하지 않고, 공동제안·자문·과제화 추진으로 표현", 1.05, 5.15, 11.0, 0.35, 16, "red", True, "center")

    # 9. PoC plan
    s = slide()
    header(s, "POC PLAN", "8주 안에 보여줄 것은 기술 과시가 아니라 현장 적용성입니다", "지하층 골조 또는 지하주차장 철근콘크리트 공정 1개 구역을 대상으로 합니다.", 9)
    weeks = [("1-2주", "현장 조사\n통신 음영 측정"), ("3-4주", "TBM·퀴즈\n전자서명 운영"), ("5-6주", "글라스 관제\n로봇개 순찰"), ("7-8주", "효과 측정\n확장안 도출")]
    x = 0.9
    for i, (wk, body) in enumerate(weeks):
        shape(s, x + i * 3.0, 2.75, 2.35, 2.2, "white", ["blue", "green", "violet", "amber"][i], True)
        text(s, wk, x + i * 3.0 + 0.2, 3.05, 1.9, 0.32, 18, ["blue", "green", "violet", "amber"][i], True, "center")
        text(s, body, x + i * 3.0 + 0.2, 3.62, 1.9, 0.72, 18, "ink", True, "center")
    text(s, "측정: TBM 참여율, 이해도, 작업중지 처리시간, 통신 커버리지, 보고서 작성시간 절감", 1.0, 5.65, 11.0, 0.36, 17, "ink", True, "center")

    # 10. Expected value
    s = slide()
    header(s, "VALUE", "대우건설이 얻는 것은 ‘신기술’보다 현장 운영 개선입니다", "평가 기준인 현업 적용성, 기술 차별성, PoC 명확성에 직접 대응합니다.", 10)
    rows = [
        ["평가 기준", "제안서에서 보여줄 답"],
        ["현업 적용성", "지하 철근콘크리트 공정 1개 구역에서 바로 PoC 가능"],
        ["기술 우수성", "AI-RAN + 글라스 + 로봇개 + SAFE-LINK 증빙 결합"],
        ["수용성", "근로자는 QR/NFC, 관리자는 기존 TBM·위험성평가 흐름 사용"],
        ["경제성", "수기 TBM·보고서·미처리 신고 누락을 줄여 관리시간 절감"],
    ]
    table(s, 1.05, 2.55, rows, [2.5, 8.35], 0.66)

    # 11. Collaboration
    s = slide()
    header(s, "COLLABORATION", "협업 구조는 단순해야 실행됩니다", "각 주체가 무엇을 제공하는지 명확히 나누어 대우건설의 부담을 줄입니다.", 11)
    note(s, "서원토건", "철근콘크리트 공정 위험 라이브러리\nSAFE-LINK 운영 및 PoC 시나리오", 0.8, 2.55, 3.45, 1.65, "blue")
    note(s, "대우건설", "PoC 현장 제공\nTBM·위험성평가 양식 및 안전관리자 피드백", 4.85, 2.55, 3.45, 1.65, "green")
    note(s, "협력 파트너", "경희대 AI-RAN 추진\n로키드 글라스\n모빌리오 또는 상용 로봇개", 8.9, 2.55, 3.45, 1.65, "violet")
    text(s, "확정된 것과 추진 중인 것을 구분해 쓰면 신뢰도가 올라갑니다.", 1.0, 5.35, 11.2, 0.42, 22, "ink", True, "center")

    # 12. Closing
    s = slide()
    header(s, "CLOSE", "SAFE-LINK UNDERGROUND는 지하 골조 안전의 빈칸을 겨냥합니다", "", 12)
    add_image(s, report, 0.9, 2.25, 4.4, 2.5)
    shape(s, 6.05, 2.28, 5.9, 2.35, "white", "line", True)
    text(s, "제안의 결론", 6.45, 2.62, 2.0, 0.3, 16, "blue", True)
    text(s, "대우건설의 지상 스마트 안전과\n서원토건의 지하 작업면 데이터를 연결하면,\n통신·시야·언어·증빙 사각을 동시에 줄일 수 있습니다.", 6.45, 3.05, 5.05, 1.0, 22, "ink", True)
    text(s, "첫 PoC는 지하 철근콘크리트 1개 구역이면 충분합니다.", 1.0, 5.65, 11.0, 0.4, 23, "green", True, "center")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def font(size=26, bold=False):
    paths = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for p in paths:
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def preview():
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    slides = [
        ("지하 골조 안전 사각을\n작업면 데이터로 메우다", "SAFE-LINK UNDERGROUND"),
        ("지하 철근콘크리트 현장은\n안전관리의 마지막 사각입니다", "통신 · 시야 · 언어"),
        ("전문건설업자가 가장 잘 아는 위험", "철근 · 거푸집 · 동바리 · 타설 · PIT"),
        ("통신·시야·순찰·증빙을\n하나의 안전 루프로", "AI-RAN + Glass + Robot + SAFE-LINK"),
        ("SAFE-LINK는 이미\n현장 운영 화면을 갖고 있습니다", "기존 플랫폼 기반 확장"),
        ("심사위원이 기억할\n세 개의 운영 개념", "Safety Passport / Pour Gate / Stop-Work"),
        ("로키드 글라스와 로봇개는\n보여주기 장비가 아닙니다", "시야와 선행탐지를 데이터로 연결"),
        ("AI-RAN은 안전 인프라로 제안", "경희대 MOU 및 공동제안 추진"),
        ("8주 안에 보여줄 것은\n현장 적용성입니다", "조사 → 운영 → 관제/순찰 → 측정"),
        ("대우건설이 얻는 것은\n현장 운영 개선입니다", "평가 기준 직접 대응"),
        ("협업 구조는 단순해야 실행됩니다", "서원토건 / 대우건설 / 파트너"),
        ("지하 골조 안전의\n빈칸을 겨냥합니다", "첫 PoC는 1개 구역이면 충분"),
    ]
    paths = []
    for i, (head, sub) in enumerate(slides, 1):
        img = Image.new("RGB", (PX_W, PX_H), "#" + PALETTE["paper"])
        d = ImageDraw.Draw(img)
        d.text((80, 60), "SAFE-LINK UNDERGROUND", font=font(24, True), fill="#" + PALETTE["blue"])
        d.multiline_text((80, 155), head, font=font(62, True), fill="#" + PALETTE["ink"], spacing=16)
        d.text((84, 410), sub, font=font(32, True), fill="#" + PALETTE["muted"])
        d.line((80, 510, 1400, 510), fill="#" + PALETTE["line"], width=3)
        for k, c in enumerate(["blue", "green", "violet", "amber"]):
            x = 100 + k * 330
            d.rounded_rectangle((x, 575, x + 260, 700), 18, fill="#FFFFFF", outline="#" + PALETTE[c], width=3)
        d.text((80, 815), "서원토건 / 대우건설 Hyper Safety & AI", font=font(20), fill="#" + PALETTE["muted"])
        d.text((1460, 815), f"{i:02d}", font=font(20), fill="#" + PALETTE["muted"])
        out = PREVIEW_DIR / f"slide_{i:02d}.png"
        img.save(out)
        paths.append(out)
    sheet = Image.new("RGB", (PX_W, 3 * 270), "white")
    d = ImageDraw.Draw(sheet)
    for i, p in enumerate(paths):
        im = Image.open(p).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 260
        sheet.paste(im, (x, y))
        d.text((x, y + 210), f"{i+1:02d}", font=font(18, True), fill="#" + PALETTE["ink"])
    sheet.save(PREVIEW_DIR / "contact_sheet.png")


if __name__ == "__main__":
    build_deck()
    preview()
    print(PPTX_PATH)
    print(PREVIEW_DIR / "contact_sheet.png")
