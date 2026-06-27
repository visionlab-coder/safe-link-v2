from pathlib import Path

from pptx import Presentation


p = Path("docs/generated/SAFE-LINK_관리자_교육용_사용설명서_최신업데이트_20260601.pptx")
prs = Presentation(p)
print(f"exists={p.exists()} bytes={p.stat().st_size} slides={len(prs.slides)}")
for i, slide in enumerate(prs.slides, 1):
    texts = []
    pictures = 0
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == 13:
            pictures += 1
        if getattr(shape, "has_text_frame", False):
            text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            if text:
                texts.append(text[:90])
    title = texts[0] if texts else ""
    print(f"{i:02d} pictures={pictures} text_shapes={len(texts)} title={title}")
