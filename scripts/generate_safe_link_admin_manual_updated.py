# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
SCREEN_DIR = OUT / "real-screens"
PREVIEW_DIR = OUT / "previews-updated"
PPTX = OUT / "SAFE-LINK_admin_training_manual_updated.pptx"
PPTX_KO = OUT / "SAFE-LINK_관리자_교육용_사용설명서_업데이트본.pptx"
CONTACT = PREVIEW_DIR / "SAFE-LINK_admin_training_manual_updated_contact.png"

INK = "111827"
MUTED = "4B5563"
PAPER = "F8F5EF"
WHITE = "FFFFFF"
LINE = "D9D2C7"
BLUE = "2563EB"
GREEN = "16A34A"
AMBER = "F59E0B"
RED = "DC2626"
SLATE = "263447"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=True, alpha=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if alpha is not None:
        shape.fill.transparency = alpha
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, text, x, y, w, h, size=24, color=INK, bold=True, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullet_list(slide, items: Iterable[str], x, y, w, h, size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(INK)
    return box


def add_image(slide, file_name: str, x=0.55, y=1.35, w=8.25, h=5.55):
    path = SCREEN_DIR / file_name
    if not path.exists():
        rect(slide, x, y, w, h, "F3F4F6", LINE)
        textbox(slide, f"화면 캡처 없음\n{file_name}", x + 0.4, y + 2.1, w - 0.8, 0.7, 24, MUTED, True, PP_ALIGN.CENTER)
        return
    with Image.open(path) as img:
        iw, ih = img.size
    target = w / h
    source = iw / ih
    if source >= target:
        ph = h
        pw = h * source
    else:
        pw = w
        ph = w / source
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(pw), height=Inches(ph))
    pic.crop_left = max(0, (pw - w) / pw / 2)
    pic.crop_right = max(0, (pw - w) / pw / 2)
    pic.crop_top = max(0, (ph - h) / ph / 2)
    pic.crop_bottom = max(0, (ph - h) / ph / 2)


def base_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    textbox(slide, title, 0.55, 0.35, 8.3, 0.55, 25, INK, True)
    if subtitle:
        textbox(slide, subtitle, 0.55, 0.88, 8.5, 0.35, 15, MUTED, True)
    textbox(slide, "SAFE-LINK v2.0 관리자 교육", 9.55, 0.45, 3.1, 0.32, 14, MUTED, True, PP_ALIGN.RIGHT)
    return slide


def screenshot_slide(prs, file_name, title, subtitle, bullets, export_note=None):
    slide = base_slide(prs, title, subtitle)
    rect(slide, 0.48, 1.28, 8.4, 5.72, WHITE, LINE, True)
    add_image(slide, file_name)
    textbox(slide, "초보자 핵심操作", 9.25, 1.3, 3.3, 0.35, 17, BLUE, True)
    bullet_list(slide, bullets, 9.25, 1.82, 3.35, 3.75, 17)
    if export_note:
        rect(slide, 9.2, 5.85, 3.35, 0.82, "EAF2FF", "B9D4FF", True)
        textbox(slide, export_note, 9.4, 6.02, 2.95, 0.45, 15, BLUE, True, PP_ALIGN.CENTER)
    return slide


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(SLATE)
    add_image(slide, "02-admin-dashboard.png", 6.0, 0.65, 6.55, 5.65)
    rect(slide, 0.55, 0.55, 5.3, 5.7, "F8F5EF", None, False)
    textbox(slide, "SAFE-LINK v2.0", 0.9, 1.1, 4.6, 0.55, 26, BLUE, True)
    textbox(slide, "관리자 교육용\n사용 설명서", 0.9, 1.75, 4.35, 1.35, 36, INK, True)
    textbox(slide, "실제 화면 캡처 기반 · 더미 데이터 포함 · TBM/ESG/내보내기 보강", 0.92, 3.55, 4.4, 0.65, 17, INK, True)
    textbox(slide, "업데이트: 2026-05-18", 0.92, 5.45, 4.1, 0.38, 15, MUTED, True)


def workflow_slide(prs):
    slide = base_slide(prs, "관리자 하루 운영 흐름", "회원가입부터 리포트 내보내기까지 이 순서로 교육합니다.")
    steps = [
        ("1", "로그인", "관리자 계정 접속"),
        ("2", "작업자 등록", "현장/언어/권한 확인"),
        ("3", "TBM 발송", "위험요인·작업내용 작성"),
        ("4", "서명 확인", "미서명자 재안내"),
        ("5", "ESG 집계", "TBM·서명·퀴즈 자동 통계"),
        ("6", "파일 내보내기", "PDF/Excel/Word/HWP 저장"),
    ]
    for i, (num, name, desc) in enumerate(steps):
        x = 0.7 + (i % 3) * 4.0
        y = 1.55 + (i // 3) * 2.2
        rect(slide, x, y, 3.35, 1.35, WHITE, LINE)
        textbox(slide, num, x + 0.2, y + 0.25, 0.45, 0.45, 22, BLUE, True, PP_ALIGN.CENTER)
        textbox(slide, name, x + 0.75, y + 0.2, 2.3, 0.35, 21, INK, True)
        textbox(slide, desc, x + 0.75, y + 0.68, 2.3, 0.32, 15, MUTED, True)


def tbm_flow_slide(prs):
    slide = base_slide(prs, "TBM 작동방법: 작성 → 발송 → 서명 → 집계", "관리자가 가장 자주 쓰는 TBM 운영 절차입니다.")
    phases = [
        ("작성", "작업명·위험요인·예방대책 입력\n예: 지하 2층 배관 작업, 추락/끼임 위험"),
        ("발송", "현장 작업자에게 TBM 공지 전송\n예: 서울 A현장 32명 대상"),
        ("서명", "작업자가 모바일에서 확인 후 서명\n예: 서명 28명, 미서명 4명"),
        ("집계", "서명률과 미확인자를 관리자 화면에서 확인\n예: 서명률 87.5%"),
    ]
    colors = [BLUE, GREEN, AMBER, RED]
    for i, (name, desc) in enumerate(phases):
        x = 0.65 + i * 3.05
        rect(slide, x, 1.45, 2.55, 3.8, WHITE, LINE)
        textbox(slide, f"{i+1}. {name}", x + 0.22, 1.78, 2.1, 0.38, 24, colors[i], True)
        textbox(slide, desc, x + 0.25, 2.45, 2.05, 1.4, 18, INK, True)
    textbox(slide, "교육 포인트: TBM은 작성 화면만 보는 것이 아니라, 반드시 서명 현황과 ESG 집계까지 이어서 확인합니다.", 0.75, 6.05, 11.8, 0.5, 18, BLUE, True, PP_ALIGN.CENTER)


def export_overview_slide(prs):
    slide = base_slide(prs, "내보내기 기능 확장", "데이터가 생성되는 주요 페이지에서 파일을 바로 만들 수 있도록 정리했습니다.")
    formats = [("PDF", "보고/공유용"), ("Excel", "원본 데이터 분석"), ("Word", "교육자료 편집"), ("HWP", "한글 문서 제출")]
    for i, (fmt, desc) in enumerate(formats):
        x = 0.8 + i * 3.05
        rect(slide, x, 1.45, 2.45, 1.3, "EAF2FF" if i % 2 == 0 else "ECFDF5", LINE)
        textbox(slide, fmt, x + 0.25, 1.75, 1.95, 0.38, 25, BLUE if i % 2 == 0 else GREEN, True, PP_ALIGN.CENTER)
        textbox(slide, desc, x + 0.2, 2.2, 2.05, 0.3, 15, MUTED, True, PP_ALIGN.CENTER)
    bullet_list(
        slide,
        [
            "ESG 리포트: 숫자 요약, TBM/퀴즈/서명 집계, 벤다이어그램 해석",
            "TBM 현황: 작업자별 서명 여부와 미서명자 목록",
            "작업자/NFC/퀴즈/인센티브: 운영 데이터 목록 저장",
            "채팅/라이브/용어집: 교육 기록과 번역 데이터를 문서화",
        ],
        1.25,
        3.35,
        10.8,
        2.1,
        22,
    )
    textbox(slide, "주의: HWP는 한글에서 열 수 있는 HTML 문서 방식으로 생성됩니다.", 1.25, 6.05, 10.8, 0.35, 16, RED, True, PP_ALIGN.CENTER)


def esg_venn_slide(prs):
    slide = base_slide(prs, "ESG 안전리포트 자동 집계 예시", "더미 데이터 기준: TBM·서명·퀴즈가 겹치는 작업자 그룹을 한눈에 확인합니다.")
    textbox(slide, "오늘 현장 요약", 0.9, 1.35, 3.3, 0.4, 22, BLUE, True)
    metrics = [("TBM 발송", "18건"), ("서명률", "92%"), ("퀴즈 평균", "84점"), ("고위험 알림", "3건")]
    for i, (label, value) in enumerate(metrics):
        y = 1.9 + i * 0.95
        textbox(slide, value, 0.95, y, 1.2, 0.36, 24, INK, True)
        textbox(slide, label, 2.15, y + 0.05, 2.2, 0.25, 14, MUTED, True)
    cx, cy = 7.5, 3.45
    for x, y, color, text in [
        (5.75, 2.25, BLUE, "TBM\n참여 31"),
        (7.0, 2.25, GREEN, "서명\n완료 29"),
        (6.38, 3.45, AMBER, "퀴즈\n통과 26"),
    ]:
        oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(2.45), Inches(2.45))
        oval.fill.solid()
        oval.fill.fore_color.rgb = rgb(color)
        oval.fill.transparency = 42
        oval.line.color.rgb = rgb(color)
        textbox(slide, text, x + 0.45, y + 0.75, 1.55, 0.8, 20, INK, True, PP_ALIGN.CENTER)
    textbox(slide, "중복 달성 24명", cx - 1.05, cy + 0.15, 2.1, 0.35, 24, RED, True, PP_ALIGN.CENTER)
    textbox(slide, "관리자는 PDF/Excel/Word/HWP로 내보내어 월간 안전회의 자료로 사용합니다.", 1.1, 6.2, 11.0, 0.35, 17, BLUE, True, PP_ALIGN.CENTER)


def checklist_slide(prs):
    slide = base_slide(prs, "관리자 운영 체크리스트", "교육 후 현장에서 바로 따라 할 수 있는 최소 절차입니다.")
    items = [
        "출근 전: 현장 작업자 명단과 언어 정보를 확인합니다.",
        "작업 전: TBM을 작성하고 대상 현장에 발송합니다.",
        "작업 시작 전: 서명 현황에서 미서명자를 확인해 재안내합니다.",
        "작업 중: 채팅/라이브 통역으로 외국인 작업자 문의를 처리합니다.",
        "작업 후: ESG 안전리포트를 열어 숫자와 벤다이어그램을 확인합니다.",
        "보고 전: 필요한 페이지에서 PDF/Excel/Word/HWP로 내보냅니다.",
    ]
    bullet_list(slide, items, 1.0, 1.35, 11.1, 4.7, 23)


def build() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    workflow_slide(prs)
    tbm_flow_slide(prs)
    export_overview_slide(prs)
    screenshot_slide(prs, "01-auth-admin-login.png", "회원가입/로그인", "관리자 교육은 로그인 화면부터 시작합니다.", ["관리자 이메일과 비밀번호 입력", "권한이 없으면 관리자 메뉴 접근 불가", "로그인 후 대시보드로 이동"])
    screenshot_slide(prs, "03-profile-setup.png", "프로필 설정", "현장과 역할 정보가 맞아야 데이터가 정확히 보입니다.", ["이름, 소속, 현장 정보 확인", "관리자 권한과 현장 매핑 확인", "외국인 작업자는 언어 정보까지 점검"])
    screenshot_slide(prs, "02-admin-dashboard.png", "관리자 대시보드", "오늘의 안전 운영 상태를 한 화면에서 확인합니다.", ["TBM, 작업자, 퀴즈, NFC 메뉴로 이동", "위험 알림과 미처리 항목 확인", "초보자는 왼쪽 메뉴부터 순서대로 교육"])
    screenshot_slide(prs, "04-tbm-create.png", "TBM 작성", "작업 전 위험요인을 짧고 명확하게 작성합니다.", ["작업명/장소/위험요인 입력", "예방대책과 작업자 안내문 작성", "대상 현장 선택 후 발송 준비"])
    screenshot_slide(prs, "04a-tbm-filled.png", "TBM 더미 데이터 입력 예시", "교육 중에는 실제 작업처럼 예시 데이터를 넣어 설명합니다.", ["예: 지하 2층 배관 설치", "위험: 추락, 협착, 화재", "대책: 안전대, 신호수, 화기감시자"])
    screenshot_slide(prs, "05-tbm-status.png", "TBM 서명 현황", "발송 후에는 반드시 서명률과 미서명자를 확인합니다.", ["서명 완료/미완료 인원 확인", "미서명자에게 재안내", "현황 데이터를 보고서로 저장"], "PDF / Excel / Word / HWP 내보내기")
    screenshot_slide(prs, "06-chat.png", "1:1 AI 번역 채팅", "외국인 작업자와 안전 안내를 주고받습니다.", ["작업자 선택 후 메시지 입력", "원문/번역문/발음 힌트 확인", "중요 대화는 교육 기록으로 저장"], "채팅 기록 내보내기")
    screenshot_slide(prs, "07-workers.png", "작업자 관리", "현장 작업자 명단은 모든 통계의 기준 데이터입니다.", ["작업자 이름/현장/언어 확인", "퇴사자 또는 현장 이동자 정리", "명단을 Excel로 내려받아 점검"], "작업자 목록 내보내기")
    screenshot_slide(prs, "08-workers-enroll.png", "작업자 등록", "신규 작업자는 현장 투입 전에 등록합니다.", ["이름, 연락처, 언어, 현장 입력", "QR/NFC와 연결될 작업자 정보 확인", "등록 후 TBM 대상에 포함되는지 확인"])
    screenshot_slide(prs, "09-nfc.png", "NFC 일일 로그", "태그 기록으로 작업자 출입과 안전 활동을 확인합니다.", ["NFC 태그 이력 확인", "작업자별 시간/현장 점검", "일일 로그를 파일로 저장"], "NFC 로그 내보내기")
    screenshot_slide(prs, "10-qr-code.png", "QR 코드", "작업자가 모바일로 빠르게 접속하도록 QR을 활용합니다.", ["현장별 QR을 확인", "교육장/게시판에 부착", "접속이 안 되면 현장 선택을 재확인"])
    screenshot_slide(prs, "11-quiz.png", "안전 퀴즈", "TBM 이해도를 점검하고 점수를 관리합니다.", ["퀴즈 목록과 응답 현황 확인", "오답이 많은 문항은 재교육", "결과를 보고서로 저장"], "퀴즈 결과 내보내기")
    screenshot_slide(prs, "12-esg.png", "ESG 안전리포트 화면", "안전 활동 데이터가 자동으로 모여 리포트가 됩니다.", ["TBM, 서명, 퀴즈, 위험 알림 집계", "숫자 카드와 그래프로 현황 파악", "회의용 문서로 바로 저장"], "PDF / Excel / Word / HWP")
    esg_venn_slide(prs)
    screenshot_slide(prs, "13-glossary.png", "안전 용어집", "현장 용어를 다국어로 정리해 오해를 줄입니다.", ["용어 검색 및 언어별 확인", "신규 용어는 교육 후 보강", "용어 목록을 파일로 공유"], "용어집 내보내기")
    screenshot_slide(prs, "14-live.png", "라이브 통역", "현장 브리핑이나 긴급 안내를 실시간으로 전달합니다.", ["언어 선택 후 통역 시작", "중요 발화는 기록으로 남김", "회의 후 기록 파일 저장"], "통역 기록 내보내기")
    checklist_slide(prs)

    prs.save(PPTX)
    prs.save(PPTX_KO)
    make_contact_sheet()
    print(PPTX)
    print(PPTX_KO)
    print(CONTACT)


def make_contact_sheet() -> None:
    thumbs = []
    for path in sorted(SCREEN_DIR.glob("*.png")):
        with Image.open(path) as img:
            thumb = img.convert("RGB")
            thumb.thumbnail((260, 178))
            canvas = Image.new("RGB", (300, 225), "white")
            canvas.paste(thumb, ((300 - thumb.width) // 2, 10))
            draw = ImageDraw.Draw(canvas)
            try:
                font = ImageFont.truetype("malgun.ttf", 14)
            except OSError:
                font = ImageFont.load_default()
            draw.text((12, 195), path.name, fill=(17, 24, 39), font=font)
            thumbs.append(canvas)
    cols = 5
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 300, rows * 225), (248, 245, 239))
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 300, (idx // cols) * 225))
    sheet.save(CONTACT)


if __name__ == "__main__":
    build()
