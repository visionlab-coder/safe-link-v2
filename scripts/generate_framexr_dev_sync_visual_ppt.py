from pathlib import Path
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
GEN = Path.home() / ".codex" / "generated_images" / "019e7e61-97bf-7ad0-9e60-9a1de84c4656"
BASE = ROOT / "docs" / "generated" / "framexr-dev-sync-20260607"
ASSET = BASE / "assets"
PREVIEW = BASE / "previews-v2"
OUT = BASE / "FRAMEXR_DEV_LOG_TODO_SYNC_이미지보강_V2_20260607.pptx"

W, H = 13.333, 7.5
COL = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "navy": "0B1220",
    "white": "FFFFFF",
    "soft": "F2F4F7",
}


def rgb(key):
    v = COL.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, trans=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = trans
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if radius:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def text(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, "paper", "paper", False)
    return s


def footer(s, n):
    text(s, "FRAMEXR Dev Sync / DEV_LOG + TODO", 0.55, 7.05, 4.8, 0.24, 10, "muted")
    text(s, f"{n:02d}", 12.3, 7.05, 0.45, 0.24, 10, "muted", align="right")


def header(s, label, title, sub, n):
    text(s, label, 0.62, 0.34, 5.8, 0.3, 12, "blue", True)
    text(s, title, 0.62, 0.74, 11.9, 0.72, 28, "ink", True)
    if sub:
        text(s, sub, 0.64, 1.5, 11.4, 0.38, 15, "muted")
    rule(s, 0.62, 2.02, 12.05)
    footer(s, n)


def pic(s, filename, x, y, w, h):
    p = ASSET / filename
    if p.exists():
        s.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(s, x, y, w, h, "soft", "line", True)
        text(s, filename, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, 13, "muted", True, "center")


def callout(s, head, body, x, y, w, h, color="blue"):
    rect(s, x, y, w, h, "white", "line", True)
    rect(s, x, y, 0.08, h, color, color, False)
    text(s, head, x + 0.22, y + 0.14, w - 0.42, 0.3, 15, color, True)
    text(s, body, x + 0.22, y + 0.53, w - 0.42, h - 0.62, 13.2, "ink")


def table(s, x, y, rows, widths, row_h=0.47):
    for r, row in enumerate(rows):
        cx = x
        for c, val in enumerate(row):
            fill = "navy" if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            rect(s, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            text(s, val, cx + 0.08, y + r * row_h + 0.08, widths[c] - 0.16, row_h - 0.1, 11.5 if r == 0 else 12.2, color, r == 0)
            cx += widths[c]


def prepare_assets():
    ASSET.mkdir(parents=True, exist_ok=True)
    latest = sorted(GEN.glob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)
    initial_names = [
        "04-qa-report-board.png",
        "03-bim-measurement-laser.png",
        "02-pour-gate-checkpoints.png",
        "01-dev-sync-hero.png",
    ]
    if not all((ASSET / name).exists() for name in initial_names):
        for src, name in zip(latest[:4], initial_names):
            shutil.copy2(src, ASSET / name)

    extra_names = [
        "10-commit-ready.png",
        "09-caution-wording.png",
        "08-open-items.png",
        "07-document-roles.png",
        "06-commit-hash.png",
        "05-document-scope.png",
    ]
    for src, name in zip(latest[:6], extra_names):
        shutil.copy2(src, ASSET / name)


def build():
    prepare_assets()
    BASE.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    s = slide(prs)
    pic(s, "01-dev-sync-hero.png", 6.55, 0, 6.78, 7.5)
    text(s, "FRAMEXR V0.3 DEVELOPMENT SYNC", 0.72, 0.68, 4.9, 0.3, 12, "blue", True)
    text(s, "DEV_LOG / TODO\n동기화 완료 보고", 0.72, 1.22, 5.35, 1.35, 33, "ink", True)
    text(s, "Pour 3중 게이트 완성 이후\n문서 6종 정합성 검증", 0.74, 3.02, 5.2, 0.78, 20, "ink", True)
    rule(s, 0.74, 4.12, 4.3, "blue")
    text(s, "코드 변경 0 / 문서 6개 / main==origin/main==e82d121 기준", 0.74, 4.45, 5.3, 0.68, 15, "muted")
    text(s, "2026.06.07", 0.74, 6.78, 2.8, 0.25, 11, "muted")

    s = slide(prs)
    header(s, "01 SCOPE", "이번 작업은 코드 패치가 아니라 문서 동기화입니다", "개발 완료 상태, 남은 범위, 데모 주의문구를 6개 문서에 맞춰 정렬했습니다.", 2)
    pic(s, "05-document-scope.png", 0.75, 2.35, 5.75, 3.55)
    table(s, 6.85, 2.32, [
        ["구분", "내용"],
        ["변경 파일", "마크다운 6개"],
        ["변경 성격", "마크다운 문서 6개 한정"],
        ["코드 영향", "apps / web / api / package.json 변경 0"],
        ["브랜치", "docs/sync-after-pour-gate-completion"],
        ["커밋 메시지", "docs: sync development log after pour gate completion"],
    ], [2.05, 3.95], 0.55)

    s = slide(prs)
    header(s, "02 POUR GATE", "Pour 3중 게이트가 현재 완료 축입니다", "marker + QA + issue category 조건을 묶어 데모 승인 흐름을 완성했습니다.", 3)
    pic(s, "02-pour-gate-checkpoints.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "Marker Gate", "타설 위치·마커 기준으로 작업 맥락을 고정", 7.0, 2.35, 4.95, 0.85, "blue")
    callout(s, "QA Gate", "QA/QC taxonomy와 체크리스트 UI를 통해 확인 항목을 묶음", 7.0, 3.38, 4.95, 0.85, "green")
    callout(s, "Issue Gate", "Field issue category UI와 C1 조건을 통해 이슈 분류를 연결", 7.0, 4.41, 4.95, 0.85, "orange")
    callout(s, "주의", "실제 현장 타설 승인 시스템이 아니라 public demo gate로 명시", 7.0, 5.44, 4.95, 0.85, "red")

    s = slide(prs)
    header(s, "03 COMPLETED", "완료 항목은 커밋 해시 기준으로 대조했습니다", "기능 이름만 나열하지 않고 실제 main 히스토리와 맞춰 문서에 반영했습니다.", 4)
    pic(s, "06-commit-hash.png", 0.75, 2.28, 5.9, 3.6)
    table(s, 6.88, 2.35, [
        ["완료 묶음", "대표 커밋"],
        ["Snapshot Issue", "d1b57bc"],
        ["Measurement Tool + 한국어 명령", "c090274 + d8bc08e"],
        ["Remote Laser + Floating Toolbar", "83a628c + 4649c93"],
        ["Demo Camera Presets + Guided Demo", "1355d1e + 6528e9a"],
        ["QA/QC Taxonomy + Checklist UI", "ccfcdf5 + c67fe37"],
        ["Pour Hard Gate + QA Preset + Issue C1", "644f406 + 67d7fb5 + e82d121"],
    ], [2.55, 3.3], 0.46)
    text(s, "문서상 완료 처리는 e82d121 기준이며, main==origin/main 상태를 명시했습니다.", 1.0, 6.05, 11.1, 0.32, 17, "blue", True, "center")

    s = slide(prs)
    header(s, "04 CAPABILITY", "XYZ Capability Audit는 과장 표현을 줄이는 방향으로 정리했습니다", "MVP와 미구현 항목을 분리해 외부 공유 시 오해 소지를 낮췄습니다.", 5)
    pic(s, "03-bim-measurement-laser.png", 0.75, 2.25, 5.9, 3.65)
    table(s, 6.95, 2.38, [
        ["항목", "상태"],
        ["Remote Laser", "0→4 / MVP"],
        ["Measurement", "1→3 / 모델 내부 거리 측정"],
        ["Floating Toolbar", "0→4 / MVP"],
        ["Snapshot Review", "0→3 / 로컬 저장"],
        ["QA Checklist", "4→5 / 데모 게이트"],
    ], [2.3, 3.7], 0.5)

    s = slide(prs)
    header(s, "05 DOCUMENTS", "6개 문서의 역할을 다르게 가져갔습니다", "같은 내용을 복붙하지 않고 로그, 할 일, 백로그, 진행도, 감사표로 분리했습니다.", 6)
    pic(s, "07-document-roles.png", 0.72, 2.25, 5.85, 3.6)
    table(s, 6.8, 2.25, [
        ["문서", "반영 핵심"],
        ["DEV_LOG.md", "Log 005, 17커밋 범위, 완료·미구현·주의문구"],
        ["TODO.md", "v0.3-03 완료, 현재 우선순위 7~13 완료"],
        ["V03_TASK_BACKLOG.md", "Floating Toolbar 완료, Demo Scenario 일부완료"],
        ["PUBLIC_SAMPLE_BACKLOG.md", "PS-12, PS-14, 공개 샘플 주의문구"],
        ["PROGRESS_MAP.md", "Pour 3중 게이트 현재 위치와 mermaid 갱신"],
        ["XYZ_CAPABILITY_AUDIT.md", "능력 점수·MVP·미구현 상태 갱신"],
    ], [2.5, 3.45], 0.5)

    s = slide(prs)
    header(s, "06 OPEN ITEMS", "미완료 항목은 계속 미완료로 남겼습니다", "실제 도면, 실측, AR 정합, 서버 기능은 데모 범위를 벗어난 것으로 유지했습니다.", 7)
    pic(s, "08-open-items.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "도면·모델", "실 DWG/DXF/PDF/IFC parser\n실 회사 도면 모델", 7.0, 2.35, 4.95, 0.82, "blue")
    callout(s, "현장 정합", "AR 카메라 정합\n실측 정확도 검증\nField Camera Preview", 7.0, 3.35, 4.95, 0.82, "green")
    callout(s, "서버 기능", "서버/DB/로그인\n프로젝트 관리\nBCFzip export/import", 7.0, 4.35, 4.95, 0.82, "orange")
    callout(s, "2차 확장", "Viewpoint Capture 정식\nLaser+Measurement 연동\nFloating Toolbar 2차 확장", 7.0, 5.35, 4.95, 0.82, "red")

    s = slide(prs)
    header(s, "07 SAFETY WORDING", "실제 현장 기능처럼 보이는 표현을 차단했습니다", "데모와 실제 운용의 경계를 문서 전체에 일관되게 넣었습니다.", 8)
    pic(s, "09-caution-wording.png", 0.75, 2.25, 5.9, 3.65)
    table(s, 6.85, 2.35, [
        ["주의문구", "의미"],
        ["public IFC-derived demo", "실제 회사 도면이 아님"],
        ["데모 게이트", "실제 현장 타설 승인 시스템이 아님"],
        ["Measurement=모델 내부 거리 측정", "실측 또는 mm급 정밀도 아님"],
        ["Snapshot=로컬 메모리 증빙", "서버/DB 저장 아님"],
        ["Private 유지·원본 미커밋", "외부 배포 범위 통제"],
    ], [2.9, 3.1], 0.52)

    s = slide(prs)
    header(s, "08 VERIFICATION", "문서 품질 검증까지 끝낸 상태입니다", "Mermaid fence, 변경 범위, 원본 파일 미포함 여부를 따로 확인했습니다.", 9)
    pic(s, "04-qa-report-board.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "Fence 균형", "6개 문서 모두 ``` 개수 짝수\nPROGRESS_MAP mermaid 2블록 정상", 7.0, 2.35, 4.95, 0.95, "blue")
    callout(s, "Git Status", "변경 6개 마크다운 한정\n코드·package 영향 0", 7.0, 3.55, 4.95, 0.95, "green")
    callout(s, "원본 제외", "PDF/DWG/DXF/IFC/BCF/ZIP/RAR/7Z/tmp/png/mjs staged 0", 7.0, 4.75, 4.95, 0.95, "orange")

    s = slide(prs)
    header(s, "09 COMMIT READY", "선택적 add 후 문서 커밋만 진행 가능한 상태입니다", "범위가 명확하므로 코드 리뷰 부담 없이 문서 동기화 커밋으로 분리할 수 있습니다.", 10)
    pic(s, "10-commit-ready.png", 0.75, 2.25, 5.9, 3.65)
    text(s, "예정 커밋", 6.95, 2.42, 2.0, 0.3, 13, "blue", True)
    text(s, "docs: sync development log after pour gate completion", 6.95, 2.82, 5.25, 0.72, 20, "ink", True)
    callout(s, "Selective Add", "DEV_LOG.md\nTODO.md\n4개 docs/product·ops 문서", 6.95, 3.85, 4.95, 0.82, "blue")
    callout(s, "Commit / Push", "문서 6개 한정\n코드 변경 0\nmain ff-only 머지 후 push", 6.95, 4.86, 4.95, 0.82, "green")
    text(s, "현재 상태: 커밋 승인 대기", 6.95, 6.05, 4.95, 0.35, 18, "blue", True, "center")

    prs.save(OUT)


def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles = [
        "DEV_LOG / TODO 동기화 완료",
        "문서 동기화 범위",
        "Pour 3중 게이트",
        "커밋 해시 대조",
        "Capability Audit",
        "6개 문서 역할",
        "미완료 유지",
        "주의문구",
        "검증 완료",
        "커밋 준비",
    ]
    image_map = {
        1: "01-dev-sync-hero.png",
        2: "05-document-scope.png",
        3: "02-pour-gate-checkpoints.png",
        4: "06-commit-hash.png",
        5: "03-bim-measurement-laser.png",
        6: "07-document-roles.png",
        7: "08-open-items.png",
        8: "09-caution-wording.png",
        9: "04-qa-report-board.png",
        10: "10-commit-ready.png",
    }
    try:
        f1 = ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf", 38)
        f2 = ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf", 20)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    paths = []
    for i, title in enumerate(titles, 1):
        canvas = Image.new("RGB", (1600, 900), "#" + COL["paper"])
        d = ImageDraw.Draw(canvas)
        d.text((80, 60), "FRAMEXR Dev Sync", font=f2, fill="#" + COL["blue"])
        d.text((80, 135), title, font=f1, fill="#" + COL["ink"])
        d.line((80, 235, 1460, 235), fill="#" + COL["line"], width=3)
        if i in image_map and (ASSET / image_map[i]).exists():
            im = Image.open(ASSET / image_map[i]).convert("RGB")
            im.thumbnail((760, 430))
            canvas.paste(im, (80, 285))
        for k, color in enumerate(("blue", "green", "orange", "red")):
            d.rounded_rectangle((980, 310 + k * 105, 1390, 370 + k * 105), 16, outline="#" + COL[color], width=3, fill="white")
        d.text((80, 830), f"{i:02d}", font=f2, fill="#" + COL["muted"])
        out = PREVIEW / f"slide_{i:02d}.png"
        canvas.save(out)
        paths.append(out)

    sheet = Image.new("RGB", (1600, 1040), "white")
    d = ImageDraw.Draw(sheet)
    for i, path in enumerate(paths):
        im = Image.open(path).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 320
        sheet.paste(im, (x, y))
        d.text((x, y + 212), f"{i + 1:02d} {titles[i]}", font=f2, fill="#" + COL["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build()
    preview()
    print(OUT)
    print(PREVIEW / "contact_sheet.png")
