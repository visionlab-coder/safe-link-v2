from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-hyper-safety-business-20260609"
ASSET = BASE / "assets"
PREVIEW = BASE / "previews"
OUT = BASE / "대우건설_HyperSafety_AI_SQ-LINK_사업제안서_PPT_v1_20260609.pptx"

SOURCE_ASSETS = [
    ROOT / "docs/generated/daewoo-drone-patrol-20260606/assets_micro_drone_nest/01-micro-drone-nest-hero.png",
    ROOT / "docs/generated/daewoo-drone-patrol-20260606/assets_micro_drone_nest/02-hard-to-reach-space.png",
    ROOT / "docs/generated/daewoo-drone-patrol-20260606/assets_micro_drone_nest/04-control-center-micro-feeds.png",
    ROOT / "docs/generated/daewoo-drone-patrol-20260606/assets_micro_drone_nest/05-safelink-action-report.png",
    ROOT / "docs/generated/daewoo-drone-patrol-20260606/assets_micro_drone_nest/06-mobile-robot-dog-nest.png",
]

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "navy": "0B1220",
    "white": "FFFFFF",
    "soft": "F2F4F7",
    "yellow": "FFF7D6",
}


def rgb(key):
    value = C.get(key, key).strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, trans=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = trans
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if radius:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def text(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, "paper", "paper", False)
    return s


def footer(s, n):
    text(s, "SQ-LINK Underground / 대우건설 Hyper Safety & AI 제안", 0.58, 7.05, 6.2, 0.24, 10, "muted")
    text(s, f"{n:02d}", 12.25, 7.05, 0.45, 0.24, 10, "muted", align="right")


def header(s, label, title, sub, n):
    text(s, label, 0.62, 0.34, 5.8, 0.3, 12, "blue", True)
    text(s, title, 0.62, 0.74, 11.9, 0.72, 28, "ink", True)
    if sub:
        text(s, sub, 0.64, 1.5, 11.4, 0.38, 15, "muted")
    rule(s, 0.62, 2.02, 12.05)
    footer(s, n)


def pic(s, name, x, y, w, h):
    path = ASSET / name
    if path.exists():
        s.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(s, x, y, w, h, "soft", "line", True)
        text(s, name, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, 13, "muted", True, "center")


def callout(s, head, body, x, y, w, h, color="blue"):
    rect(s, x, y, w, h, "white", "line", True)
    rect(s, x, y, 0.08, h, color, color, False)
    text(s, head, x + 0.22, y + 0.14, w - 0.42, 0.32, 15, color, True)
    text(s, body, x + 0.22, y + 0.55, w - 0.42, h - 0.64, 13.2, "ink")


def table(s, x, y, rows, widths, row_h=0.48):
    for r, row in enumerate(rows):
        cx = x
        for c, val in enumerate(row):
            fill = "navy" if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            rect(s, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            text(s, val, cx + 0.08, y + r * row_h + 0.08, widths[c] - 0.16, row_h - 0.1, 11.5 if r == 0 else 12.3, color, r == 0)
            cx += widths[c]


def prepare_assets():
    ASSET.mkdir(parents=True, exist_ok=True)
    names = [
        "01-underground-hero.png",
        "02-hard-to-reach-space.png",
        "03-control-center.png",
        "04-safelink-report.png",
        "05-mobile-nest.png",
    ]
    for src, name in zip(SOURCE_ASSETS, names):
        if src.exists():
            shutil.copy2(src, ASSET / name)

    # Native report diagram assets for slides where a generated photo is not ideal.
    make_layer_diagram(ASSET / "06-layer-architecture.png")
    make_poc_roadmap(ASSET / "07-poc-roadmap.png")
    make_brand_map(ASSET / "08-brand-map.png")


def load_font(size=34, bold=False):
    try:
        path = r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf"
        return ImageFont.truetype(path, size)
    except Exception:
        return ImageFont.load_default()


def make_layer_diagram(path):
    im = Image.new("RGB", (1600, 900), "#" + C["paper"])
    d = ImageDraw.Draw(im)
    f1, f2 = load_font(44, True), load_font(25)
    d.text((80, 70), "SQ-LINK Layer Architecture", font=f1, fill="#" + C["ink"])
    layers = [
        ("AI 현장 운영 엔진", "품질·안전·공정·원가 판단", C["blue"]),
        ("SAFE-LINK 작업데이터", "TBM·작업중지·서명·보고서", C["green"]),
        ("BIM / 3D 좌표", "층·높이·깊이·작업구역 매핑", C["orange"]),
        ("AI-RAN 통신망", "지하 통신 음영 해소", C["red"]),
        ("현장 장비", "스마트글라스·로봇개·드론·근로자 단말", C["navy"]),
    ]
    for i, (head, body, color) in enumerate(layers):
        y = 200 + i * 115
        d.rounded_rectangle((130, y, 1470, y + 76), 18, fill="white", outline="#" + color, width=4)
        d.rectangle((130, y, 148, y + 76), fill="#" + color)
        d.text((185, y + 13), head, font=f2, fill="#" + color)
        d.text((600, y + 15), body, font=f2, fill="#" + C["ink"])
    im.save(path)


def make_poc_roadmap(path):
    im = Image.new("RGB", (1600, 900), "#" + C["paper"])
    d = ImageDraw.Draw(im)
    f1, f2 = load_font(44, True), load_font(24)
    d.text((80, 70), "8주 PoC Roadmap", font=f1, fill="#" + C["ink"])
    steps = [
        ("1-2주", "현장 선정·통신 음영 측정"),
        ("3-4주", "QR/NFC·SAFE-LINK 작업데이터 적용"),
        ("5-6주", "스마트글라스·로봇/드론 시나리오"),
        ("7-8주", "성과측정·자동보고·확산안"),
    ]
    x0, y0 = 130, 330
    for i, (week, body) in enumerate(steps):
        x = x0 + i * 360
        d.rounded_rectangle((x, y0, x + 270, y0 + 180), 22, fill="white", outline="#" + [C["blue"], C["green"], C["orange"], C["red"]][i], width=4)
        d.text((x + 35, y0 + 28), week, font=f1, fill="#" + [C["blue"], C["green"], C["orange"], C["red"]][i])
        d.text((x + 28, y0 + 105), body, font=f2, fill="#" + C["ink"])
        if i < 3:
            d.line((x + 285, y0 + 90, x + 345, y0 + 90), fill="#" + C["muted"], width=5)
            d.polygon([(x + 345, y0 + 90), (x + 325, y0 + 78), (x + 325, y0 + 102)], fill="#" + C["muted"])
    im.save(path)


def make_brand_map(path):
    im = Image.new("RGB", (1600, 900), "#" + C["paper"])
    d = ImageDraw.Draw(im)
    f1, f2, f3 = load_font(46, True), load_font(28, True), load_font(23)
    d.text((80, 70), "SAFE-LINK에서 SQ-LINK로", font=f1, fill="#" + C["ink"])
    boxes = [
        ("SAFE-LINK", "안전\nTBM·작업중지·증거", C["blue"]),
        ("SQ-LINK", "안전 + 품질\n현장 AI 운영체계", C["green"]),
        ("Company AI", "원가·공정·날씨·법규\n디지털트윈 확장", C["orange"]),
    ]
    for i, (head, body, color) in enumerate(boxes):
        x = 150 + i * 470
        d.rounded_rectangle((x, 285, x + 350, 570), 28, fill="white", outline="#" + color, width=5)
        d.text((x + 55, 335), head, font=f2, fill="#" + color)
        d.multiline_text((x + 55, 425), body, font=f3, fill="#" + C["ink"], spacing=8)
        if i < 2:
            d.line((x + 370, 425, x + 445, 425), fill="#" + C["muted"], width=6)
            d.polygon([(x + 445, 425), (x + 420, 410), (x + 420, 440)], fill="#" + C["muted"])
    im.save(path)


def build():
    prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    s = slide(prs)
    pic(s, "01-underground-hero.png", 6.55, 0, 6.78, 7.5)
    text(s, "DAEWOO HYPER SAFETY & AI", 0.72, 0.68, 4.8, 0.3, 12, "blue", True)
    text(s, "SQ-LINK\nUnderground", 0.72, 1.22, 5.25, 1.35, 35, "ink", True)
    text(s, "지하·골조 현장 3D 좌표 통신망 기반\n안전·품질 AI 운영 플랫폼", 0.74, 3.03, 5.4, 0.8, 20, "ink", True)
    rule(s, 0.74, 4.15, 4.35, "blue")
    text(s, "서원토건 미래전략TF / 사업보고서 v1 기반 제출용 PPT v1", 0.74, 4.47, 5.25, 0.55, 14, "muted")
    text(s, "2026.06.09", 0.74, 6.78, 2.6, 0.24, 11, "muted")

    s = slide(prs)
    header(s, "01 EXECUTIVE SUMMARY", "제안의 중심을 다시 세웁니다", "전무님 발언 기준: 장비가 아니라 통신·좌표·작업데이터 운영체계가 핵심입니다.", 2)
    callout(s, "기존 문제", "SAFE-LINK, 로봇, 드론, 스마트글라스, AI-RAN이 병렬로 섞여 핵심이 흐려짐", 0.85, 2.35, 3.6, 1.1, "red")
    callout(s, "새 중심", "지하·골조 현장 3D 좌표 통신망과 SAFE-LINK 작업데이터를 연결", 4.85, 2.35, 3.6, 1.1, "blue")
    callout(s, "최종 방향", "안전앱에서 안전·품질·원가·공정까지 보는 회사 AI로 확장", 8.85, 2.35, 3.6, 1.1, "green")
    text(s, "우리는 로봇을 만드는 회사가 아니라, 로봇과 사람이 현장에서 함께 일할 수 있는 기반을 만드는 전문건설사입니다.", 1.05, 5.1, 11.3, 0.55, 24, "ink", True, "center")

    s = slide(prs)
    header(s, "02 REAL PROBLEM", "지하·골조 현장은 가장 위험하지만 가장 연결되지 않습니다", "통신, GPS, 수직 위치, 시야, 작업데이터가 동시에 끊기는 곳입니다.", 3)
    pic(s, "01-underground-hero.png", 0.75, 2.25, 5.9, 3.65)
    table(s, 6.9, 2.35, [
        ["문제", "현장 영향"],
        ["통신 음영", "신고·영상·장비 제어 단절"],
        ["GPS 불가", "로봇/작업자 위치 추적 한계"],
        ["수직 좌표 부재", "층·높이·깊이 판단 어려움"],
        ["CCTV 사각", "PIT·코어·기계실 확인 한계"],
        ["데이터 분리", "안전·품질·공정이 따로 관리"],
    ], [2.2, 3.65], 0.52)

    s = slide(prs)
    header(s, "03 STRATEGIC SHIFT", "로봇·드론은 핵심 기술이 아니라 응용 장비입니다", "전문 업체 장비를 현장에서 작동하게 만드는 운영 기반이 서원토건의 제안입니다.", 4)
    pic(s, "05-mobile-nest.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "하지 않을 말", "우리가 로봇개·드론을 직접 개발한다", 7.0, 2.35, 4.95, 0.8, "red")
    callout(s, "해야 할 말", "로봇·드론·스마트글라스가 현장에서 작동할 통신·좌표·데이터 기반을 만든다", 7.0, 3.35, 4.95, 1.0, "blue")
    callout(s, "대우 제안 포인트", "장비 구매가 아니라 지하 골조 현장 PoC 공동 실증", 7.0, 4.65, 4.95, 0.8, "green")

    s = slide(prs)
    header(s, "04 WHAT WE HAVE", "서원토건은 이미 현장 데이터의 출발점을 갖고 있습니다", "SAFE-LINK와 철근콘크리트 공정 지식이 제안의 출발점입니다.", 5)
    table(s, 0.85, 2.25, [
        ["보유 기반", "설명"],
        ["SAFE-LINK", "다국어 TBM·위험성평가·서명·작업중지·보고서"],
        ["QR/NFC", "작업구역 접근과 이력 생성"],
        ["철근콘크리트 전문성", "철근·거푸집·동바리·타설 전후 리스크 이해"],
        ["스마트글라스 계획", "Rokid 기반 영상·열화상·품질확인 확장"],
        ["로봇/드론 시나리오", "선행순찰·사각공간 촬영·증거수집"],
        ["특허/협업", "특허출원 준비, AI-RAN 산학협력 추진"],
    ], [3.05, 8.45], 0.55)

    s = slide(prs)
    header(s, "05 CORE TECHNOLOGY", "핵심은 지하 3D 좌표 통신망입니다", "지하에서도 층·높이·깊이를 포함해 장비와 작업자를 같은 좌표계에 묶어야 합니다.", 6)
    pic(s, "06-layer-architecture.png", 0.75, 2.18, 6.1, 3.85)
    callout(s, "AI-RAN/O-RAN", "지하 통신 음영 해소 PoC", 7.1, 2.35, 4.85, 0.75, "blue")
    callout(s, "3D 좌표", "평면이 아니라 수직 위치까지 포함", 7.1, 3.28, 4.85, 0.75, "green")
    callout(s, "BIM 연결", "도면에서 구역을 선택하면 장비가 확인", 7.1, 4.21, 4.85, 0.75, "orange")
    callout(s, "표현 제한", "cm급 구현 완료가 아니라 PoC 검증", 7.1, 5.14, 4.85, 0.75, "red")

    s = slide(prs)
    header(s, "06 SMART GLASS", "스마트글라스는 안전뿐 아니라 품질 장비입니다", "Rokid 기반 영상·열화상 기능을 현장 작업면 데이터와 연결합니다.", 7)
    pic(s, "03-control-center.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "안전", "개구부·PIT·코어·동바리 위험 확인", 7.0, 2.35, 4.95, 0.8, "blue")
    callout(s, "품질", "철근 배근, 거푸집, 타설 전 체크리스트 보조", 7.0, 3.35, 4.95, 0.8, "green")
    callout(s, "관제", "관리자 시야를 원격 관제와 보고서 증거로 연결", 7.0, 4.35, 4.95, 0.8, "orange")

    s = slide(prs)
    header(s, "07 ROBOT & DRONE", "로봇개와 드론은 이동형 현장 센서가 됩니다", "사람이 먼저 들어가기 어려운 구역을 선행 확인하고 SQ-LINK 증거 흐름에 연결합니다.", 8)
    pic(s, "05-mobile-nest.png", 0.75, 2.25, 5.9, 3.65)
    table(s, 6.9, 2.38, [
        ["응용", "역할"],
        ["로봇개", "지하 통로·PIT·코어 주변 선행 순찰"],
        ["초소형 드론", "협소·상부·사각공간 촬영"],
        ["BIM 지시", "도면상 특정 구역 선택 후 확인"],
        ["SAFE-LINK", "사진·위치·작업조·조치 보고 연결"],
        ["협업 방식", "전문 업체/상용 장비 활용"],
    ], [2.2, 3.65], 0.52)

    s = slide(prs)
    header(s, "08 BUSINESS EXPANSION", "SAFE-LINK에서 SQ-LINK로 확장합니다", "안전 중심 앱을 안전+품질+회사 AI 플랫폼으로 확장하는 브랜드 전략입니다.", 9)
    pic(s, "08-brand-map.png", 0.75, 2.25, 6.0, 3.6)
    callout(s, "브랜드 권장", "SQ-LINK: Safety + Quality + Link", 7.05, 2.4, 4.85, 0.8, "blue")
    callout(s, "확장 범위", "안전, 품질, 원가, 공정, 날씨, 법규, 장비", 7.05, 3.4, 4.85, 0.8, "green")
    callout(s, "전무님 의도", "짧고 강한 이름으로 회사 AI 브랜드화", 7.05, 4.4, 4.85, 0.8, "orange")

    s = slide(prs)
    header(s, "09 POC PROPOSAL", "대우건설과 8주 PoC를 제안합니다", "구매 요청이 아니라 현장 제공, 피드백, 공동 실증 파트너십입니다.", 10)
    pic(s, "07-poc-roadmap.png", 0.75, 2.25, 6.0, 3.6)
    callout(s, "대상", "지하층 골조 또는 지하주차장 철근콘크리트 구역 1개소", 7.05, 2.35, 4.85, 0.8, "blue")
    callout(s, "범위", "통신 음영, QR/NFC, 스마트글라스, 로봇/드론 시나리오", 7.05, 3.35, 4.85, 0.8, "green")
    callout(s, "성과", "안전+품질 체크리스트, 처리시간, 현장 수용성, 자동보고", 7.05, 4.35, 4.85, 0.8, "orange")

    s = slide(prs)
    header(s, "10 COLLABORATION", "협업 구조는 컨소시엄·공동과제형으로 잡습니다", "서원토건이 모든 비용과 기술을 떠안는 구조가 아니라 역할을 나눕니다.", 11)
    table(s, 0.85, 2.3, [
        ["주체", "역할"],
        ["대우건설", "PoC 현장, 안전/품질팀 피드백, 기존 체계 비교검증"],
        ["서원토건", "철근콘크리트 공정 지식, SAFE/SQ-LINK 운영, 시나리오 설계"],
        ["경희대", "AI-RAN/O-RAN 자문 및 공동과제 추진"],
        ["로봇/드론 업체", "하드웨어, 운용 기술, 현장 적용 협력"],
        ["Rokid", "스마트글라스 영상·열화상 기반 관제 검토"],
    ], [2.7, 8.8], 0.58)

    s = slide(prs)
    header(s, "11 RISK WORDING", "과장 표현은 피하고 PoC 검증 언어로 갑니다", "심사 설득력은 높이고, 기술 확정 표현의 리스크는 줄입니다.", 12)
    table(s, 0.9, 2.35, [
        ["피해야 할 표현", "대체 표현"],
        ["로봇개를 직접 개발", "전문 업체와 협업 또는 상용 장비 활용"],
        ["cm급 위치정밀도 구현 완료", "현장 자동화에 필요한 정밀 위치 인식 가능성 검증"],
        ["AI-RAN 구축 확정", "AI-RAN/O-RAN 기반 공동과제·PoC 추진"],
        ["타설 승인 자동화 완료", "타설 전 판단 보조 게이트 PoC"],
        ["스마트글라스 품질 자동화 완료", "품질검측 보조 및 증거화 시나리오"],
    ], [4.0, 7.3], 0.58)

    s = slide(prs)
    header(s, "12 NEXT ACTION", "다음 회의 전 준비할 것", "기술보다 먼저 제안의 흐름과 현장 적용 가능성을 확실히 정리해야 합니다.", 13)
    callout(s, "1. 기능 확인", "SAFE-LINK 최신 시연\nNFC/QR·TBM·작업중지 가능 범위 확인", 0.9, 2.45, 3.55, 1.1, "blue")
    callout(s, "2. 브랜드", "SQ-LINK 포함 이름 후보 5개 이상\n짧고 강한 통합 브랜드", 4.9, 2.45, 3.55, 1.1, "green")
    callout(s, "3. 제출물", "10장 내외 PPT\n사업보고서 v1\nPoC 예산안", 8.9, 2.45, 3.55, 1.1, "orange")
    callout(s, "4. 협업", "경희대 AI-RAN 문구\n로봇/스마트글라스 후보\n대우 요청사항", 2.9, 4.35, 3.65, 1.1, "red")
    callout(s, "5. 핵심 메시지", "장비가 아니라 통신·좌표·작업데이터 운영체계", 6.95, 4.35, 3.65, 1.1, "blue")

    BASE.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles = [
        "SQ-LINK Underground",
        "제안의 중심 재정의",
        "지하·골조 문제",
        "장비에서 기반으로",
        "우리가 가진 것",
        "3D 좌표 통신망",
        "스마트글라스",
        "로봇·드론",
        "브랜드 확장",
        "8주 PoC",
        "협업 구조",
        "표현 리스크",
        "다음 액션",
    ]
    image_map = {
        1: "01-underground-hero.png",
        3: "01-underground-hero.png",
        4: "05-mobile-nest.png",
        6: "06-layer-architecture.png",
        7: "03-control-center.png",
        8: "05-mobile-nest.png",
        9: "08-brand-map.png",
        10: "07-poc-roadmap.png",
    }
    f1, f2 = load_font(36, True), load_font(19)
    paths = []
    for i, title in enumerate(titles, 1):
        im = Image.new("RGB", (1600, 900), "#" + C["paper"])
        d = ImageDraw.Draw(im)
        d.text((80, 60), "대우건설 Hyper Safety & AI", font=f2, fill="#" + C["blue"])
        d.text((80, 135), title, font=f1, fill="#" + C["ink"])
        d.line((80, 235, 1460, 235), fill="#" + C["line"], width=3)
        if i in image_map and (ASSET / image_map[i]).exists():
            photo = Image.open(ASSET / image_map[i]).convert("RGB")
            photo.thumbnail((760, 430))
            im.paste(photo, (80, 285))
        for k, color in enumerate(("blue", "green", "orange", "red")):
            d.rounded_rectangle((980, 310 + k * 105, 1390, 370 + k * 105), 16, outline="#" + C[color], width=3, fill="white")
        d.text((80, 830), f"{i:02d}", font=f2, fill="#" + C["muted"])
        p = PREVIEW / f"slide_{i:02d}.png"
        im.save(p)
        paths.append(p)

    sheet = Image.new("RGB", (1600, 1350), "white")
    d = ImageDraw.Draw(sheet)
    for i, p in enumerate(paths):
        im = Image.open(p).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 310
        sheet.paste(im, (x, y))
        d.text((x, y + 212), f"{i + 1:02d} {titles[i]}", font=f2, fill="#" + C["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build()
    preview()
    print(OUT)
    print(PREVIEW / "contact_sheet.png")
