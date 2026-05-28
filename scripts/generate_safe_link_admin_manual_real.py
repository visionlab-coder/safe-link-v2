from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
SCREEN_DIR = OUT / "real-screens"
PREVIEWS = OUT / "previews-real"
PPTX = OUT / "SAFE-LINK_admin_training_manual_real_enriched.pptx"

INK = "101827"
MUTED = "657083"
PAPER = "F7F4EE"
PANEL = "FFFFFF"
LINE = "DDD6CC"
BLUE = "2F75FF"
PINK = "EF3D9A"
AMBER = "F5B82E"
GREEN = "19C37D"
NAVY = "121827"


SLIDES = [
    ("01-auth-admin-login.png", "회원가입 / 로그인", "언어 선택 후 관리자 역할로 진입하고 이메일·비밀번호로 로그인합니다.", [("언어", "처음 접속하면 사용할 언어를 선택합니다."), ("역할", "관리자 교육에서는 관리자 버튼을 선택합니다."), ("계정", "기존 계정은 로그인, 신규 관리자는 회원가입을 진행합니다.")]),
    ("03-profile-setup.png", "최초 프로필 설정", "권한, 이름, 현장 정보를 맞춰야 관리자 메뉴와 데이터가 정확히 연결됩니다.", [("이름", "보고서와 채팅에 표시될 이름입니다."), ("역할", "안전관리자 등 실제 담당 권한을 확인합니다."), ("현장", "담당 현장 코드가 맞는지 반드시 확인합니다.")]),
    ("02-admin-dashboard.png", "관리자 대시보드", "로그인 후 가장 먼저 보는 화면입니다. 시스템 상태와 핵심 메뉴 진입점을 확인합니다.", [("상태", "DB, 번역, 음성 서비스 준비 상태를 확인합니다."), ("메뉴", "TBM, 채팅, 작업자, QR/NFC 등으로 이동합니다."), ("알림", "미확인 작업자나 신규 메시지를 확인합니다.")]),
    ("04-tbm-create.png", "TBM 작성 / 발송", "작업 내용과 위험요인을 입력하고 AI 가이드 또는 기초교육 라이브러리를 활용합니다.", [("초안", "오늘 작업과 위험요인을 짧게 입력합니다."), ("AI", "AI 가이드 생성으로 안전 문장을 보강합니다."), ("발송", "대상 작업자에게 TBM을 전파합니다.")]),
    ("04a-tbm-filled.png", "TBM 예시 입력 화면", "더미 시나리오를 입력한 실제 화면입니다. 발송 버튼을 누르면 작업자에게 TBM이 전달됩니다.", [("예시", "3층 철근 양중 / 고소작업 TBM을 입력합니다."), ("확인", "작업 반경, 안전고리, 신호수, 통제선을 확인합니다."), ("발송", "TBM 브로드캐스트로 대상 작업자에게 전송합니다.")]),
    ("05-tbm-status.png", "TBM 서명 현황", "발송 후에는 누가 확인했고 누가 미확인 상태인지 추적합니다.", [("완료율", "서명 완료율과 전체 인원을 확인합니다."), ("미확인자", "서명하지 않은 작업자를 찾아 재안내합니다."), ("기록", "서명 시간과 이력은 보고 자료가 됩니다.")]),
    ("06-chat.png", "1:1 AI 번역 채팅", "외국인 작업자와 짧고 명확하게 소통하는 화면입니다.", [("선택", "대화할 작업자를 선택합니다."), ("입력", "한 문장에 하나의 요청만 작성합니다."), ("확인", "중요 지시는 번역 의미를 한 번 더 확인합니다.")]),
    ("07-workers.png", "작업자 관리", "작업자 명단, 언어, 소속 현장을 관리합니다.", [("명단", "등록된 작업자 목록을 확인합니다."), ("언어", "TBM 번역과 음성 안내의 기준입니다."), ("수정", "전출·퇴사자는 즉시 정리합니다.")]),
    ("08-workers-enroll.png", "작업자 등록", "신규 작업자의 이름, 연락처, 언어, 현장 정보를 입력합니다.", [("필수값", "이름과 연락처는 정확히 입력합니다."), ("언어", "작업자가 이해 가능한 언어를 선택합니다."), ("현장", "현장 코드 연결을 확인합니다.")]),
    ("09-nfc.png", "NFC 관리", "NFC 스티커와 현장 태그 기록을 관리합니다.", [("발급", "작업자별 NFC 스티커를 발급합니다."), ("태그", "출입·TBM 확인 태그 기록을 봅니다."), ("교체", "분실 시 기존 토큰을 비활성화합니다.")]),
    ("10-qr-code.png", "QR 코드", "작업자 접속 또는 현장 입장용 QR을 생성하고 공유합니다.", [("생성", "필요한 목적에 맞는 QR을 만듭니다."), ("공유", "현장 게시 또는 작업자 안내에 사용합니다."), ("확인", "스캔 후 이동 경로가 맞는지 확인합니다.")]),
    ("11-quiz.png", "안전 퀴즈", "TBM 이해도를 짧게 확인하고 참여 현황을 관리합니다.", [("생성", "교육 내용 기반으로 문항을 만듭니다."), ("참여", "작업자별 응답 여부를 확인합니다."), ("보상", "필요 시 인센티브와 연결합니다.")]),
    ("12-esg.png", "ESG 리포트", "교육, 서명, 참여 데이터를 관리 보고용으로 정리합니다.", [("수집", "TBM·서명·퀴즈 데이터를 확인합니다."), ("정리", "보고서에 필요한 지표를 모읍니다."), ("공유", "현장/본사 보고 자료로 활용합니다.")]),
    ("13-glossary.png", "현장 용어집", "현장 은어와 표준 표현을 관리해 번역 품질을 높입니다.", [("등록", "은어, 표준어, 설명을 함께 저장합니다."), ("검수", "위험 작업 용어는 관리자가 확인합니다."), ("반영", "TBM 작성과 번역 품질 개선에 활용합니다.")]),
    ("14-live.png", "라이브 통역", "즉시 소통이 필요한 상황에서 실시간 음성·텍스트 통역을 사용합니다.", [("시작", "언어와 모드를 확인하고 통역을 시작합니다."), ("짧게", "짧은 문장으로 말해야 정확도가 높습니다."), ("안전", "위험 지시는 반드시 재확인합니다.")]),
]


def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=True, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    return box


def add_chip(slide, text, x, y, w, color=BLUE):
    s = add_rect(slide, x, y, w, 0.46, color, radius=True)
    tf = s.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = rgb("FFFFFF")


def add_screen(slide, path: Path, x=0.62, y=1.75, w=10.15):
    add_rect(slide, x - 0.04, y - 0.04, w + 0.08, 6.95, "FFFFFF", LINE, radius=True)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(6.86))


def title(slide, heading, subtitle, idx):
    add_text(slide, heading, 0.62, 0.34, 11.2, 0.64, 32, INK, True)
    add_text(slide, subtitle, 0.64, 1.05, 11.4, 0.58, 17, MUTED, True)
    add_text(slide, f"SAFE-LINK V2.0 관리자 교육 | 실제 화면 캡처 | {idx:02d}", 0.65, 8.48, 5.2, 0.22, 7.5, MUTED)


def normalize_deck_text(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Malgun Gothic"
                    run.font.bold = True
                    if run.font.size is None or run.font.size < Pt(15):
                        run.font.size = Pt(15)


def build():
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 16, 9, NAVY, radius=False)
    cover_img = SCREEN_DIR / "02-admin-dashboard.png"
    s.shapes.add_picture(str(cover_img), Inches(7.2), Inches(0.75), width=Inches(8.0), height=Inches(5.44))
    add_text(s, "SAFE-LINK V2.0", 0.75, 0.72, 4.1, 0.36, 18, "FFFFFF", True)
    add_text(s, "관리자 교육용\n사용 설명서", 0.75, 2.0, 6.4, 1.55, 42, "FFFFFF", True)
    add_text(s, "실제 SAFE-LINK 화면을 직접 캡처해 구성한 초보 관리자용 운영 가이드", 0.82, 4.02, 6.1, 0.55, 15, "D8DFEA")
    add_chip(s, "Actual Product Screens", 0.82, 5.55, 2.25, PINK)
    add_text(s, "회원가입 · 로그인 · 대시보드 · TBM · 서명 · 채팅 · 작업자 · NFC/QR · 퀴즈 · ESG · 용어집 · 라이브 통역", 0.82, 6.55, 6.6, 0.65, 13, "D8DFEA")

    # Flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 16, 9, PAPER, radius=False)
    title(s, "처음 관리자가 기억할 전체 흐름", "화면을 모두 외우는 것이 아니라, 매일 반복되는 운영 순서를 먼저 익힙니다.", 2)
    flow = ["로그인", "프로필 설정", "대시보드 확인", "TBM 작성", "서명 추적", "문의 대응", "기록 관리"]
    colors = [BLUE, PINK, AMBER, GREEN, BLUE, PINK, AMBER]
    for i, name in enumerate(flow):
        x = 0.85 + i * 2.12
        add_rect(s, x, 3.15, 1.55, 1.35, "FFFFFF", LINE)
        add_chip(s, str(i + 1), x + 0.5, 3.43, 0.55, colors[i])
        add_text(s, name, x + 0.1, 4.05, 1.35, 0.25, 12, INK, True, PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(s, "→", x + 1.58, 3.67, 0.38, 0.25, 20, PINK, True, PP_ALIGN.CENTER)
    add_rect(s, 1.1, 6.1, 13.8, 0.84, "FFFFFF", LINE)
    add_text(s, "교육 핵심: TBM을 보내고, 서명 여부를 확인하고, 미확인자에게 연락하고, 기록을 남기는 흐름입니다.", 1.42, 6.37, 12.5, 0.28, 16, INK, True)

    # TBM workflow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 16, 9, PAPER, radius=False)
    title(s, "TBM 작동방법: 관리자와 작업자가 같은 기록을 봅니다", "TBM은 단순 공지가 아니라 발송, 번역, 확인, 서명, 보고까지 이어지는 운영 흐름입니다.", 3)
    steps = [
        ("1", "작성", "관리자가 작업·위험요인·보호구를 입력"),
        ("2", "보강", "AI 가이드와 용어집으로 문장을 정리"),
        ("3", "발송", "현장 작업자에게 언어별 안내 전송"),
        ("4", "확인", "작업자가 읽고 서명 또는 NFC/QR 확인"),
        ("5", "조치", "미확인자 재알림, 채팅, 보고서 기록"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.82 + i * 3.02
        add_rect(s, x, 2.15, 2.45, 2.25, "FFFFFF", LINE)
        add_chip(s, num, x + 0.22, 2.45, 0.5, [BLUE, PINK, AMBER, GREEN, BLUE][i])
        add_text(s, head, x + 0.82, 2.44, 1.25, 0.26, 17, INK, True)
        add_text(s, body, x + 0.25, 3.18, 1.9, 0.58, 11, MUTED)
        if i < 4:
            add_text(s, "→", x + 2.48, 3.02, 0.35, 0.25, 22, PINK, True, PP_ALIGN.CENTER)
    add_rect(s, 1.0, 5.25, 14.0, 1.55, "FFFFFF", LINE)
    add_text(s, "더미 TBM 예시", 1.35, 5.58, 2.0, 0.28, 15, BLUE, True)
    add_text(s, "3층 철근 양중 작업: 크레인 작업 반경 5m 접근 금지, 안전고리 체결, 신호수 지시 준수, 낙하물 위험구역 통제선 설치", 3.0, 5.58, 10.8, 0.46, 14, INK, True)
    add_text(s, "관리자는 이 문장을 입력하고, 발송 후 서명 현황에서 Nguyen An / Somchai / Rustam 등 작업자의 확인 여부를 추적합니다.", 1.35, 6.22, 12.2, 0.34, 12, MUTED)

    # Dummy data pack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 16, 9, PAPER, radius=False)
    title(s, "교육용 더미 데이터 구성", "이번 PPT의 실제 화면은 아래 더미 데이터를 넣은 교육 현장을 기준으로 캡처했습니다.", 4)
    columns = [0.9, 4.25, 7.6, 10.95]
    groups = [
        ("현장", [("현장명", "SAFE-LINK 교육 현장"), ("코드", "SL-000001-0101"), ("담당", "교육용 관리자")]),
        ("작업자", [("김민수", "KR / 철근 / 서명 완료"), ("Nguyen An", "VN / 형틀 / 채팅 예시"), ("Somchai", "TH / 비계 / 미확인")]),
        ("TBM", [("주제", "3층 철근 양중 작업"), ("위험", "낙하물·고소작업·동선"), ("조치", "안전고리·통제선·신호수")]),
        ("운영 기록", [("채팅", "한국어↔베트남어"), ("용어집", "양중·안전고리·통제선"), ("ESG", "TBM/서명/장비 지급 집계")]),
    ]
    for gi, (head, rows) in enumerate(groups):
        x = columns[gi]
        add_rect(s, x, 2.05, 2.85, 4.5, "FFFFFF", LINE)
        add_chip(s, head, x + 0.25, 2.35, 1.1, [BLUE, PINK, AMBER, GREEN][gi])
        for ri, (k, v) in enumerate(rows):
            y = 3.1 + ri * 0.92
            add_text(s, k, x + 0.28, y, 0.95, 0.24, 11.5, INK, True)
            add_text(s, v, x + 1.15, y, 1.38, 0.38, 9.5, MUTED)
    add_text(s, "이 데이터는 교육용이며, 실제 운영 시에는 현장명·작업자·TBM 내용·서명 기록이 운영 DB 기준으로 표시됩니다.", 1.0, 7.15, 12.8, 0.38, 13, INK, True)

    for i, (filename, heading, subtitle, notes) in enumerate(SLIDES, start=3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_rect(s, 0, 0, 16, 9, PAPER, radius=False)
        title(s, heading, subtitle, i + 2)
        add_screen(s, SCREEN_DIR / filename)
        add_rect(s, 11.08, 1.75, 4.28, 6.86, PANEL, LINE)
        add_chip(s, "따라 하기", 11.42, 2.12, 1.15, BLUE)
        for n, (head, body) in enumerate(notes):
            y = 2.82 + n * 1.45
            add_chip(s, str(n + 1), 11.45, y, 0.42, [BLUE, PINK, AMBER][n % 3])
            add_text(s, head, 12.0, y - 0.06, 2.8, 0.34, 18, INK, True)
            add_text(s, body, 12.0, y + 0.38, 3.0, 0.72, 15, MUTED, True)

    # Checklist
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 16, 9, PAPER, radius=False)
    title(s, "관리자 일일 운영 체크리스트", "실제 화면 기준으로 매일 확인할 항목입니다.", len(SLIDES) + 5)
    checks = [
        ("출근 직후", "대시보드에서 서비스 상태와 담당 현장을 확인"),
        ("작업 전", "TBM 작성 화면에서 오늘 교육을 작성하고 발송"),
        ("작업 시작 전", "서명 현황에서 미확인자에게 재안내"),
        ("작업 중", "1:1 채팅 또는 라이브 통역으로 문의 처리"),
        ("퇴근 전", "퀴즈, ESG, 서명 기록을 보고 자료로 정리"),
    ]
    for i, (head, body) in enumerate(checks):
        y = 1.95 + i * 1.08
        add_rect(s, 1.0, y, 14.0, 0.75, "FFFFFF", LINE)
        add_chip(s, str(i + 1), 1.32, y + 0.2, 0.46, [BLUE, PINK, AMBER, GREEN, BLUE][i])
        add_text(s, head, 2.0, y + 0.14, 2.2, 0.36, 18, INK, True)
        add_text(s, body, 4.2, y + 0.14, 10.0, 0.36, 16, MUTED, True)

    normalize_deck_text(prs)
    prs.save(PPTX)
    make_previews()
    print(f"pptx={PPTX}")
    print(f"screens={SCREEN_DIR}")
    print(f"previews={PREVIEWS}")


def fnt(size, bold=False):
    p = r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf"
    return ImageFont.truetype(p, size=size) if Path(p).exists() else ImageFont.load_default()


def make_previews():
    thumbs = []
    for idx, (filename, heading, _, _) in enumerate(SLIDES, start=1):
        img = Image.new("RGB", (640, 360), "#" + PAPER)
        d = ImageDraw.Draw(img)
        d.text((26, 18), heading, font=fnt(22, True), fill="#" + INK)
        screen = Image.open(SCREEN_DIR / filename).convert("RGB")
        screen.thumbnail((410, 280))
        img.paste(screen, (28, 68))
        d.rounded_rectangle((458, 86, 612, 286), radius=12, fill="#FFFFFF", outline="#" + LINE, width=2)
        d.text((476, 116), "실제 화면", font=fnt(18, True), fill="#" + BLUE)
        d.text((476, 154), "따라 하기", font=fnt(16, True), fill="#" + INK)
        out = PREVIEWS / f"real_slide_{idx:02d}.png"
        img.save(out)
        thumbs.append(img.resize((320, 180)))
    sheet = Image.new("RGB", (4 * 320, ((len(thumbs) + 3) // 4) * 180), "#FFFFFF")
    for i, im in enumerate(thumbs):
        sheet.paste(im, ((i % 4) * 320, (i // 4) * 180))
    sheet.save(PREVIEWS / "contact_sheet.png")


if __name__ == "__main__":
    build()
