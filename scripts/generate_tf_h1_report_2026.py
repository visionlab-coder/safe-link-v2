# -*- coding: utf-8 -*-
"""
미래전략TF 2026 상반기 보고 (발표: 2026-06-02, 김무빈 차장 / 팀원 공유용)
근거(1차 자료): GS본사_브리핑자료_20260528_v3.pdf(15p, 마스터), 품의서_모빌리오_PoC실비_20260529_v3,
                과천G타운_CCTV_연구현황, seowon H2 선전 박람회 계획, 특허리뷰_20260526.
원칙: 기술 과장 방지 · 확정/진행/준비/검토 명확 구분 · 검증 가능 표현 우선.
디자인: 다크 + 레드/시안 + Nano Banana Pro 이미지(표지+섹션).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ASSET = 'docs/generated/report_assets'
OUTFILE = 'docs/generated/미래전략TF_2026_상반기보고_20260602.pptx'

BG=RGBColor(0x08,0x0A,0x12); CARD=RGBColor(0x12,0x16,0x24); CARD2=RGBColor(0x18,0x1E,0x30)
RED=RGBColor(0xC0,0x39,0x2B); CYAN=RGBColor(0x2D,0xD4,0xE8); GOLD=RGBColor(0xE5,0xB8,0x55)
GREEN=RGBColor(0x3F,0xC3,0x80); WHITE=RGBColor(0xF4,0xF6,0xFA); GRAY=RGBColor(0xA8,0xB0,0xC0)
DGRAY=RGBColor(0x6B,0x74,0x86); LINE=RGBColor(0x2A,0x31,0x45); FONT='맑은 고딕'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; W=13.333

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s
def rect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def set_alpha(sp,a):
    srgb=sp._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'),{'val':str(int(a*1000))}))
def txt(s,x,y,w,h,runs,size=18,bold=False,color=WHITE,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=1.0,italic=False):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs,str): runs=[(runs,color,bold,size)]
    first=True
    for it in runs:
        t,c,b,sz=(it+(color,bold,size))[:4] if isinstance(it,tuple) else (it,color,bold,size)
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.line_spacing=spacing
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=italic
        r.font.color.rgb=c; r.font.name=FONT
    return tb
def bullets(s,x,y,w,h,items,size=14,color=GRAY,gap=1.3,mk='·  ',mkcolor=RED):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,t in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=gap; p.space_after=Pt(2)
        rm=p.add_run(); rm.text=mk; rm.font.size=Pt(size); rm.font.bold=True; rm.font.color.rgb=mkcolor; rm.font.name=FONT
        rt=p.add_run(); rt.text=t; rt.font.size=Pt(size); rt.font.color.rgb=color; rt.font.name=FONT
    return tb
def header(s,kicker,title,sub=None):
    rect(s,0,0,0.18,7.5,RED)
    txt(s,0.7,0.42,11,0.35,kicker,12,bold=True,color=CYAN)
    txt(s,0.7,0.74,12,0.7,title,26,bold=True,color=WHITE)
    if sub: txt(s,0.7,1.42,12,0.4,sub,13,color=GRAY)
    rect(s,0.7,1.92,11.93,0.018,LINE)
def footer(s,page):
    txt(s,0.7,7.04,9,0.3,"서원토건 미래전략TF  ·  2026 상반기 보고  ·  검증 가능 표현 기준",9,color=DGRAY)
    txt(s,11.6,7.04,1,0.3,str(page),9,color=DGRAY,align=PP_ALIGN.RIGHT)
def card(s,x,y,w,h,color=CARD,accent=None):
    c=rect(s,x,y,w,h,color)
    if accent: rect(s,x,y,w,0.07,accent)
    return c
def chip(s,x,y,label,color):
    w=0.12+len(label)*0.135
    rect(s,x,y,w,0.32,color)
    txt(s,x,y,w,0.32,label,10.5,bold=True,color=BG,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    return w
def divider(img,kicker,title,sub):
    s=slide()
    s.shapes.add_picture(f'{ASSET}/{img}',0,0,width=Inches(W),height=Inches(7.5))
    ov=rect(s,0,0,W,7.5,RGBColor(0x05,0x06,0x0C)); set_alpha(ov,42)
    left=rect(s,0,0,7.2,7.5,RGBColor(0x05,0x06,0x0C)); set_alpha(left,32)
    rect(s,0.7,3.0,0.9,0.06,RED)
    txt(s,0.7,3.2,8,0.4,kicker,14,bold=True,color=CYAN)
    txt(s,0.7,3.62,8.2,1.4,title,38,bold=True,color=WHITE,spacing=1.05)
    txt(s,0.72,5.05,7.6,0.8,sub,14,color=GRAY,spacing=1.2)

# ── S1 표지 ──────────────────────────────────────────
s=slide()
s.shapes.add_picture(f'{ASSET}/cover.png',0,0,width=Inches(W),height=Inches(7.5))
ov=rect(s,0,0,W,7.5,RGBColor(0x05,0x06,0x0C)); set_alpha(ov,30)
ld=rect(s,0,0,8.4,7.5,RGBColor(0x04,0x05,0x0A)); set_alpha(ld,30)
rect(s,0,0,W,0.12,RED); rect(s,0,7.38,W,0.12,RED)
txt(s,0.85,1.45,9,0.45,"(주)서원토건 · 미래전략TF",15,bold=True,color=CYAN)
txt(s,0.85,2.3,11,1.6,[("2026 상반기 보고",WHITE,True,52)],spacing=1.0)
rect(s,0.9,3.5,1.1,0.06,RED)
txt(s,0.85,3.72,11.5,0.5,"로봇 도입 · AI 안전 커뮤니케이션 · 산학협력 · 특허",18,bold=True,color=WHITE)
txt(s,0.85,4.35,11,0.45,"로봇·AI·안전을 현장 적용 사업 과제로 — 단계별 검증 중심",13,color=GRAY,italic=True)
txt(s,0.85,6.2,9,0.4,"발표  김무빈 차장 (미래전략TF 팀장)   |   2026. 06. 02",13,bold=True,color=WHITE)
txt(s,0.85,6.68,11,0.35,"※ 본 보고는 확정/진행/준비/검토를 구분한 검증 가능 진척 중심 (과장 표현 배제)",10.5,color=DGRAY)

# ── S2 목차 ──────────────────────────────────────────
s=slide(); rect(s,0,0,0.18,7.5,RED)
txt(s,0.7,0.6,11,0.4,"AGENDA",13,bold=True,color=CYAN)
txt(s,0.7,0.98,11,0.7,"목   차",30,bold=True,color=WHITE)
rect(s,0.7,1.78,11.93,0.018,LINE)
ag=[("01","미래전략TF 추진 방향","로봇·AI·안전 6대 축과 기획 원칙"),
    ("02","로봇 도입 — 모빌리오 협업","사족보행 로봇 현장 적용 PoC (과천)"),
    ("03","SAFE-LINK 안전 커뮤니케이션","20개국어 전달·검증·기록 — 현재 PoC 단계"),
    ("04","산학협력 · R&D · 특허","성균관대 MOU · 경희대 협의 · 특허 출원"),
    ("05","대외활동 · 박람회","WIS·AI EXPO·선전 하이테크 박람회"),
    ("06","하반기 로드맵 및 관리 과제","6~12월 실행 과제 · 리스크 관리")]
y=2.05
for no,t,d in ag:
    rect(s,0.7,y,0.78,0.78,CARD)
    txt(s,0.7,y,0.78,0.78,no,22,bold=True,color=RED,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.7,y+0.07,10.8,0.45,t,18,bold=True,color=WHITE)
    txt(s,1.7,y+0.5,10.8,0.3,d,12,color=GRAY); y+=0.86
footer(s,2)

# ── S3 Executive Summary ─────────────────────────────
s=slide()
header(s,"EXECUTIVE SUMMARY","상반기 진척 요약 — 실행 기반 확보",
       "로봇·특허·산학·SAFE-LINK PoC 준비까지 — 확정/진행/준비/검토를 구분해 관리")
items=[("모빌리오 협업","확정",GREEN,"본사 미팅·과천 답사 완료\n6월 과천 현장 PoC 진행"),
       ("성균관대 MOU","확정",GREEN,"건설환경공학부 공동연구\n과천 G타운 연구용 CCTV 실증"),
       ("SAFE-LINK","PoC 준비",CYAN,"20개국어 안전 커뮤니케이션\n청주·과천 현장 적용 검증"),
       ("특허 출원","최종 단계",GOLD,"발명신고·변리사 검토 완료\n6월 출원 예정"),
       ("박람회·전시","참관 완료",GREEN,"WIS·AI EXPO KOREA 참관 완료\n선전 하이테크(11월) 검토")]
cw=2.28; gap=0.16; x0=0.7; y0=2.25
for i,(t,tag,ac,desc) in enumerate(items):
    x=x0+i*(cw+gap); card(s,x,y0,cw,3.15,CARD,accent=ac)
    txt(s,x+0.18,y0+0.28,cw-0.36,0.5,t,17,bold=True,color=WHITE)
    chip(s,x+0.18,y0+0.95,tag,ac)
    rect(s,x+0.18,y0+1.5,cw-0.36,0.016,LINE)
    txt(s,x+0.18,y0+1.68,cw-0.36,1.3,desc,11.5,color=GRAY,spacing=1.18)
rect(s,0.7,5.9,11.93,0.9,CARD2)
txt(s,0.95,6.02,11.4,0.7,[("관리 원칙   ",CYAN,True,13),
    ("기술 과장 표현을 배제하고, 확정·논의 중·계획 중을 명확히 구분하여 '검증 가능한 진척'만 보고합니다.",WHITE,False,13.5)],spacing=1.15)
footer(s,3)

# ── S4 TF 추진 방향 ─────────────────────────────────
s=slide()
header(s,"01  TF DIRECTION","미래전략TF 추진 방향 — 로봇·AI·안전",
       "로봇·AI·안전을 현장 적용 사업 과제로 설계하고, 단계별로 검증")
axes=[("01","AI · 플랫폼",CYAN),("02","로봇 · 자동화",RED),("03","스마트안전",GOLD),
      ("04","R&D · 특허",CYAN),("05","산학협력",RED),("06","전략 파트너십",GOLD)]
cw=1.86; x0=0.7; y0=2.3
for i,(n,t,ac) in enumerate(axes):
    x=x0+i*(cw+0.075); card(s,x,y0,cw,1.5,CARD,accent=ac)
    txt(s,x+0.16,y0+0.26,cw-0.3,0.4,n,12,bold=True,color=ac)
    txt(s,x+0.16,y0+0.7,cw-0.3,0.6,t,13.5,bold=True,color=WHITE,spacing=1.0)
panels=[("기획 원칙","현장성 · 법적 증빙 · 확장성 · 운영 편의성",CYAN),
        ("진행 현황","로봇 PoC 진행 · 특허 출원 준비 · 산학협력 추진",GREEN),
        ("관리 포인트","기술 과장 방지 · 단계별 검증 · 비용·운영 리스크 관리",GOLD)]
y=4.2
for t,d,ac in panels:
    card(s,0.7,y,11.93,0.78,CARD2); rect(s,0.7,y,0.07,0.78,ac)
    txt(s,1.0,y,2.6,0.78,t,14,bold=True,color=ac,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.6,y,8.9,0.78,d,13,color=WHITE,anchor=MSO_ANCHOR.MIDDLE); y+=0.88
footer(s,4)

# ── S5 DIV 로봇 ─────────────────────────────────────
divider("sec_robot.png","SECTION 02",
        "로봇 도입\n모빌리오 협업",
        "사족보행 로봇을 안전·생산성·디지털 전환 과제로 — 단계적 PoC 후 확산 판단")

# ── S6 로봇 도입 사업 방향 ───────────────────────────
s=slide()
header(s,"02  로봇 도입","건설 현장 로봇 도입을 핵심 과제로 추진",
       "안전·생산성·디지털 전환 관점에서 사족보행 로봇의 현장 적용 가능성 검토")
why=[("안전성 확보","사람이 접근하기 어려운 위험 구역을 로봇이 대체 점검",RED),
     ("운영 효율","반복 점검 작업의 자동화로 관리자 생산성 향상",CYAN),
     ("디지털 전환","현장 데이터 자동 수집·기록으로 안전관리 체계 고도화",GOLD)]
y0=2.3
for i,(t,d,ac) in enumerate(why):
    y=y0+i*1.05; card(s,0.7,y,7.4,0.92,CARD); rect(s,0.7,y,0.07,0.92,ac)
    txt(s,1.0,y,2.5,0.92,t,15,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.4,y,4.5,0.92,d,12,color=GRAY,anchor=MSO_ANCHOR.MIDDLE,spacing=1.12)
card(s,8.35,2.3,4.28,3.67,CARD2)
txt(s,8.6,2.5,3.8,0.4,"검토 영역",13,bold=True,color=CYAN)
rect(s,8.6,2.92,3.8,0.016,LINE)
bullets(s,8.6,3.12,3.8,2.7,["사족보행 로봇 플랫폼 검토","현장 점검 · 모니터링 적용",
    "안전관리 데이터 자동 수집","SAFE-LINK 운영체계와 연계","단계적 PoC 후 확산 판단"],size=12.5,gap=1.55,mkcolor=CYAN)
txt(s,0.7,6.2,11.93,0.5,[("관리 포인트   ",GOLD,True,12),
    ("기술 기대효과보다 '검증 가능한 업무 개선'을 우선 — 단계적 PoC로 현장 적용성을 정량 확인.",GRAY,False,12)])
footer(s,6)

# ── S7 모빌리오 협업 진행 ────────────────────────────
s=slide()
header(s,"02  로봇 도입","모빌리오 협업 — 본사 미팅·과천 답사·6월 PoC",
       "사족보행 로봇 기반 건설현장 적용성을 단계별 검증으로 확인")
card(s,0.7,2.3,4.3,3.67,CARD2)
txt(s,0.95,2.5,3.8,0.4,"PARTNER",11,bold=True,color=CYAN)
txt(s,0.95,2.82,3.8,0.5,"(주)모빌리오",19,bold=True,color=WHITE)
txt(s,0.95,3.38,3.8,0.4,"산업용 사족보행 로봇 · 디지털 트윈",12,color=GRAY)
rect(s,0.95,3.85,3.55,0.016,LINE)
bullets(s,0.95,4.05,3.8,1.9,["자율주행 사족보행 로봇 플랫폼","위험 지역 모니터링 · 현장 점검",
    "LiDAR · 카메라 · 가스 센서 모듈","CES 2026 혁신상 · 글로벌 실증"],size=11.5,gap=1.45)
steps=[("완료","본사 미팅","모빌리오 본사 방문·약 2시간 협의, 적용 가능성 사전 협의",GREEN),
       ("5.28","과천 현장 답사·1차 시연","로봇 주행환경·적용 영역 현장 확인 (완료)",RED),
       ("6월","과천 현장 PoC","현장 적용성·운영 가능성 검증 후 확산 판단",CYAN)]
y=2.3
for tag,t,d,ac in steps:
    card(s,5.2,y,7.43,1.12,CARD); rect(s,5.2,y,1.2,1.12,ac)
    txt(s,5.2,y,1.2,1.12,tag,15,bold=True,color=BG,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,6.6,y+0.16,5.8,0.4,t,14.5,bold=True,color=WHITE)
    txt(s,6.6,y+0.6,5.8,0.45,d,11.5,color=GRAY,spacing=1.1); y+=1.27
footer(s,7)

# ── S8 모빌리오 6월 PoC 상세 ────────────────────────
s=slide()
header(s,"02  로봇 도입","6월 모빌리오 PoC — 양사 경영진 참관",
       "사족보행 로봇 건설현장 도입 타당성 검증(PoC) · 거래처 (주)모빌리오")
sched=[("5.28 (목)","1차 현장 시연 — 완료",GREEN),
       ("6.05 (금)","로봇개 사전 셋팅 및 시운전",CYAN),
       ("6월 2째 주","서원토건 대표이사 · 대우건설 CSO 현장 PoC 참관",RED),
       ("6.22 (월)","대우건설 회장 · 서원토건 이강범 전무 PoC 참관",RED)]
y0=2.25
for i,(d,t,ac) in enumerate(sched):
    y=y0+i*0.82; card(s,0.7,y,7.4,0.7,CARD); rect(s,0.7,y,0.07,0.7,ac)
    txt(s,1.0,y,1.95,0.7,d,13,bold=True,color=ac,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.0,y,5.0,0.7,t,12.5,color=WHITE,anchor=MSO_ANCHOR.MIDDLE,spacing=1.05)
card(s,8.35,2.25,4.28,1.55,CARD2,accent=GOLD)
txt(s,8.6,2.42,3.8,0.4,"집행 / 조건",12,bold=True,color=GOLD)
txt(s,8.6,2.8,3.85,0.9,[("₩550,000  ",WHITE,True,16),("(예약 보증금)",GRAY,False,11),
    ("\n도입 확정 시 구매대금 100% 차감 · 미진행 시 추가 비용 없음",GRAY,False,11)],spacing=1.12)
card(s,8.35,3.95,4.28,2.02,CARD2,accent=CYAN)
txt(s,8.6,4.12,3.8,0.4,"PoC 서비스 범위",12,bold=True,color=CYAN)
bullets(s,8.6,4.5,3.85,1.4,["현장 정밀 맵핑 · 자율주행 세팅","엔지니어 현장 파견·운용 교육",
    "PoC 결과 타당성 분석 보고서"],size=11,gap=1.35,mkcolor=CYAN)
rect(s,0.7,5.62,7.4,0.78,CARD2)
txt(s,0.95,5.74,7.0,0.6,[("전략 의의   ",CYAN,True,12),
    ("양사(서원토건·대우건설) 최고경영진 참관으로 전략적 파트너십 가시화.",WHITE,False,12)],spacing=1.12)
txt(s,0.7,7.04,11.5,0.3,"견적 22605141423 · (주)모빌리오 837-88-02126 · PoC 타당성 분석 보고서 제공",8.5,color=DGRAY)
txt(s,12.0,7.04,0.6,0.3,"8",9,color=DGRAY,align=PP_ALIGN.RIGHT)

# ── S9 DIV SAFE-LINK ────────────────────────────────
divider("sec1_product.png","SECTION 03",
        "SAFE-LINK\n안전 커뮤니케이션",
        "외국인 근로자 안전 커뮤니케이션을 전달·검증·기록 가능한 운영 프로세스로 — 현재 PoC 단계")

# ── S10 SAFE-LINK 개요 ──────────────────────────────
s=slide()
header(s,"03  SAFE-LINK","현장 안전 운영체계의 디지털 전환 과제",
       "다국어 번역 기능이 아니라, 안전 커뮤니케이션의 운영 리스크를 줄이는 관리 체계")
stat=[("20","개국어","한국어 포함 다국어 대응",CYAN),
      ("1 Touch","NFC","NFC 기반 TBM 교육 가동",RED),
      ("Evidence","증빙","교육·지시·이력 자동 기록",GOLD),
      ("PoC","검증","현장 적용성 검증 중심",GREEN)]
cw=2.85; x0=0.7; y0=2.3
for i,(big,sub,d,ac) in enumerate(stat):
    x=x0+i*(cw+0.13); card(s,x,y0,cw,1.85,CARD,accent=ac)
    txt(s,x+0.2,y0+0.28,cw-0.4,0.55,big,22,bold=True,color=WHITE)
    txt(s,x+0.2,y0+0.92,cw-0.4,0.35,sub,12,bold=True,color=ac)
    txt(s,x+0.2,y0+1.3,cw-0.4,0.45,d,11,color=GRAY,spacing=1.1)
card(s,0.7,4.45,5.9,1.55,CARD2)
txt(s,0.95,4.6,5.4,0.4,"현재 단계",12,bold=True,color=CYAN)
txt(s,0.95,4.98,5.5,1.0,"서원토건 미래전략TF 내부 기획 및 PoC 준비 단계.\nPoC 현장: 청주 센텀 푸르지오자이 · 과천 G-town",12.5,color=WHITE,spacing=1.25)
card(s,6.73,4.45,5.9,1.55,CARD2)
txt(s,6.98,4.6,5.4,0.4,"해결 접근",12,bold=True,color=GOLD)
txt(s,6.98,4.98,5.5,1.0,"TBM·위험지시·작업중지권 전달의 공백을\nPTT·NFC·AI 번역·역번역 검증·안전일지 증빙으로 연결",12.5,color=WHITE,spacing=1.25)
footer(s,10)

# ── S11 SAFE-LINK 처리 흐름 ─────────────────────────
s=slide()
header(s,"03  SAFE-LINK","발화·번역·검증·전달·기록을 하나의 흐름으로",
       "핵심은 '전달'이 아니라, 전달 결과를 검증하고 기록으로 남기는 현장 운영 프로세스")
flow=[("STEP 1","관리자 발화","PTT"),("STEP 2","STT 변환","음성→텍스트"),
      ("STEP 3","AI 번역","20개국어"),("STEP 4","역번역 검증","오역 리스크 관리"),
      ("STEP 5","근로자 수신","모바일·스피커"),("STEP 6","일지·증빙 저장","이력 자동화")]
cw=1.92; x0=0.7; y0=2.5
for i,(stp,t,d) in enumerate(flow):
    x=x0+i*(cw+0.07)
    hl=(i==3)
    card(s,x,y0,cw,1.7,CARD2 if hl else CARD)
    if hl: rect(s,x,y0,cw,0.06,RED)
    txt(s,x+0.12,y0+0.22,cw-0.24,0.3,stp,10,bold=True,color=(RED if hl else DGRAY),align=PP_ALIGN.CENTER)
    txt(s,x+0.12,y0+0.6,cw-0.24,0.5,t,13,bold=True,color=WHITE,align=PP_ALIGN.CENTER,spacing=1.0)
    txt(s,x+0.12,y0+1.18,cw-0.24,0.4,d,10.5,color=GRAY,align=PP_ALIGN.CENTER)
    if i<5: txt(s,x+cw-0.05,y0+0.55,0.3,0.5,"›",22,bold=True,color=DGRAY)
card(s,0.7,4.7,11.93,1.3,CARD2)
txt(s,0.95,4.85,11.4,0.4,"핵심 차별",12,bold=True,color=CYAN)
txt(s,0.95,5.25,11.4,0.7,"역번역 검증으로 오역에 의한 안전사고 가능성을 사전 차단하고, NFC 태그 기반 TBM 교육 자동 가동과\n이력 자동 관리로 전체 흐름을 단일 운영 체계로 통합합니다.",13,color=WHITE,spacing=1.25)
footer(s,11)

# ── S12 SAFE-LINK 상반기 진척 ───────────────────────
s=slide()
header(s,"03  SAFE-LINK","SAFE-LINK 상반기 진척",
       "고도화 구조 확정 + PoC 준비 — 다음 초점은 현장 시연·검증")
prog=[("고도화 구조 확정","PTT·NFC·역번역 검증 흐름 기획 완료, 단일 운영 체계로 정리",GREEN,"확정"),
      ("PoC 준비","청주·과천 현장 적용성 검증 준비 — 관리자 교육·단말·네트워크 사전 점검",CYAN,"진행"),
      ("내부 보안·무결성","역번역 검증·이력 기록 등 증빙 신뢰성 확보 (개인정보·음성·영상 증빙 체계 검토)",GOLD,"진행"),
      ("다음 초점","6~7월 내부 시연 → 9월 2차 현장 후보 검토 → 특허·스마트글라스 연동 단계화",DGRAY,"예정")]
y0=2.3
for i,(t,d,ac,tag) in enumerate(prog):
    y=y0+i*1.08; card(s,0.7,y,11.93,0.95,CARD); rect(s,0.7,y,0.07,0.95,ac)
    txt(s,1.0,y+0.15,3.0,0.65,t,15,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    chip(s,4.05,y+0.32,tag,ac)
    txt(s,5.4,y+0.15,7.0,0.65,d,11.5,color=GRAY,anchor=MSO_ANCHOR.MIDDLE,spacing=1.12)
footer(s,12)

# ── S13 DIV 산학·R&D ────────────────────────────────
divider("sec4_academy.png","SECTION 04",
        "산학협력\nR&D · 특허",
        "확정 과제(MOU·과제수행)와 사전 협의를 명확히 구분해 관리")

# ── S14 산학협력 ────────────────────────────────────
s=slide()
header(s,"04  산학협력","성균관대 MOU 완료 · 경희대 사전 협의",
       "성균관대는 협력 기반 확보, 경희대는 통신 기술 관점의 사전 협의로 진행")
# 성균관대
card(s,0.7,2.3,5.9,3.67,CARD,accent=RED)
txt(s,0.95,2.5,4.3,0.45,"성균관대학교",18,bold=True,color=WHITE)
chip(s,5.05,2.55,"MOU 완료",GREEN)
txt(s,0.95,3.05,5.4,0.35,"건설환경공학부 연구팀",12,bold=True,color=RED)
rect(s,0.95,3.45,5.4,0.016,LINE)
bullets(s,0.95,3.62,5.4,2.3,["건설 AI 플랫폼 · 시설물 유지관리 공동연구 검토",
    "로봇 자동화 · 자율시공 영역 연구 협력","과천 G타운 연구용 CCTV 실증 (3월 1·2차 설치)",
    "향후 과제수행협약 및 R&D 연계 검토"],size=12,gap=1.5)
# 경희대
card(s,6.73,2.3,5.9,3.67,CARD,accent=CYAN)
txt(s,6.98,2.5,4.3,0.45,"경희대학교",18,bold=True,color=WHITE)
chip(s,11.1,2.55,"사전 협의",GOLD)
txt(s,6.98,3.05,5.4,0.35,"AI-RAN / 6G 네트워크 기술",12,bold=True,color=CYAN)
rect(s,6.98,3.45,5.4,0.016,LINE)
bullets(s,6.98,3.62,5.4,2.3,["AI-RAN / 6G 네트워크 기술 협력 가능성 탐색",
    "지하·터널·음영지역 통신 안정성 관점 접점","현재 단계: 아이디어 협의",
    "확정 협력은 별도 단계(MOU·과제수행)로 관리"],size=12,gap=1.5,mkcolor=CYAN)
txt(s,0.7,6.15,11.93,0.5,[("산학협력 원칙   ",GOLD,True,12),
    ("확정 협력은 MOU·과제수행 단위로, 사전 협의는 별도 단계로 분리 관리.",GRAY,False,12)])
footer(s,14)

# ── S15 정부공동과제 + 특허 ─────────────────────────
s=slide()
header(s,"04  R&D · 특허","정부공동과제 물색 · 특허 출원 최종단계",
       "산학·파트너 역량을 국책과제와 특허로 연결")
# 정부과제
card(s,0.7,2.3,5.9,3.67,CARD2,accent=CYAN)
txt(s,0.95,2.5,5.4,0.4,"정부공동과제 · 국책과제",15,bold=True,color=CYAN)
rect(s,0.95,2.95,5.4,0.016,LINE)
bullets(s,0.95,3.15,5.4,2.7,["성균관대·경희대 정부공동과제 물색",
    "모빌리오와 정부과제(로봇·안전) 연계 논의","9월 국책과제 후보 본격 검토 (하반기)",
    "건설 AI·로봇·안전 분야 R&D 연계"],size=12.5,gap=1.7,mkcolor=CYAN)
# 특허
card(s,6.73,2.3,5.9,3.67,CARD2,accent=GOLD)
s.shapes.add_picture(f'{ASSET}/sec2_patent.png',Inches(6.73),Inches(2.3),width=Inches(5.9),height=Inches(1.35))
ovp=rect(s,6.73,2.3,5.9,1.35,RGBColor(0x05,0x06,0x0C)); set_alpha(ovp,40)
txt(s,6.98,2.45,5.4,0.4,"특허 출원",15,bold=True,color=GOLD)
chip(s,11.35,2.5,"최종단계",GOLD)
txt(s,6.98,3.7,5.5,0.4,"다국어 현장 안전 커뮤니케이션의 검증 세션 및",12.5,bold=True,color=WHITE,spacing=1.05)
txt(s,6.98,3.98,5.5,0.4,"법적 증거 데이터셋 생성 시스템과 방법",12.5,bold=True,color=WHITE)
bullets(s,6.98,4.5,5.5,1.4,["발명신고서 정리·변리사 검토 완료 → 6월 출원 예정",
    "독립항 3 + 종속항 12 (총 15) · 관리번호 YDP20260105KR",
    "차별점: NFC/QR·번역·서명·해시를 하나의 세션으로 결합"],size=11,gap=1.4,mkcolor=GOLD)
footer(s,15)

# ── S16 DIV 대외활동 ────────────────────────────────
divider("sec3_business.png","SECTION 05",
        "대외활동\n박람회 · 전시",
        "국내 전시 참관 완료(WIS·AI EXPO KOREA) · 해외(선전) 참관 검토 — 기술 동향·파트너 발굴")

# ── S17 박람회·전시 ─────────────────────────────────
s=slide()
header(s,"05  대외활동","박람회 · 전시 — 국내 참관 완료",
       "WIS·AI EXPO KOREA 참관 완료로 PoC 파트너 발굴 · 해외(선전) 참관 검토")
ex=[("WIS 2026","참관 완료",GREEN,"스마트안전·통신·AI 신기술 동향 파악\n및 파트너 후보 탐색"),
    ("AI EXPO KOREA 2026","참관 완료",GREEN,"2026.05.06~08 · 코엑스 A홀 (350개사)\nSAFE-LINK 연동 파트너 발굴·PoC 협의")]
for i,(t,tag,ac,d) in enumerate(ex):
    x=0.7+i*6.03; card(s,x,2.3,5.9,1.55,CARD,accent=ac)
    txt(s,x+0.25,2.48,4.3,0.45,t,16,bold=True,color=WHITE)
    chip(s,x+4.75,2.52,tag,ac)
    txt(s,x+0.25,3.0,5.4,0.7,d,11.5,color=GRAY,spacing=1.15)
rect(s,0.7,4.05,11.93,0.55,CARD2)
txt(s,0.95,4.13,11.4,0.4,[("발굴 성과   ",GREEN,True,11.5),
    ("현장 안전 AI·스마트글라스(AR)·AI 번역·NFC·웨어러블 등 SAFE-LINK 연동 파트너 후보 다수 접촉 → P1 업체 PoC 협의 착수",GRAY,False,11.5)],anchor=MSO_ANCHOR.MIDDLE)
card(s,0.7,4.8,11.93,1.5,CARD2,accent=RED)
txt(s,0.95,4.92,11.4,0.4,"중국 하이테크 박람회 (해외)",13,bold=True,color=RED)
rect(s,0.95,5.32,11.5,0.016,LINE)
ch=[("상반기 · 광저우","WCME(자재·장비·시험·보수) + CIHIE(PC·BIM·스마트시공·건설로봇)","일정 캔슬 → 하반기 통합",GOLD),
    ("하반기 · 선전 (CHTF)","스마트건설·AI·로봇·BIM·현장 디지털화 신기술","11.26~28 · 대표·임원 동행 검토",CYAN)]
y=5.45
for t,d,st,ac in ch:
    txt(s,0.95,y,3.3,0.42,t,12,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.3,y,5.4,0.42,d,11,color=GRAY,anchor=MSO_ANCHOR.MIDDLE,spacing=1.05)
    txt(s,9.75,y,2.75,0.42,st,10.5,bold=True,color=ac,anchor=MSO_ANCHOR.MIDDLE,spacing=1.0); y+=0.45
footer(s,17)

# ── S18 DIV 하반기 ──────────────────────────────────
divider("sec5_future.png","SECTION 06",
        "하반기 로드맵\n관리 과제",
        "6~12월 실행 과제 — 로봇·특허·현장 검증·산학협력 병행 관리")

# ── S19 하반기 로드맵 (월별) ────────────────────────
s=slide()
header(s,"06  하반기 로드맵","2026 하반기 실행 과제 (6~12월)",
       "로봇·특허·현장 검증·산학협력을 월 단위로 정렬")
road=[("6월","모빌리오 과천 PoC · 특허 출원 · SAFE-LINK 기능 정의서 정리",RED),
      ("7월","모빌리오 PoC 결과 정리 · 경희대 논의 구체화 · 스마트글라스 연동 요구사항",CYAN),
      ("8월","스마트글라스 PoC 계획안 작성 · 현장 적용 리스크 검토",CYAN),
      ("9월","국책과제 후보 검토 · 2차 현장 후보군 점검",GOLD),
      ("10월","PoC 결과 프레임 정리 · 운영비·보안·법무 이슈 검토",CYAN),
      ("11월","선전 하이테크 박람회 참관 검토 · 해외 기술 탐색",RED),
      ("12월","2027년 실행계획 및 예산·조직 운영안 수립",GOLD)]
y0=2.25; rh=0.64
for i,(m,d,ac) in enumerate(road):
    y=y0+i*rh; rect(s,0.7,y,1.5,rh-0.08,CARD2); rect(s,0.7,y,0.06,rh-0.08,ac)
    txt(s,0.7,y,1.5,rh-0.08,m,14,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,2.3,y,10.33,rh-0.08,CARD)
    txt(s,2.55,y,9.9,rh-0.08,d,12.5,color=WHITE,anchor=MSO_ANCHOR.MIDDLE,spacing=1.05)
footer(s,19)

# ── S20 관리 과제 ───────────────────────────────────
s=slide()
header(s,"06  관리 과제","과장 표현보다 검증 가능한 표현을 우선",
       "확정/논의 중/계획 중을 분리하고, 법무·보안·운영 리스크를 선제 관리")
mp=[("표현 관리","확정 / 논의 중 / 계획 중을 명확히 구분하여 보고",CYAN),
    ("기술 검증","로봇·번역 정확도와 검증 기준을 정량 지표로 수치화",RED),
    ("현장 운영","관리자 교육 · 단말 배포 · 네트워크 환경 사전 점검",GOLD),
    ("법무 · 보안","개인정보·음성·영상 데이터 증빙 체계 사전 검토",GREEN)]
cw=5.9; ch2=1.5; x0=0.7; y0=2.3
for i,(t,d,ac) in enumerate(mp):
    x=x0+(i%2)*(cw+0.13); y=y0+(i//2)*(ch2+0.16)
    card(s,x,y,cw,ch2,CARD); rect(s,x,y,0.07,ch2,ac)
    txt(s,x+0.3,y+0.25,1.2,0.9,f"0{i+1}",26,bold=True,color=ac)
    txt(s,x+1.5,y+0.28,cw-1.7,0.45,t,16,bold=True,color=WHITE)
    txt(s,x+1.5,y+0.78,cw-1.7,0.6,d,12,color=GRAY,spacing=1.15)
rect(s,0.7,5.7,11.93,0.85,CARD2)
txt(s,0.95,5.85,11.4,0.6,[("핵심 방침   ",GOLD,True,13),
    ("기술 기대효과보다 '검증 가능한 업무 개선'을 우선 제시합니다.",WHITE,False,13.5)],spacing=1.1)
footer(s,20)

# ── S21 마무리/결론 ─────────────────────────────────
s=slide()
s.shapes.add_picture(f'{ASSET}/cover.png',0,0,width=Inches(W),height=Inches(7.5))
ov=rect(s,0,0,W,7.5,RGBColor(0x04,0x05,0x0A)); set_alpha(ov,55)
rect(s,0,0,W,0.12,RED); rect(s,0,7.38,W,0.12,RED)
txt(s,0,1.35,W,0.5,"CONCLUSION · 미래전략TF 2026 상반기",13,bold=True,color=CYAN,align=PP_ALIGN.CENTER)
txt(s,0,2.1,W,1.0,"로봇·AI·안전 통합 과제로 정교하게 검증",36,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,0,3.2,W,0.5,"상반기는 실행 기반 확보 — 하반기는 PoC 검증과 확산 판단의 단계",15,color=GRAY,align=PP_ALIGN.CENTER)
sums=[("상반기","실행 기반 확보 (협업·MOU·특허·PoC 준비)"),
      ("하반기","모빌리오 PoC·SAFE-LINK 시연·산학 병행 검증"),
      ("Next","과천 PoC 결과로 로봇 확산·SAFE-LINK 연계 시점 확정")]
cw=3.85; x0=(W-(cw*3+0.25*2))/2; y0=4.05
for i,(k,v) in enumerate(sums):
    x=x0+i*(cw+0.25); c=rect(s,x,y0,cw,1.25,RGBColor(0x0C,0x0E,0x18)); set_alpha(c,20)
    txt(s,x,y0+0.2,cw,0.35,k,12,bold=True,color=CYAN,align=PP_ALIGN.CENTER)
    txt(s,x+0.2,y0+0.56,cw-0.4,0.6,v,12,bold=True,color=WHITE,align=PP_ALIGN.CENTER,spacing=1.12)
txt(s,0,5.7,W,0.4,"감사합니다   |   Q & A",18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,0,6.55,W,0.35,"이강범 전무 · 김무빈 차장 · 경영지원본부 미래전략TF · visionlab@seowonenc.co.kr",11,color=GRAY,align=PP_ALIGN.CENTER)

os.makedirs('docs/generated',exist_ok=True)
prs.save(OUTFILE)
print('SAVED',OUTFILE,'slides:',len(prs.slides._sldIdLst))
