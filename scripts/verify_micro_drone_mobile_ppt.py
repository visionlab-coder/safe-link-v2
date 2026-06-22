from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606" / "SAFE-LINK_Micro_Drone_Nest_대우건설_Advanced_V2_이동형네스트_20260607.pptx"
DOC = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606" / "SAFE_LINK_Micro_Drone_Nest_이동형네스트_추가반영.md"
PREVIEW = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606" / "previews-micro-drone-advanced-v2" / "contact_sheet.png"

prs = Presentation(PPTX)
all_text = []
picture_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text:
            all_text.append(shape.text)
        if shape.shape_type == 13:
            picture_count += 1

joined = " ".join(all_text)
print(f"pptx_exists={PPTX.exists()}")
print(f"doc_exists={DOC.exists()}")
print(f"preview_exists={PREVIEW.exists()}")
print(f"slide_count={len(prs.slides)}")
print(f"picture_count={picture_count}")
print(f"has_mobile_nest={'4족보행 로봇은 이동형 네스트' in joined}")
print(f"has_fixed_mobile_split={'고정형 네스트' in joined and '이동형 네스트' in joined}")
print(f"has_safelink={'SAFE-LINK' in joined}")
print(f"has_budget={'5,500만' in joined}")
