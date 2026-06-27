# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
SCREEN_DIR = OUT / "real-screens"
WORKER_DIR = OUT / "worker-screens-20260601"
PREVIEW_DIR = OUT / "previews-admin-manual-20260601"
PPTX = OUT / "SAFE-LINK_관리자_교육용_사용설명서_최신업데이트_20260601.pptx"
PPTX_ASCII = OUT / "SAFE-LINK_admin_training_manual_latest_20260601.pptx"
CONTACT = PREVIEW_DIR / "contact_sheet.png"

INK = "111827"
MUTED = "5B6472"
PAPER = "F7F4EE"
WHITE = "FFFFFF"
LINE = "DDD6CC"
BLUE = "2563EB"
GREEN = "16A34A"
AMBER = "F59E0B"
RED = "DC2626"
NAVY = "121827"
PINK = "EF3D9A"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=True, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def bullet_list(slide, items: Iterable[str], x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(color)
    return box


def chip(slide, text, x, y, w, fill=BLUE):
    shape = add_rect(slide, x, y, w, 0.42, fill, radius=True)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = rgb(WHITE)
    return shape


def base_slide(prs: Presentation, title: str, subtitle: str | None = None, idx: int | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    add_text(slide, title, 0.62, 0.35, 11.8, 0.55, 28, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.98, 12.2, 0.42, 15.5, MUTED, True)
    if idx is not None:
        add_text(slide, f"SAFE-LINK v2.0 관리자 교육 | 최신 업데이트 반영 | {idx:02d}", 0.65, 8.55, 5.8, 0.2, 8, MUTED, True)
    return slide


def add_screen(slide, file_name: str, x=0.62, y=1.72, w=10.0, h=5.72):
    path = SCREEN_DIR / file_name
    add_rect(slide, x - 0.05, y - 0.05, w + 0.1, h + 0.1, WHITE, LINE, True)
    if not path.exists():
        add_text(slide, f"화면 캡처 없음\n{file_name}", x + 0.6, y + 2.1, w - 1.2, 0.8, 22, MUTED, True, PP_ALIGN.CENTER)
        return
    with Image.open(path) as img:
        iw, ih = img.size
    target = w / h
    source = iw / ih
    if source >= target:
        ph = h
        pw = h * source
    else:
        pw = w
        ph = w / source
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(pw), height=Inches(ph))
    pic.crop_left = max(0, (pw - w) / pw / 2)
    pic.crop_right = max(0, (pw - w) / pw / 2)
    pic.crop_top = max(0, (ph - h) / ph / 2)
    pic.crop_bottom = max(0, (ph - h) / ph / 2)


def add_image_from_path(slide, path: Path, x=0.62, y=1.45, w=5.1, h=6.3):
    add_rect(slide, x - 0.05, y - 0.05, w + 0.1, h + 0.1, WHITE, LINE, True)
    if not path.exists():
        add_text(slide, f"이미지 없음\n{path.name}", x + 0.4, y + 2.8, w - 0.8, 0.7, 20, MUTED, True, PP_ALIGN.CENTER)
        return
    with Image.open(path) as img:
        iw, ih = img.size
    target = w / h
    source = iw / ih
    if source >= target:
        ph = h
        pw = h * source
    else:
        pw = w
        ph = w / source
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(pw), height=Inches(ph))
    pic.crop_left = max(0, (pw - w) / pw / 2)
    pic.crop_right = max(0, (pw - w) / pw / 2)
    pic.crop_top = max(0, (ph - h) / ph / 2)
    pic.crop_bottom = max(0, (ph - h) / ph / 2)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_screen(slide, "02-admin-dashboard.png", 7.0, 0.72, 8.15, 5.25)
    add_text(slide, "SAFE-LINK v2.0", 0.78, 0.78, 4.5, 0.42, 22, WHITE, True)
    add_text(slide, "관리자 교육용\n사용설명서", 0.78, 1.95, 5.9, 1.25, 42, WHITE, True)
    add_text(slide, "최신 업데이트 반영본 · 2026-06-01", 0.82, 3.55, 5.5, 0.35, 17, "D8DFEA", True)
    add_text(slide, "로그인부터 TBM, 근로자 관리, NFC/QR, 퀴즈, ESG 보고서, 실시간 통역, 법적 증빙까지 실제 관리자 운영 흐름 기준으로 정리했습니다.", 0.82, 4.25, 5.7, 0.8, 15, "D8DFEA", True)
    chip(slide, "시연 전 교육자료", 0.82, 5.65, 1.8, PINK)
    chip(slide, "실제 화면 기반", 2.82, 5.65, 1.8, BLUE)
    chip(slide, "2026.06 업데이트", 4.82, 5.65, 2.0, GREEN)


def agenda(prs):
    slide = base_slide(prs, "교육 목표와 운영 흐름", "처음 접속한 관리자가 하루 현장 운영을 끝낼 때까지 필요한 순서만 익힙니다.", 2)
    steps = [
        ("1", "로그인", "관리자 계정 접속과 권한 확인"),
        ("2", "현장 설정", "프로필, 현장 코드, 근로자 언어 확인"),
        ("3", "TBM 운영", "작성, 발송, 서명 확인, 미확인자 조치"),
        ("4", "현장 소통", "1:1 채팅과 실시간 통역"),
        ("5", "검증 매체", "NFC/QR 발급과 출입·교육 확인"),
        ("6", "이해도 확인", "퀴즈 생성, 발송, 응답 확인"),
        ("7", "보고서", "ESG/법적 증빙 보고서와 파일 내보내기"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.85 + (i % 4) * 3.7
        y = 1.75 + (i // 4) * 2.6
        add_rect(slide, x, y, 3.05, 1.55, WHITE, LINE)
        chip(slide, num, x + 0.25, y + 0.28, 0.52, [BLUE, PINK, AMBER, GREEN][i % 4])
        add_text(slide, head, x + 0.92, y + 0.25, 1.8, 0.35, 19, INK, True)
        add_text(slide, body, x + 0.28, y + 0.82, 2.45, 0.36, 12.5, MUTED, True)
    add_text(slide, "운영 원칙: SAFE-LINK는 “교육 안내”가 아니라 교육 확인, 현장 소통, 증빙 보존까지 이어지는 현장 안전관리 흐름입니다.", 1.0, 7.15, 13.4, 0.35, 16, BLUE, True, PP_ALIGN.CENTER)


def update_summary(prs):
    slide = base_slide(prs, "이번 최신본에 반영한 주요 업데이트", "지난 교육자료 이후 추가된 기능과 시연 전 안정화 내용을 관리자 관점으로 반영했습니다.", 3)
    items = [
        ("실시간 통역", "관리자 발화를 한국어 원문과 언어별 번역 payload로 저장하고, 근로자 화면에서 본인 언어로 확인"),
        ("TBM 기반 퀴즈", "TBM 공지와 실시간 발화 내용을 합쳐 퀴즈 생성, 컬럼 불일치 수정으로 생성 안정화"),
        ("보고서 증빙", "SHA-256 기반 보고서 envelope, 법적 고지·보존·개인정보 문구 정리"),
        ("검증 URL", "보고서/QR 검증 흐름과 외부 확인용 verify 화면 구조 반영"),
        ("DB 안정화", "live_translations.translations JSONB 컬럼 적용 및 실제 insert/delete smoke 통과"),
        ("시연 점검", "관리자 로그인 API, 세션 확인, 관리자 화면, 퀴즈 생성, 쿠키 제거 후 인증 실패 확인"),
    ]
    for i, (head, body) in enumerate(items):
        x = 0.9 + (i % 2) * 7.0
        y = 1.65 + (i // 2) * 1.75
        chip(slide, f"UPDATE {i+1}", x, y, 1.25, [BLUE, GREEN, AMBER, PINK, BLUE, GREEN][i])
        add_text(slide, head, x + 1.45, y - 0.02, 4.8, 0.32, 19, INK, True)
        add_text(slide, body, x, y + 0.55, 5.9, 0.55, 13.5, MUTED, True)


def screen_slide(prs, idx, file_name, title, subtitle, bullets, note=None):
    slide = base_slide(prs, title, subtitle, idx)
    add_screen(slide, file_name)
    add_rect(slide, 10.95, 1.72, 4.15, 5.72, WHITE, LINE)
    chip(slide, "따라 하기", 11.28, 2.08, 1.15, BLUE)
    bullet_list(slide, bullets, 11.28, 2.78, 3.2, 3.55, 15.5)
    if note:
        add_rect(slide, 11.25, 6.5, 3.3, 0.55, "EAF2FF", "B9D4FF")
        add_text(slide, note, 11.42, 6.67, 2.95, 0.2, 11.5, BLUE, True, PP_ALIGN.CENTER)


def worker_screen_slide(prs, idx, file_name, title, subtitle, bullets):
    slide = base_slide(prs, title, subtitle, idx)
    add_image_from_path(slide, WORKER_DIR / file_name, 0.88, 1.42, 5.25, 6.25)
    add_rect(slide, 6.65, 1.72, 8.15, 5.72, WHITE, LINE)
    chip(slide, "근로자 화면 교육 포인트", 7.0, 2.08, 2.0, GREEN)
    bullet_list(slide, bullets, 7.0, 2.82, 6.9, 3.7, 18)
    add_text(slide, "관리자 교육 시 근로자 화면을 함께 보여주면, 관리자가 본인 화면의 조치가 근로자에게 어떻게 보이는지 이해할 수 있습니다.", 7.0, 6.72, 6.8, 0.38, 14.5, BLUE, True)


def live_update_slide(prs):
    slide = base_slide(prs, "실시간 통역: 최신 저장 구조", "이번 업데이트에서 시연 리스크가 컸던 live_translations 저장 구조를 보강했습니다.", 18)
    add_screen(slide, "14-live.png", 0.62, 1.72, 8.35, 5.72)
    add_rect(slide, 9.35, 1.72, 5.7, 5.72, WHITE, LINE)
    chip(slide, "DB 구조", 9.72, 2.05, 1.05, GREEN)
    bullet_list(
        slide,
        [
            "text_ko: 관리자 한국어 발화 원문 저장",
            "translations: 언어 코드별 사전 번역 JSON 저장",
            "근로자 화면은 row.translations[내 언어]를 우선 표시",
            "저장 실패 시 관리자 화면에 [저장 실패] 메시지 표시",
            "시연 전 insert/delete smoke 테스트 통과",
        ],
        9.72,
        2.65,
        4.6,
        3.4,
        15,
    )
    add_text(slide, "교육 포인트: 라이브 통역은 마이크 권한과 네트워크가 필요합니다. 시연 전 브라우저 마이크 권한을 반드시 허용해 둡니다.", 9.72, 6.55, 4.8, 0.5, 13.5, RED, True)


def report_update_slide(prs):
    slide = base_slide(prs, "보고서와 법적 증빙: 최신 반영 사항", "관리자 교육에서 보고서는 단순 출력물이 아니라 사후 증빙자료라는 점을 강조합니다.", 21)
    items = [
        ("무결성", "보고서 생성 시 SHA-256 해시와 생성자·범위·원천 테이블 정보를 함께 보존"),
        ("법적 고지", "전자문서 효력, 보존 정책, 개인정보 최소 수집 문구를 보고서 envelope에 반영"),
        ("내보내기", "PDF, Excel, Word, HWP, JSON 등 운영 목적별 출력 형식 제공"),
        ("검증", "verify URL/QR 흐름으로 보고서 진위 확인 시나리오 설명 가능"),
    ]
    for i, (head, body) in enumerate(items):
        x = 1.0 + (i % 2) * 7.1
        y = 1.75 + (i // 2) * 2.1
        add_rect(slide, x, y, 5.9, 1.4, WHITE, LINE)
        chip(slide, head, x + 0.28, y + 0.28, 1.05, [BLUE, GREEN, AMBER, PINK][i])
        add_text(slide, body, x + 1.55, y + 0.24, 3.85, 0.55, 14.5, INK, True)
    add_rect(slide, 1.0, 6.35, 13.1, 0.75, "EAF2FF", "B9D4FF")
    add_text(slide, "교육 멘트: 사고가 난 뒤 자료를 찾는 시스템이 아니라, 매일의 TBM·서명·퀴즈·소통 이력이 자동으로 증빙화되는 구조입니다.", 1.28, 6.58, 12.2, 0.24, 15, BLUE, True, PP_ALIGN.CENTER)


def demo_checklist(prs):
    slide = base_slide(prs, "내일 시연 전 최종 체크리스트", "미래전략TF 미팅 전에는 기능 설명보다 실제 흐름 안정성이 더 중요합니다.", 23)
    checks = [
        ("로그인", "관리자 계정으로 로그인 후 /admin 진입 확인"),
        ("DB", "live_translations.translations = jsonb / NOT NULL / default '{}' 확인"),
        ("TBM", "공지 작성, 발송, 상태 화면에서 TBM 표시 확인"),
        ("라이브", "브라우저 마이크 권한 허용, 발화 저장 오류 없는지 확인"),
        ("퀴즈", "TBM 또는 fallback 퀴즈 생성 결과 확인"),
        ("보고서", "ESG/법적 보고서 내보내기 버튼과 생성 결과 확인"),
        ("근로자", "근로자 화면에서 TBM/퀴즈/라이브 수신 경로 확인"),
    ]
    for i, (head, body) in enumerate(checks):
        y = 1.45 + i * 0.82
        chip(slide, str(i + 1), 1.05, y, 0.45, [BLUE, GREEN, AMBER, PINK][i % 4])
        add_text(slide, head, 1.75, y - 0.02, 1.7, 0.25, 17, INK, True)
        add_text(slide, body, 3.7, y - 0.01, 9.8, 0.25, 15, MUTED, True)
    add_text(slide, "확인 완료 기준: build 성공 + smoke 통과 + 실제 브라우저에서 로그인/주요 화면/로그아웃 리허설 완료", 1.05, 7.45, 12.8, 0.35, 15.5, RED, True, PP_ALIGN.CENTER)


def closing(prs):
    slide = base_slide(prs, "관리자 교육 핵심 요약", "현장 관리자는 모든 기능을 외우는 것이 아니라, 하루 운영 흐름 안에서 필요한 화면을 정확히 쓰면 됩니다.", 24)
    add_text(slide, "1", 1.25, 1.85, 0.5, 0.45, 28, BLUE, True)
    add_text(slide, "TBM은 작성보다 확인이 중요합니다.", 1.9, 1.88, 10.5, 0.35, 24, INK, True)
    add_text(slide, "2", 1.25, 3.05, 0.5, 0.45, 28, GREEN, True)
    add_text(slide, "근로자별 언어·서명·퀴즈·소통 이력은 나중에 증빙이 됩니다.", 1.9, 3.08, 11.2, 0.35, 24, INK, True)
    add_text(slide, "3", 1.25, 4.25, 0.5, 0.45, 28, AMBER, True)
    add_text(slide, "보고서는 버튼 하나로 뽑되, 원천 데이터는 매일 쌓여야 합니다.", 1.9, 4.28, 11.2, 0.35, 24, INK, True)
    add_text(slide, "4", 1.25, 5.45, 0.5, 0.45, 28, PINK, True)
    add_text(slide, "시연 전에는 로그인, 마이크, DB 컬럼, 퀴즈 생성, 로그아웃을 반드시 확인합니다.", 1.9, 5.48, 11.4, 0.35, 24, INK, True)


def normalize(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Malgun Gothic"
                    run.font.bold = True


def build_deck():
    OUT.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    cover(prs)
    agenda(prs)
    update_summary(prs)
    screen_slide(prs, 4, "01-auth-admin-login.png", "1. 관리자 로그인", "언어 선택 후 관리자 역할로 접속합니다.", ["한국어 선택 후 관리자 역할을 선택합니다.", "이메일과 비밀번호를 입력합니다.", "로그인 후 관리자 대시보드로 이동되는지 확인합니다."])
    screen_slide(prs, 5, "03-profile-setup.png", "2. 최초 프로필 설정", "이름, 권한, 현장 코드를 정확히 맞춥니다.", ["현장 코드와 담당 현장이 맞는지 확인합니다.", "관리자 표시 이름은 보고서와 채팅에 사용됩니다.", "권한이 맞지 않으면 관리자 메뉴 접근이 제한됩니다."])
    screen_slide(prs, 6, "02-admin-dashboard.png", "3. 관리자 대시보드", "하루 운영은 대시보드에서 시작합니다.", ["DB, 번역, 음성 기능 상태를 확인합니다.", "TBM, 근로자, NFC/QR, 퀴즈, ESG 메뉴로 이동합니다.", "시연 전에는 주요 카드가 정상 표시되는지 확인합니다."])
    screen_slide(prs, 7, "04-tbm-create.png", "4. TBM 작성", "작업 내용과 위험요인을 입력하고 AI 가이드로 문장을 보강합니다.", ["작업 위치, 작업조, 위험요인을 구체적으로 입력합니다.", "안전고리, 낙하물, 이동 동선 등 현장 키워드를 포함합니다.", "발송 전 문장이 현장 근로자에게 이해 가능한지 확인합니다."])
    screen_slide(prs, 8, "04a-tbm-filled.png", "5. TBM 입력 예시", "실제 교육용 더미 데이터 기준의 작성 예시입니다.", ["오늘 작업과 위험요인을 한 화면에서 확인합니다.", "발송 전 오탈자와 현장명, 대상 근로자를 확인합니다.", "TBM은 사후 증빙이 되므로 간결하지만 구체적으로 작성합니다."])
    screen_slide(prs, 9, "05-tbm-status.png", "6. TBM 서명 현황", "발송 후에는 서명률과 미확인자를 추적합니다.", ["서명 완료자와 미확인자를 구분합니다.", "미확인자는 작업 전 재안내합니다.", "서명 시간은 보고서 증빙 자료로 활용됩니다."])
    screen_slide(prs, 10, "07-workers.png", "7. 근로자 관리", "근로자 명단, 언어, 현장 배정을 관리합니다.", ["근로자별 선호 언어를 확인합니다.", "현장 이동 또는 퇴사자는 즉시 정리합니다.", "연락처와 현장 배정 오류는 TBM/퀴즈 수신 문제로 이어집니다."])
    screen_slide(prs, 11, "08-workers-enroll.png", "8. 근로자 등록", "신규 근로자의 기본 정보와 언어를 등록합니다.", ["이름, 연락처, 국적, 언어를 정확히 입력합니다.", "현장 코드를 확인합니다.", "등록 후 NFC/QR 또는 로그인 경로를 안내합니다."])
    screen_slide(prs, 12, "09-nfc.png", "9. NFC 관리", "NFC 스티커와 현장 태깅 기록을 관리합니다.", ["근로자별 NFC 매체를 발급합니다.", "분실 또는 재발급 시 기존 매체를 폐기 처리합니다.", "출입, TBM 확인, 교육 참여 기록과 연결됩니다."])
    screen_slide(prs, 13, "10-qr-code.png", "10. QR 코드", "NFC가 어려운 상황에서는 QR/서명 URL을 fallback으로 사용합니다.", ["QR 목적과 대상 현장을 확인합니다.", "현장 게시 또는 근로자 안내에 사용합니다.", "검증 URL과 연결되는 경우 외부 확인 흐름을 설명합니다."])
    screen_slide(prs, 14, "06-chat.png", "11. 1:1 다국어 채팅", "외국인 근로자와 짧고 명확한 문장으로 소통합니다.", ["대상 근로자를 선택합니다.", "한 번에 하나의 지시만 보냅니다.", "위험·작업중지 표현은 별도 확인과 조치가 필요합니다."])
    screen_slide(prs, 15, "11-quiz.png", "12. 안전 퀴즈", "TBM 이해도를 확인하고 보충교육 대상을 찾습니다.", ["TBM 내용 기반으로 퀴즈를 생성합니다.", "응답률과 점수를 확인합니다.", "미응답자 또는 낮은 점수자는 재안내합니다."], "최신: text_ko 기반 생성 안정화")
    screen_slide(prs, 16, "12-esg.png", "13. ESG 안전 리포트", "교육, 서명, 퀴즈, 현장 기록을 보고서로 정리합니다.", ["기간과 현장을 선택합니다.", "TBM/서명/퀴즈/안전활동 지표를 확인합니다.", "필요한 형식으로 내보냅니다."])
    screen_slide(prs, 17, "13-glossary.png", "14. 현장 용어집", "건설 현장 표현과 표준 용어를 관리해 번역 품질을 높입니다.", ["현장 은어와 표준 표현을 등록합니다.", "위험 작업 키워드를 정리합니다.", "TBM, 채팅, 실시간 통역 품질 개선에 사용됩니다."])
    live_update_slide(prs)
    screen_slide(prs, 19, "14-live.png", "15. 실시간 통역 운영", "관리자 발화를 근로자 언어로 전달합니다.", ["마이크 권한을 허용합니다.", "짧은 문장으로 천천히 말합니다.", "저장 실패 메시지가 뜨면 네트워크와 DB 상태를 확인합니다."], "최신: translations JSONB 적용")
    worker_screen_slide(prs, 20, "worker_01_login_entry.png", "근로자 화면 1. 접속과 본인 확인", "QR/NFC 또는 전화번호로 현장 교육에 들어오는 화면입니다.", ["근로자는 본인 언어를 선택합니다.", "NFC/QR이 우선이고, 필요 시 전화번호 기반 fallback을 사용합니다.", "현장 배정이 맞아야 TBM과 퀴즈를 받을 수 있습니다."])
    worker_screen_slide(prs, 21, "worker_02_tbm_confirm.png", "근로자 화면 2. TBM 확인과 서명", "관리자가 발송한 TBM은 근로자 언어로 표시되고 서명 증빙으로 남습니다.", ["작업 내용, 위험요인, 보호구 항목을 읽습니다.", "확인 후 전자서명을 남깁니다.", "서명 시각과 세션 정보는 보고서 증빙에 사용됩니다."])
    worker_screen_slide(prs, 22, "worker_03_quiz.png", "근로자 화면 3. 안전 퀴즈", "TBM 이해도를 확인하고 재교육 대상을 구분합니다.", ["문항은 TBM 내용 기반으로 생성됩니다.", "근로자는 객관식 또는 OX 문항에 응답합니다.", "점수와 응답 여부는 관리자 퀴즈 화면과 보고서에 연결됩니다."])
    worker_screen_slide(prs, 23, "worker_04_chat_live.png", "근로자 화면 4. 채팅과 실시간 통역", "관리자와의 1:1 대화 및 라이브 방송을 본인 언어로 확인합니다.", ["짧은 문장 단위로 번역 품질을 높입니다.", "관리자 발화는 text_ko와 translations JSON으로 저장됩니다.", "위험 표현은 향후 C15 위험신호 라우팅과 연결됩니다."])
    worker_screen_slide(prs, 24, "worker_05_stop_work.png", "근로자 화면 5. 작업중지와 안전신고", "위험을 발견하면 근로자 언어로 신고하고 조치 이력을 남깁니다.", ["위험 위치와 상황을 입력합니다.", "관리자에게 라우팅되고 조치 상태를 추적합니다.", "작업중지 기록은 감사 이벤트와 보고서 증빙으로 활용됩니다."])
    report_update_slide(prs)
    demo_checklist(prs)
    closing(prs)

    normalize(prs)
    prs.save(PPTX)
    prs.save(PPTX_ASCII)


def font(size, bold=False):
    path = Path(r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf")
    return ImageFont.truetype(str(path), size=size) if path.exists() else ImageFont.load_default()


def make_previews():
    slides = [
        ("표지", "SAFE-LINK v2.0 관리자 교육용 사용설명서"),
        ("흐름", "교육 목표와 운영 흐름"),
        ("업데이트", "이번 최신본에 반영한 주요 업데이트"),
        ("로그인", "관리자 로그인"),
        ("프로필", "최초 프로필 설정"),
        ("대시보드", "관리자 대시보드"),
        ("TBM", "TBM 작성"),
        ("TBM 예시", "TBM 입력 예시"),
        ("서명", "TBM 서명 현황"),
        ("근로자", "근로자 관리"),
        ("등록", "근로자 등록"),
        ("NFC", "NFC 관리"),
        ("QR", "QR 코드"),
        ("채팅", "1:1 다국어 채팅"),
        ("퀴즈", "안전 퀴즈"),
        ("ESG", "ESG 안전 리포트"),
        ("용어집", "현장 용어집"),
        ("라이브 구조", "실시간 통역 저장 구조"),
        ("라이브 운영", "실시간 통역 운영"),
        ("근로자", "근로자 화면 1. 접속과 본인 확인"),
        ("근로자", "근로자 화면 2. TBM 확인과 서명"),
        ("근로자", "근로자 화면 3. 안전 퀴즈"),
        ("근로자", "근로자 화면 4. 채팅과 실시간 통역"),
        ("근로자", "근로자 화면 5. 작업중지와 안전신고"),
        ("보고서", "보고서와 법적 증빙"),
        ("체크리스트", "내일 시연 전 최종 체크리스트"),
        ("요약", "관리자 교육 핵심 요약"),
    ]
    thumbs = []
    for idx, (label, title) in enumerate(slides, 1):
        img = Image.new("RGB", (640, 360), "#" + PAPER)
        d = ImageDraw.Draw(img)
        d.text((28, 24), f"{idx:02d}", font=font(20, True), fill="#" + BLUE)
        d.text((88, 24), title, font=font(23, True), fill="#" + INK)
        d.text((88, 62), label, font=font(15, True), fill="#" + MUTED)
        screen_map = {
            4: "01-auth-admin-login.png", 5: "03-profile-setup.png", 6: "02-admin-dashboard.png",
            7: "04-tbm-create.png", 8: "04a-tbm-filled.png", 9: "05-tbm-status.png",
            10: "07-workers.png", 11: "08-workers-enroll.png", 12: "09-nfc.png",
            13: "10-qr-code.png", 14: "06-chat.png", 15: "11-quiz.png",
            16: "12-esg.png", 17: "13-glossary.png", 18: "14-live.png", 19: "14-live.png",
        }
        worker_map = {
            20: "worker_01_login_entry.png",
            21: "worker_02_tbm_confirm.png",
            22: "worker_03_quiz.png",
            23: "worker_04_chat_live.png",
            24: "worker_05_stop_work.png",
        }
        screen_path = None
        if idx in screen_map and (SCREEN_DIR / screen_map[idx]).exists():
            screen_path = SCREEN_DIR / screen_map[idx]
        elif idx in worker_map and (WORKER_DIR / worker_map[idx]).exists():
            screen_path = WORKER_DIR / worker_map[idx]

        if screen_path:
            screen = Image.open(screen_path).convert("RGB")
            screen.thumbnail((390, 235))
            img.paste(screen, (34, 105))
            d.rounded_rectangle((455, 110, 604, 280), radius=12, fill="#FFFFFF", outline="#" + LINE, width=2)
            d.text((474, 140), "화면 이미지", font=font(18, True), fill="#" + BLUE)
            d.text((474, 178), "따라 하기", font=font(16, True), fill="#" + INK)
        else:
            d.rounded_rectangle((80, 120, 560, 275), radius=14, fill="#FFFFFF", outline="#" + LINE, width=2)
            d.text((115, 170), "최신 업데이트 반영", font=font(24, True), fill="#" + BLUE)
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        thumbs.append(img.resize((256, 144)))

    rows = (len(thumbs) + 3) // 4
    sheet = Image.new("RGB", (256 * 4, 144 * rows), "#FFFFFF")
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % 4) * 256, (idx // 4) * 144))
    sheet.save(CONTACT)


if __name__ == "__main__":
    build_deck()
    make_previews()
    print(f"pptx={PPTX}")
    print(f"pptx_ascii={PPTX_ASCII}")
    print(f"preview_contact={CONTACT}")
