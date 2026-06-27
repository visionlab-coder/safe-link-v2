from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-magbrake-fall-arrest-20260607"
ASSET = BASE / "assets"
OUT = BASE / "SAFE-LINK_MagBrake_Fall_Arrest_대우현대국가과제_제안서_V3_용어주석_20260607.pptx"
PREVIEW = BASE / "previews-magbrake-ppt-v3"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8", "ink": "101828", "muted": "667085", "line": "D0D5DD",
    "blue": "175CD3", "green": "027A48", "orange": "B54708", "red": "B42318",
    "white": "FFFFFF", "soft": "F2F4F7", "navy": "0B1220", "copper": "B86B2B",
    "yellow": "F2C94C"
}

def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))

def rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    if radius:
        try: shp.adjustments[0] = 0.05
        except Exception: pass
    return shp

def txt(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = body; r.font.name = "맑은 고딕"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box

def rule(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color); shp.line.fill.background()

def footer(slide, n):
    txt(slide, "SAFE-LINK MagBrake Fall Arrest / 서원토건·서원이앤에이", .55, 7.04, 5.7, .24, 9.5, "muted")
    txt(slide, f"{n:02d}", 12.25, 7.04, .45, .24, 10, "muted", align="right")

def footnote(slide, body):
    rect(slide, .62, 6.52, 11.6, .4, "soft", "line", True)
    txt(slide, "용어 주석: " + body, .78, 6.61, 11.25, .22, 9.2, "muted")

def header(slide, label, title, sub, n):
    txt(slide, label, .62, .34, 6.0, .3, 12, "blue", True)
    txt(slide, title, .62, .74, 11.8, .72, 27, "ink", True)
    if sub: txt(slide, sub, .64, 1.5, 11.2, .4, 14.5, "muted")
    rule(slide, .62, 2.02, 12.05)
    footer(slide, n)

def img(slide, name, x, y, w, h):
    p = ASSET / name
    if p.exists(): slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        txt(slide, name, x, y+h/2-.15, w, .3, 13, "muted", True, "center")

def note(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True); rect(slide, x, y, .08, h, color, color, False)
    txt(slide, head, x+.22, y+.15, w-.42, .32, 14.5, color, True)
    txt(slide, body, x+.22, y+.55, w-.42, h-.66, 13.2, "ink")

def table(slide, x, y, rows, widths, row_h=.48, head="navy"):
    for r,row in enumerate(rows):
        cx=x
        for c,value in enumerate(row):
            fill=head if r==0 else ("white" if r%2 else "soft")
            color="white" if r==0 else "ink"
            rect(slide,cx,y+r*row_h,widths[c],row_h,fill,"line",False)
            txt(slide,value,cx+.08,y+r*row_h+.085,widths[c]-.16,row_h-.12,10.7 if r==0 else 11.8,color,r==0)
            cx += widths[c]

def build():
    prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H); blank = prs.slide_layouts[6]
    def slide():
        s=prs.slides.add_slide(blank); rect(s,0,0,W,H,"paper","paper",False); return s

    # 1
    s=slide(); img(s,"01-site-fall-hazard-srl.png",6.45,0,6.88,7.5)
    txt(s,"DAEWOO · HYUNDAI E&C · NATIONAL R&D",.72,.64,5,.28,11.5,"blue",True)
    txt(s,"SAFE-LINK\nMagBrake\nFall Arrest",.72,1.16,5.4,1.82,34,"ink",True)
    txt(s,"와전류 자기제동 기반 스마트 추락방지 및 구조알림 시스템",.74,3.35,5.25,.62,20,"ink",True)
    rule(s,.74,4.18,4.3,"blue")
    txt(s,"기존 안전대·SRL의 한계를 보완하고, 추락 충격 완화·즉시 구조·보고 자동화를 결합하는 현실적 R&D 제안",.74,4.52,5.25,.8,15,"muted")
    txt(s,"서원토건 / 서원이앤에이 / SAFE-LINK",.74,6.78,3.9,.25,11,"muted")
    footnote(s, "SRL(Self-Retracting Lifeline)=자동으로 줄이 감기고 추락 시 잠기는 안전블록/자기수축식 생명줄")

    # 2
    s=slide(); header(s,"01 PROBLEM","추락방지 장비는 있어도 사고 후 충격과 구조 지연은 남습니다","안전대와 SRL은 필수 장비지만, 충격 완화·위치 식별·구조 워크플로우는 별도 문제입니다.",2)
    img(s,"01-site-fall-hazard-srl.png",.75,2.35,5.75,3.55)
    note(s,"충격", "추락 시 큰 체포하중이 신체와 앵커에 전달될 수 있음",7.0,2.35,4.85,.88,"red")
    note(s,"매달림", "구조가 늦어지면 현수 상태 자체가 2차 위험이 됨",7.0,3.45,4.85,.88,"orange")
    note(s,"위치", "층, 구역, 앵커, 작업자 식별이 지연될 수 있음",7.0,4.55,4.85,.88,"blue")
    note(s,"보고", "조치 이력과 재발방지 보고가 수작업으로 남음",7.0,5.65,4.85,.88,"green")
    footnote(s, "체포하중=추락이 멈추는 순간 몸과 장비에 걸리는 충격 하중 / 앵커=안전줄을 고정하는 구조물")

    # 3 Existing devices comparison
    s=slide(); header(s,"02 EXISTING DEVICES","기존 장비와 경쟁하지 않고, 빈틈을 메웁니다","안전대·SRL·풍선형 조끼는 각각 역할이 다르고, MagBrake는 감속·감지·구조알림을 보강합니다.",3)
    rows=[["구분","강점","한계","MagBrake 차별점"],
          ["안전대/안전벨트","기본 추락방지 PPE","충격·위치·구조알림 없음","센서·위치·알림 추가"],
          ["SRL/안전블록","추락 시 라인 잠김","급정지 충격·상태로그 부족","와전류 감속 보조"],
          ["충격흡수 랜야드","웨빙 찢김으로 충격 흡수","1회성·작동 후 교체","반복 시험 가능한 감속데이터"],
          ["풍선형 조끼","충돌/낙상 완충 가능성","고소 추락 정지장치 아님","SRL 계통과 직접 연동"],
          ["MagBrake","감속+감지+구조알림","인증 전 R&D 필요","기존 장비 고도화"]]
    table(s,.7,2.35,rows,[2.0,2.7,3.05,3.65],.53)
    footnote(s, "랜야드=안전대와 고정점을 연결하는 줄 / 충격흡수 랜야드=추락 시 일부가 찢어지며 충격을 줄이는 연결줄")

    # 4 Advantages
    s=slide(); header(s,"03 BEST ADVANTAGE","최고 장점은 생명을 구할 가능성을 높이면서 관리비용을 낮추는 것입니다","장비 하나가 아니라 사고 대응 시간을 줄이는 안전관리 시스템으로 봐야 합니다.",4)
    note(s,"충격 완화", "와전류 감속과 기존 충격흡수장치를 병행해 급정지 충격 저감을 목표",.85,2.55,3.65,1.45,"blue")
    note(s,"즉시 구조", "추락 이벤트, 작업자, 위치, 생체/움직임 신호를 즉시 전송",4.85,2.55,3.65,1.45,"red")
    note(s,"비용 절감", "구조 지연, 보고 누락, 사고조사 시간, 재발방지 문서 작성 부담 감소",8.85,2.55,3.65,1.45,"green")
    txt(s,"한 명의 생명, 한 건의 중대재해, 한 번의 구조 지연을 줄이는 것만으로도 PoC 가치는 충분합니다.",1.0,5.15,11.2,.55,24,"ink",True,"center")
    footnote(s, "PoC(Proof of Concept)=정식 제품화 전 현장에서 가능성과 효과를 검증하는 실증")

    # 5 Principle deep
    s=slide(); header(s,"04 PRINCIPLE","작동원리: 낙하 에너지를 회전으로 바꾸고, 와전류로 감속합니다","자석이 사람을 들어 올리는 것이 아니라, 회전 드럼의 속도를 비접촉 방식으로 낮추는 원리입니다.",5)
    img(s,"02-magbrake-device-cutaway.png",.75,2.25,5.75,3.55)
    steps=[("1 낙하", "라이프라인 급출"),("2 회전", "드럼/디스크 고속회전"),("3 유도전류", "도체에 와전류 발생"),("4 반대자기장", "렌츠 법칙에 따라 운동 저항"),("5 열·센서", "에너지 일부 열로 전환·로그 저장")]
    for i,(a,b) in enumerate(steps):
        y=2.42+i*.72
        rect(s,6.95,y,4.9,.52,"white",["red","orange","copper","blue","green"][i],True)
        txt(s,a,7.12,y+.1,1.2,.22,11.5,["red","orange","copper","blue","green"][i],True)
        txt(s,b,8.3,y+.1,3.25,.22,11.5,"ink",True)
    txt(s,"빠르게 풀릴수록 제동 저항이 커지는 특성을 이용하되, 최종 정지는 기존 SRL 잠김 구조와 함께 수행합니다.",6.95,6.15,4.9,.42,13.5,"blue",True,"center")
    footnote(s, "와전류=움직이는 자석과 도체 사이에서 생기는 유도전류 / 렌츠 법칙=유도전류가 원래 움직임을 방해하는 방향으로 작용한다는 법칙")

    # 6 Device structure
    s=slide(); header(s,"05 DEVICE STRUCTURE","MagBrake는 SRL 내부 또는 외부 보조 모듈로 설계합니다","1차 개발은 실제 인증제품 개조가 아니라 실험용 인터페이스와 벤치모듈로 시작합니다.",6)
    img(s,"02-magbrake-device-cutaway.png",.75,2.3,5.85,3.55)
    rows=[["구성","역할"],["영구자석 배열","도체 디스크에 자기장 제공"],["전도체 디스크/드럼","회전하며 와전류 발생"],["엔코더","라인 풀림 속도와 회전수 측정"],["온도센서","발열·반복사용 위험 감시"],["하중센서","체포하중과 앵커 상태 기록"],["MCU/통신","SAFE-LINK 이벤트 전송"]]
    table(s,6.9,2.35,rows,[2.2,3.55],.48)
    footnote(s, "엔코더=회전속도·회전량 측정 센서 / MCU=소형 제어칩 / 도체=전기가 흐르는 금속, 예: 구리·알루미늄")

    # 7 Cost model
    s=slide(); header(s,"06 COST MODEL","목표는 고가 장비 대체가 아니라 ‘추가 안전가치 대비 낮은 도입비’입니다","가격은 인증·양산 전 가정치이며, PoC에서는 개발비와 장비비를 분리합니다.",7)
    rows=[["항목","예상 범위","설명"],["기존 안전대","수만~수십만 원","작업자 기본 PPE"],["일반 SRL/안전블록","수십만~100만 원+","길이·브랜드·인증별 편차"],["MagBrake 보조모듈 목표","50만~150만 원/대","양산 시 SRL 추가 모듈 목표"],["스마트 센서 태그","5만~20만 원/개","IMU/BLE/UWB 구성별 차이"],["SAFE-LINK 관제 연동","현장 단위 구독/구축","알림·보고·데이터 저장"],["1차 PoC","3천만~7천만 원","시제품, 더미시험, 관제 개발 포함"]]
    table(s,.85,2.35,rows,[2.7,2.6,5.7],.55)
    txt(s,"비용 논리는 ‘장비 단가’가 아니라 구조시간 단축, 사고조사 비용 감소, 보고 자동화, 중대재해 리스크 저감입니다.",1.0,6.35,11.2,.34,15.5,"green",True,"center")
    footnote(s, "도입비용은 인증·양산 전 추정치이며, 실제 비용은 제조사 견적·시험범위·현장 구축범위에 따라 달라짐")

    # 8 System architecture
    s=slide(); header(s,"07 SYSTEM","감속장치만으로는 부족합니다. 센서·위치·관제가 붙어야 제품 가치가 생깁니다","기계장치와 SAFE-LINK 플랫폼을 결합해 구조 워크플로우까지 닫습니다.",8)
    nodes=[("Harness","IMU 추락감지","blue"),("SRL","잠김 구조","green"),("MagBrake","와전류 감속","orange"),("Location","UWB/BLE/NFC","blue"),("SAFE-LINK","알림·조치보고","red")]
    for i,(a,b,c) in enumerate(nodes):
        x=.75+i*2.45
        rect(s,x,3.0,1.85,1.25,"white",c,True)
        txt(s,a,x+.12,3.18,1.6,.28,14,c,True,"center"); txt(s,b,x+.12,3.64,1.6,.28,11.5,"ink",True,"center")
        if i<4: txt(s,"→",x+1.86,3.4,.35,.3,20,"muted",True,"center")
    txt(s,"추락 이벤트 → 위치/작업자 식별 → 긴급 알림 → 구조 상태 관리 → 초기 보고서 자동 생성",1.0,5.4,11.2,.45,20,"ink",True,"center")
    footnote(s, "IMU=가속도·자이로 센서로 움직임을 감지하는 장치 / UWB·BLE·NFC=근거리 위치·식별 통신 기술")

    # 9 SAFE-LINK dashboard using user image
    s=slide(); header(s,"08 SAFE-LINK CONTROL","관제센터는 추락 감지 후 즉시 구조 체크리스트를 실행합니다","사용자가 제공한 관제 화면 콘셉트를 반영해 위치·상태·구조 연락망·보고서 생성을 연결합니다.",9)
    img(s,"06-user-safelink-fall-control-dashboard.png",.75,2.25,7.0,3.95)
    note(s,"이벤트 상세","작업자, 발생 위치, 심각도, 생존신호 확인",8.05,2.35,3.8,.9,"red")
    note(s,"조치 체크리스트","현장확인, 근로자 상태, 구조팀 호출, 위험통제",8.05,3.45,3.8,.9,"orange")
    note(s,"보고 자동화","이벤트 보고서, PDF, 이메일 발송까지 연결",8.05,4.55,3.8,.9,"blue")
    footnote(s, "UWB=정밀 위치 측정용 초광대역 통신 / BLE=저전력 블루투스 / NFC=태그 접촉식 식별 기술")

    # 10 Test
    s=slide(); header(s,"09 TEST STRATEGY","성능 주장은 더미 낙하시험 데이터로만 말해야 합니다","사람 대상이 아니라 40/60/80kg 더미, 벤치시험, 반복시험으로 검증합니다.",10)
    img(s,"03-dummy-drop-test-lab.png",.75,2.35,5.85,3.55)
    rows=[["시험","측정 데이터"],["벤치 회전시험","회전속도, 제동력, 발열"],["단계하중","감속거리, 온도, 반복성"],["더미낙하","최대하중, 정지시간, 충격량"],["SRL 비교","SRL 단독 vs MagBrake 보조"],["반복시험","30회 이상 데이터 축적"]]
    table(s,6.85,2.42,rows,[2.45,3.35],.53)
    footnote(s, "더미 낙하시험=사람 대신 시험용 인체모형을 떨어뜨려 하중·정지거리·충격량을 측정하는 시험")

    # 11 Standards
    s=slide(); header(s,"10 CERTIFICATION","안전장비는 인증을 우회할 수 없습니다","대외 제안에서는 ‘현장 즉시 사용’이 아니라 ‘시험기관 검증을 전제로 한 R&D’로 표현합니다.",11)
    note(s,"기준","개인추락방지시스템, SRL, 충격흡수, 앵커, 최대 체포하중 기준 검토",.9,2.45,3.55,1.35,"blue")
    note(s,"시험기관","KOSHA/KCL/KTR 등 시험·인증 가능성 사전 협의",4.9,2.45,3.55,1.35,"green")
    note(s,"위험관리","발열, 자석 이탈, 과하중, 반복사용, 센서 오탐을 시험 항목화",8.9,2.45,3.55,1.35,"orange")
    rows=[["표현 원칙","제안서 문구"],["대체품 아님","기존 안전대·SRL 보조 시스템"],["현장 즉시사용 아님","더미시험·시험기관 검증 후 단계 적용"],["성능 과장 금지","감속 보조·충격완화·구조알림 중심"],["SAFE-LINK 가치","위치·알림·조치·보고 자동화"]]
    table(s,1.45,4.35,rows,[3.4,6.8],.45)
    footnote(s, "KOSHA=한국산업안전보건공단 / KCL·KTR=국내 시험·인증기관 / 인증 전에는 현장 실사용 장비로 표현하면 안 됨")

    # 12 Roadmap
    s=slide(); header(s,"11 ROADMAP","18개월 개발, 3~6개월 내 보여줄 수 있는 PoC", "장기 인증 로드맵과 단기 데모를 분리해야 건설사 제안이 현실적입니다.",12)
    rows=[["단계","기간","목표","산출물"],["0","0~1개월","기술정의·IP·협력구조","개념서/NDA/역할표"],["1","1~3개월","센서형 하네스 PoC","IMU 감지·SAFE-LINK 알림"],["2","3~6개월","와전류 벤치 실험","감속모듈·데이터"],["3","6~9개월","더미 낙하 테스트","80kg 더미·하중데이터"],["4","9~12개월","SRL 통합 시제품","통합 구조·반복시험"],["5","12~18개월","현장 비작업 테스트","시험기관·인증 검토"]]
    table(s,.75,2.35,rows,[1.0,1.45,3.75,4.75],.49)
    footnote(s, "벤치 실험=사람·현장 투입 전 실험대에서 부품 성능을 측정하는 시험 / SRL 통합=안전블록 구조와 결합 검증")

    # 13 Builder/National
    s=slide(); header(s,"12 DEPLOYMENT","대우건설·현대엔지니어링·국가과제는 서로 다른 언어로 제안합니다","같은 기술이라도 건설사는 PoC, 국책과제는 R&D 검증성과를 봅니다.",13)
    img(s,"05-rnd-roadmap-collaboration.png",.75,2.35,5.85,3.55)
    rows=[["대상","제안 포인트"],["대우건설","Hyper Safety 공모전, 골조현장 PoC, SAFE-LINK 보고 자동화"],["현대엔지니어링","스마트건설 안전기술, 교육장/모형현장 실증"],["국가과제","스마트 PPE, 중대재해 예방, 더미시험 데이터"],["산학협력","와전류 해석, 센서 알고리즘, 시험데이터 분석"]]
    table(s,6.85,2.55,rows,[2.0,3.8],.58)
    footnote(s, "스마트 PPE=센서·통신 기능이 들어간 개인보호구 / R&D=연구개발 / SaaS=구독형 소프트웨어 서비스")

    # 14 IP/Business
    s=slide(); header(s,"13 IP & BUSINESS","특허는 장치 단독보다 ‘장치+센서+구조알림+보고’ 결합으로 잡습니다","비즈니스는 완제품 판매, 보조모듈, SAFE-LINK 관제 구독으로 나눌 수 있습니다.",14)
    rows=[["IP 축","내용"],["기계 구조","SRL 드럼/디스크 결합 와전류 감속 모듈"],["센서 이벤트","IMU, 하중, 온도, 앵커 상태 기반 추락 판정"],["위치 식별","UWB/BLE/NFC 작업자·구역·앵커 매칭"],["구조 알림","안전관리자·반장·본사 알림 및 상태 추적"],["보고 자동화","초기보고서·재발방지 보고서 생성"]]
    table(s,.9,2.35,rows,[2.35,4.6],.54)
    rows2=[["사업모델","설명"],["시제품 공동개발","건설사/제조사/연구기관 과제"],["보조모듈 판매","기존 SRL 고도화 옵션"],["SAFE-LINK 구독","관제·보고·데이터 저장"],["시험데이터/IP","국가과제·특허·인증 자산"]]
    table(s,7.35,2.35,rows2,[2.0,3.1],.54)
    footnote(s, "IP=지식재산권, 특허·디자인·영업비밀 등 / 청구항=특허에서 보호받고자 하는 기술 범위")

    # 15 Collaboration
    s=slide(); header(s,"14 COLLABORATION","협업 요청은 역할을 명확히 나눠야 실행됩니다","건설사, 안전장비 제조사, 연구기관, 시험기관, SAFE-LINK의 역할을 분리합니다.",15)
    rows=[["주체","역할"],["건설사","PoC 현장/교육장 제공, 안전관리자 피드백"],["서원토건","골조공정 위험 시나리오, 전문건설업 현장성"],["서원이앤에이/SAFE-LINK","알림·조치·보고 플랫폼, 대시보드"],["SRL/안전장비 협력사","기구 설계, 제품화, 인증 대응"],["대학/연구기관","와전류 해석, 센서 알고리즘, 시험 데이터 분석"],["시험기관","더미 낙하시험, 기준 검토, 인증 로드맵"]]
    table(s,.9,2.35,rows,[2.6,8.5],.53)
    footnote(s, "시험기관 협의는 성능 보증이 아니라 향후 인증 가능성을 확인하기 위한 사전 검토 단계")

    # 16 Closing
    s=slide(); header(s,"15 CLOSING","MagBrake는 추락을 ‘없앤다’가 아니라 구조 가능한 시간과 데이터를 만듭니다","사람의 목숨을 구할 가능성을 높이고, 사고 대응·보고·재발방지 비용을 낮추는 현실적 R&D입니다.",16)
    note(s,"1단계","센서형 하네스 + SAFE-LINK 긴급알림\n1~3개월 내 데모 가능",.9,2.55,3.55,1.55,"blue")
    note(s,"2단계","와전류 자기제동 벤치 실험\n3~6개월 내 데이터 확보",4.9,2.55,3.55,1.55,"orange")
    note(s,"3단계","더미 낙하·SRL 통합·시험기관 검토\n6~18개월 R&D",8.9,2.55,3.55,1.55,"green")
    txt(s,"SAFE-LINK MagBrake Fall Arrest System",1.0,5.15,11.2,.42,28,"ink",True,"center")
    txt(s,"충격을 줄이고, 위치를 알리고, 구조를 빠르게 하고, 보고를 자동화하는 스마트 추락방지 시스템",1.3,5.78,10.6,.42,18,"muted",True,"center")
    footnote(s, "본 제안은 기존 법정 안전장비를 대체하는 것이 아니라 시험·인증을 거쳐 보조 안전기능을 추가하는 R&D 과제")

    BASE.mkdir(parents=True, exist_ok=True); prs.save(OUT)

def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles=["Cover","Problem","Existing Devices","Best Advantage","Principle","Device","Cost Model","System","Control","Test","Certification","Roadmap","Deployment","IP/Business","Collaboration","Closing"]
    try:
        f1=ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf",31); f2=ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf",17)
    except Exception: f1=f2=ImageFont.load_default()
    image_map={1:"01-site-fall-hazard-srl.png",2:"01-site-fall-hazard-srl.png",5:"02-magbrake-device-cutaway.png",6:"02-magbrake-device-cutaway.png",9:"06-user-safelink-fall-control-dashboard.png",10:"03-dummy-drop-test-lab.png",13:"05-rnd-roadmap-collaboration.png"}
    paths=[]
    for i,t in enumerate(titles,1):
        canvas=Image.new("RGB",(1600,900),"#"+C["paper"]); d=ImageDraw.Draw(canvas)
        d.text((80,60),"SAFE-LINK MagBrake Fall Arrest",font=f2,fill="#"+C["blue"]); d.text((80,128),t,font=f1,fill="#"+C["ink"])
        d.line((80,225,1460,225),fill="#"+C["line"],width=3)
        if i in image_map and (ASSET/image_map[i]).exists():
            im=Image.open(ASSET/image_map[i]).convert("RGB"); im.thumbnail((760,430)); canvas.paste(im,(80,285))
        for k,col in enumerate(["blue","green","orange","red"]):
            d.rounded_rectangle((980,310+k*105,1390,370+k*105),16,outline="#"+C[col],width=3,fill="white")
        d.text((80,830),f"{i:02d}",font=f2,fill="#"+C["muted"])
        p=PREVIEW/f"slide_{i:02d}.png"; canvas.save(p); paths.append(p)
    sheet=Image.new("RGB",(1600,1260),"white"); d=ImageDraw.Draw(sheet)
    for i,p in enumerate(paths):
        im=Image.open(p).resize((300,169)); x=25+(i%5)*315; y=25+(i//5)*300
        sheet.paste(im,(x,y)); d.text((x,y+176),f"{i+1:02d} {titles[i]}",font=f2,fill="#"+C["ink"])
    sheet.save(PREVIEW/"contact_sheet.png")

if __name__=="__main__":
    build(); preview(); print(OUT); print(PREVIEW/"contact_sheet.png")
