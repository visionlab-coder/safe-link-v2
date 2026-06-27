from collections import OrderedDict
from pathlib import Path
import re
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
SRC = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_v3_가독성개선_20260611.pptx"
DST = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_v4_페이지압축_20260613.pptx"

NAVY = RGBColor(27, 52, 78)
BLUE = RGBColor(42, 86, 127)
GRAY = RGBColor(82, 82, 82)
BLACK = RGBColor(30, 30, 30)
LIGHT_BLUE = RGBColor(234, 241, 248)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(184, 196, 208)


def shape_text(shape):
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""


def normalize_title(title):
    title = re.sub(r"\s*\(\d+/\d+\)\s*$", "", title.strip())
    return title


def tx(slide, x, y, w, h, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT, fill=None):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = Inches(0.025)
    tf.margin_right = Inches(0.025)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def header(slide, page):
    slide.shapes.add_shape(1, Inches(0.42), Inches(0.50), Inches(10.85), Inches(0.01)).line.color.rgb = LINE
    tx(slide, 0.45, 0.20, 7.0, 0.25, "대우건설 Hyper Safety & AI Open Innovation 제출본 v4", 9.5, color=GRAY)
    tx(slide, 8.6, 0.20, 2.65, 0.25, "CONFIDENTIAL | 서원토건", 9.5, color=GRAY, align=PP_ALIGN.RIGHT)
    slide.shapes.add_shape(1, Inches(0.42), Inches(7.68), Inches(10.85), Inches(0.01)).line.color.rgb = LINE
    tx(slide, 0.45, 7.75, 4.5, 0.20, "공모전 제출 및 PoC 협의용", 8.5, color=GRAY)
    tx(slide, 10.55, 7.75, 0.7, 0.20, f"{page:02d}", 8.5, color=GRAY, align=PP_ALIGN.RIGHT)


def table_ppt(slide, x, y, w, h, headers, rows, widths, font=12):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    total = sum(widths)
    for c, cw in enumerate(widths):
        table.columns[c].width = Inches(w * cw / total)
    row_h = h / (len(rows) + 1)
    for r in range(len(rows) + 1):
        table.rows[r].height = Inches(row_h)
        for c in range(len(headers)):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.018)
            cell.margin_right = Inches(0.018)
            cell.margin_top = Inches(0.006)
            cell.margin_bottom = Inches(0.006)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if r == 0 else WHITE
            cell.text = headers[c] if r == 0 else rows[r - 1][c]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 or c == 0 else PP_ALIGN.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                for rr in p.runs:
                    rr.font.name = "Malgun Gothic"
                    rr.font.size = Pt(font)
                    rr.font.bold = r == 0
                    rr.font.color.rgb = NAVY if r == 0 else BLACK
    return shape


def extract_sections():
    prs = Presentation(SRC)
    sections = OrderedDict()
    final_text = None
    cover_meta = []
    for idx, slide in enumerate(prs.slides):
        tables = [sh.table for sh in slide.shapes if sh.has_table]
        texts = [shape_text(sh) for sh in slide.shapes if shape_text(sh)]
        if idx == 0:
            if tables:
                cover_meta = [
                    [tables[0].cell(r, c).text.strip() for c in range(len(tables[0].columns))]
                    for r in range(1, len(tables[0].rows))
                ]
            continue
        title = next((t for t in texts if re.match(r"^\d+\.", t)), "")
        if title.startswith("21."):
            final_text = "\n".join(t for t in texts if not t.startswith("대우건설") and "CONFIDENTIAL" not in t and not re.fullmatch(r"\d+", t))
            continue
        if not title or not tables:
            continue
        base = normalize_title(title)
        lead = ""
        for t in texts:
            if t != title and not t.startswith("대우건설") and not t.startswith("CONFIDENTIAL") and not t.startswith("공모전 제출") and not re.fullmatch(r"\d+", t):
                lead = t
                break
        table = tables[0]
        headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
        rows = [
            [table.cell(r, c).text.strip() for c in range(len(table.columns))]
            for r in range(1, len(table.rows))
        ]
        if base not in sections:
            sections[base] = {"lead": lead, "headers": headers, "rows": []}
        sections[base]["rows"].extend(rows)
        if not sections[base]["lead"] and lead:
            sections[base]["lead"] = lead
    return cover_meta, sections, final_text


def widths_for(headers):
    n = len(headers)
    if n == 2:
        return [2.4, 8.6]
    if n == 3:
        return [2.4, 4.3, 4.3]
    if n == 4:
        return [2.1, 2.0, 4.5, 3.4]
    return [1] * n


def chunk_size(headers):
    # Keep 12pt body text, but remove previous excessive blank space by allowing more rows per page.
    if len(headers) == 4:
        return 4
    if len(headers) == 3:
        return 4
    return 5


def build():
    cover_meta, sections, final_text = extract_sections()
    out = Presentation()
    out.slide_width = Inches(11.69)
    out.slide_height = Inches(8.27)
    blank = out.slide_layouts[6]

    slide = out.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_shape(1, Inches(0.7), Inches(0.85), Inches(10.3), Inches(0.02)).line.color.rgb = NAVY
    tx(slide, 1.0, 1.25, 9.7, 0.38, "공문서형 제출 제안서 v4 페이지압축본", 14, color=GRAY, align=PP_ALIGN.CENTER)
    tx(slide, 1.0, 1.9, 9.7, 0.95, "대우건설 Hyper Safety & AI\nOpen Innovation", 30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tx(slide, 1.1, 3.15, 9.5, 0.75, "SQ-LINK Underground\n지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼", 16, color=BLUE, align=PP_ALIGN.CENTER)
    rows = cover_meta or [["문서번호", "SW-DW-HSAI-20260611-04"], ["작성일", "2026. 06. 13."], ["제안기관", "서원토건"], ["참고자료", "중기부 R&D HWPX 2종 + 검토의견 MD"]]
    table_ppt(slide, 1.65, 4.62, 8.4, 1.48, ["항목", "내용"], rows[:4], [2.1, 6.3], 11)
    tx(slide, 0.7, 7.35, 10.3, 0.28, "CONFIDENTIAL - 공모전 제출 및 PoC 협의용", 10, color=GRAY, align=PP_ALIGN.CENTER)

    page = 1
    for title, sec in sections.items():
        headers = sec["headers"]
        rows = sec["rows"]
        cs = chunk_size(headers)
        chunks = [rows[i : i + cs] for i in range(0, len(rows), cs)]
        for ci, chunk in enumerate(chunks, 1):
            page += 1
            slide = out.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = WHITE
            header(slide, page - 1)
            suffix = f" ({ci}/{len(chunks)})" if len(chunks) > 1 else ""
            tx(slide, 0.55, 0.68, 10.7, 0.42, title + suffix, 18, bold=True, color=NAVY)
            if sec["lead"]:
                tx(slide, 0.58, 1.14, 10.65, 0.34, sec["lead"], 12, color=GRAY)
            table_ppt(slide, 0.52, 1.60, 10.75, 5.85, headers, chunk, widths_for(headers), 12)

    page += 1
    slide = out.slides.add_slide(blank)
    header(slide, page - 1)
    tx(slide, 0.6, 0.75, 10.6, 0.42, "21. 최종 제안 문장", 18, bold=True, color=NAVY)
    if final_text:
        lines = [l for l in final_text.splitlines() if l and not l.startswith("21.") and "현장 PoC" not in l]
        main = lines[0] if lines else "지하·골조 현장 자동화의 병목은 장비 부족이 아니라 통신, 3D 좌표, 작업데이터의 단절입니다."
        sub = "\n".join(lines[1:]) if len(lines) > 1 else ""
    else:
        main = "지하·골조 현장 자동화의 병목은 장비 부족이 아니라\n통신, 3D 좌표, 작업데이터의 단절입니다."
        sub = "서원토건은 SAFE-LINK에서 출발해 안전과 품질을 연결하는 SQ-LINK Underground 현장 AI 운영체계로 확장하고,\n대우건설과 함께 그 가능성을 실제 현장에서 PoC로 검증하고자 합니다."
    tx(slide, 1.0, 2.0, 9.7, 1.55, main, 22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tx(slide, 1.25, 4.2, 9.2, 1.35, sub, 14, color=BLUE, align=PP_ALIGN.CENTER)
    tx(slide, 3.0, 6.15, 5.7, 0.5, "현장 PoC 및 공동 실증 파트너십 제안", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)

    out.save(DST)
    with ZipFile(DST) as z:
        bad = z.testzip()
        if bad:
            raise RuntimeError(bad)
    print(SRC)
    print(DST)
    print("slides", len(out.slides))


if __name__ == "__main__":
    build()
