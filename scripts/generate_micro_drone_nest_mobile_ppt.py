from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606"
ASSET = BASE / "assets_micro_drone_nest"
OUT = BASE / "SAFE-LINK_Micro_Drone_Nest_대우건설_Advanced_V2_이동형네스트_20260607.pptx"
PREVIEW = BASE / "previews-micro-drone-advanced-v2"
DOC = BASE / "SAFE_LINK_Micro_Drone_Nest_이동형네스트_추가반영.md"

W, H = 13.333, 7.5
COL = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "navy": "0B1220",
    "white": "FFFFFF",
    "soft": "F2F4F7",
}


def rgb(key):
    value = COL.get(key, key).strip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, trans=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = trans
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    if radius:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def text(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def footer(slide, n):
    text(slide, "SAFE-LINK Micro Drone Nest / 서원토건", 0.58, 7.04, 5.0, 0.26, 10, "muted")
    text(slide, f"{n:02d}", 12.28, 7.04, 0.45, 0.26, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    text(slide, label, 0.62, 0.34, 5.8, 0.3, 12, "blue", True)
    text(slide, title, 0.62, 0.74, 11.9, 0.72, 28, "ink", True)
    if sub:
        text(slide, sub, 0.64, 1.5, 11.4, 0.38, 15, "muted")
    rule(slide, 0.62, 2.02, 12.05)
    footer(slide, n)


def pic(slide, filename, x, y, w, h):
    path = ASSET / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        text(slide, filename, x + 0.15, y + h / 2 - 0.15, w - 0.3, 0.3, 13, "muted", True, "center")


def callout(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True)
    rect(slide, x, y, 0.08, h, color, color, False)
    text(slide, head, x + 0.22, y + 0.15, w - 0.42, 0.3, 15, color, True)
    text(slide, body, x + 0.22, y + 0.55, w - 0.42, h - 0.65, 13.5, "ink")


def table(slide, x, y, rows, widths, row_h=0.48):
    for r, row in enumerate(rows):
        cx = x
        for c, value in enumerate(row):
            fill = "navy" if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            rect(slide, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            text(slide, value, cx + 0.08, y + r * row_h + 0.08, widths[c] - 0.16, row_h - 0.1, 11.5 if r == 0 else 12.4, color, r == 0)
            cx += widths[c]


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, "paper", "paper", False)
    return slide


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    s = new_slide(prs)
    pic(s, "01-micro-drone-nest-hero.png", 6.58, 0, 6.75, 7.5)
    text(s, "DAEWOO HYPER SAFETY & AI", 0.72, 0.68, 4.8, 0.3, 12, "blue", True)
    text(s, "SAFE-LINK\nMicro Drone Nest", 0.72, 1.28, 5.25, 1.38, 35, "ink", True)
    text(s, "초소형 드론 네스트 기반\n건설현장 사각공간 자동정찰 플랫폼", 0.74, 3.08, 5.2, 0.78, 20, "ink", True)
    rule(s, 0.74, 4.12, 4.4, "blue")
    text(s, "고정형 네스트 + 이동형 4족보행 로봇 네스트 + SAFE-LINK 조치보고 연동", 0.74, 4.45, 5.25, 0.7, 15, "muted")
    text(s, "서원토건 / 2026.06.07", 0.74, 6.78, 3.4, 0.25, 11, "muted")

    s = new_slide(prs)
    header(s, "01 DIRECTION", "드론이 아니라, 현장의 작은 안전 감각기관입니다", "손바닥보다 작은 초소형 드론이 네스트에서 출발하고 복귀하며 반복 정찰합니다.", 2)
    table(s, 0.9, 2.45, [
        ["구분", "일반 드론 순찰", "Micro Drone Nest"],
        ["크기", "중형·대형 드론 중심", "손바닥보다 작은 초소형 기체"],
        ["공간", "외부·상부·넓은 구역", "PIT, 샤프트, 점검구, 협소공간"],
        ["운용", "사람 조종 중심", "정기 순찰 + 관제센터 명령"],
        ["복귀", "수동 회수·충전", "네스트 자동 복귀·충전"],
        ["연동", "영상 확인 중심", "위험후보 → 조치 → 보고"],
    ], [1.55, 4.4, 5.2], 0.58)
    text(s, "핵심은 ‘날아다니는 카메라’가 아니라 현장 곳곳에 배치되는 자동 안전 감각망입니다.", 1.0, 6.1, 11.1, 0.36, 19, "blue", True, "center")

    s = new_slide(prs)
    header(s, "02 TARGET", "진짜 문제는 넓은 하늘이 아니라 좁은 사각공간입니다", "사람이 들어가기 어렵고 CCTV가 보지 못하는 공간을 반복 확인합니다.", 3)
    pic(s, "02-hard-to-reach-space.png", 0.75, 2.35, 5.75, 3.6)
    table(s, 6.85, 2.45, [
        ["사각공간", "주요 위험"],
        ["엘리베이터 샤프트", "추락, 동하중, 시야 손실"],
        ["지하 PIT·집수정", "침수, 조도 부족, 유해가스"],
        ["천장·배관 주변", "누수, 임시전기, 화재 위험"],
        ["동바리·거푸집 내부", "지지상태 이상, 발판 불량"],
        ["개구부 하부", "추락 전 사전 확인 필요"],
    ], [2.35, 3.45], 0.53)

    s = new_slide(prs)
    header(s, "03 FIXED NEST", "고정형 네스트는 반복 위험구역의 상시 거점입니다", "샤프트, PIT, 점검구처럼 매일 확인해야 하는 곳에 분산 설치합니다.", 4)
    pic(s, "03-distributed-nests.png", 0.75, 2.3, 5.85, 3.55)
    callout(s, "보관·충전", "초소형 드론을 보호하고 접점식 또는 무선충전으로 대기", 6.95, 2.35, 4.95, 0.95, "blue")
    callout(s, "상태 전송", "배터리, 기체 상태, 네스트 통신 상태를 SAFE-LINK로 전송", 6.95, 3.55, 4.95, 0.95, "green")
    callout(s, "미션 대기", "정기 순찰 또는 관제센터 긴급 호출을 기다리는 현장 거점", 6.95, 4.75, 4.95, 0.95, "orange")

    s = new_slide(prs)
    header(s, "04 MOBILE NEST", "4족보행 로봇은 이동형 네스트가 됩니다", "고정형 네스트가 닿지 않는 임시 위험구역까지 초소형 드론 정찰 범위를 확장합니다.", 5)
    pic(s, "06-mobile-robot-dog-nest.png", 0.7, 2.25, 6.05, 3.8)
    callout(s, "이동형 거점", "4족보행 로봇이 변화하는 작업면, 지하층, 통로, 개구부 주변까지 이동", 7.0, 2.3, 4.95, 0.85, "blue")
    callout(s, "드론 전개", "로봇 위 네스트에서 초소형 드론 이륙 → 상부·내부·협소부 촬영", 7.0, 3.35, 4.95, 0.85, "green")
    callout(s, "복귀·충전", "정찰 후 로봇 네스트로 복귀해 충전하거나 다음 지점으로 이동", 7.0, 4.4, 4.95, 0.85, "orange")
    callout(s, "안전용 전환", "군·보안 분야의 정찰 운용 개념을 건설 안전 점검용으로 전환", 7.0, 5.45, 4.95, 0.85, "red")
    text(s, "고정형 네스트가 상시 감각기관이라면, 4족보행 로봇 네스트는 이동형 감각기관입니다.", 0.9, 6.32, 11.5, 0.35, 17, "blue", True, "center")

    s = new_slide(prs)
    header(s, "05 OPERATION", "정기 자동순찰과 긴급 호출을 함께 운용합니다", "자동순찰은 반복성을 만들고, 관제센터 명령은 현장 대응성을 만듭니다.", 6)
    steps = [("스케줄", "정기 순찰"), ("네스트", "상태 확인"), ("드론", "근접 촬영"), ("AI", "위험후보"), ("관리자", "조치 요청"), ("복귀", "자동충전")]
    for i, (a, b) in enumerate(steps):
        x = 0.55 + i * 2.05
        color = ["blue", "green", "orange", "red", "blue", "green"][i]
        rect(s, x, 2.75, 1.55, 1.12, "white", color, True)
        text(s, a, x + 0.1, 2.92, 1.35, 0.26, 13.5, color, True, "center")
        text(s, b, x + 0.1, 3.38, 1.35, 0.24, 12, "ink", True, "center")
        if i < 5:
            text(s, ">", x + 1.58, 3.16, 0.28, 0.24, 17, "muted", True, "center")
    callout(s, "고정형 미션", "반복 위험구역을 매일 같은 시간에 자동 확인해 누락을 줄입니다.", 1.0, 4.8, 5.5, 1.0, "blue")
    callout(s, "이동형 미션", "작업면이 바뀌면 4족보행 로봇이 네스트를 싣고 근처까지 이동한 뒤 드론을 띄웁니다.", 6.9, 4.8, 5.3, 1.0, "green")

    s = new_slide(prs)
    header(s, "06 CONTROL CENTER", "관제센터는 작은 드론들의 실시간 눈을 갖게 됩니다", "네스트 위치, 영상 피드, 위험후보, 조치상태를 한 화면에서 관리합니다.", 7)
    pic(s, "04-control-center-micro-feeds.png", 0.75, 2.35, 5.85, 3.55)
    table(s, 6.85, 2.45, [
        ["기능", "설명"],
        ["네스트 상태", "드론 유무, 배터리, 충전, 통신"],
        ["실시간 영상", "사각공간별 최근 촬영·라이브 피드"],
        ["원격 명령", "재촬영, 이동, 긴급복귀, 미션중지"],
        ["조치관리", "담당자, 기한, 완료사진, 확정 상태"],
        ["보고 자동화", "일일·주간 안전점검 리포트 생성"],
    ], [2.1, 3.65], 0.53)

    s = new_slide(prs)
    header(s, "07 TECH & AI", "1차 PoC는 완전자율보다 위험후보 식별부터 시작합니다", "안전관리자가 최종 확인하는 human-in-the-loop 구조가 현실적입니다.", 8)
    table(s, 0.85, 2.35, [
        ["초소형 드론", "요구사항"],
        ["크기·안전", "손바닥 이하, 보호가드, 저속 운용"],
        ["시야", "근접 카메라, LED 라이트"],
        ["위치", "Optical Flow, AprilTag, UWB"],
        ["복귀", "네스트 마커 기반 자동복귀"],
    ], [2.0, 3.75], 0.55)
    table(s, 6.75, 2.35, [
        ["AI 식별", "위험후보"],
        ["샤프트", "개구부 방치, 안전난간 손실"],
        ["PIT", "침수, 조도 부족, 이물질"],
        ["천장", "누수, 전선 노출, 화재위험"],
        ["협소부", "안전망 손실, 동하중 위험"],
    ], [2.0, 3.75], 0.55)
    text(s, "촬영 → 프레임 추출 → 위험후보 식별 → 관리자 확인 → 조치요청 → 데이터셋 축적", 1.0, 6.0, 11.2, 0.34, 17, "blue", True, "center")

    s = new_slide(prs)
    header(s, "08 POC & BUDGET", "3~5개 사각공간이면 8주 PoC가 가능합니다", "대우건설 PoC 지원금 기준으로 고정형+이동형 네스트 검증안을 제안합니다.", 9)
    table(s, 0.85, 2.35, [
        ["구성", "내용"],
        ["장비", "초소형 드론 2~5대, 고정형 네스트 3~5개"],
        ["이동형", "4족보행 로봇 1대 + 탑재형 네스트 1식"],
        ["대상", "샤프트, PIT, 점검구, 지하 통로, 개구부 하부"],
        ["운영", "일 3회 자동순찰 + 관제센터 수동 호출"],
        ["기간", "8주 PoC"],
    ], [2.0, 4.9], 0.52)
    callout(s, "1차 PoC 예산", "초소형 드론·고정형 네스트 2,000만\n이동형 네스트 연동 1,500만\nSAFE-LINK 관제·보고 1,000만\n운영·분석·예비비 1,000만\n합계 5,500만 원형", 8.05, 2.42, 3.75, 2.7, "green")
    callout(s, "선행 검증", "제출 전 100만~500만 원 규모로 기체·네스트 목업과 관제 화면을 먼저 검증할 수 있습니다.", 8.05, 5.45, 3.75, 0.82, "orange")

    s = new_slide(prs)
    header(s, "09 KPI & RISK", "성과지표는 안전효용 기준으로 제시해야 설득력이 생깁니다", "초소형 드론은 안전성, 개인정보, 충전, 통신 리스크를 처음부터 관리해야 합니다.", 10)
    table(s, 0.85, 2.35, [
        ["정량 KPI", "목표"],
        ["사각공간·네스트", "3~5개"],
        ["드론 운용", "2~5대, 일 3회"],
        ["자동복귀·충전", "80% 이상"],
        ["영상 전송", "90% 이상"],
        ["위험후보 식별", "70% 이상"],
        ["조치요청", "위험 발견 후 1분 이내"],
    ], [2.8, 2.4], 0.48)
    table(s, 6.45, 2.35, [
        ["리스크", "대응"],
        ["충돌·접촉", "보호가드, 저속 운용, 통제시간 운용"],
        ["추락", "경량 기체, 비상정지, 회수 프로토콜"],
        ["개인정보", "촬영 고지, 얼굴 비식별, 접근권한 제한"],
        ["통신장애", "자동복귀, 수동 회수 프로토콜"],
        ["배터리", "과열 차단, 네스트 온도감시"],
    ], [2.15, 3.75], 0.48)

    s = new_slide(prs)
    header(s, "10 VALUE", "대우건설이 얻는 것은 장비가 아니라 사각공간의 상시 감각입니다", "현장 하나에서 검증하고 대우건설형 Hyper Safety 패키지로 확장합니다.", 11)
    pic(s, "05-safelink-action-report.png", 0.75, 2.25, 5.75, 3.55)
    callout(s, "대우건설 요청", "PoC 현장 1개소\n사각공간 지정\n안전관리자 피드백\n비행·촬영 시간 협의", 6.95, 2.35, 4.95, 1.1, "blue")
    callout(s, "제공 가치", "작업자 진입 전 위험 확인\n조치보고 자동화\n위치·시간·영상 이력 추적", 6.95, 3.75, 4.95, 1.1, "green")
    callout(s, "확장안", "현장별 네스트 추가\n로봇개·AI-RAN 연계\n정부과제·공동개발 연결", 6.95, 5.15, 4.95, 1.1, "orange")

    BASE.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


def write_doc():
    DOC.write_text(
        """# SAFE-LINK Micro Drone Nest V2 추가 반영

작성일: 2026-06-07

## 추가 컨셉: 4족보행 로봇 탑재형 이동형 네스트

기존 Micro Drone Nest는 샤프트, PIT, 점검구, 천장 배관 주변처럼 반복적으로 확인해야 하는 위험구역에 고정형 네스트를 설치하는 구조였다. 여기에 4족보행 로봇을 결합하면, 네스트 자체가 현장을 이동할 수 있는 이동형 정찰 거점이 된다.

핵심 문장:

> 고정형 네스트가 상시 감각기관이라면, 4족보행 로봇 네스트는 이동형 감각기관이다.

## 운용 방식

1. 4족보행 로봇이 지하층, 통로, 개구부 주변, 작업면 변경 구역으로 이동한다.
2. 로봇 상부의 네스트에서 초소형 드론이 이륙한다.
3. 드론은 사람이 들어가기 어렵거나 로봇 카메라가 직접 볼 수 없는 상부·내부·협소공간을 촬영한다.
4. 촬영 영상과 위험후보는 SAFE-LINK 관제센터로 전송된다.
5. 필요 시 관제센터가 재촬영, 긴급복귀, 다음 지점 이동을 명령한다.
6. 드론은 로봇 네스트로 복귀해 충전하거나 다음 정찰을 대기한다.

## 고정형 네스트와 이동형 네스트의 역할 분담

| 구분 | 고정형 네스트 | 이동형 4족보행 로봇 네스트 |
|---|---|---|
| 설치 위치 | 샤프트, PIT, 점검구 등 반복 위험구역 | 매일 바뀌는 작업면, 임시 위험구역 |
| 장점 | 상시 대기, 반복 순찰, 관리 단순 | 현장 변화 대응, 넓은 지하층 커버, 긴급 출동 |
| 드론 역할 | 고정 위치 주변 반복 촬영 | 로봇 도달 지점 주변의 협소·상부·내부 촬영 |
| 관제 방식 | 정기 스케줄 중심 | 관제센터 명령 + 상황 대응 중심 |
| SAFE-LINK 연동 | 정기 보고, 위험후보 알림 | 긴급 알림, 조치요청, 이동형 증거 수집 |

## 대우건설 제안서 반영 포인트

- 군·보안 분야의 4족보행 로봇+정찰 드론 운용 개념을 건설 안전 점검용으로 전환한다.
- 무장, 군사용 감시, 보안 위협 대응이 아니라 사각공간 안전 확인과 중대재해 예방 목적임을 명확히 한다.
- 로봇개는 드론을 대체하는 장비가 아니라 초소형 드론의 이동형 네스트다.
- 대우건설 PoC에서는 고정형 네스트 3~5개와 이동형 네스트 1식을 함께 검증한다.
- 모빌리오 협업 또는 중국 로봇 도입 여부는 추후 조달 방식으로 열어두고, 제안서에는 기술 구성과 PoC 목표 중심으로 표현한다.
""",
        encoding="utf-8",
    )


def make_preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles = [
        "SAFE-LINK Micro Drone Nest",
        "작은 안전 감각기관",
        "좁은 사각공간",
        "고정형 네스트",
        "4족보행 이동형 네스트",
        "운영 시나리오",
        "관제센터",
        "기술과 AI",
        "PoC와 예산",
        "KPI와 리스크",
        "대우건설 확장 가치",
    ]
    image_map = {
        1: "01-micro-drone-nest-hero.png",
        3: "02-hard-to-reach-space.png",
        4: "03-distributed-nests.png",
        5: "06-mobile-robot-dog-nest.png",
        7: "04-control-center-micro-feeds.png",
        11: "05-safelink-action-report.png",
    }
    try:
        f1 = ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf", 38)
        f2 = ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf", 20)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    paths = []
    for idx, title in enumerate(titles, 1):
        canvas = Image.new("RGB", (1600, 900), "#" + COL["paper"])
        draw = ImageDraw.Draw(canvas)
        draw.text((80, 60), "SAFE-LINK Micro Drone Nest V2", font=f2, fill="#" + COL["blue"])
        draw.text((80, 135), title, font=f1, fill="#" + COL["ink"])
        draw.line((80, 235, 1460, 235), fill="#" + COL["line"], width=3)
        if idx in image_map and (ASSET / image_map[idx]).exists():
            im = Image.open(ASSET / image_map[idx]).convert("RGB")
            im.thumbnail((760, 430))
            canvas.paste(im, (80, 285))
        for k, color in enumerate(["blue", "green", "orange", "red"]):
            draw.rounded_rectangle((980, 310 + k * 105, 1390, 370 + k * 105), 16, outline="#" + COL[color], width=3, fill="white")
        draw.text((80, 830), f"{idx:02d}", font=f2, fill="#" + COL["muted"])
        out = PREVIEW / f"slide_{idx:02d}.png"
        canvas.save(out)
        paths.append(out)

    sheet = Image.new("RGB", (1600, 1100), "white")
    draw = ImageDraw.Draw(sheet)
    for i, path in enumerate(paths):
        im = Image.open(path).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 345
        sheet.paste(im, (x, y))
        draw.text((x, y + 214), f"{i + 1:02d} {titles[i]}", font=f2, fill="#" + COL["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build_ppt()
    write_doc()
    make_preview()
    print(OUT)
    print(DOC)
    print(PREVIEW / "contact_sheet.png")
