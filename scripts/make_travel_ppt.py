from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RED   = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0xAA, 0xAA, 0xAA)
DARK  = RGBColor(0x07, 0x07, 0x0E)
LGRAY = RGBColor(0x55, 0x55, 0x66)
BLUE  = RGBColor(0x52, 0x98, 0xDB)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
CARD  = RGBColor(0x10, 0x10, 0x1C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide():
    sl = prs.slides.add_slide(blank)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK
    return sl

def rect(sl, x, y, w, h, color):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def label(sl, x, y, w, h, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

# ── SLIDE 1: 표지 ──────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 0.08, RED)
rect(sl, 0, 7.42, 13.33, 0.08, RED)
label(sl, 0, 0.6, 13.33, 0.55, "SAFE-LINK", 13, bold=True, color=RED, align=PP_ALIGN.CENTER)
label(sl, 0, 1.2, 13.33, 1.8, "Travel Talk", 72, color=WHITE, align=PP_ALIGN.CENTER)
label(sl, 0, 3.0, 13.33, 0.55, "언어의 벽을 넘어, 세계와 대화하세요", 20, color=GRAY, align=PP_ALIGN.CENTER)
label(sl, 0, 3.55, 13.33, 0.45, "言葉の壁を越えて · Break the language barrier", 13, color=LGRAY, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 4.3, 6.33, 1.8, CARD)
label(sl, 3.5, 4.45, 6.33, 0.45, "앱 설치 불필요  ·  アプリ不要", 14, color=WHITE, align=PP_ALIGN.CENTER)
label(sl, 3.5, 4.95, 6.33, 0.45, "QR 스캔만으로 즉시 연결", 14, color=WHITE, align=PP_ALIGN.CENTER)
label(sl, 3.5, 5.45, 6.33, 0.45, "5개 언어 실시간 번역 · 통역", 14, color=WHITE, align=PP_ALIGN.CENTER)
label(sl, 0, 6.85, 13.33, 0.4, "서원토건 미래전략TF  ·  2026", 11, color=LGRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: 지원 언어 & 시나리오 ─────────────────────────
sl = add_slide()
rect(sl, 0, 0, 0.08, 7.5, RED)
label(sl, 0.4, 0.35, 12, 0.55, "지원 언어 및 사용 시나리오", 28, bold=True)
label(sl, 0.4, 0.9, 12, 0.38, "5개 언어 간 양방향 실시간 번역", 14, color=GRAY)
langs = [("KR","한국어"),("JP","日本語"),("US","English"),("CN","中文"),("VN","Viet")]
flags = ["🇰🇷","🇯🇵","🇺🇸","🇨🇳","🇻🇳"]
for i,(code,name) in enumerate(langs):
    x = 0.4 + i*2.5
    rect(sl, x, 1.5, 2.2, 2.0, CARD)
    label(sl, x, 1.6, 2.2, 0.75, flags[i], 34, align=PP_ALIGN.CENTER)
    label(sl, x, 2.35, 2.2, 0.45, name, 15, bold=True, align=PP_ALIGN.CENTER)
label(sl, 0.4, 3.75, 12, 0.45, "주요 활용 장면", 18, bold=True, color=RED)
scenes = [
    ("✈️  해외여행", "일본·중국·베트남\n현지인과 즉석 대화"),
    ("🛎️  비즈니스", "외국 바이어·파트너\n실시간 협의"),
    ("🏗️  현장 소통", "외국인 근로자\n업무 지시 전달"),
    ("🏥  긴급 상황", "병원·관공서\n언어 장벽 돌파"),
]
for i,(title,desc) in enumerate(scenes):
    x = 0.4 + i*3.1
    rect(sl, x, 4.3, 2.85, 2.5, CARD)
    label(sl, x+0.1, 4.45, 2.65, 0.55, title, 14, bold=True)
    label(sl, x+0.1, 5.05, 2.65, 1.5, desc, 12, color=GRAY)

# ── SLIDE 3: 호스트 시작 방법 ─────────────────────────────
sl = add_slide()
label(sl, 0.5, 0.3, 12, 0.55, "호스트 시작 방법 (한국인)", 28, bold=True)
label(sl, 0.5, 0.85, 12, 0.38, "3단계로 대화 시작 — 1분이면 충분합니다", 14, color=GRAY)
steps = [
    ("1","언어 선택","'한국어' 국기를 선택합니다","🇰🇷"),
    ("2","새 대화 시작","'새 대화 시작' 버튼을 누릅니다\nQR 코드가 자동 생성됩니다","▶"),
    ("3","화면 보여주기","QR 대기 화면을 외국인에게 보여줍니다\n화면에 4개 언어로 안내가 표시됩니다\n상대방이 스캔하면 자동 연결!","📷"),
]
for i,(num,title,desc,icon) in enumerate(steps):
    x = 0.5 + i*4.1
    rect(sl, x, 1.55, 0.5, 0.5, RED)
    label(sl, x, 1.57, 0.5, 0.45, num, 20, bold=True, align=PP_ALIGN.CENTER)
    rect(sl, x, 2.15, 3.65, 3.6, CARD)
    label(sl, x+0.15, 2.3, 3.35, 0.55, icon+"  "+title, 17, bold=True)
    label(sl, x+0.15, 2.95, 3.35, 2.5, desc, 13, color=GRAY)
rect(sl, 0.5, 5.95, 12.3, 1.0, RGBColor(0x1A,0x08,0x08))
label(sl, 0.7, 6.05, 12.0, 0.45, "💡  QR 대기 화면에 일·영·중·베트남어로 '스캔해주세요' 안내 자동 표시", 13, color=RED)
label(sl, 0.7, 6.52, 12.0, 0.38, "      말 한마디 없이 화면만 보여줘도 외국인이 바로 이해하고 스캔합니다", 12, color=GRAY, italic=True)

# ── SLIDE 4: 게스트(외국인) 참여 방법 ────────────────────
sl = add_slide()
label(sl, 0.5, 0.3, 12, 0.55, "게스트 참여 방법 (외국인)", 28, bold=True)
label(sl, 0.5, 0.85, 12, 0.38, "앱 설치 불필요  ·  스마트폰 카메라만 있으면 OK", 14, color=GRAY)
steps2 = [
    ("1","QR 스캔","QRコードをスキャン\nScan the QR code\n请扫描二维码\nQuet ma QR"),
    ("2","언어 선택","자국 국기를 탭합니다\nTap your flag\n点击您的国旗\nNhan vao co ban"),
    ("3","자동 연결!","자동으로 채팅방 입장\nAuto-connected!\n自动进入聊天室\nTu dong ket noi!"),
]
for i,(num,title,desc) in enumerate(steps2):
    x = 0.5 + i*4.1
    rect(sl, x, 1.55, 0.5, 0.5, RED)
    label(sl, x, 1.57, 0.5, 0.45, num, 20, bold=True, align=PP_ALIGN.CENTER)
    rect(sl, x, 2.15, 3.65, 4.2, CARD)
    label(sl, x+0.15, 2.3, 3.35, 0.55, title, 17, bold=True)
    label(sl, x+0.15, 2.95, 3.35, 3.2, desc, 12, color=GRAY)
label(sl, 0.5, 6.55, 12.3, 0.5, "브라우저만 있으면 OK — 카카오톡 · 라인 · 위챗 인앱브라우저에서도 작동합니다", 13, color=GRAY)

# ── SLIDE 5: 대화 모드 ────────────────────────────────────
sl = add_slide()
label(sl, 0.5, 0.3, 12, 0.55, "대화 모드 선택", 28, bold=True)
label(sl, 0.5, 0.85, 12, 0.38, "호스트가 상황에 맞게 선택합니다", 14, color=GRAY)
rect(sl, 0.5, 1.45, 5.85, 5.1, CARD)
label(sl, 0.7, 1.6, 5.5, 0.55, "💬  일반 대화 모드", 19, bold=True)
label(sl, 0.7, 2.2, 5.5, 0.38, "버튼을 눌러 말하고 놓으면 번역", 13, color=GRAY)
items_a = [
    "긴 문장·복잡한 내용에 적합",
    "발음 + 역번역 학습 기능 사용 가능",
    "오인식 확인 후 재전송 가능",
    "번역 속도: ~1초 내외",
]
for j,it in enumerate(items_a):
    label(sl, 0.7, 2.75+j*0.62, 5.5, 0.5, ("✅  " if j<3 else "⏱  ")+it, 13, color=WHITE if j<3 else GRAY)
rect(sl, 6.95, 1.45, 5.85, 5.1, RGBColor(0x16,0x05,0x05))
label(sl, 7.15, 1.6, 5.5, 0.55, "🎙  동시통역 모드", 19, bold=True, color=RED)
label(sl, 7.15, 2.2, 5.5, 0.38, "말하면 자동으로 실시간 번역·전송", 13, color=GRAY)
items_b = [
    "빠른 현장 지시·짧은 문장에 최적",
    "말하는 즉시 상대 화면에 표시",
    "마이크 On 상태 유지",
    "번역 속도: ~0.8초 (800ms)",
]
for j,it in enumerate(items_b):
    label(sl, 7.15, 2.75+j*0.62, 5.5, 0.5, ("✅  " if j<3 else "⚡  ")+it, 13, color=WHITE if j<3 else GRAY)
label(sl, 0.5, 6.75, 12.3, 0.45, "긴 설명·계약 내용 → 일반 대화  |  현장 지시·짧은 교환 → 동시통역 추천", 13, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 6: 추가 기능 ───────────────────────────────────
sl = add_slide()
label(sl, 0.5, 0.3, 12, 0.55, "추가 기능", 28, bold=True)
rect(sl, 0.5, 1.1, 5.85, 4.2, CARD)
label(sl, 0.7, 1.25, 5.5, 0.55, "📖  학습 모드", 18, bold=True, color=BLUE)
label(sl, 0.7, 1.85, 5.5, 0.38, "한국인 호스트 전용", 12, color=GRAY)
label(sl, 0.7, 2.3, 5.5, 2.8, "ON 시 번역 + 한글 발음 + 역번역 표시\n\n예)  日本語です\n      [발음]  니혼고데스\n      [역번역]  일본어입니다\n\n상대 언어를 따라 읽고 싶을 때 사용", 12, color=GRAY)
rect(sl, 6.95, 1.1, 5.85, 4.2, CARD)
label(sl, 7.15, 1.25, 5.5, 0.55, "🔊  음성 설정", 18, bold=True, color=GREEN)
label(sl, 7.15, 1.85, 5.5, 3.8, "음성 ON/OFF  — 무음 환경에서 끄기\n\n남자/여자 음성 전환\n→ 상황·취향에 맞게 선택\n\n▶  재생 버튼\n→ 수신 메시지 다시 듣기\n\nGoogle Neural2 AI 음성 사용", 12, color=GRAY)
rect(sl, 0.5, 5.5, 12.3, 1.5, RGBColor(0x0E,0x0E,0x1A))
label(sl, 0.7, 5.65, 12.0, 0.45, "내부 기술 스택", 13, bold=True, color=GRAY)
label(sl, 0.7, 6.1, 12.0, 0.7, "번역: Papago(네이버)  ·  음성합성(TTS): Google Neural2  ·  음성인식(STT): Google Cloud STT", 12, color=LGRAY)

# ── SLIDE 7: FAQ ─────────────────────────────────────────
sl = add_slide()
label(sl, 0.5, 0.3, 12, 0.55, "자주 묻는 질문 (FAQ)", 28, bold=True)
faqs = [
    ("Q. 앱을 설치해야 하나요?",
     "아니요. 스마트폰 브라우저만 있으면 됩니다. 카카오톡 인앱브라우저도 OK."),
    ("Q. 외국인이 QR을 어떻게 스캔하나요?",
     "QR 대기 화면에 일·영·중·베 4개 언어로 안내가 표시됩니다. 화면을 보여주기만 하면 됩니다."),
    ("Q. 인터넷이 필요한가요?",
     "예. 호스트와 게스트 모두 모바일 데이터 또는 Wi-Fi가 필요합니다."),
    ("Q. 동시에 여러 명이 참여할 수 있나요?",
     "현재는 1:1 대화 방식입니다. 추후 그룹 모드 추가 예정입니다."),
    ("Q. 번역이 얼마나 정확한가요?",
     "Papago(네이버) 기반으로 여행 회화 수준의 정확도를 제공합니다. 학습 모드의 역번역으로 추가 확인 가능합니다."),
]
for i,(q,a) in enumerate(faqs):
    y = 1.05 + i*1.2
    rect(sl, 0.5, y, 12.3, 1.1, CARD)
    label(sl, 0.7, y+0.06, 12.0, 0.4, q, 13, bold=True, color=RED)
    label(sl, 0.7, y+0.5, 11.8, 0.5, a, 11, color=GRAY)

# ── SLIDE 8: 마무리 ──────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 0.08, RED)
rect(sl, 0, 7.42, 13.33, 0.08, RED)
label(sl, 0, 1.6, 13.33, 0.65, "SAFE-LINK Travel Talk", 40, color=WHITE, align=PP_ALIGN.CENTER)
label(sl, 0, 2.35, 13.33, 0.5, "이제 언어는 더 이상 장벽이 아닙니다", 18, color=GRAY, align=PP_ALIGN.CENTER)
label(sl, 0, 2.9, 13.33, 0.4, "言語の壁を越えて、世界と話しましょう", 13, color=LGRAY, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 3.65, 6.33, 2.3, CARD)
label(sl, 3.5, 3.82, 6.33, 0.45, "접속 주소", 13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
label(sl, 3.5, 4.3, 6.33, 0.6, "safe-link-v2.vercel.app/travel", 17, bold=True, color=RED, align=PP_ALIGN.CENTER)
label(sl, 3.5, 4.92, 6.33, 0.45, "또는 SAFE-LINK 앱 -> Travel Talk 메뉴", 12, color=LGRAY, align=PP_ALIGN.CENTER)
label(sl, 0, 6.35, 13.33, 0.4, "서원토건 미래전략TF  |  기술지원본부  |  2026", 11, color=LGRAY, align=PP_ALIGN.CENTER)

out = "C:/Users/seowo/Documents/dev/seowon-projects/SAFE-LINK-V2/SAFE-LINK_TravelTalk_사용설명서.pptx"
prs.save(out)
print("saved:", out)
