from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "generated" / "framexr-snapshot-viewpoint-20260607" / "FRAMEXR_Snapshot_Viewpoint_Capture_MVP_이미지보강_20260607.pptx"
PREVIEW = ROOT / "docs" / "generated" / "framexr-snapshot-viewpoint-20260607" / "previews" / "contact_sheet.png"
prs = Presentation(PPTX)
all_text = " ".join(
    shape.text
    for slide in prs.slides
    for shape in slide.shapes
    if getattr(shape, "has_text_frame", False) and shape.text
)
pictures = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
print(f"pptx_exists={PPTX.exists()}")
print(f"preview_exists={PREVIEW.exists()}")
print(f"slide_count={len(prs.slides)}")
print(f"picture_count={pictures}")
print(f"has_embed={'A안 임베드' in all_text}")
print(f"has_no_restore={'Restore 버튼' in all_text}")
print(f"has_coordinate_basis={'public-ifc-derived-demo' in all_text}")
print(f"has_no_server={'서버 저장' in all_text}")
print(f"has_no_bcf={'BCF' in all_text}")
