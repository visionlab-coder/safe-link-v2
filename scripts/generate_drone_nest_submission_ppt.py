from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606"
ASSET = BASE / "assets"
OUT = BASE / "SAFE-LINK_Drone_Nest_대우건설_제출용_PPT_실사이미지_20260606.pptx"
PREVIEW = BASE / "previews-submission-ppt"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "white": "FFFFFF",
    "soft": "F2F4F7",
    "navy": "0B1220",
}


def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, transparency=0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.05
        except Exception:
            pass
    return shp


def txt(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def line(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.01))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()


def footer(slide, n):
    txt(slide, "SAFE-LINK Drone Nest / 서원토건", 0.55, 7.05, 4.2, 0.22, 10, "muted")
    txt(slide, f"{n:02d}", 12.25, 7.05, 0.45, 0.22, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    txt(slide, label, 0.62, 0.34, 5.0, 0.28, 11, "blue", True)
    txt(slide, title, 0.62, 0.72, 11.5, 0.7, 28, "ink", True)
    if sub:
        txt(slide, sub, 0.64, 1.46, 11.0, 0.38, 14, "muted")
    line(slide, 0.62, 1.96, 12.05)
    footer(slide, n)


def image(slide, name, x, y, w, h):
    p = ASSET / name
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        txt(slide, name, x, y + h / 2 - 0.15, w, 0.3, 14, "muted", True, "center")


def note(slide, title, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True)
    rect(slide, x, y, 0.08, h, color, color, False)
    txt(slide, title, x + 0.22, y + 0.18, w - 0.4, 0.32, 15, color, True)
    txt(slide, body, x + 0.22, y + 0.62, w - 0.38, h - 0.76, 15, "ink")


def table(slide, x, y, rows, widths, row_h=0.56):
    for r, row in enumerate(rows):
        cx = x
        for c, value in enumerate(row):
            fill = "navy" if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            bold = r == 0
            rect(slide, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            txt(slide, value, cx + 0.08, y + r * row_h + 0.1, widths[c] - 0.16, row_h - 0.14, 12 if r == 0 else 13, color, bold)
            cx += widths[c]


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide(bg="paper"):
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, W, H, bg, bg, False)
        return s

    # 1 cover
    s = slide()
    image(s, "01-cover-drone-construction-hero.png", 0, 0, W, H)
    rect(s, 0, 0, 6.25, H, "navy", "navy", False, 10)
    txt(s, "DAEWOO HYPER SAFETY & AI", 0.7, 0.65, 4.8, 0.3, 12, "white", True)
    txt(s, "SAFE-LINK\nDrone Nest", 0.7, 1.35, 4.95, 1.28, 38, "white", True)
    txt(s, "AI 자율드론 기반 건설현장 위험구역 자동순찰 및 중대재해 예방 플랫폼", 0.72, 3.05, 4.65, 0.78, 18, "white", True)
    line(s, 0.72, 4.05, 4.1, "white")
    txt(s, "드론 순찰 → AI 위험판독 → SAFE-LINK 알림·조치·보고", 0.72, 4.35, 4.75, 0.55, 15, "white")
    txt(s, "서원토건 / 2026.06.06", 0.72, 6.78, 3.0, 0.24, 11, "white")

    # 2 problem
    s = slide()
    header(s, "01 PROBLEM", "위험구역은 매일 바뀌지만, 순찰 방식은 그대로입니다", "수동 순찰과 고정 CCTV만으로는 고층 골조현장의 반복 위험을 놓치기 쉽습니다.", 2)
    image(s, "02-problem-risk-zones.png", 0.75, 2.35, 5.9, 3.55)
    note(s, "수동 순찰", "넓은 현장 전체를 같은 기준으로 반복 확인하기 어렵고 보고가 수작업입니다.", 7.05, 2.35, 4.9, 0.95, "blue")
    note(s, "CCTV 사각", "고정형 카메라는 층별 공정 변화와 자재 이동을 따라가지 못합니다.", 7.05, 3.55, 4.9, 0.95, "orange")
    note(s, "골조 반복위험", "개구부, 단부, 자재 적치, 타워크레인 반경이 매일 새로 생깁니다.", 7.05, 4.75, 4.9, 0.95, "red")

    # 3 positioning
    s = slide()
    header(s, "02 POSITIONING", "단순 드론 촬영이 아니라 안전관리 프로세스 자동화입니다", "촬영 파일을 남기는 것이 아니라 위험 발견부터 조치 완료 보고까지 연결합니다.", 3)
    rows = [
        ["일반 드론 촬영", "SAFE-LINK Drone Nest"],
        ["조종자가 필요할 때 촬영", "위험 스팟 10~20개 반복순찰"],
        ["영상 파일 제공", "AI 위험후보 자동 태깅"],
        ["사람이 별도 확인", "SAFE-LINK 담당자 알림"],
        ["보고서 수작업", "조치 전후 리포트 자동화"],
    ]
    table(s, 1.0, 2.45, rows, [5.3, 5.3], 0.66)
    txt(s, "핵심 차별성: 드론 데이터를 안전관리 업무 흐름에 직접 연결", 1.0, 6.05, 11.0, 0.38, 21, "blue", True, "center")

    # 4 solution flow
    s = slide()
    header(s, "03 SOLUTION", "Drone Nest 운영 루프", "현장에 상주하는 순찰 시스템처럼 작동하는 것이 목표입니다.", 4)
    labels = [("위험구역 등록", "개구부·단부·자재"), ("웨이포인트 순찰", "일 2~3회 반복"), ("AI 위험판독", "위험후보 태깅"), ("SAFE-LINK 조치", "알림·담당자·보고")]
    for i, (a, b) in enumerate(labels):
        x = 0.8 + i * 3.05
        rect(s, x, 3.0, 2.35, 1.35, "white", ["blue", "green", "orange", "red"][i], True)
        txt(s, a, x + 0.18, 3.22, 1.95, 0.3, 16, ["blue", "green", "orange", "red"][i], True, "center")
        txt(s, b, x + 0.18, 3.72, 1.95, 0.3, 14, "ink", True, "center")
        if i < 3:
            txt(s, "→", x + 2.42, 3.38, 0.4, 0.32, 24, "muted", True, "center")
    txt(s, "위험 발견은 시작이고, 조치 완료 보고까지 닫혀야 안전관리 데이터가 됩니다.", 1.0, 5.45, 11.0, 0.42, 20, "ink", True, "center")

    # 5 drone nest device
    s = slide()
    header(s, "04 DRONE NEST", "정밀복귀·충전까지 가야 반복순찰이 됩니다", "1차 PoC는 반자동부터 시작하고, 확장 단계에서 자동충전 도크를 붙입니다.", 5)
    image(s, "03-drone-docking-nest.png", 0.75, 2.25, 5.9, 3.55)
    note(s, "1단계", "저가형 오픈소스 드론으로 웨이포인트 순찰과 촬영 업로드 검증", 7.05, 2.35, 4.8, 0.95, "blue")
    note(s, "2단계", "AI 위험태그와 SAFE-LINK 조치 리포트 연결", 7.05, 3.55, 4.8, 0.95, "green")
    note(s, "3단계", "정밀착륙, 자동충전 도크, 대우건설형 패키지화", 7.05, 4.75, 4.8, 0.95, "orange")

    # 6 AI risk targets
    s = slide()
    header(s, "05 AI RISK DETECTION", "1차 PoC는 ‘잘 보이는 위험’부터 정확히 잡습니다", "처음부터 모든 위험을 맞히기보다, 시각 판독 가능한 항목을 안전하게 시작합니다.", 6)
    rows = [
        ["우선순위", "감지 대상", "이유"],
        ["1", "안전난간 미설치", "추락 중대위험"],
        ["2", "개구부 덮개 미설치", "골조현장 반복 위험"],
        ["3", "자재 통로 침범", "동선·전도 위험"],
        ["4", "위험구역 접근", "통제구역 관리"],
        ["5", "보호구 미착용", "AI 비전 확장 항목"],
    ]
    table(s, 1.0, 2.35, rows, [1.45, 4.0, 5.2], 0.58)
    txt(s, "AI는 자동 확정이 아니라 ‘위험 후보’를 만들고, 안전관리자가 최종 확인합니다.", 1.0, 6.05, 11.0, 0.36, 18, "red", True, "center")

    # 7 dashboard/report
    s = slide()
    header(s, "06 SAFE-LINK LINKAGE", "드론이 본 위험은 SAFE-LINK에서 조치로 닫힙니다", "관리자는 사진, 위치, 위험태그, 담당자, 조치 상태를 한 화면에서 봅니다.", 7)
    image(s, "04-safelink-ai-dashboard-field.png", 0.75, 2.28, 5.55, 3.35)
    image(s, "05-before-after-ai-report.png", 6.75, 2.28, 5.55, 3.35)
    txt(s, "좌: 드론 순찰 관제 / 우: 조치 전후 리포트", 1.0, 5.88, 11.0, 0.3, 15, "muted", True, "center")

    # 8 PoC
    s = slide()
    header(s, "07 POC PLAN", "8주 PoC로 실증 가능한 범위만 제안합니다", "대우건설 현장 1개소, 위험스팟 10~20개, 일 2~3회 순찰 기준입니다.", 8)
    weeks = [("1-2주", "위험구역 선정\n비행계획 수립"), ("3-4주", "반자동 순찰\n촬영 업로드"), ("5-6주", "AI 위험태그\nSAFE-LINK 알림"), ("7-8주", "KPI 측정\n확장안 도출")]
    for i, (a, b) in enumerate(weeks):
        x = 0.8 + i * 3.0
        rect(s, x, 2.8, 2.35, 2.15, "white", ["blue", "green", "orange", "red"][i], True)
        txt(s, a, x + 0.15, 3.08, 2.05, 0.32, 18, ["blue", "green", "orange", "red"][i], True, "center")
        txt(s, b, x + 0.15, 3.68, 2.05, 0.6, 17, "ink", True, "center")
    txt(s, "초기 PoC는 조종자 상시 대기와 수동전환 가능 구조로 운용", 1.0, 5.8, 11.0, 0.34, 17, "ink", True, "center")

    # 9 KPI/safety operation
    s = slide()
    header(s, "08 KPI & SAFETY", "성과와 안전 운용 기준을 함께 제시합니다", "드론 PoC는 기술보다 현장 안전운용 신뢰가 중요합니다.", 9)
    rows = [
        ["KPI", "목표"],
        ["순찰 스팟", "10~20개"],
        ["반복 촬영 성공률", "90% 이상"],
        ["SAFE-LINK 업로드율", "95% 이상"],
        ["위험태그 정확도", "1차 PoC 70% 이상"],
        ["수동 점검 시간 절감", "20~30%"],
        ["일일 리포트 생성", "5분 이내"],
    ]
    table(s, 0.85, 2.35, rows, [4.9, 2.2], 0.48)
    note(s, "운용 원칙", "저고도·통제구역 중심\n조종자 상시 대기\n작업자 고지·얼굴 비식별\n보험·승인 필요 시 사전 검토", 8.35, 2.45, 3.7, 2.45, "green")

    # 10 closing
    s = slide()
    header(s, "09 COLLABORATION VALUE", "대우건설에는 안전관리 자동화 PoC, 서원토건에는 공동상품화 기회", "현장 하나에서 검증하고, 대우건설형 스마트 안전 패키지로 확장합니다.", 10)
    note(s, "대우건설 요청", "PoC 현장 제공\n안전관리자 피드백\n위험구역 데이터\n운용 안전 기준 협의", 0.9, 2.55, 3.45, 2.0, "blue")
    note(s, "서원토건 제공", "골조현장 위험 정의\nSAFE-LINK 연동 설계\nPoC 운영·성과분석\n후속 공동개발안", 4.95, 2.55, 3.45, 2.0, "green")
    note(s, "최종 가치", "순찰 사각 감소\n조치보고 자동화\n중대재해 예방\n스마트 안전 브랜드 강화", 9.0, 2.55, 3.45, 2.0, "orange")
    txt(s, "SAFE-LINK Drone Nest는 ‘드론을 날리는 제안’이 아니라\n위험 발견부터 조치 완료까지 닫는 안전관리 자동순찰 플랫폼입니다.", 1.05, 5.55, 11.2, 0.75, 24, "ink", True, "center")

    BASE.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(OUT))
    titles = [
        "SAFE-LINK Drone Nest",
        "위험구역은 매일 바뀝니다",
        "촬영이 아니라 프로세스 자동화",
        "Drone Nest 운영 루프",
        "정밀복귀·충전까지",
        "AI 위험감지",
        "SAFE-LINK 조치 연동",
        "8주 PoC",
        "KPI & Safety",
        "협업 가치",
    ]
    try:
        f1 = ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf", 42)
        f2 = ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf", 22)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    paths = []
    for i, title in enumerate(titles, 1):
        img = Image.new("RGB", (1600, 900), "#" + C["paper"])
        d = ImageDraw.Draw(img)
        d.text((90, 70), "SAFE-LINK Drone Nest", font=f2, fill="#" + C["blue"])
        d.text((90, 150), title, font=f1, fill="#" + C["ink"])
        d.line((90, 250, 1450, 250), fill="#" + C["line"], width=3)
        if i in [1, 2, 5, 7]:
            names = {1: "01-cover-drone-construction-hero.png", 2: "02-problem-risk-zones.png", 5: "03-drone-docking-nest.png", 7: "04-safelink-ai-dashboard-field.png"}
            p = ASSET / names[i]
            if p.exists():
                im = Image.open(p).convert("RGB")
                im.thumbnail((920, 520))
                img.paste(im, (90, 300))
        for k, color in enumerate(["blue", "green", "orange", "red"]):
            d.rounded_rectangle((980, 330 + k * 105, 1390, 390 + k * 105), 16, outline="#" + C[color], width=3, fill="white")
        d.text((90, 830), f"{i:02d}", font=f2, fill="#" + C["muted"])
        out = PREVIEW / f"slide_{i:02d}.png"
        img.save(out)
        paths.append(out)
    sheet = Image.new("RGB", (1600, 840), "white")
    d = ImageDraw.Draw(sheet)
    for i, p in enumerate(paths):
        im = Image.open(p).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 265
        sheet.paste(im, (x, y))
        d.text((x, y + 210), f"{i+1:02d}", font=f2, fill="#" + C["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build()
    preview()
    print(OUT)
    print(PREVIEW / "contact_sheet.png")
