from pathlib import Path
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
GEN = Path.home() / ".codex" / "generated_images" / "019e7e61-97bf-7ad0-9e60-9a1de84c4656"
BASE = ROOT / "docs" / "generated" / "framexr-snapshot-viewpoint-20260607"
ASSET = BASE / "assets"
PREVIEW = BASE / "previews"
OUT = BASE / "FRAMEXR_Snapshot_Viewpoint_Capture_MVP_이미지보강_20260607.pptx"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "navy": "0B1220",
    "white": "FFFFFF",
    "soft": "F2F4F7",
}


def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.04
        except Exception:
            pass
    return shp


def text(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()


def pic(slide, name, x, y, w, h):
    p = ASSET / name
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        text(slide, name, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, 13, "muted", True, "center")


def callout(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True)
    rect(slide, x, y, 0.08, h, color, color, False)
    text(slide, head, x + 0.22, y + 0.14, w - 0.42, 0.3, 15, color, True)
    text(slide, body, x + 0.22, y + 0.54, w - 0.42, h - 0.62, 13.2, "ink")


def footer(slide, n):
    text(slide, "FRAMEXR Snapshot Viewpoint Capture MVP", 0.55, 7.05, 5.4, 0.24, 10, "muted")
    text(slide, f"{n:02d}", 12.28, 7.05, 0.45, 0.24, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    text(slide, label, 0.62, 0.34, 5.8, 0.3, 12, "blue", True)
    text(slide, title, 0.62, 0.74, 11.9, 0.72, 28, "ink", True)
    if sub:
        text(slide, sub, 0.64, 1.5, 11.4, 0.38, 15, "muted")
    rule(slide, 0.62, 2.02, 12.05)
    footer(slide, n)


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, "paper", "paper", False)
    return s


def table(slide, x, y, rows, widths, row_h=0.5):
    for r, row in enumerate(rows):
        cx = x
        for c, value in enumerate(row):
            fill = "navy" if r == 0 else ("white" if r % 2 else "soft")
            color = "white" if r == 0 else "ink"
            rect(slide, cx, y + r * row_h, widths[c], row_h, fill, "line", False)
            text(slide, value, cx + 0.08, y + r * row_h + 0.08, widths[c] - 0.16, row_h - 0.1, 11.5 if r == 0 else 12.3, color, r == 0)
            cx += widths[c]


def prepare_assets():
    ASSET.mkdir(parents=True, exist_ok=True)
    latest = sorted(GEN.glob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)[:4]
    names = [
        "04-demo-coordinate-basis.png",
        "03-issue-card-viewpoint-badge.png",
        "02-embedded-evidence-object.png",
        "01-snapshot-viewpoint-hero.png",
    ]
    for src, name in zip(latest, names):
        shutil.copy2(src, ASSET / name)


def build():
    prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    s = slide(prs)
    pic(s, "01-snapshot-viewpoint-hero.png", 6.5, 0, 6.83, 7.5)
    text(s, "FRAMEXR V0.3 MVP", 0.72, 0.68, 4.8, 0.3, 12, "blue", True)
    text(s, "Snapshot Attach 시\nViewpoint 자동 저장", 0.72, 1.22, 5.35, 1.35, 32, "ink", True)
    text(s, "A안 임베드 방식\nimageDataUrl + viewpoint metadata", 0.74, 3.0, 5.2, 0.78, 20, "ink", True)
    rule(s, 0.74, 4.08, 4.3, "blue")
    text(s, "Restore / BCF export / 서버 저장 없이, evidence 객체 안에 시점 정보를 함께 저장하는 1차 MVP", 0.74, 4.42, 5.2, 0.78, 14.5, "muted")
    text(s, "2026.06.07", 0.74, 6.78, 2.8, 0.25, 11, "muted")

    s = slide(prs)
    header(s, "01 WHY", "스냅샷만 있으면 ‘어디서 본 장면인지’가 약합니다", "이미지에 카메라 위치와 대상점이 함께 저장되어야 리뷰자가 같은 맥락을 이해할 수 있습니다.", 2)
    pic(s, "01-snapshot-viewpoint-hero.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "기존 증거", "이미지 썸네일과 카운트는 남지만, 당시 카메라 시점은 별도로 보존되지 않음", 7.0, 2.35, 4.95, 0.95, "orange")
    callout(s, "MVP 보강", "SnapshotEvidence 안에 viewpoint metadata를 임베드하여 이미지와 시점을 한 객체로 저장", 7.0, 3.55, 4.95, 0.95, "blue")
    callout(s, "사용자 효과", "동일한 이슈 스냅샷이 어느 방향, 어떤 모델 기준에서 찍혔는지 설명 가능", 7.0, 4.75, 4.95, 0.95, "green")

    s = slide(prs)
    header(s, "02 DATA MODEL", "A안은 evidence 객체 안에 시점 데이터를 함께 넣습니다", "별도 viewpointByEvidence map 없이 SnapshotEvidence.viewpoint 하나로 묶습니다.", 3)
    pic(s, "02-embedded-evidence-object.png", 0.75, 2.25, 5.9, 3.65)
    table(s, 6.85, 2.35, [
        ["필드", "저장 내용"],
        ["cameraPosition", "카메라 위치 x/y/z"],
        ["cameraTarget", "controls target x/y/z"],
        ["cameraUp", "카메라 up vector"],
        ["fov", "Perspective camera 시야각"],
        ["renderMode / sceneSource", "현재 뷰어 상태"],
        ["marker / issue / component", "선택 맥락 id"],
    ], [2.5, 3.45], 0.5)

    s = slide(prs)
    header(s, "03 CAPTURE FLOW", "handleAttachSnapshot 성공 후 viewpoint도 같이 생성합니다", "captureViewpointRef가 null이면 기존 스냅샷만 저장해 회귀를 막습니다.", 4)
    callout(s, "1. dataURL 생성", "기존 snapshot imageDataUrl 생성 흐름 유지", 0.9, 2.45, 3.45, 0.95, "blue")
    callout(s, "2. camera/controls 읽기", "scene effect에서 camera position, target, up, fov를 캡처 가능하게 구성", 4.9, 2.45, 3.45, 0.95, "green")
    callout(s, "3. evidence 임베드", "SnapshotEvidence.viewpoint에 메타데이터를 넣고 동일 evidence 객체에 저장", 8.9, 2.45, 3.45, 0.95, "orange")
    callout(s, "4. cleanup", "viewer cleanup 시 captureViewpointRef.current = null 처리", 2.9, 4.45, 3.65, 0.95, "red")
    callout(s, "5. fallback", "mock scene 또는 ref null이면 guid/viewpoint 없이 기존 동작 유지", 6.95, 4.45, 3.65, 0.95, "blue")

    s = slide(prs)
    header(s, "04 ISSUE CARD", "카드에는 ‘시점 포함’ 배지만 표시합니다", "Restore 버튼, 전체 메타데이터 표시는 이번 MVP 범위에서 제외합니다.", 5)
    pic(s, "03-issue-card-viewpoint-badge.png", 0.75, 2.25, 5.9, 3.65)
    callout(s, "보이는 것", "thumbnail 영역 또는 evidence 영역에 짧은 시점 포함 배지", 7.0, 2.35, 4.95, 0.95, "blue")
    callout(s, "유지할 것", "thumbnail, count, snapshot local-only 고지", 7.0, 3.55, 4.95, 0.95, "green")
    callout(s, "금지할 것", "Viewpoint Restore 버튼, 상세 메타데이터 표, 다운로드 기능", 7.0, 4.75, 4.95, 0.95, "red")

    s = slide(prs)
    header(s, "05 COORDINATE BASIS", "좌표 표현은 반드시 데모 기준으로 제한합니다", "public IFC-derived demo world-m 좌표일 뿐, 실제 측량/AR/회사도면 좌표가 아닙니다.", 6)
    pic(s, "04-demo-coordinate-basis.png", 0.75, 2.15, 6.05, 3.85)
    callout(s, "worldUnit", '"m"', 7.05, 2.35, 4.95, 0.65, "blue")
    callout(s, "coordinateBasis", '"public-ifc-derived-demo"', 7.05, 3.18, 4.95, 0.65, "green")
    callout(s, "selectedComponentGuid", "demo id 또는 public IFC-derived id일 수 있음. 실제 IfcGlobalId로 단정 금지", 7.05, 4.01, 4.95, 0.95, "orange")
    callout(s, "명확한 제외", "실제 측량 좌표, 실제 회사 도면 좌표, 실제 AR 정합 좌표 아님", 7.05, 5.15, 4.95, 0.95, "red")

    s = slide(prs)
    header(s, "06 IMPLEMENTATION SCOPE", "코드 변경은 4개 파일로 제한합니다", "ViewerPage 내부 최소 구현으로 진행하고, 헬퍼 분리는 후속 필요성만 보고합니다.", 7)
    table(s, 1.0, 2.3, [
        ["파일", "변경 내용"],
        ["apps/web/src/types/issues.ts", "SnapshotViewpoint 타입 + SnapshotEvidence.viewpoint"],
        ["apps/web/src/pages/ViewerPage.tsx", "captureViewpointRef + attach 시 viewpoint 생성"],
        ["apps/web/src/components/issues/LinkedFieldIssuesCard.tsx", "시점 포함 배지 표시"],
        ["apps/web/src/i18n/dictionary.ts", "snapshot.viewpointBadge ko/en 추가"],
    ], [4.25, 7.0], 0.62)
    text(s, "utils/viewpoint.ts 신규 헬퍼, apps/api, 외부 라이브러리, 서버 저장, BCF export는 제외합니다.", 1.0, 5.55, 11.2, 0.36, 17, "red", True, "center")

    s = slide(prs)
    header(s, "07 VALIDATION", "검증은 기능 저장과 회귀 방지를 함께 봅니다", "스냅샷 첨부, 시점 필드 생성, 배지 표시, 기존 기능 회귀 없음이 핵심입니다.", 8)
    table(s, 0.95, 2.3, [
        ["검증 축", "확인 내용"],
        ["Build", "apps/web npm run build"],
        ["Snapshot", "첨부 정상 + viewpoint 생성 정상"],
        ["Camera", "position / target / up / fov 저장"],
        ["Context", "renderMode / sceneSource / marker / issue / component id 저장"],
        ["Fallback", "public IFC Demo id 저장, mock scene undefined fallback"],
        ["Regression", "thumbnail/count/local-only/QA/Pour/Measurement/Toolbar 회귀 없음"],
        ["Excluded", "Restore, BCFzip/JSON, 서버/API, 원본 파일 staged 0"],
    ], [2.7, 8.65], 0.5)

    prs.save(OUT)


def preview():
    PREVIEW.mkdir(parents=True, exist_ok=True)
    titles = [
        "Snapshot Viewpoint Capture",
        "왜 필요한가",
        "A안 임베드 데이터 모델",
        "Capture Flow",
        "Issue Card 배지",
        "Demo Coordinate Basis",
        "구현 범위",
        "검증 항목",
    ]
    image_map = {
        1: "01-snapshot-viewpoint-hero.png",
        2: "01-snapshot-viewpoint-hero.png",
        3: "02-embedded-evidence-object.png",
        5: "03-issue-card-viewpoint-badge.png",
        6: "04-demo-coordinate-basis.png",
    }
    try:
        f1 = ImageFont.truetype(r"C:\Windows\Fonts\malgunbd.ttf", 38)
        f2 = ImageFont.truetype(r"C:\Windows\Fonts\malgun.ttf", 20)
    except Exception:
        f1 = f2 = ImageFont.load_default()
    paths = []
    for i, title in enumerate(titles, 1):
        canvas = Image.new("RGB", (1600, 900), "#" + C["paper"])
        d = ImageDraw.Draw(canvas)
        d.text((80, 60), "FRAMEXR Snapshot Viewpoint MVP", font=f2, fill="#" + C["blue"])
        d.text((80, 135), title, font=f1, fill="#" + C["ink"])
        d.line((80, 235, 1460, 235), fill="#" + C["line"], width=3)
        if i in image_map and (ASSET / image_map[i]).exists():
            im = Image.open(ASSET / image_map[i]).convert("RGB")
            im.thumbnail((760, 430))
            canvas.paste(im, (80, 285))
        for k, color in enumerate(("blue", "green", "orange", "red")):
            d.rounded_rectangle((980, 310 + k * 105, 1390, 370 + k * 105), 16, outline="#" + C[color], width=3, fill="white")
        d.text((80, 830), f"{i:02d}", font=f2, fill="#" + C["muted"])
        out = PREVIEW / f"slide_{i:02d}.png"
        canvas.save(out)
        paths.append(out)

    sheet = Image.new("RGB", (1600, 760), "white")
    d = ImageDraw.Draw(sheet)
    for i, path in enumerate(paths):
        im = Image.open(path).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 340
        sheet.paste(im, (x, y))
        d.text((x, y + 212), f"{i + 1:02d} {titles[i]}", font=f2, fill="#" + C["ink"])
    sheet.save(PREVIEW / "contact_sheet.png")


if __name__ == "__main__":
    build()
    preview()
    print(OUT)
    print(PREVIEW / "contact_sheet.png")
