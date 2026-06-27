from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "generated" / "daewoo-hyper-safety-20260604"
PREVIEW_DIR = OUT_DIR / "previews-safe-link-underground-ppt"
PPTX_PATH = OUT_DIR / "SAFE-LINK_UNDERGROUND_대우건설_HyperSafety_AI_제안서_가독성개선_20260604.pptx"

W, H = 13.333, 7.5
PX_W, PX_H = 1600, 900

COLORS = {
    "bg": "F7F6F3",
    "white": "FFFFFF",
    "ink": "111827",
    "muted": "6B7280",
    "line": "E5E7EB",
    "blue": "2563EB",
    "green": "059669",
    "yellow": "F59E0B",
    "red": "DC2626",
    "purple": "7C3AED",
    "slate": "374151",
}


def rgb(hex_value):
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def add_text(slide, text, x, y, w, h, size=22, color="ink", bold=False, align="left", font="맑은 고딕"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_rect(slide, x, y, w, h, fill="white", line="line", radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.color.rgb = rgb(COLORS[line])
    shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def add_rule(slide, x, y, w, color="line"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS[color])
    line.line.fill.background()
    return line


def add_pill(slide, text, x, y, w, color="blue"):
    add_rect(slide, x, y, w, 0.34, fill="white", line=color, radius=True)
    add_text(slide, text, x + 0.08, y + 0.04, w - 0.16, 0.26, size=11, color=color, bold=True, align="center")


def add_footer(slide, num):
    add_text(slide, "SAFE-LINK UNDERGROUND | 서원토건", 0.55, 7.0, 4.8, 0.26, size=11, color="muted")
    add_text(slide, f"{num:02d}", 12.22, 7.0, 0.5, 0.26, size=11, color="muted", align="right")


def title(slide, eyebrow, heading, sub, num=None):
    add_text(slide, eyebrow, 0.65, 0.4, 6.0, 0.34, size=14, color="blue", bold=True)
    add_text(slide, heading, 0.65, 0.82, 11.6, 0.95, size=32, color="ink", bold=True)
    if sub:
        add_text(slide, sub, 0.67, 1.78, 11.4, 0.5, size=16, color="muted")
    add_rule(slide, 0.65, 2.28, 12.0)
    if num:
        add_footer(slide, num)


def card(slide, x, y, w, h, head, body, accent="blue"):
    add_rect(slide, x, y, w, h, fill="white", line="line", radius=True)
    add_rect(slide, x, y, 0.08, h, fill=accent, line=accent)
    add_text(slide, head, x + 0.22, y + 0.18, w - 0.38, 0.38, size=18, color="ink", bold=True)
    add_text(slide, body, x + 0.22, y + 0.76, w - 0.38, h - 0.88, size=15, color="slate")


def flow_node(slide, x, y, w, label, body, color="blue"):
    add_rect(slide, x, y, w, 1.1, fill="white", line=color, radius=True)
    add_text(slide, label, x + 0.16, y + 0.12, w - 0.32, 0.3, size=14, color=color, bold=True, align="center")
    add_text(slide, body, x + 0.14, y + 0.48, w - 0.28, 0.52, size=14, color="ink", align="center")


def table(slide, x, y, cols, rows, widths, row_h=0.48):
    cx = x
    for i, c in enumerate(cols):
        add_rect(slide, cx, y, widths[i], row_h, fill="ink", line="ink")
        add_text(slide, c, cx + 0.08, y + 0.1, widths[i] - 0.16, 0.32, size=14, color="white", bold=True, align="center")
        cx += widths[i]
    for r, row in enumerate(rows):
        cx = x
        fill = "white" if r % 2 == 0 else "bg"
        for i, cell in enumerate(row):
            add_rect(slide, cx, y + row_h * (r + 1), widths[i], row_h, fill=fill, line="line")
            add_text(slide, cell, cx + 0.1, y + row_h * (r + 1) + 0.1, widths[i] - 0.2, row_h - 0.14, size=14, color="ink")
            cx += widths[i]


def image_if_exists(slide, path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def new_slide(bg="bg"):
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, W, H, fill=bg, line=bg)
        return s

    # 1 cover
    s = new_slide()
    add_text(s, "DAEWOO HYPER SAFETY & AI OPEN INNOVATION", 0.7, 0.63, 8.8, 0.34, size=14, color="blue", bold=True)
    add_text(s, "SAFE-LINK\nUNDERGROUND", 0.7, 1.22, 8.0, 1.75, size=44, color="ink", bold=True)
    add_text(s, "AI-RAN 기반 철근콘크리트 지하 골조 Hyper Safety 플랫폼", 0.75, 3.12, 8.1, 0.42, size=18, color="slate")
    add_rule(s, 0.75, 3.62, 4.1, "blue")
    add_text(s, "지하 음영지역 · 로키드 스마트글라스 · AI 로봇개 · 다국어 TBM · 작업중지 증빙", 0.75, 3.92, 7.5, 0.66, size=17, color="ink")
    card(s, 8.55, 0.85, 3.95, 1.15, "전장", "지하층, 코어, PIT, 동바리, 타설구간", "blue")
    card(s, 8.55, 2.25, 3.95, 1.15, "방식", "통신·시야·이동탐지·증빙의 결합", "green")
    card(s, 8.55, 3.65, 3.95, 1.15, "목표", "대우건설 스마트 안전의 지하 사각 보완", "purple")
    add_text(s, "서원토건", 0.75, 6.72, 2.5, 0.34, size=14, color="muted", bold=True)
    add_text(s, "2026.06.04", 10.65, 6.72, 1.8, 0.34, size=14, color="muted", align="right")

    # 2
    s = new_slide()
    title(s, "PROBLEM", "가장 위험한 곳이 가장 연결되지 않습니다", "철근콘크리트 지하 골조는 위험·언어·통신 사각이 동시에 발생하는 작업면입니다.", 2)
    card(s, 0.75, 2.55, 3.75, 2.7, "통신 음영", "지하층·PIT·코어 내부에서 모바일 앱, CCTV, 센서 관제가 끊깁니다.", "red")
    card(s, 4.82, 2.55, 3.75, 2.7, "공정 위험", "동바리, 거푸집, 철근, 타설, 개구부 위험이 하루 단위로 바뀝니다.", "yellow")
    card(s, 8.9, 2.55, 3.75, 2.7, "외국인 근로자", "TBM을 했는지보다 이해했는지, 위험을 말할 수 있는지가 중요합니다.", "blue")

    # 3
    s = new_slide()
    title(s, "POSITIONING", "대우건설 기존 체계를 대체하지 않고 지하 사각을 완성합니다", "지상 관제는 원청 시스템, 지하 작업면 실행은 SAFE-LINK UNDERGROUND가 맡습니다.", 3)
    flow_node(s, 0.85, 2.75, 2.2, "지상 관제", "대우건설 기존 스마트 안전", "slate")
    flow_node(s, 3.45, 2.75, 2.2, "통신 버블", "AI-RAN / O-RAN", "blue")
    flow_node(s, 6.05, 2.75, 2.2, "작업면 시야", "로키드 스마트글라스", "green")
    flow_node(s, 8.65, 2.75, 2.2, "선행 순찰", "AI 로봇개", "purple")
    flow_node(s, 11.05, 2.75, 1.55, "증빙", "SAFE-LINK", "yellow")
    add_text(s, "지상 → 지하 → 작업자 → 관리자 → 증빙 보고서", 1.65, 4.5, 10.0, 0.55, size=25, color="ink", bold=True, align="center")
    add_text(s, "핵심은 ‘앱 하나’가 아니라 통신, 착용형 관제, 이동형 센서, 다국어 안전기록을 하나의 현장 데이터 흐름으로 묶는 것입니다.", 1.05, 5.25, 11.2, 0.7, size=16, color="slate", align="center")

    # 4
    s = new_slide()
    title(s, "WHY REBAR CONCRETE", "철근콘크리트 전문건설업 현장성이 제안의 방어력입니다", "서원토건은 지하 골조 공정의 실제 작업순서와 사고 포인트를 알고 있습니다.", 4)
    rows = [
        ["철근 배근", "찔림, 협착, 철근 다발 전도, 양중 낙하"],
        ["거푸집·동바리", "전도, 미고정, 수평재 누락, 잭베이스 불량"],
        ["콘크리트 타설", "호스 휘핑, 측압 증가, 동선 충돌, 미끄럼"],
        ["지하층·PIT", "산소결핍, 조도 부족, 통신 단절, 추락"],
    ]
    table(s, 1.05, 2.55, ["공정", "현장 위험"], rows, [2.4, 8.6], row_h=0.62)
    add_text(s, "원청 관제가 놓치는 ‘작업면 단위 실행’을 전문건설업 데이터로 보강", 0.85, 5.72, 11.4, 0.48, size=21, color="blue", bold=True, align="center")

    # 5
    s = new_slide()
    title(s, "ARCHITECTURE", "네 개의 기술을 하나의 안전 루프로 연결합니다", "AI-RAN, 스마트글라스, 로봇개, SAFE-LINK가 각자 다른 사각을 맡습니다.", 5)
    card(s, 0.75, 2.45, 2.85, 2.35, "AI-RAN", "지하 음영구간에 안전 통신 버블을 만들고 통신 품질을 안전지표로 관리", "blue")
    card(s, 3.9, 2.45, 2.85, 2.35, "로키드 글라스", "열화상·영상 촬영·안드로이드 기반 착용형 현장 관제", "green")
    card(s, 7.05, 2.45, 2.85, 2.35, "AI 로봇개", "작업 전 선행 순찰, 위험구역 촬영, 센서 확장 노드", "purple")
    card(s, 10.2, 2.45, 2.45, 2.35, "SAFE-LINK", "TBM, 퀴즈, 신고, 조치, 해시 증빙", "yellow")
    add_text(s, "결과: 지하 작업면의 ‘연결 불가, 확인 불가, 설명 불가’를 제거", 0.9, 5.52, 11.6, 0.5, size=23, color="ink", bold=True, align="center")

    # 6
    s = new_slide()
    title(s, "AI-RAN", "지하 작업구역 단위 안전 통신 버블", "통신망 공사가 아니라, 위험 작업구역에 필요한 만큼 빠르게 펼치는 안전 인프라입니다.", 6)
    add_text(s, "추진 방식", 0.9, 2.52, 2.8, 0.36, size=19, color="ink", bold=True)
    card(s, 0.9, 3.05, 3.4, 1.1, "경희대 홍인기 교수", "MOU 체결 추진 및 AI-RAN/O-RAN 자문·공동제안 검토", "blue")
    card(s, 4.85, 3.05, 3.4, 1.1, "국책과제", "지하 건설현장 안전통신 과제화 추진", "green")
    card(s, 8.8, 3.05, 3.4, 1.1, "대우 PoC", "통신 음영 측정과 현장 실증 연계", "purple")
    add_text(s, "표현 수위: MOU 전에는 ‘확정 협력기관’이 아니라 ‘MOU 및 공동제안 추진 예정’으로 기재", 0.9, 5.03, 11.3, 0.58, size=16, color="red", bold=True, align="center")

    # 7
    s = new_slide()
    title(s, "SMART GLASS", "로키드 스마트글라스는 움직이는 작업면 관제 장치입니다", "CCTV가 못 보는 곳을 작업조장 시야로 보고, 열화상과 영상으로 기록합니다.", 7)
    card(s, 0.8, 2.55, 3.65, 2.45, "기능", "안드로이드 기반\n영상 촬영\n열화상 카메라\n음성 지시", "green")
    card(s, 4.85, 2.55, 3.65, 2.45, "철콘 적용", "타설·양생 온도 이상\n동바리·거푸집 촬영\n개구부 위험 확인", "blue")
    card(s, 8.9, 2.55, 3.65, 2.45, "SAFE-LINK 연동", "다국어 안전지시\n위험 캡처 저장\n작업중지 이벤트 연결", "purple")
    add_text(s, "착용형 관제의 가치는 ‘보는 것’이 아니라 보고 즉시 지시하고 증빙으로 남기는 것입니다.", 0.9, 5.62, 11.5, 0.46, size=19, color="ink", bold=True, align="center")

    # 8
    s = new_slide()
    title(s, "ROBOT DOG", "로봇개는 지하 작업면의 이동형 위험탐지 노드입니다", "공급사는 PoC에서 확정하되, 제안 구조에서는 필수 구성으로 포함합니다.", 8)
    rows = [
        ["확보 방식", "모빌리오 협업 또는 중국산 상용 로봇 도입 검토"],
        ["역할", "지하 선행 순찰, 통신 품질 확인, 위험구역 촬영"],
        ["센서 확장", "산소, 가스, 조도, 열화상, 위치 데이터"],
        ["SAFE-LINK 연결", "순찰 이벤트를 TBM·작업중지·보고서와 연결"],
    ]
    table(s, 1.0, 2.45, ["구분", "내용"], rows, [2.1, 8.9], row_h=0.62)
    add_text(s, "단독 로봇 시연이 아니라 ‘안전 데이터 흐름의 이동형 센서’로 제안", 0.85, 5.72, 11.5, 0.46, size=21, color="purple", bold=True, align="center")

    # 9
    s = new_slide()
    title(s, "SAFETY PASSPORT", "Rebar-Zone Safety Passport", "근로자 개인 이수만 보지 않고, 작업구역별 안전허가 상태를 관리합니다.", 9)
    steps = [
        ("QR/NFC", "구역 접속"),
        ("TBM", "공종별 교육"),
        ("QUIZ", "언어별 이해 확인"),
        ("SIGN", "안전약속 전자서명"),
        ("ROBOT", "선행 순찰 완료"),
        ("GO", "작업 허가"),
    ]
    x = 0.65
    for i, (a, b) in enumerate(steps):
        flow_node(s, x + i * 2.05, 3.0, 1.65, a, b, ["blue", "green", "purple", "yellow", "slate", "red"][i])
    add_text(s, "작업구역이 바뀌면 안전여권도 바뀝니다. 이것이 철근콘크리트 현장에 맞는 방식입니다.", 0.9, 5.12, 11.5, 0.58, size=19, color="ink", bold=True, align="center")

    # 10
    s = new_slide()
    title(s, "POUR GATE", "Pour-Go / No-Go AI 타설 승인 게이트", "타설 전 조건이 하나라도 미달이면 시스템이 보류 상태를 표시합니다.", 10)
    rows = [
        ["구조", "동바리·거푸집 점검 완료"],
        ["근로자", "TBM 이수율·퀴즈 통과·전자서명"],
        ["현장", "개구부 방호·레미콘 동선·펌프카 위험"],
        ["기술", "통신 버블 정상·글라스 관제 준비·로봇 순찰 이상 없음"],
        ["이슈", "작업중지 미처리 건 없음"],
    ]
    table(s, 1.0, 2.42, ["확인 영역", "Go 조건"], rows, [2.0, 9.0], row_h=0.62)
    add_text(s, "타설은 골조 공정의 핵심 이벤트입니다. 타설 승인 게이트 하나로 안전과 생산성을 동시에 설명할 수 있습니다.", 0.9, 5.72, 11.5, 0.52, size=16, color="slate", align="center")

    # 11
    s = new_slide()
    title(s, "MULTILINGUAL STOP-WORK", "외국어 위험신고를 관리자 조치와 법적 증빙까지 연결합니다", "작업중지권을 제도가 아니라 실제 작동하는 기능으로 만듭니다.", 11)
    flow_node(s, 0.9, 2.9, 2.0, "근로자", "모국어 신고", "green")
    flow_node(s, 3.25, 2.9, 2.0, "AI", "정규화·위험분류", "blue")
    flow_node(s, 5.6, 2.9, 2.0, "관리자", "한국어 알림", "yellow")
    flow_node(s, 7.95, 2.9, 2.0, "조치", "현장 확인·처리", "purple")
    flow_node(s, 10.3, 2.9, 2.0, "증빙", "해시체인 저장", "red")
    add_text(s, "베트남어로 말해도, 관리자는 한국어로 받고, 근로자는 모국어로 결과를 받습니다.", 0.9, 5.08, 11.5, 0.52, size=21, color="ink", bold=True, align="center")

    # 12
    s = new_slide()
    title(s, "TRIANGULATION", "Robot · Glass · SAFE-LINK 삼각 증빙", "서로 다른 세 데이터가 같은 작업구역에 묶이면 사고 전후 설명력이 커집니다.", 12)
    card(s, 1.0, 2.45, 3.25, 2.35, "로봇개", "사람이 들어가기 전\n위험구역·통신·센서 확인", "purple")
    card(s, 5.05, 2.45, 3.25, 2.35, "스마트글라스", "작업조장 시야\n열화상·영상 기록", "green")
    card(s, 9.1, 2.45, 3.25, 2.35, "SAFE-LINK", "TBM·신고·서명·조치\n해시 보고서", "blue")
    add_text(s, "Evidence-First Safety Twin: 3D 모델보다 먼저 필요한 것은 설명 가능한 안전 이력입니다.", 0.9, 5.52, 11.5, 0.52, size=20, color="ink", bold=True, align="center")

    # 13
    s = new_slide()
    title(s, "POC", "8주 PoC로 현장 적용성을 증명합니다", "대우건설 지하층 골조 또는 지하주차장 철근콘크리트 공정 1개 구역을 대상으로 합니다.", 13)
    rows = [
        ["1-2주", "현장 조사, 통신 음영 측정, TBM·위험성평가 양식 반영"],
        ["3-4주", "NFC/QR TBM, 퀴즈, 전자서명, 다국어 안내 운영"],
        ["5-6주", "로키드 글라스 관제, 로봇개 순찰, 작업중지 신고 운영"],
        ["7-8주", "효과 측정, 보고서 자동화, 대우건설 확장 적용안 도출"],
    ]
    table(s, 1.0, 2.45, ["기간", "실증 내용"], rows, [1.7, 9.3], row_h=0.7)

    # 14
    s = new_slide()
    title(s, "METRICS", "성과는 안전관리자가 바로 이해하는 지표로 측정합니다", "기술 시연보다 현장 효과가 평가의 핵심입니다.", 14)
    metrics = [
        ("TBM 참여율", "근로자별 이수"),
        ("이해도", "퀴즈·재교육"),
        ("신고 처리시간", "접수→조치"),
        ("통신 커버리지", "음영 히트맵"),
        ("보고시간 절감", "수기 대비"),
        ("증빙 완결성", "해시 이력"),
    ]
    for i, (m, d) in enumerate(metrics):
        x = 0.9 + (i % 3) * 4.05
        y = 2.45 + (i // 3) * 1.55
        card(s, x, y, 3.45, 1.15, m, d, ["blue", "green", "purple", "yellow", "red", "slate"][i])

    # 15
    s = new_slide()
    title(s, "COLLABORATION", "협업 구조는 명확하고, 표현 수위는 안전하게 가져갑니다", "확정된 것은 확정으로, 추진 중인 것은 추진으로 써야 평가 리스크가 없습니다.", 15)
    rows = [
        ["서원토건", "철근콘크리트 공정 데이터, SAFE-LINK 운영, 현장 시나리오"],
        ["대우건설", "PoC 현장, TBM·위험성평가 양식, 안전관리자 피드백"],
        ["경희대 홍인기 교수", "AI-RAN/O-RAN MOU 및 국책과제·공동제안 추진"],
        ["로봇개 파트너", "모빌리오 협업 또는 상용 로봇 도입 검토"],
        ["로키드", "실제 스마트글라스 제품 기반 현장 적용"],
    ]
    table(s, 0.9, 2.32, ["주체", "역할"], rows, [2.25, 9.0], row_h=0.62)

    # 16
    s = new_slide()
    title(s, "LANDING", "대우건설의 지하 안전 사각을 서원토건의 작업면 데이터로 메웁니다", "SAFE-LINK UNDERGROUND는 원청 관제와 전문건설업 실행 사이의 빈 공간을 연결합니다.", 16)
    add_text(s, "지하 통신", 1.0, 2.72, 2.45, 0.36, size=20, color="blue", bold=True)
    add_text(s, "+", 3.15, 2.65, 0.4, 0.4, size=28, color="muted", bold=True, align="center")
    add_text(s, "작업면 시야", 3.7, 2.72, 2.55, 0.36, size=20, color="green", bold=True)
    add_text(s, "+", 5.85, 2.65, 0.4, 0.4, size=28, color="muted", bold=True, align="center")
    add_text(s, "선행 순찰", 6.45, 2.72, 2.45, 0.36, size=20, color="purple", bold=True)
    add_text(s, "+", 8.55, 2.65, 0.4, 0.4, size=28, color="muted", bold=True, align="center")
    add_text(s, "다국어 증빙", 9.1, 2.72, 2.9, 0.36, size=20, color="yellow", bold=True)
    add_text(s, "철근콘크리트 지하 골조 Hyper Safety의 새로운 표준", 1.0, 4.0, 11.2, 0.6, size=30, color="ink", bold=True, align="center")
    add_text(s, "AI-RAN · 로키드 스마트글라스 · AI 로봇개 · SAFE-LINK", 1.6, 5.15, 10.2, 0.42, size=18, color="slate", align="center")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def font(size=28, bold=False):
    candidates = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for c in candidates:
        if Path(c).exists():
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()


def draw_wrapped(draw, text, xy, max_w, fnt, fill, line_gap=8, align="left"):
    words = []
    for part in text.split("\n"):
        if words:
            words.append("\n")
        buf = ""
        for ch in part:
            if draw.textlength(buf + ch, font=fnt) <= max_w:
                buf += ch
            else:
                if buf:
                    words.append(buf)
                buf = ch
        if buf:
            words.append(buf)
    x, y = xy
    for line in words:
        if line == "\n":
            y += fnt.size + line_gap
            continue
        tw = draw.textlength(line, font=fnt)
        dx = {"left": 0, "center": (max_w - tw) / 2, "right": max_w - tw}[align]
        draw.text((x + dx, y), line, font=fnt, fill=fill)
        y += fnt.size + line_gap


def create_preview():
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    slides = [
        ("SAFE-LINK\nUNDERGROUND", "AI-RAN 기반 철근콘크리트 지하 골조 Hyper Safety 플랫폼"),
        ("가장 위험한 곳이\n가장 연결되지 않습니다", "통신 음영 · 공정 위험 · 외국인 근로자"),
        ("대우건설 기존 체계의\n지하 사각을 완성", "지상 관제 + 지하 작업면 실행"),
        ("철근콘크리트 현장성이\n제안의 방어력", "배근 · 거푸집 · 동바리 · 타설 · PIT"),
        ("4개 기술을\n하나의 안전 루프로", "AI-RAN · 로키드 글라스 · 로봇개 · SAFE-LINK"),
        ("AI-RAN 안전 통신 버블", "경희대 홍인기 교수 MOU 및 공동제안 추진"),
        ("로키드 스마트글라스", "열화상 · 영상 · 안드로이드 기반 착용형 관제"),
        ("AI 로봇개", "지하 선행 순찰과 이동형 위험탐지 노드"),
        ("Rebar-Zone\nSafety Passport", "작업구역별 안전허가 상태 관리"),
        ("Pour-Go / No-Go AI", "콘크리트 타설 전 승인 게이트"),
        ("Multilingual Stop-Work", "외국어 위험신고를 관리자 조치와 증빙으로 연결"),
        ("Robot · Glass · SAFE-LINK", "삼각 데이터로 Evidence-First Safety Twin 구성"),
        ("8주 PoC", "조사 → 운영 → 관제/순찰 → 효과 측정"),
        ("측정 지표", "참여율 · 이해도 · 처리시간 · 커버리지 · 보고절감"),
        ("협업 구조", "서원토건 · 대우건설 · 경희대 · 로봇개 파트너 · 로키드"),
        ("지하 안전 사각을\n작업면 데이터로 메웁니다", "철근콘크리트 지하 골조 Hyper Safety의 새로운 표준"),
    ]
    thumbs = []
    for idx, (head, sub) in enumerate(slides, 1):
        img = Image.new("RGB", (PX_W, PX_H), "#" + COLORS["bg"])
        d = ImageDraw.Draw(img)
        d.rectangle([0, 0, PX_W, PX_H], fill="#" + COLORS["bg"])
        d.text((90, 70), "SAFE-LINK UNDERGROUND", font=font(24, True), fill="#" + COLORS["blue"])
        draw_wrapped(d, head, (90, 165), 1100, font(62, True), "#" + COLORS["ink"], line_gap=12)
        draw_wrapped(d, sub, (94, 390), 1120, font(30), "#" + COLORS["slate"], line_gap=10)
        d.line([90, 515, 1370, 515], fill="#" + COLORS["line"], width=3)
        colors = ["blue", "green", "purple", "yellow"]
        for i, c in enumerate(colors):
            x = 110 + i * 330
            d.rounded_rectangle([x, 575, x + 260, 710], radius=24, fill="#FFFFFF", outline="#" + COLORS["line"], width=2)
            d.rectangle([x, 575, x + 12, 710], fill="#" + COLORS[c])
        d.text((90, 820), "서원토건 | 대우건설 Hyper Safety & AI Open Innovation", font=font(20), fill="#" + COLORS["muted"])
        d.text((1450, 820), f"{idx:02d}", font=font(20), fill="#" + COLORS["muted"])
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        img.save(out)
        thumbs.append(out)

    sheet = Image.new("RGB", (PX_W, 4 * 260), "white")
    for i, p in enumerate(thumbs):
        im = Image.open(p).resize((360, 203))
        x = 30 + (i % 4) * 390
        y = 25 + (i // 4) * 250
        sheet.paste(im, (x, y))
        ImageDraw.Draw(sheet).text((x, y + 208), f"{i+1:02d}", font=font(18, True), fill="#" + COLORS["ink"])
    sheet.save(PREVIEW_DIR / "contact_sheet.png")


if __name__ == "__main__":
    create_deck()
    create_preview()
    print(PPTX_PATH)
    print(PREVIEW_DIR / "contact_sheet.png")
