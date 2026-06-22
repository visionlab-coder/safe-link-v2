from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE


ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "docs" / "generated" / "daewoo-drone-patrol-20260606"
ASSET = BASE / "assets_micro_drone_nest"
OUT = BASE / "SAFE-LINK_Micro_Drone_Nest_대우건설_정정본_20260606.pptx"

W, H = 13.333, 7.5
C = {
    "paper": "FAFAF8",
    "ink": "101828",
    "muted": "667085",
    "line": "D0D5DD",
    "blue": "175CD3",
    "green": "027A48",
    "orange": "B54708",
    "red": "B42318",
    "white": "FFFFFF",
    "soft": "F2F4F7",
    "navy": "0B1220",
    "note": "FFF7D6",
}


def rgb(key):
    v = C.get(key, key).strip("#")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, transparency=0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.05
        except Exception:
            pass
    return shp


def txt(slide, body, x, y, w, h, size=18, color="ink", bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = body
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def rule(slide, x, y, w, color="line"):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(color)
    s.line.fill.background()


def footer(slide, n):
    txt(slide, "SAFE-LINK Micro Drone Nest / 서원토건", 0.55, 7.04, 4.6, 0.24, 10, "muted")
    txt(slide, f"{n:02d}", 12.25, 7.04, 0.45, 0.24, 10, "muted", align="right")


def header(slide, label, title, sub, n):
    txt(slide, label, 0.62, 0.34, 5.6, 0.3, 12, "blue", True)
    txt(slide, title, 0.62, 0.74, 11.7, 0.72, 29, "ink", True)
    if sub:
        txt(slide, sub, 0.64, 1.5, 11.2, 0.38, 15, "muted")
    rule(slide, 0.62, 2.02, 12.05)
    footer(slide, n)


def image_or_placeholder(slide, filename, x, y, w, h, title):
    p = ASSET / filename
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        rect(slide, x, y, w, h, "soft", "line", True)
        txt(slide, title, x + 0.25, y + h / 2 - 0.42, w - 0.5, 0.36, 18, "ink", True, "center")
        txt(slide, "이미지 2.0 실사풍 컷 삽입 위치", x + 0.25, y + h / 2 + 0.04, w - 0.5, 0.28, 13, "muted", align="center")


def note(slide, head, body, x, y, w, h, color="blue"):
    rect(slide, x, y, w, h, "white", "line", True)
    rect(slide, x, y, 0.08, h, color, color, False)
    txt(slide, head, x + 0.22, y + 0.17, w - 0.42, 0.32, 15, color, True)
    txt(slide, body, x + 0.22, y + 0.58, w - 0.42, h - 0.72, 15, "ink")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        rect(s, 0, 0, W, H, "paper", "paper", False)
        return s

    # 1
    s = slide()
    image_or_placeholder(s, "01-micro-drone-nest-hero.png", 6.55, 0, 6.78, 7.5, "초소형 드론 네스트")
    txt(s, "DAEWOO HYPER SAFETY & AI", 0.72, 0.7, 4.8, 0.3, 12, "blue", True)
    txt(s, "SAFE-LINK\nMicro Drone Nest", 0.72, 1.35, 5.2, 1.3, 35, "ink", True)
    txt(s, "손바닥보다 작은 초소형 드론이 사각공간을 정찰하고 스스로 둥지로 복귀·충전하는 분산형 안전관제 플랫폼", 0.74, 3.05, 5.15, 1.0, 18, "ink", True)
    rule(s, 0.74, 4.35, 4.35, "blue")
    txt(s, "초소형 드론 · 분산 네스트 · 자동충전 · 실시간 관제 · SAFE-LINK 조치보고", 0.74, 4.66, 5.1, 0.68, 15, "muted")
    txt(s, "서원토건 / 2026.06.06", 0.74, 6.78, 3.2, 0.25, 11, "muted")

    # 2
    s = slide()
    header(s, "01 CORE CORRECTION", "핵심은 대형 드론 순찰이 아닙니다", "건설현장 곳곳의 작은 둥지에서 초소형 드론들이 자동으로 나가고 돌아오는 구조입니다.", 2)
    note(s, "크기", "손바닥보다 작은 초소형 드론\n보호가드 또는 덕트형 프로펠러", 0.9, 2.65, 3.5, 1.4, "blue")
    note(s, "설치", "샤프트, PIT, 천장 점검구, 구조물 틈 주변에 소형 네스트 설치", 4.9, 2.65, 3.5, 1.4, "green")
    note(s, "운용", "자동 이륙·촬영·복귀·충전\n관제센터 원격 명령 가능", 8.9, 2.65, 3.5, 1.4, "orange")
    txt(s, "이 제안은 ‘드론을 날리는 것’이 아니라\n현장 사각공간에 작은 안전 감각기관을 심는 것입니다.", 1.1, 5.1, 11.0, 0.75, 25, "ink", True, "center")

    # 3
    s = slide()
    header(s, "02 BLIND SPOTS", "사람 손이 닿지 않는 공간이 진짜 타깃입니다", "CCTV도, 일반 순찰도, 대형 드론도 접근하기 어려운 좁은 공간을 맡습니다.", 3)
    image_or_placeholder(s, "02-hard-to-reach-space.png", 0.75, 2.35, 5.75, 3.6, "사각공간 정찰")
    note(s, "대상 공간", "엘리베이터 샤프트\n지하 PIT·집수정\n천장 내부·배관 주변\n거푸집·동바리 틈\n개구부 하부", 7.0, 2.45, 4.9, 2.3, "blue")
    txt(s, "넓은 상공보다 좁고 위험한 내부 공간을 반복 확인하는 것이 차별점입니다.", 7.0, 5.25, 4.9, 0.55, 18, "green", True, "center")

    # 4
    s = slide()
    header(s, "03 NEST NETWORK", "네스트는 드론의 보관함이 아니라 현장 센서 거점입니다", "각 네스트는 드론 상태, 배터리, 통신, 촬영 스케줄을 관제센터와 연결합니다.", 4)
    image_or_placeholder(s, "03-distributed-nests.png", 0.75, 2.35, 5.9, 3.55, "분산형 네스트")
    note(s, "소형 설치", "기둥, 벽면, 점검구, 난간 주변에 설치 가능한 작은 둥지", 7.05, 2.42, 4.85, 0.95, "blue")
    note(s, "자동 복귀", "정찰 후 네스트로 복귀해 접점식 또는 무선 방식 충전", 7.05, 3.62, 4.85, 0.95, "green")
    note(s, "상태 관제", "네스트별 드론 위치, 배터리, 통신, 영상 상태 확인", 7.05, 4.82, 4.85, 0.95, "orange")

    # 5
    s = slide()
    header(s, "04 OPERATING LOOP", "자동순찰과 관제센터 명령이 함께 작동합니다", "정기 순찰은 자동으로, 특이상황은 관제센터 명령으로 대응합니다.", 5)
    steps = [("스케줄", "정기 순찰"), ("이륙", "네스트 출발"), ("촬영", "근접 확인"), ("전송", "실시간 관제"), ("복귀", "자동충전")]
    for i, (a, b) in enumerate(steps):
        x = 0.72 + i * 2.45
        rect(s, x, 3.05, 1.85, 1.25, "white", ["blue", "green", "orange", "red", "blue"][i], True)
        txt(s, a, x + 0.12, 3.25, 1.6, 0.3, 16, ["blue", "green", "orange", "red", "blue"][i], True, "center")
        txt(s, b, x + 0.12, 3.73, 1.6, 0.26, 13, "ink", True, "center")
        if i < 4:
            txt(s, "→", x + 1.88, 3.43, 0.35, 0.3, 22, "muted", True, "center")
    txt(s, "관제센터는 특정 네스트의 드론을 즉시 호출해 추가 촬영·재정찰을 명령할 수 있습니다.", 1.0, 5.25, 11.1, 0.45, 20, "ink", True, "center")

    # 6
    s = slide()
    header(s, "05 CONTROL CENTER", "관제센터는 작은 드론들의 실시간 눈을 갖게 됩니다", "현장 여러 곳의 네스트 영상이 SAFE-LINK 관제 화면으로 모입니다.", 6)
    image_or_placeholder(s, "04-control-center-micro-feeds.png", 0.75, 2.35, 5.9, 3.55, "관제센터")
    note(s, "실시간 영상", "네스트별 카메라 피드와 위험 후보 확인", 7.05, 2.42, 4.85, 0.95, "blue")
    note(s, "원격 명령", "재촬영, 특정 구역 이동, 긴급 복귀 명령", 7.05, 3.62, 4.85, 0.95, "green")
    note(s, "SAFE-LINK 연동", "알림, 조치, 작업중지, 보고서 자동 연결", 7.05, 4.82, 4.85, 0.95, "orange")

    # 7
    s = slide()
    header(s, "06 POC", "1차 PoC는 3~5개 사각공간이면 충분합니다", "처음부터 현장 전체가 아니라 사람이 보기 어려운 핵심 사각공간을 검증합니다.", 7)
    note(s, "장비", "초소형 드론 2~5대\n소형 네스트 3~5개", 0.9, 2.65, 3.5, 1.4, "blue")
    note(s, "운영", "일 3회 자동 순찰\n관제센터 수동 호출", 4.9, 2.65, 3.5, 1.4, "green")
    note(s, "검증", "자동복귀·충전\n실시간 영상\n위험후보 태깅", 8.9, 2.65, 3.5, 1.4, "orange")
    txt(s, "PoC 대상: 샤프트, PIT, 점검구, 동바리 틈, 개구부 하부 등 3~5개 위치", 1.0, 5.25, 11.2, 0.42, 19, "ink", True, "center")

    # 8
    s = slide()
    header(s, "07 VALUE", "대우건설이 얻는 것은 사각공간의 상시 감각입니다", "고정 CCTV와 일반 드론 사이의 빈 공간을 초소형 드론 네스트가 메웁니다.", 8)
    image_or_placeholder(s, "05-safelink-action-report.png", 0.75, 2.35, 5.75, 3.55, "SAFE-LINK 조치보고")
    note(s, "안전", "작업자 접근 전 위험 확인\n좁은 공간 반복 점검", 7.0, 2.42, 4.9, 0.95, "blue")
    note(s, "효율", "사람이 직접 들어가기 어려운 공간의 확인 시간 절감", 7.0, 3.62, 4.9, 0.95, "green")
    note(s, "데이터", "위치·시간·영상·조치 이력을 SAFE-LINK에 누적", 7.0, 4.82, 4.9, 0.95, "orange")
    txt(s, "초소형 드론 네스트는 건설현장에 ‘작은 눈’을 여러 개 심는 방식입니다.", 1.0, 6.23, 11.2, 0.4, 22, "ink", True, "center")

    BASE.mkdir(parents=True, exist_ok=True)
    ASSET.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
