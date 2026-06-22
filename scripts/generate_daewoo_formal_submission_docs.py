from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZipFile, ZIP_DEFLATED
import xml.etree.ElementTree as ET

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt as PPTPt
from pptx.dml.color import RGBColor as PPTRGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
OUT.mkdir(parents=True, exist_ok=True)

DOCX_PATH = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_20260611.docx"
PPTX_PATH = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_20260611.pptx"

NAVY = RGBColor(31, 58, 88)
BLUE = RGBColor(48, 88, 128)
GRAY = RGBColor(89, 89, 89)
LIGHT_BLUE = "EAF1F8"
LIGHT_GRAY = "F2F4F7"
LINE = "B8C4D0"
WHITE = "FFFFFF"

PPT_NAVY = PPTRGBColor(31, 58, 88)
PPT_BLUE = PPTRGBColor(48, 88, 128)
PPT_GRAY = PPTRGBColor(89, 89, 89)
PPT_LIGHT_BLUE = PPTRGBColor(234, 241, 248)
PPT_LIGHT_GRAY = PPTRGBColor(242, 244, 247)
PPT_LINE = PPTRGBColor(184, 196, 208)
PPT_WHITE = PPTRGBColor(255, 255, 255)


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_border(cell, color=LINE, size="6"):
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        tag = "w:{}".format(edge)
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), size)
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), color)


def set_cell_text(cell, text, bold=False, size=9, color=None, align=None):
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    p = cell.paragraphs[0]
    p.alignment = align or WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.line_spacing = 1.15
    run = p.add_run(text)
    run.bold = bold
    run.font.name = "Malgun Gothic"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def add_docx_heading(doc, text, level=1):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(10 if level == 1 else 6)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    r.bold = True
    r.font.name = "Malgun Gothic"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    r.font.size = Pt(14 if level == 1 else 11)
    r.font.color.rgb = NAVY if level == 1 else BLUE
    return p


def add_docx_para(doc, text, bold=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.line_spacing = 1.18
    r = p.add_run(text)
    r.bold = bold
    r.font.name = "Malgun Gothic"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    r.font.size = Pt(9.5)
    r.font.color.rgb = RGBColor(35, 35, 35)
    return p


def add_docx_table(doc, headers, rows, widths=None, font_size=8.5):
    table = doc.add_table(rows=1, cols=len(headers))
    table.autofit = False
    for idx, h in enumerate(headers):
        cell = table.rows[0].cells[idx]
        if widths:
            cell.width = Cm(widths[idx])
        set_cell_shading(cell, LIGHT_BLUE)
        set_cell_border(cell)
        set_cell_text(cell, h, bold=True, size=font_size, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER)
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            if widths:
                cells[idx].width = Cm(widths[idx])
            set_cell_border(cells[idx])
            set_cell_text(cells[idx], value, size=font_size, align=WD_ALIGN_PARAGRAPH.CENTER if idx == 0 else WD_ALIGN_PARAGRAPH.LEFT)
    doc.add_paragraph().paragraph_format.space_after = Pt(2)
    return table


def build_docx():
    doc = Document()
    section = doc.sections[0]
    section.page_width = Cm(21.0)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(1.7)
    section.bottom_margin = Cm(1.6)
    section.left_margin = Cm(1.7)
    section.right_margin = Cm(1.7)

    styles = doc.styles
    styles["Normal"].font.name = "Malgun Gothic"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    styles["Normal"].font.size = Pt(9.5)

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    hr = header.add_run("대우건설 Hyper Safety & AI Open Innovation 제출본 | 서원토건")
    hr.font.name = "Malgun Gothic"
    hr._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    hr.font.size = Pt(8)
    hr.font.color.rgb = GRAY

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer.add_run("CONFIDENTIAL - PoC 및 공동실증 협의용")
    fr.font.name = "Malgun Gothic"
    fr._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    fr.font.size = Pt(8)
    fr.font.color.rgb = GRAY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(70)
    r = p.add_run("대우건설 Hyper Safety & AI Open Innovation\n공문서형 제출 제안서")
    r.bold = True
    r.font.name = "Malgun Gothic"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    r.font.size = Pt(22)
    r.font.color.rgb = NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("SQ-LINK Underground\n지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼")
    r.font.name = "Malgun Gothic"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    r.font.size = Pt(12)
    r.font.color.rgb = BLUE

    table = doc.add_table(rows=5, cols=2)
    table.autofit = False
    meta = [
        ("문서번호", "SW-DW-HSAI-20260611-01"),
        ("작성일", "2026. 06. 11."),
        ("제안기관", "서원토건"),
        ("제안구분", "현장 PoC 및 공동 실증 파트너십 제안"),
        ("문서성격", "공모전 제출용 제안서 / 대외 협의용"),
    ]
    for i, (k, v) in enumerate(meta):
        for j, text in enumerate((k, v)):
            cell = table.rows[i].cells[j]
            cell.width = Cm(5 if j == 0 else 11)
            set_cell_border(cell)
            set_cell_shading(cell, LIGHT_BLUE if j == 0 else WHITE)
            set_cell_text(cell, text, bold=(j == 0), size=9, color=NAVY if j == 0 else None, align=WD_ALIGN_PARAGRAPH.CENTER if j == 0 else WD_ALIGN_PARAGRAPH.LEFT)

    doc.add_page_break()

    add_docx_heading(doc, "1. 제안 개요")
    add_docx_para(
        doc,
        "본 제안은 지하층, PIT, 코어, 지하주차장, 기계실 및 골조 구간에서 반복되는 통신·위치·작업데이터 단절 문제를 "
        "대우건설 현장에서 PoC로 검증하기 위한 공모전 제출 문서이다. 제안의 핵심은 장비 도입이 아니라 장비가 실제 현장에서 "
        "작동할 수 있도록 통신, 3D 좌표, BIM, SAFE-LINK 작업데이터를 하나의 안전·품질 AI 운영체계로 연결하는 것이다.",
    )
    add_docx_table(
        doc,
        ["구분", "내용"],
        [
            ("제안명", "SQ-LINK Underground"),
            ("목표", "지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 PoC"),
            ("대상", "대우건설 지하층 또는 지하주차장 철근콘크리트 작업구역 1개소"),
            ("기간", "8주"),
            ("산출물", "현장 진단표, 작업데이터 샘플, 영상·사진 증거, PoC 결과보고서, 후속과제 제안"),
        ],
        widths=[4, 12],
    )

    add_docx_heading(doc, "2. 현장 문제 정의")
    add_docx_table(
        doc,
        ["영역", "현장 문제", "운영 리스크"],
        [
            ("통신", "지하층·PIT·코어에서 Wi-Fi/LTE/GPS 불안정", "영상 관제, 위치 확인, 전자증거 수집 단절"),
            ("위치", "평면 좌표만으로 층·높이·깊이 구분 곤란", "작업자·장비 위치와 BIM/도면 매칭 실패"),
            ("데이터", "안전, 품질, 공정, 작업허가 데이터 분리", "사고·하자·작업중지 원인 추적 어려움"),
            ("장비", "로봇·드론·스마트글라스가 단독 장비로 운영", "촬영은 가능하나 판단·보고 자동화 미흡"),
            ("실증", "장비 시연 중심 PoC에 치우침", "실무자 수용성과 지속 운영성 검증 부족"),
        ],
        widths=[2.3, 6.5, 7.2],
    )

    add_docx_heading(doc, "3. 제안 솔루션 체계")
    add_docx_para(doc, "SQ-LINK Underground는 SAFE-LINK를 안전앱에서 안전·품질 AI 운영 플랫폼으로 확장하는 구조이다.")
    add_docx_table(
        doc,
        ["계층", "구성", "역할"],
        [
            ("Layer 1", "AI-RAN/O-RAN 검토, 이음5G 특화망, UWB/BLE/Wi-Fi RTT, 기준점 비콘", "지하 통신 음영 측정 및 3D 위치 인식 가능성 검증"),
            ("Layer 2", "SAFE-LINK, QR/NFC, TBM, 위험성평가, 안전서약, 작업중지", "근로자·작업조·작업구역·조치 이력의 전자증거화"),
            ("Layer 3", "스마트글라스, 모바일 영상, 사진 증거", "원격 관제, 품질검측, 위험구역 확인"),
            ("Layer 4", "로봇개, 드론, 이동형 카메라", "사람이 접근하기 어려운 구역의 선행 촬영 및 점검"),
            ("Layer 5", "BIM/도면, 작업데이터, 통신상태, 품질 체크리스트, 공정·날씨·법규", "점검 대상 추천, 보고서 자동화, 안전·품질 판단 보조"),
        ],
        widths=[2.2, 7.0, 6.8],
        font_size=8,
    )

    add_docx_heading(doc, "4. 중기부 R&D와 대우건설 공모전 역할 분담")
    add_docx_table(
        doc,
        ["구분", "중기부 R&D", "대우건설 공모전"],
        [
            ("목적", "지하 특화망·로봇·웨어러블 통합 제어 원천기술 개발", "대우건설 현장에서 안전·품질 PoC 및 실증 검증"),
            ("중심", "통신 성능, 로봇 제어, 하이브리드 통신 모듈", "현장 업무 적용성, 실무 수용성, 안전·품질 데이터 연결"),
            ("산출물", "기술개발 결과, 성능시험, TRL 향상", "PoC 결과보고서, 현장 피드백, 후속 공동과제 기획"),
            ("리스크 관리", "국가 R&D 중복지원 규정 확인", "현장 실증 범위로 한정하고 원천기술 개발비와 구분"),
        ],
        widths=[3, 6.5, 6.5],
    )

    add_docx_heading(doc, "5. PoC 추진 계획")
    add_docx_table(
        doc,
        ["단계", "기간", "주요 내용", "산출물"],
        [
            ("1단계", "1~2주", "지하 통신 음영, 작업구역, 위험·품질 체크리스트 조사", "현장 진단표, PoC 상세계획"),
            ("2단계", "3~4주", "QR/NFC, TBM, 위험성평가, 작업중지, 품질 체크 항목 연동", "작업데이터 샘플, 체크리스트"),
            ("3단계", "5~6주", "스마트글라스 영상, 로봇개/드론 선행점검 시나리오 검증", "사진·영상 증거, 운영 로그"),
            ("4단계", "7~8주", "실무자 피드백, 성과지표 측정, 후속과제 도출", "PoC 결과보고서, 개선안"),
        ],
        widths=[2.3, 2.4, 7.0, 4.3],
    )

    add_docx_heading(doc, "6. 정량 성과지표")
    add_docx_table(
        doc,
        ["지표", "PoC 목표", "측정 방법"],
        [
            ("통신 음영 지도화", "PoC 구역 내 음영 구간 식별 완료", "현장 측정 로그 및 위치별 수신 상태 기록"),
            ("작업구역 식별", "층·구역 단위 식별 가능성 확인", "QR/NFC, 기준점, 작업구역 매핑 비교"),
            ("작업데이터 생성률", "대상 작업의 80% 이상 데이터화", "TBM, 위험성평가, 체크리스트 등록 건수"),
            ("스마트글라스 활용성", "품질·안전 확인 시나리오 3건 이상 수행", "영상 캡처, 관리자 피드백"),
            ("로봇/드론 시나리오", "선행점검 후보 구역 2건 이상 검토", "이동 가능성, 촬영 가능성, 위험요인 기록"),
            ("실무자 수용성", "안전·공사·품질 담당자 피드백 확보", "인터뷰 및 설문"),
        ],
        widths=[4, 5.5, 6.5],
        font_size=8.2,
    )

    add_docx_heading(doc, "7. TRL 및 기술개발 목표")
    add_docx_table(
        doc,
        ["항목", "착수 수준", "PoC 종료 목표"],
        [
            ("SAFE-LINK 작업데이터", "현장 기능 시연 가능", "대우 PoC 구역 작업데이터 샘플 생성"),
            ("지하 통신·좌표", "후보 기술 검토 단계", "음영지도 및 3D 위치 인식 가능성 검증"),
            ("스마트글라스", "장비 적용 검토", "안전·품질 영상 관제 시나리오 검증"),
            ("로봇/드론", "적용 시나리오 검토", "선행점검 후보 구역과 운용 조건 도출"),
            ("통합 보고서", "SAFE-LINK 보고 기능 기반", "안전+품질+위치+증거 통합 보고서 초안"),
        ],
        widths=[4, 5.5, 6.5],
    )
    add_docx_para(doc, "TRL 기준으로는 현재 TRL 4~5 수준의 구성요소를 대우 현장 PoC를 통해 TRL 5~6 수준의 현장 검증 단계로 끌어올리는 것을 목표로 한다.")

    add_docx_heading(doc, "8. 컨소시엄 및 역할")
    add_docx_table(
        doc,
        ["주체", "역할"],
        [
            ("서원토건", "철근콘크리트 공정 지식, SAFE-LINK 운영, PoC 시나리오, 작업데이터 설계"),
            ("대우건설", "PoC 현장 제공, 안전·공사·품질 실무 검증, 기존 체계와 비교 평가"),
            ("경희대 AI-RAN/O-RAN 연구팀", "지하 통신, AI-RAN/O-RAN, 이음5G 특화망 자문 및 후속 공동과제 검토"),
            ("로봇·드론 협력사", "상용 장비 또는 협업 장비 기반 선행점검 시나리오 검토"),
            ("스마트글라스 협력사", "현장 영상 관제, 품질검측 보조, 착용형 디바이스 적용성 검토"),
        ],
        widths=[5, 11],
    )

    add_docx_heading(doc, "9. 지식재산권·데이터·보안 관리")
    add_docx_table(
        doc,
        ["구분", "관리 방안"],
        [
            ("IP 방향", "검증 이벤트 기반 작업증거 데이터셋, BIM/도면 위치와 작업증거 매핑, 안전+품질 체크 자동화 중심 권리화 검토"),
            ("데이터 원칙", "대우건설 현장 데이터는 PoC 목적 범위에서만 사용하며, 민감 정보와 도면은 접근 권한을 분리"),
            ("보안", "자료 열람·수정·반출 이력 관리, 외부 발표 전 대우건설 검토, 원본 데이터 보관·폐기 범위 협약"),
            ("공동성과", "PoC 결과물, 공동특허, 후속과제 자료의 권리 범위를 사전 합의"),
        ],
        widths=[4, 12],
    )

    add_docx_heading(doc, "10. 대우건설 협력 요청 사항")
    add_docx_table(
        doc,
        ["번호", "요청 사항"],
        [
            ("1", "지하층 또는 골조 작업구역 1개소 PoC 현장 제공"),
            ("2", "안전관리자, 공사팀, 품질팀 담당자 인터뷰 및 피드백 협조"),
            ("3", "현장 통신 음영 측정과 장비 운용 시간 협의"),
            ("4", "기존 스마트안전·품질관리 체계와 비교 검증"),
            ("5", "PoC 결과보고서 공동 검토"),
            ("6", "성과 확인 시 후속 공동과제, 공동특허, 확대 실증 검토"),
        ],
        widths=[2, 14],
    )

    add_docx_heading(doc, "11. 제출 전 확정 필요 항목")
    add_docx_table(
        doc,
        ["항목", "현재 상태", "보강 방향"],
        [
            ("PoC 예산", "미확정", "8주 기준 인력, 장비, 스마트글라스, 로봇/드론 대여, 통신 측정비로 분리 산정"),
            ("현장 범위", "미확정", "지하층, 지하주차장, PIT, 코어, 기계실 중 1개 구역으로 한정"),
            ("협력사", "후보 단계", "스마트글라스, 로봇/드론, 통신 자문 기관을 역할별로 명시"),
            ("정량 성과", "PoC 지표 제시", "작업데이터 생성률, 체크리스트 누락 감소, 처리시간 단축 등 현장 지표 확정"),
            ("중복지원 리스크", "관리 필요", "중기부 R&D는 원천기술, 대우 공모는 현장 실증으로 문구 분리"),
        ],
        widths=[3.3, 3.7, 9],
        font_size=8,
    )

    add_docx_heading(doc, "12. 최종 제안 문장")
    add_docx_para(
        doc,
        "지하·골조 현장 자동화의 병목은 장비 부족이 아니라 통신, 3D 좌표, 작업데이터의 단절입니다. "
        "서원토건은 SAFE-LINK에서 출발해 안전과 품질을 연결하는 SQ-LINK Underground 현장 AI 운영체계로 확장하고, "
        "대우건설과 함께 그 가능성을 실제 현장에서 PoC로 검증하고자 합니다.",
        bold=True,
    )

    doc.save(DOCX_PATH)


def ppt_textbox(slide, x, y, w, h, text, size=12, bold=False, color=PPT_NAVY, align=PP_ALIGN.LEFT, fill=None, line=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = PPTPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def ppt_header_footer(slide, page, title="대우건설 Hyper Safety & AI Open Innovation 제출본"):
    left = 0.45
    top = 0.25
    width = 10.8
    slide.shapes.add_shape(1, Inches(left), Inches(0.55), Inches(width), Inches(0.01)).line.color.rgb = PPT_LINE
    ppt_textbox(slide, left, top, 6.7, 0.25, title, 8.5, color=PPT_GRAY)
    ppt_textbox(slide, 8.9, top, 2.35, 0.25, "CONFIDENTIAL | 서원토건", 8.5, color=PPT_GRAY, align=PP_ALIGN.RIGHT)
    slide.shapes.add_shape(1, Inches(left), Inches(7.65), Inches(width), Inches(0.01)).line.color.rgb = PPT_LINE
    ppt_textbox(slide, left, 7.72, 4.0, 0.2, "PoC 및 공동실증 협의용", 7.2, color=PPT_GRAY)
    ppt_textbox(slide, 10.55, 7.72, 0.7, 0.2, f"{page:02d}", 7.2, color=PPT_GRAY, align=PP_ALIGN.RIGHT)


def ppt_title(slide, text):
    ppt_textbox(slide, 0.6, 0.78, 10.5, 0.38, text, size=16, bold=True, color=PPT_NAVY)


def ppt_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=8.2):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = Inches(w * cw / total)
    for r_idx in range(len(rows) + 1):
        for c_idx in range(len(headers)):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPT_LIGHT_BLUE if r_idx == 0 else PPT_WHITE
            cell.text = headers[c_idx] if r_idx == 0 else rows[r_idx - 1][c_idx]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c_idx == 0 or r_idx == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Malgun Gothic"
                    run.font.size = PPTPt(font_size)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = PPT_NAVY if r_idx == 0 else PPTRGBColor(35, 35, 35)
    return table_shape


def ppt_bullets(slide, x, y, w, h, items, size=9):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "· " + item
        p.font.name = "Malgun Gothic"
        p.font.size = PPTPt(size)
        p.font.color.rgb = PPTRGBColor(45, 45, 45)
        p.space_after = PPTPt(3)
    return shape


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)
    blank = prs.slide_layouts[6]

    # Cover: formal document cover, not a presentation title slide.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPT_WHITE
    slide.shapes.add_shape(1, Inches(0.6), Inches(0.85), Inches(10.45), Inches(0.02)).line.color.rgb = PPT_NAVY
    ppt_textbox(slide, 0.65, 1.15, 10.1, 0.35, "공모전 제출 제안서", 12, color=PPT_GRAY, align=PP_ALIGN.CENTER)
    ppt_textbox(slide, 1.0, 1.85, 9.65, 0.95, "대우건설 Hyper Safety & AI\nOpen Innovation", 28, bold=True, color=PPT_NAVY, align=PP_ALIGN.CENTER)
    ppt_textbox(slide, 1.2, 3.05, 9.25, 0.7, "SQ-LINK Underground\n지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼", 14, color=PPT_BLUE, align=PP_ALIGN.CENTER)
    ppt_table(
        slide,
        2.0,
        4.55,
        7.7,
        1.45,
        ["항목", "내용"],
        [
            ("문서번호", "SW-DW-HSAI-20260611-01"),
            ("작성일", "2026. 06. 11."),
            ("제안기관", "서원토건"),
            ("제안구분", "현장 PoC 및 공동 실증 파트너십 제안"),
        ],
        col_widths=[2.1, 5.6],
        font_size=8.8,
    )
    ppt_textbox(slide, 0.65, 7.35, 10.4, 0.25, "CONFIDENTIAL - PoC 및 공동실증 협의용", 8.5, color=PPT_GRAY, align=PP_ALIGN.CENTER)

    slide_specs = [
        (
            "1. 제안 개요",
            "본 제안은 지하·골조 현장의 통신·위치·작업데이터 단절을 대우건설 현장에서 PoC로 검증하기 위한 공모전 제출 문서이다.",
            ["구분", "내용"],
            [
                ("제안명", "SQ-LINK Underground"),
                ("목표", "지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 PoC"),
                ("대상", "대우건설 지하층 또는 지하주차장 철근콘크리트 작업구역 1개소"),
                ("기간", "8주"),
                ("산출물", "현장 진단표, 작업데이터 샘플, 영상·사진 증거, PoC 결과보고서"),
            ],
            [2.5, 8.0],
        ),
        (
            "2. 현장 문제 정의",
            "장비는 늘었지만 지하·골조 구간에서는 연결 기반이 부족해 현장 자동화가 끊긴다.",
            ["영역", "현장 문제", "운영 리스크"],
            [
                ("통신", "지하층·PIT·코어에서 Wi-Fi/LTE/GPS 불안정", "영상 관제, 위치 확인, 전자증거 수집 단절"),
                ("위치", "평면 좌표만으로 층·높이·깊이 구분 곤란", "작업자·장비 위치와 BIM/도면 매칭 실패"),
                ("데이터", "안전, 품질, 공정, 작업허가 데이터 분리", "사고·하자·작업중지 원인 추적 어려움"),
                ("장비", "로봇·드론·스마트글라스가 단독 장비로 운영", "촬영은 가능하나 판단·보고 자동화 미흡"),
                ("실증", "장비 시연 중심 PoC에 치우침", "실무자 수용성과 지속 운영성 검증 부족"),
            ],
            [1.5, 4.8, 4.7],
        ),
        (
            "3. 제안 솔루션 체계",
            "SAFE-LINK를 안전앱에서 안전·품질 AI 운영 플랫폼으로 확장한다.",
            ["계층", "구성", "역할"],
            [
                ("Layer 1", "AI-RAN/O-RAN 검토, 이음5G 특화망, UWB/BLE/Wi-Fi RTT, 기준점 비콘", "지하 통신 음영 측정 및 3D 위치 인식 가능성 검증"),
                ("Layer 2", "SAFE-LINK, QR/NFC, TBM, 위험성평가, 안전서약, 작업중지", "근로자·작업조·작업구역·조치 이력의 전자증거화"),
                ("Layer 3", "스마트글라스, 모바일 영상, 사진 증거", "원격 관제, 품질검측, 위험구역 확인"),
                ("Layer 4", "로봇개, 드론, 이동형 카메라", "접근 곤란 구역의 선행 촬영 및 점검"),
                ("Layer 5", "BIM/도면, 작업데이터, 통신상태, 품질 체크리스트, 공정·날씨·법규", "점검 대상 추천, 보고서 자동화, 안전·품질 판단 보조"),
            ],
            [1.5, 5.0, 4.5],
        ),
        (
            "4. R&D와 공모전 역할 분담",
            "중기부 R&D는 원천기술 개발, 대우건설 공모전은 현장 PoC·실증으로 명확히 분리한다.",
            ["구분", "중기부 R&D", "대우건설 공모전"],
            [
                ("목적", "지하 특화망·로봇·웨어러블 통합 제어 원천기술 개발", "대우건설 현장에서 안전·품질 PoC 및 실증 검증"),
                ("중심", "통신 성능, 로봇 제어, 하이브리드 통신 모듈", "현장 업무 적용성, 실무 수용성, 안전·품질 데이터 연결"),
                ("산출물", "기술개발 결과, 성능시험, TRL 향상", "PoC 결과보고서, 현장 피드백, 후속 공동과제 기획"),
                ("리스크", "국가 R&D 중복지원 규정 확인", "현장 실증 범위로 한정하고 원천기술 개발비와 구분"),
            ],
            [1.7, 4.65, 4.65],
        ),
        (
            "5. PoC 추진 계획",
            "8주 동안 진단, 데이터 연결, 장비 시나리오, 보고·평가를 단계적으로 수행한다.",
            ["단계", "기간", "주요 내용", "산출물"],
            [
                ("1단계", "1~2주", "지하 통신 음영, 작업구역, 위험·품질 체크리스트 조사", "현장 진단표, PoC 상세계획"),
                ("2단계", "3~4주", "QR/NFC, TBM, 위험성평가, 작업중지, 품질 체크 항목 연동", "작업데이터 샘플, 체크리스트"),
                ("3단계", "5~6주", "스마트글라스 영상, 로봇개/드론 선행점검 시나리오 검증", "사진·영상 증거, 운영 로그"),
                ("4단계", "7~8주", "실무자 피드백, 성과지표 측정, 후속과제 도출", "PoC 결과보고서, 개선안"),
            ],
            [1.6, 1.4, 5.0, 3.0],
        ),
        (
            "6. 정량 성과지표",
            "성과는 시연 여부가 아니라 현장 데이터 생성, 위치·통신 검증, 실무 수용성으로 판단한다.",
            ["지표", "PoC 목표", "측정 방법"],
            [
                ("통신 음영 지도화", "PoC 구역 내 음영 구간 식별 완료", "현장 측정 로그 및 위치별 수신 상태 기록"),
                ("작업구역 식별", "층·구역 단위 식별 가능성 확인", "QR/NFC, 기준점, 작업구역 매핑 비교"),
                ("작업데이터 생성률", "대상 작업의 80% 이상 데이터화", "TBM, 위험성평가, 체크리스트 등록 건수"),
                ("스마트글라스 활용성", "품질·안전 확인 시나리오 3건 이상 수행", "영상 캡처, 관리자 피드백"),
                ("로봇/드론 시나리오", "선행점검 후보 구역 2건 이상 검토", "이동 가능성, 촬영 가능성, 위험요인 기록"),
                ("실무자 수용성", "안전·공사·품질 담당자 피드백 확보", "인터뷰 및 설문"),
            ],
            [2.5, 4.1, 4.4],
        ),
        (
            "7. TRL 및 기술개발 목표",
            "대우건설 PoC는 완성품 납품이 아니라 TRL 5~6 수준의 현장 검증 단계로 설정한다.",
            ["항목", "착수 수준", "PoC 종료 목표"],
            [
                ("SAFE-LINK 작업데이터", "현장 기능 시연 가능", "대우 PoC 구역 작업데이터 샘플 생성"),
                ("지하 통신·좌표", "후보 기술 검토 단계", "음영지도 및 3D 위치 인식 가능성 검증"),
                ("스마트글라스", "장비 적용 검토", "안전·품질 영상 관제 시나리오 검증"),
                ("로봇/드론", "적용 시나리오 검토", "선행점검 후보 구역과 운용 조건 도출"),
                ("통합 보고서", "SAFE-LINK 보고 기능 기반", "안전+품질+위치+증거 통합 보고서 초안"),
            ],
            [3.2, 3.8, 4.0],
        ),
        (
            "8. 컨소시엄 및 역할",
            "서원토건은 현장 공정 지식과 작업데이터를 맡고, 대우건설은 현장 검증 체계를 제공한다.",
            ["주체", "역할"],
            [
                ("서원토건", "철근콘크리트 공정 지식, SAFE-LINK 운영, PoC 시나리오, 작업데이터 설계"),
                ("대우건설", "PoC 현장 제공, 안전·공사·품질 실무 검증, 기존 체계와 비교 평가"),
                ("경희대 AI-RAN/O-RAN 연구팀", "지하 통신, AI-RAN/O-RAN, 이음5G 특화망 자문 및 후속 공동과제 검토"),
                ("로봇·드론 협력사", "상용 장비 또는 협업 장비 기반 선행점검 시나리오 검토"),
                ("스마트글라스 협력사", "현장 영상 관제, 품질검측 보조, 착용형 디바이스 적용성 검토"),
            ],
            [3.4, 7.6],
        ),
        (
            "9. 지식재산권·데이터·보안 관리",
            "현장 도면, 영상, 위치 로그, 작업자 데이터는 PoC 목적 범위에서만 사용한다.",
            ["구분", "관리 방안"],
            [
                ("IP 방향", "검증 이벤트 기반 작업증거 데이터셋, BIM/도면 위치와 작업증거 매핑, 안전+품질 체크 자동화 중심 권리화 검토"),
                ("데이터 원칙", "대우건설 현장 데이터는 PoC 목적 범위에서만 사용하며, 민감 정보와 도면은 접근 권한을 분리"),
                ("보안", "자료 열람·수정·반출 이력 관리, 외부 발표 전 대우건설 검토, 원본 데이터 보관·폐기 범위 협약"),
                ("공동성과", "PoC 결과물, 공동특허, 후속과제 자료의 권리 범위를 사전 합의"),
            ],
            [2.3, 8.7],
        ),
        (
            "10. 대우건설 협력 요청 사항",
            "본 제안은 구매 요청이 아니라 현장 PoC와 공동 실증 파트너십 요청이다.",
            ["번호", "요청 사항"],
            [
                ("1", "지하층 또는 골조 작업구역 1개소 PoC 현장 제공"),
                ("2", "안전관리자, 공사팀, 품질팀 담당자 인터뷰 및 피드백 협조"),
                ("3", "현장 통신 음영 측정과 장비 운용 시간 협의"),
                ("4", "기존 스마트안전·품질관리 체계와 비교 검증"),
                ("5", "PoC 결과보고서 공동 검토"),
                ("6", "성과 확인 시 후속 공동과제, 공동특허, 확대 실증 검토"),
            ],
            [1.4, 9.6],
        ),
        (
            "11. 제출 전 확정 필요 항목",
            "최종 제출 전 예산, 현장 범위, 협력사, 정량성과, 중복지원 리스크 문구를 확정한다.",
            ["항목", "현재 상태", "보강 방향"],
            [
                ("PoC 예산", "미확정", "8주 기준 인력, 장비, 스마트글라스, 로봇/드론 대여, 통신 측정비로 분리 산정"),
                ("현장 범위", "미확정", "지하층, 지하주차장, PIT, 코어, 기계실 중 1개 구역으로 한정"),
                ("협력사", "후보 단계", "스마트글라스, 로봇/드론, 통신 자문 기관을 역할별로 명시"),
                ("정량 성과", "PoC 지표 제시", "작업데이터 생성률, 체크리스트 누락 감소, 처리시간 단축 등 현장 지표 확정"),
                ("중복지원 리스크", "관리 필요", "중기부 R&D는 원천기술, 대우 공모는 현장 실증으로 문구 분리"),
            ],
            [2.3, 2.3, 6.4],
        ),
    ]

    for idx, (title, lead, headers, rows, widths) in enumerate(slide_specs, start=2):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PPT_WHITE
        ppt_header_footer(slide, idx - 1)
        ppt_title(slide, title)
        ppt_textbox(slide, 0.65, 1.25, 10.5, 0.35, lead, 9.3, color=PPT_GRAY)
        ppt_table(slide, 0.65, 1.85, 10.55, 5.25, headers, rows, col_widths=widths, font_size=7.8 if len(rows) >= 6 else 8.3)

    slide = prs.slides.add_slide(blank)
    ppt_header_footer(slide, 12)
    ppt_title(slide, "12. 최종 제안 문장")
    ppt_textbox(
        slide,
        1.05,
        2.0,
        9.6,
        1.7,
        "지하·골조 현장 자동화의 병목은 장비 부족이 아니라\n통신, 3D 좌표, 작업데이터의 단절입니다.",
        20,
        bold=True,
        color=PPT_NAVY,
        align=PP_ALIGN.CENTER,
    )
    ppt_textbox(
        slide,
        1.3,
        4.1,
        9.1,
        1.2,
        "서원토건은 SAFE-LINK에서 출발해 안전과 품질을 연결하는 SQ-LINK Underground 현장 AI 운영체계로 확장하고,\n대우건설과 함께 그 가능성을 실제 현장에서 PoC로 검증하고자 합니다.",
        13,
        color=PPT_BLUE,
        align=PP_ALIGN.CENTER,
    )
    ppt_textbox(slide, 3.2, 6.1, 5.3, 0.45, "현장 PoC 및 공동 실증 파트너십 제안", 12, bold=True, color=PPT_WHITE, fill=PPT_NAVY, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)
    scrub_pptx_slide_number_placeholders(PPTX_PATH)


def scrub_pptx_slide_number_placeholders(path):
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    ET.register_namespace("p", ns["p"])
    ET.register_namespace("a", ns["a"])
    tmp = NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.close()
    with ZipFile(path, "r") as src, ZipFile(tmp.name, "w", ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename.startswith("ppt/slide") and item.filename.endswith(".xml"):
                try:
                    root = ET.fromstring(data)
                    changed = False
                    for sp_tree in root.findall(".//p:spTree", ns):
                        for sp in list(sp_tree.findall("p:sp", ns)):
                            has_sld_num = sp.find(".//p:ph[@type='sldNum']", ns) is not None
                            has_slide_number_text = any((t.text or "") == "Slide Number" for t in sp.findall(".//a:t", ns))
                            if has_sld_num or has_slide_number_text:
                                sp_tree.remove(sp)
                                changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                except ET.ParseError:
                    pass
            dst.writestr(item, data)
    Path(tmp.name).replace(path)


def main():
    build_docx()
    build_pptx()
    print(DOCX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
