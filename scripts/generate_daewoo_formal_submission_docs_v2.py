from pathlib import Path
from zipfile import ZipFile

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt as PPTPt
from pptx.dml.color import RGBColor as PPTRGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "generated"
DOCX_PATH = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_v3_가독성개선_20260611.docx"
PPTX_PATH = OUT / "대우건설_HyperSafety_AI_공모전_공문서형_제출본_v3_가독성개선_20260611.pptx"

NAVY = RGBColor(27, 52, 78)
BLUE = RGBColor(42, 86, 127)
GRAY = RGBColor(82, 82, 82)
BLACK = RGBColor(30, 30, 30)
LIGHT_BLUE = "EAF1F8"
LIGHT_GRAY = "F5F6F8"
LINE = "B8C4D0"

PPT_NAVY = PPTRGBColor(27, 52, 78)
PPT_BLUE = PPTRGBColor(42, 86, 127)
PPT_GRAY = PPTRGBColor(82, 82, 82)
PPT_BLACK = PPTRGBColor(30, 30, 30)
PPT_LIGHT_BLUE = PPTRGBColor(234, 241, 248)
PPT_WHITE = PPTRGBColor(255, 255, 255)
PPT_LINE = PPTRGBColor(184, 196, 208)


def shade(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def border(cell, color=LINE):
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        el = borders.find(qn("w:" + edge))
        if el is None:
            el = OxmlElement("w:" + edge)
            borders.append(el)
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "5")
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), color)


def run_font(run, size=9, bold=False, color=BLACK):
    run.font.name = "Malgun Gothic"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def cell_text(cell, text, size=11, bold=False, color=BLACK, align=None):
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    p = cell.paragraphs[0]
    p.alignment = align or WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.line_spacing = 1.12
    r = p.add_run(str(text))
    run_font(r, size=size, bold=bold, color=color)


def heading(doc, text, level=1):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(11 if level == 1 else 6)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    run_font(r, size=16 if level == 1 else 13.5, bold=True, color=NAVY if level == 1 else BLUE)


def para(doc, text, bold=False):
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = 1.18
    p.paragraph_format.space_after = Pt(3)
    r = p.add_run(text)
    run_font(r, size=12, bold=bold, color=BLACK)


def table_doc(doc, headers, rows, widths=None, font=10.5):
    t = doc.add_table(rows=1, cols=len(headers))
    t.autofit = False
    for c, h in enumerate(headers):
        cell = t.rows[0].cells[c]
        if widths:
            cell.width = Cm(widths[c])
        shade(cell, LIGHT_BLUE)
        border(cell)
        cell_text(cell, h, size=font, bold=True, color=NAVY, align=WD_ALIGN_PARAGRAPH.CENTER)
    for row in rows:
        cells = t.add_row().cells
        for c, val in enumerate(row):
            if widths:
                cells[c].width = Cm(widths[c])
            border(cells[c])
            cell_text(cells[c], val, size=font, align=WD_ALIGN_PARAGRAPH.CENTER if c == 0 else WD_ALIGN_PARAGRAPH.LEFT)
    doc.add_paragraph().paragraph_format.space_after = Pt(1)
    return t


sections = [
    {
        "title": "1. 제출 요약문",
        "lead": "중기부 R&D 문서의 기술개발 구조를 대우건설 공모전의 현장 PoC 논리로 전환한 제출본이다.",
        "headers": ["항목", "상세 내용"],
        "rows": [
            ("제안명", "SQ-LINK Underground: 지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼"),
            ("제안 목적", "지하 통신 음영, 3D 위치, 작업데이터, 스마트글라스, 로봇·드론 점검 가능성을 대우건설 현장에서 검증"),
            ("핵심 전환", "중기부 R&D는 원천기술 개발, 대우건설 공모전은 현장 적용성·실무 수용성·공동 실증 검증으로 역할 분리"),
            ("주요 산출물", "현장 음영지도, 작업데이터 샘플, 스마트글라스 영상 증거, 로봇/드론 점검 조건, PoC 결과보고서"),
            ("최종 메시지", "장비 자체가 아니라 통신·좌표·작업데이터 기반을 만들어 안전·품질 자동화의 작동 조건을 검증"),
        ],
        "widths": [3.0, 13.0],
    },
    {
        "title": "2. 기술의 정의 및 필요성",
        "lead": "HWPX의 과제 정의를 대우건설 현장 문제에 맞춰 재정의한다.",
        "headers": ["구분", "정부 R&D 원문 구조", "대우건설 제출 반영"],
        "rows": [
            ("기술 정의", "지하 밀폐 환경 내 특화망 기반 웨어러블 및 사족보행 로봇 통합 제어 시스템", "지하·골조 현장 내 작업자, 스마트글라스, 로봇, 드론, BIM, 작업데이터 통합 운영체계"),
            ("필요성", "Wi-Fi/LTE 신호 감쇄, GPS 미약, 원격 제어 한계, 작업자 안전 사각지대", "지하층·PIT·코어·기계실에서 안전·품질 확인과 증거 수집이 끊기는 문제 해결"),
            ("트렌드", "지하 공간 활용 증대, 산업 현장 무인화, 안정 통신 인프라와 로봇 제어 융합", "대형 건설현장의 스마트 안전·품질, BIM, 디지털트윈, 현장 AI 전환과 직접 연계"),
            ("정책 정합성", "12대 국가전략기술 중 첨단 모빌리티 및 차세대 통신 분야와 부합", "대우건설 Hyper Safety & AI 주제에 안전·품질·AI·현장 자동화를 결합"),
        ],
        "widths": [2.5, 6.5, 7.0],
    },
    {
        "title": "3. 현장 문제 및 한계 분석",
        "lead": "중기부 계획서의 문제점·한계 항목을 건설현장 심사자가 이해할 수 있는 언어로 구체화한다.",
        "headers": ["문제 영역", "현장 한계", "PoC에서 확인할 사항"],
        "rows": [
            ("통신 단절", "지하층·PIT·코어·지하주차장에서 기존 Wi-Fi/LTE 신호 감쇄와 단절 빈번", "위치별 수신 상태, 영상 전송 가능성, 작업데이터 업로드 지연 확인"),
            ("GPS 부재", "지하에서는 평면 위치뿐 아니라 층·높이·깊이 식별이 어려움", "QR/NFC, 기준점, UWB/BLE/Wi-Fi RTT 후보의 구역 식별 가능성 검토"),
            ("장비 고립", "로봇·드론·스마트글라스가 BIM·작업허가·품질 체크와 분리되어 단독 시연에 머묾", "장비가 수집한 사진·영상이 작업구역과 보고서에 연결되는지 검증"),
            ("안전·품질 분리", "철근 배근, 거푸집, 동바리, 타설 전 점검은 안전과 품질이 같은 작업면에서 발생", "TBM·위험성평가·품질 체크·작업중지 이력을 하나의 증거 패키지로 구성"),
            ("경제 손실", "통신 두절로 작업 중단·재개, 재확인, 재촬영, 보고 누락 비용 발생", "대우 PoC에서는 정량 손실액 산정보다 손실 발생 이벤트와 처리 시간을 우선 계측"),
        ],
        "widths": [2.5, 6.8, 6.7],
    },
    {
        "title": "4. 핵심 기술 요소",
        "lead": "HWPX의 3대 핵심기술을 대우건설 현장 PoC 구조로 확장한다.",
        "headers": ["핵심기술", "정부 R&D 기술 내용", "대우 PoC 적용"],
        "rows": [
            ("지하 특화망 구축", "지하 환경에 최적화된 이음5G 주파수 설계 및 효율적인 기지국 배치", "대우 현장 음영지도 작성, 통신 취약 구간 도출, 후속 특화망 설계 근거 확보"),
            ("초저지연 통신 최적화", "URLLC, 네트워크 슬라이싱, 빔포밍, 스마트 안테나 기반 지연·손실 최소화", "스마트글라스 영상, 작업데이터, 로봇/드론 점검 로그의 실시간성 요구조건 정리"),
            ("로봇-웨어러블 통합 제어", "사족보행 로봇과 착용형 기기를 단일 플랫폼에서 통합 제어", "로봇·드론·스마트글라스는 장비 개발이 아니라 작업구역·BIM·보고서와 연결하는 응용 계층으로 검증"),
            ("현장 AI 운영 엔진", "엣지 컴퓨팅, 통합 제어 플랫폼, 다중 디바이스 연동", "안전·품질 체크리스트, 작업중지, 사진·영상 증거, 담당자 조치 이력을 자동 보고서로 구성"),
        ],
        "widths": [3.0, 6.5, 6.5],
    },
    {
        "title": "5. 주요 성능지표 및 PoC 검증지표",
        "lead": "중기부 R&D의 정량 성능지표를 대우건설 현장 검증지표로 변환한다.",
        "headers": ["지표", "R&D 최종목표", "세계수준", "대우 PoC 검증 방식"],
        "rows": [
            ("통신 지연 시간", "10ms 이하", "5ms 이하", "PoC에서는 구간별 지연 로그와 영상/데이터 전송 지연을 계측하고, 특화망 적용 필요 구간을 도출"),
            ("패킷 손실률", "0.1% 미만", "0.01% 미만", "지하 작업구역별 업로드 실패, 재전송, 영상 끊김 이벤트 기록"),
            ("로봇 동시 제어", "5대 이상", "10대 이상", "대우 PoC에서는 다중 장비 동시 제어 완성이 아니라 로봇/드론/글라스 동시 데이터 흐름을 검증"),
            ("통신 도달 거리", "500m 이상", "1,000m 이상", "PoC 구역 내 거리·층·구획별 수신 상태와 음영 분포를 지도화"),
            ("웨어러블 전송속도", "100Mbps 이상", "200Mbps 이상", "스마트글라스 영상, 센서, 사진 증거 전송 가능 조건을 확인"),
        ],
        "widths": [2.6, 2.4, 2.4, 8.6],
    },
    {
        "title": "6. 상세 개발 방법",
        "lead": "HWPX 상세계획서의 수행내용을 현장 실증 중심으로 풀어낸다.",
        "headers": ["기술 항목", "세부 방법", "대우건설 현장 산출물"],
        "rows": [
            ("전파 특성 분석", "3.7GHz 대역 5G 특화망 주파수의 지하 전파 특성 분석 및 전파 모델링", "구역별 음영지도, 수신상태 로그, 통신 취약구간 목록"),
            ("신호 집중/간섭 저감", "빔포밍 및 스마트 안테나 적용 검토", "현장 구조물, 코어, PIT, 벽체 등에 따른 신호 감쇄 원인 정리"),
            ("저지연 경로 최적화", "초저지연 통신 프로토콜 및 네트워크 슬라이싱 설계", "영상·작업데이터·비상신고별 우선순위 기준 수립"),
            ("하이브리드 통신", "5G 특화망 + Wi-Fi 6E/UWB 보조 통신 자동 전환", "신호 취약 구간에서 보조 통신 후보와 전환 조건 검토"),
            ("로봇 프레임워크", "ROS 기반 로봇 제어 프레임워크 연동", "상용 로봇/드론 협력사 연동 조건, 접근 가능 구역, 촬영 각도 조건 도출"),
            ("엣지 플랫폼", "현장 엣지 컴퓨팅 기반 실시간 데이터 처리", "현장 내 즉시 보고서, 조치 알림, 증거 패키지 생성 가능성 검증"),
        ],
        "widths": [3.0, 6.5, 6.5],
    },
    {
        "title": "7. 단계별 추진 로드맵",
        "lead": "중기부 24개월 개발 로드맵과 대우 8주 PoC를 병행 가능한 구조로 분리한다.",
        "headers": ["구분", "기간", "주요 활동", "성과물"],
        "rows": [
            ("대우 PoC 1단계", "1~2주", "현장 진단, 통신 음영, 작업구역, 위험·품질 체크리스트 조사", "현장 진단표, PoC 상세계획"),
            ("대우 PoC 2단계", "3~4주", "QR/NFC, TBM, 위험성평가, 작업중지, 품질 체크 항목 연동", "작업데이터 샘플, 체크리스트"),
            ("대우 PoC 3단계", "5~6주", "스마트글라스 영상, 로봇/드론 선행점검 시나리오 검증", "영상·사진 증거, 운영 로그"),
            ("대우 PoC 4단계", "7~8주", "실무자 피드백, 성과지표 측정, 후속과제 도출", "PoC 결과보고서, 개선안"),
            ("후속 R&D", "24개월", "전파 모델링, 하이브리드 모듈, 통합 제어 플랫폼, 지하 테스트베드 검증", "성능검증 보고서, 공인시험기관 성적서, TRL 6 달성"),
        ],
        "widths": [3.0, 2.2, 6.5, 4.3],
    },
    {
        "title": "8. TRL 단계 및 산출물",
        "lead": "대우건설 공모전은 완성 납품이 아니라 TRL 5~6 현장검증으로 가는 교량 역할이다.",
        "headers": ["TRL", "정부 R&D 활동", "대우건설 PoC와 연결되는 산출물"],
        "rows": [
            ("TRL 4", "지하 환경 전파 특성 분석, 5G 특화망 주파수 설계, 하이브리드 통신 모듈 개념 설계", "대우 현장 음영지도와 구역별 통신 요구조건"),
            ("TRL 5", "하이브리드 통신 모듈 시제품, 통합 제어 플랫폼 알파 버전, 실험실 성능 검증", "SAFE-LINK 작업데이터와 스마트글라스·로봇/드론 데이터 흐름 샘플"),
            ("TRL 6", "지하 실증 테스트베드 구축, 종합 성능 검증, 시스템 최적화", "대우 현장 PoC 결과보고서와 후속 공동실증 설계 근거"),
        ],
        "widths": [2.0, 7.0, 7.0],
    },
    {
        "title": "9. 선행연구 및 보유기반 활용",
        "lead": "HWPX의 선행연구 활용 구조를 SAFE-LINK와 서원토건의 현장 기반으로 재작성한다.",
        "headers": ["선행 기반", "HWPX 활용 논리", "대우 제출본 보강 내용"],
        "rows": [
            ("지하 전파 분석", "스마트 광산/지하 환경 전파 특성 분석 결과를 주파수 설계에 활용", "대우 현장의 실제 지하층, 코어, PIT, 지하주차장을 실증 데이터셋으로 확보"),
            ("로봇 원격제어", "재난 현장 로봇 원격 제어 기술을 하이브리드 통신 모듈에 적용", "로봇개/드론은 자체개발 주장이 아니라 상용 장비의 현장 운용 조건 검증으로 제한"),
            ("웨어러블 통신", "웨어러블 데이터 전송 속도 100Mbps 이상 목표 달성을 위한 기반", "스마트글라스 영상 관제와 품질검측 보조를 대우 PoC의 핵심 장면으로 설정"),
            ("SAFE-LINK", "작업자 안전 데이터와 전자증거를 생성하는 현장 운영 기반", "TBM, 위험성평가, 작업중지, 품질 체크, 관리자 조치 이력을 하나의 증거 패키지로 연결"),
        ],
        "widths": [3.0, 6.5, 6.5],
    },
    {
        "title": "10. 기술적 애로사항 및 해결 과제",
        "lead": "심사자가 중요하게 보는 기술 난제를 먼저 제시하고 해결 방향을 붙인다.",
        "headers": ["애로사항", "원인", "해결 방향"],
        "rows": [
            ("전파 감쇄 및 간섭", "지하 터널, 벽체, 철근, 장비, 다중 경로 간섭으로 통신 품질 저하", "전파 모델링, 특화망 주파수 설계, 빔포밍/스마트 안테나 적용 후보 검토"),
            ("초저지연 제어", "로봇, 스마트글라스, 작업데이터가 동시에 연결될 때 지연과 손실 발생", "네트워크 슬라이싱, QoS, 데이터 우선순위, 엣지 처리 기준 설계"),
            ("이종 디바이스 연동", "스마트글라스, 로봇, 드론, 모바일, BIM의 프로토콜과 데이터 형식 상이", "개방형 API, ROS 연동, 작업구역 ID, 증거 데이터 표준 구조 설계"),
            ("위치 인식", "GPS 부재와 수직 위치 식별 곤란", "QR/NFC 기준점, UWB/BLE/Wi-Fi RTT 후보, BIM 좌표 매핑을 조합한 PoC"),
            ("현장 수용성", "작업자가 번거롭다고 느끼면 데이터가 누락됨", "TBM-위험성평가-품질체크-보고서 흐름을 현장 작업순서에 맞게 최소 클릭으로 설계"),
        ],
        "widths": [3.0, 5.7, 7.3],
    },
    {
        "title": "11. 사업화 애로사항 및 해결 과제",
        "lead": "기술이 좋아도 현장 도입이 되지 않는 문제를 별도 관리한다.",
        "headers": ["사업화 리스크", "상세 내용", "대응 전략"],
        "rows": [
            ("초기 시장 진입", "건설현장은 신규 장비·시스템 도입에 보수적이고 안전 규제 부담이 큼", "대우건설 PoC를 통해 성능보다 현장 적용성, 실무 피드백, 운영 절차를 먼저 검증"),
            ("경쟁 기술", "기존 Wi-Fi/LTE, CCTV, 단일 안전앱, 단일 스마트글라스 솔루션과 비교됨", "통신·좌표·작업데이터·BIM·장비를 묶는 운영체계 차별성을 전면화"),
            ("표준화 부족", "지하 특화망 기반 로봇/웨어러블 제어의 표준 가이드라인 부족", "대우 현장 결과를 기반으로 현장 표준 운영 절차와 데이터 항목 정의"),
            ("양산/공급망", "통신 모듈, 제어 보드, 장비 협력사, 유지보수 체계 필요", "하드웨어 자체 개발을 주장하지 않고 전문 협력사/상용 장비 도입으로 리스크 축소"),
        ],
        "widths": [3.0, 6.0, 7.0],
    },
    {
        "title": "12. 지식재산권 확보 및 회피 전략",
        "lead": "HWPX의 선행특허 분석 구조를 제출본에 맞게 반영한다.",
        "headers": ["구분", "검토 대상", "대우 제출본 적용"],
        "rows": [
            ("선행특허 검토", "이음5G 프로비저닝, 사족보행 로봇 핸들링/한손 제어, 햅틱 기반 웨어러블 제어", "본 제안은 장비 자체보다 지하·골조 현장 작업 이벤트와 증거 데이터 연결에 초점"),
            ("회피 전략", "선행특허와 기술 범위 및 핵심 사상을 분리", "통신망·로봇 자체 구현 주장을 줄이고 현장 업무 흐름, BIM 매핑, 안전+품질 증거화로 차별화"),
            ("출원 방향 1", "지하 환경 전파 특성 분석 및 주파수 최적화", "대우 PoC 음영지도와 구역별 통신 요구조건을 IP 근거 데이터로 활용"),
            ("출원 방향 2", "로봇 및 웨어러블 통합 제어 플랫폼", "스마트글라스/로봇/드론 촬영 결과를 작업구역·체크리스트·보고서에 자동 편입하는 구조"),
            ("출원 방향 3", "저지연/고신뢰 통신 최적화", "위험신고, 작업중지, 영상관제 등 현장 이벤트별 QoS 우선순위 설계"),
        ],
        "widths": [3.0, 6.0, 7.0],
    },
    {
        "title": "13. 제품화·양산 및 후속 확장",
        "lead": "대우 PoC 이후 제품화와 양산으로 이어지는 경로를 제시한다.",
        "headers": ["단계", "기간", "주요 활동", "산출물"],
        "rows": [
            ("적용설계", "2026.07~2026.10", "이음5G 통신 모듈 소형화, 웨어러블/로봇 인터페이스 통합 설계", "제품화 상세 설계 보고서"),
            ("시작품 제작", "2026.11~2027.02", "통합 제어 시스템 시작품 제작, 기능·성능 테스트", "시작품, 내부 테스트 결과"),
            ("현장시험", "2027.03~2027.06", "지하 실증 테스트베드 운용, 지연·손실률 현장 측정", "현장 실증 테스트 결과"),
            ("인증/표준화", "2027.07~2027.10", "국내외 안전·성능 인증 준비, 표준화 기술 제안", "공인 인증 계획서, 표준화 제안서"),
            ("양산 준비", "2027.11~2028.02", "생산 라인 구축, 공급망 확보, 시험 생산", "초기 양산 제품, 품질관리 기준"),
        ],
        "widths": [2.5, 3.0, 7.0, 3.5],
    },
    {
        "title": "14. 시장성 및 매출·고용 목표",
        "lead": "HWPX의 시장·매출·고용 논리를 대우건설 후속 사업화 관점으로 정리한다.",
        "headers": ["항목", "HWPX 근거", "대우 제출본 활용"],
        "rows": [
            ("시장 성장", "글로벌 특화망 기반 로봇 통합 제어 시장은 연평균 26~34% 성장 가정", "대우 PoC를 초기 레퍼런스로 확보해 국내 건설·지하 시설 시장 진입 근거로 활용"),
            ("매출 목표", "연구매출 2027년 10억, 2028년 30억, 2029년 80억, 2030년 150억, 2031년 250억", "대우 실증 성공 시 건설사 PoC 패키지, 안전+품질 보고 SaaS, 현장 통신진단 서비스로 분리"),
            ("고용 창출", "개발 완료 후 5년간 15명 이상, 상세계획서 기준 직접 20명·간접 50명 효과", "통신, 로봇 SW, 필드 엔지니어, 데이터 분석, 현장 CS 인력 채용 계획과 연결"),
            ("경제효과", "생산성 30% 향상, 안전사고 90% 이상 감소 기대", "대우 PoC에서는 해당 수치를 확정 주장하지 않고, 체크 누락·처리시간·증거 생성률 개선 지표로 선행 검증"),
        ],
        "widths": [3.0, 6.5, 6.5],
    },
    {
        "title": "15. SWOT 및 판로 전략",
        "lead": "상세계획서의 SWOT를 대우건설 공모전 관점으로 압축 반영한다.",
        "headers": ["구분", "주요 내용", "전략"],
        "rows": [
            ("Strength", "이음5G 기반 저지연 통신, 로봇/웨어러블 연동, 지하 전파 특성 분석, SAFE-LINK 현장 데이터", "대우 PoC에서 전문건설사 기반 현장 문제 정의와 데이터 생성력을 강조"),
            ("Weakness", "초기 레퍼런스 부족, 양산 체계 미비, 대규모 마케팅 한계", "대우 현장 1개 구역 실증으로 레퍼런스 확보 후 국가과제·공동특허로 확장"),
            ("Opportunity", "스마트 건설, 지하 공간 개발, 재난 안전, 5G 특화망 확산, AI·디지털트윈 융합", "안전앱이 아니라 안전·품질·공정·원가로 확장 가능한 현장 AI로 포지셔닝"),
            ("Threat", "글로벌 선도기업 진입, 보안/프라이버시 이슈, 초기 투자 부담, 표준화 지연", "하드웨어 직접 개발보다 협력사 기반 생태계와 데이터/IP 차별화로 리스크 완화"),
        ],
        "widths": [2.2, 6.8, 7.0],
    },
    {
        "title": "16. 응용분야 및 파급효과",
        "lead": "지하 건설 현장을 시작으로 광산, 터널, 시설물 점검, 재난 구조까지 확장 가능하다.",
        "headers": ["응용 분야", "적용 시나리오", "기대 효과"],
        "rows": [
            ("대우 지하·골조 현장", "지하주차장, PIT, 코어, 기계실, 타설 전 점검, 개구부 관리", "안전·품질 데이터 통합, 보고 누락 감소, 실무자 피드백 기반 표준화"),
            ("터널 및 지하 구조물", "굴착 장비 원격 운용, 작업자 위치, 통신 음영 진단", "위험 작업 구역의 선행 점검과 무인화 기반 확보"),
            ("지하 시설물 점검", "전력구, 통신구, 배관, 유지보수 로봇 점검", "영상·위치·작업기록의 증거화와 점검 보고 자동화"),
            ("방재 및 재난 구조", "지하 매몰 현장 수색, 구조대원 실시간 정보 공유", "위험구역 접근 전 정보 수집과 상황 대응 속도 개선"),
            ("스마트 팩토리/물류", "이음5G 기반 저지연 통신과 이동형 장비 관제", "건설 외 산업용 통신·로봇 제어 시장으로 기술 확산"),
        ],
        "widths": [3.2, 6.4, 6.4],
    },
    {
        "title": "17. 보안관리 및 연구윤리 준수",
        "lead": "대우건설 현장 데이터 보호와 기술유출 방지 체계를 제출본에 명시한다.",
        "headers": ["관리 영역", "세부 방안"],
        "rows": [
            ("자료 등급", "현장 도면, 영상, 위치 로그, 작업자 데이터, PoC 결과보고서를 보안 등급별로 분류"),
            ("접근 통제", "참여 인력별 계정 부여, 최소 권한 원칙, 열람·수정·반출 이력 기록"),
            ("개인정보", "근로자 얼굴, 음성, 위치, 연락처 등은 마스킹·익명화 기준을 대우건설과 사전 합의"),
            ("외부 발표", "사진, 영상, 도면, 성과자료는 외부 발표 전 대우건설 검토 후 사용"),
            ("연구윤리", "원본 데이터 보존, 보고서 근거자료 관리, 결과 조작·과장 방지, 참여자 보안교육 시행"),
        ],
        "widths": [3.2, 12.8],
    },
    {
        "title": "18. 컨소시엄 및 역할분담",
        "lead": "HWPX 검토의견에서 지적된 통신·로봇·현장 3각 구성을 명확히 반영한다.",
        "headers": ["주체", "역할", "대우 PoC 기여"],
        "rows": [
            ("서원토건", "철근콘크리트 공정 지식, SAFE-LINK 운영, 작업데이터 설계", "실제 현장 업무흐름 기반 PoC 시나리오 설계"),
            ("대우건설", "PoC 현장 제공, 안전·공사·품질 실무 검증, 기존 체계와 비교", "대형 건설사 현장 검증성과 실무 피드백 확보"),
            ("경희대 AI-RAN/O-RAN 연구팀", "지하 통신, AI-RAN/O-RAN, 이음5G 특화망 자문", "통신 음영 분석과 후속 공동과제 기술 근거 제공"),
            ("로봇·드론 협력사", "상용 장비 기반 선행점검 시나리오 검토", "장비 자체 개발 리스크 없이 현장 운용 조건 검증"),
            ("스마트글라스 협력사", "영상 관제, 품질검측 보조, 착용형 디바이스 적용성 검토", "작업자 시야 공유와 안전·품질 증거 수집 검증"),
        ],
        "widths": [3.2, 6.4, 6.4],
    },
    {
        "title": "19. 대우건설에 요청하는 협력 사항",
        "lead": "본 제안은 구매 요청이 아니라 공동 실증 파트너십 요청이다.",
        "headers": ["번호", "요청 사항", "목적"],
        "rows": [
            ("1", "지하층 또는 골조 작업구역 1개소 PoC 현장 제공", "실제 통신 음영과 작업 흐름 검증"),
            ("2", "안전관리자, 공사팀, 품질팀 담당자 인터뷰 및 피드백", "실무 수용성과 기존 체계 대비 개선점 확인"),
            ("3", "현장 통신 음영 측정과 장비 운용 시간 협의", "측정 가능 시간대와 작업 간섭 최소화"),
            ("4", "기존 스마트안전·품질관리 체계와 비교 검증", "대우건설 내부 기준에 맞는 성과 판단"),
            ("5", "PoC 결과보고서 공동 검토", "후속 확대 실증과 공동과제 기획 근거 확보"),
            ("6", "성과 확인 시 공동특허, 국가과제, 확대 실증 검토", "R&D와 현장 사업화의 선순환 구조 구축"),
        ],
        "widths": [1.5, 7.0, 7.5],
    },
    {
        "title": "20. 제출 전 확정 필요 항목",
        "lead": "HWPX의 [보완 필요] 항목을 대우건설 제출본용 체크리스트로 전환한다.",
        "headers": ["항목", "현재 상태", "보강 방향"],
        "rows": [
            ("PoC 예산", "미확정", "8주 기준 인력, 장비, 스마트글라스, 로봇/드론 대여, 통신 측정비로 분리 산정"),
            ("현장 범위", "미확정", "지하층, 지하주차장, PIT, 코어, 기계실 중 1개 구역으로 한정"),
            ("협력사", "후보 단계", "통신 자문, 스마트글라스, 로봇/드론 협력사를 역할별로 명시"),
            ("정량 성과", "PoC 지표 제시", "작업데이터 생성률, 체크리스트 누락 감소, 신고 처리시간 등 현장 지표 확정"),
            ("TRL", "4~5 착수, 5~6 종료 목표", "PoC는 현장 검증 근거, R&D는 원천기술 고도화로 분리"),
            ("IP", "방향 제시", "SAFE-LINK 증거 데이터셋, BIM 매핑, 안전+품질 체크 자동화 중심으로 명세화"),
            ("중복지원", "관리 필요", "중기부 R&D=원천기술, 대우 공모=현장 실증 문구를 모든 문서에 일관 적용"),
        ],
        "widths": [3.0, 4.0, 9.0],
    },
]


def build_docx():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Cm(1.55)
    sec.bottom_margin = Cm(1.45)
    sec.left_margin = Cm(1.55)
    sec.right_margin = Cm(1.55)
    doc.styles["Normal"].font.name = "Malgun Gothic"
    doc.styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    doc.styles["Normal"].font.size = Pt(12)

    hp = sec.header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r = hp.add_run("대우건설 Hyper Safety & AI Open Innovation 제출본 v3 | 서원토건")
    run_font(r, size=9, color=GRAY)
    fp = sec.footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = fp.add_run("CONFIDENTIAL - 공모전 제출 및 PoC 협의용")
    run_font(r, size=9, color=GRAY)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(70)
    r = p.add_run("대우건설 Hyper Safety & AI Open Innovation\n공문서형 제출 제안서 v3 가독성개선본")
    run_font(r, size=21, bold=True, color=NAVY)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("SQ-LINK Underground\n지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼")
    run_font(r, size=12, color=BLUE)
    table_doc(
        doc,
        ["항목", "내용"],
        [
            ("문서번호", "SW-DW-HSAI-20260611-03"),
            ("작성일", "2026. 06. 11."),
            ("제안기관", "서원토건"),
            ("참고자료", "중기부 R&D 지원사업 HWPX 2종, 검토의견 MD, 기존 대우건설 제안서"),
            ("문서성격", "공모전 제출용 / 대우건설 PoC 및 공동실증 협의용"),
        ],
        widths=[4.0, 12.0],
        font=11,
    )
    doc.add_page_break()

    for s in sections:
        heading(doc, s["title"])
        para(doc, s["lead"])
        table_doc(doc, s["headers"], s["rows"], widths=s["widths"], font=10 if len(s["headers"]) >= 4 else 10.5)

    heading(doc, "21. 최종 제안 문장")
    para(doc, "지하·골조 현장 자동화의 병목은 장비 부족이 아니라 통신, 3D 좌표, 작업데이터의 단절입니다. 서원토건은 SAFE-LINK에서 출발해 안전과 품질을 연결하는 SQ-LINK Underground 현장 AI 운영체계로 확장하고, 대우건설과 함께 그 가능성을 실제 현장에서 PoC로 검증하고자 합니다.", bold=True)
    doc.save(DOCX_PATH)


def tx(slide, x, y, w, h, text, size=12, bold=False, color=PPT_BLACK, align=PP_ALIGN.LEFT, fill=None):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Malgun Gothic"
    r.font.size = PPTPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def header(slide, page):
    slide.shapes.add_shape(1, Inches(0.45), Inches(0.55), Inches(10.8), Inches(0.01)).line.color.rgb = PPT_LINE
    tx(slide, 0.5, 0.25, 6.8, 0.28, "대우건설 Hyper Safety & AI Open Innovation 제출본 v3", 10, color=PPT_GRAY)
    tx(slide, 8.5, 0.25, 2.7, 0.28, "CONFIDENTIAL | 서원토건", 10, color=PPT_GRAY, align=PP_ALIGN.RIGHT)
    slide.shapes.add_shape(1, Inches(0.45), Inches(7.65), Inches(10.8), Inches(0.01)).line.color.rgb = PPT_LINE
    tx(slide, 0.5, 7.72, 4.5, 0.24, "공모전 제출 및 PoC 협의용", 9, color=PPT_GRAY)
    tx(slide, 10.55, 7.72, 0.65, 0.24, f"{page:02d}", 9, color=PPT_GRAY, align=PP_ALIGN.RIGHT)


def table_ppt(slide, x, y, w, h, headers, rows, widths, font=12):
    sh = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    t = sh.table
    total = sum(widths)
    for c, cw in enumerate(widths):
        t.columns[c].width = Inches(w * cw / total)
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = t.cell(r, c)
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPT_LIGHT_BLUE if r == 0 else PPT_WHITE
            cell.text = headers[c] if r == 0 else str(rows[r - 1][c])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r == 0 or c == 0 else PP_ALIGN.LEFT
                for rr in p.runs:
                    rr.font.name = "Malgun Gothic"
                    rr.font.size = PPTPt(font)
                    rr.font.bold = r == 0
                    rr.font.color.rgb = PPT_NAVY if r == 0 else PPT_BLACK
    return sh


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPT_WHITE
    slide.shapes.add_shape(1, Inches(0.7), Inches(0.9), Inches(10.3), Inches(0.02)).line.color.rgb = PPT_NAVY
    tx(slide, 1.0, 1.3, 9.7, 0.38, "공문서형 제출 제안서 v3 가독성개선본", 14, color=PPT_GRAY, align=PP_ALIGN.CENTER)
    tx(slide, 1.0, 1.95, 9.7, 0.95, "대우건설 Hyper Safety & AI\nOpen Innovation", 30, bold=True, color=PPT_NAVY, align=PP_ALIGN.CENTER)
    tx(slide, 1.1, 3.18, 9.5, 0.75, "SQ-LINK Underground\n지하·골조 현장 3D 좌표 통신망 기반 안전·품질 AI 운영 플랫폼", 16, color=PPT_BLUE, align=PP_ALIGN.CENTER)
    table_ppt(slide, 1.7, 4.65, 8.3, 1.65, ["항목", "내용"], [("문서번호", "SW-DW-HSAI-20260611-03"), ("작성일", "2026. 06. 11."), ("제안기관", "서원토건"), ("참고자료", "중기부 R&D HWPX 2종 + 검토의견 MD")], [2.0, 5.7], 11)
    tx(slide, 0.7, 7.35, 10.3, 0.28, "CONFIDENTIAL - 공모전 제출 및 PoC 협의용", 10, color=PPT_GRAY, align=PP_ALIGN.CENTER)

    page = 1
    for s in sections:
        chunk_size = 2 if len(s["headers"]) >= 4 else 3
        chunks = [s["rows"][i:i + chunk_size] for i in range(0, len(s["rows"]), chunk_size)]
        for ci, chunk in enumerate(chunks, start=1):
            page += 1
            slide = prs.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = PPT_WHITE
            header(slide, page - 1)
            suffix = f" ({ci}/{len(chunks)})" if len(chunks) > 1 else ""
            tx(slide, 0.6, 0.78, 10.6, 0.42, s["title"] + suffix, 18, bold=True, color=PPT_NAVY)
            tx(slide, 0.65, 1.28, 10.5, 0.42, s["lead"], 12, color=PPT_GRAY)
            table_ppt(slide, 0.6, 1.9, 10.6, 5.05, s["headers"], chunk, s["widths"], 12)

    slide = prs.slides.add_slide(blank)
    page += 1
    header(slide, page - 1)
    tx(slide, 0.6, 0.78, 10.6, 0.42, "21. 최종 제안 문장", 18, bold=True, color=PPT_NAVY)
    tx(slide, 1.0, 2.05, 9.7, 1.55, "지하·골조 현장 자동화의 병목은 장비 부족이 아니라\n통신, 3D 좌표, 작업데이터의 단절입니다.", 22, bold=True, color=PPT_NAVY, align=PP_ALIGN.CENTER)
    tx(slide, 1.25, 4.25, 9.2, 1.35, "서원토건은 SAFE-LINK에서 출발해 안전과 품질을 연결하는 SQ-LINK Underground 현장 AI 운영체계로 확장하고,\n대우건설과 함께 그 가능성을 실제 현장에서 PoC로 검증하고자 합니다.", 14, color=PPT_BLUE, align=PP_ALIGN.CENTER)
    tx(slide, 3.0, 6.15, 5.7, 0.5, "현장 PoC 및 공동 실증 파트너십 제안", 14, bold=True, color=PPT_WHITE, align=PP_ALIGN.CENTER, fill=PPT_NAVY)
    prs.save(PPTX_PATH)


def validate():
    for path in (DOCX_PATH, PPTX_PATH):
        with ZipFile(path) as z:
            bad = z.testzip()
            if bad:
                raise RuntimeError(f"corrupt zip member: {bad}")


def main():
    build_docx()
    build_pptx()
    validate()
    print(DOCX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
