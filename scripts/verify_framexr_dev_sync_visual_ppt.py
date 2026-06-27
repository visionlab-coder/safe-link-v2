from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "generated" / "framexr-dev-sync-20260607" / "FRAMEXR_DEV_LOG_TODO_SYNC_이미지보강_V2_20260607.pptx"
PREVIEW = ROOT / "docs" / "generated" / "framexr-dev-sync-20260607" / "previews-v2" / "contact_sheet.png"
prs = Presentation(PPTX)
text = " ".join(
    shape.text
    for slide in prs.slides
    for shape in slide.shapes
    if getattr(shape, "has_text_frame", False) and shape.text
)
pictures = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
print(f"pptx_exists={PPTX.exists()}")
print(f"preview_exists={PREVIEW.exists()}")
print(f"slide_count={len(prs.slides)}")
print(f"picture_count={pictures}")
print(f"has_pour_gate={'Pour 3중 게이트' in text}")
print(f"has_code_zero={'코드 변경 0' in text}")
print(f"has_commit_message={'docs: sync development log after pour gate completion' in text}")
print(f"has_caution={'실제 현장 타설 승인 시스템이 아님' in text}")
