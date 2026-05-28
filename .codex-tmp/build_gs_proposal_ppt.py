from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path.cwd()
ASSET_DIR = ROOT / "public" / "presentation-assets" / "gs-proposal"
OUT = ROOT / "docs" / "generated" / "SAFE-LINK_GS건설_제안서_이미지적용_20260527.pptx"


W = Inches(13.333)
H = Inches(7.5)


def set_run(run, size=24, bold=False, color=(255, 255, 255)):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def textbox(slide, x, y, w, h, text, size=24, bold=False, color=(255, 255, 255), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size, bold, color)
    return box


def bullet_list(slide, x, y, w, h, items, size=19, color=(35, 45, 55), bullet_color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*color)
        p.space_after = Pt(8)
    return box


def rect(slide, x, y, w, h, fill, line=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = RGBColor(*line)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle=None, dark=False):
    color = (23, 33, 43) if not dark else (255, 255, 255)
    textbox(slide, 0.65, 0.45, 8.7, 0.55, title, 25, True, color)
    if subtitle:
        textbox(slide, 0.68, 1.03, 9.5, 0.36, subtitle, 11, False, (95, 105, 116) if not dark else (210, 218, 226))
    rect(slide, 0.65, 1.35, 1.3, 0.03, (245, 178, 35))


def add_footer(slide, n):
    textbox(slide, 0.65, 7.08, 5.2, 0.18, "SAFE-LINK × GS건설 제안 초안 | 서원토건", 7.5, False, (118, 126, 134))
    textbox(slide, 12.45, 7.08, 0.32, 0.18, str(n), 8, False, (118, 126, 134), PP_ALIGN.RIGHT)


def metric(slide, x, y, num, label):
    textbox(slide, x, y, 1.7, 0.4, num, 24, True, (245, 178, 35))
    textbox(slide, x, y + 0.48, 2.3, 0.35, label, 9.5, False, (230, 236, 242))


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# 1 Cover
s = prs.slides.add_slide(blank)
s.shapes.add_picture(str(ASSET_DIR / "safe-link-cover.png"), 0, 0, width=W, height=H)
rect(s, 0, 0, 5.55, 7.5, (7, 29, 52), transparency=8)
rect(s, 0, 0, 13.333, 7.5, (0, 0, 0), transparency=72)
textbox(s, 0.68, 0.72, 2.2, 0.28, "PROPOSAL DRAFT", 10, True, (245, 178, 35))
textbox(s, 0.65, 1.28, 4.9, 1.5, "SAFE-LINK", 44, True, (255, 255, 255))
textbox(s, 0.68, 2.32, 4.5, 1.2, "외국인 근로자 안전 커뮤니케이션을\n증빙 가능한 데이터로 바꾸는 현장 OS", 20, True, (255, 255, 255))
textbox(s, 0.7, 4.58, 4.55, 0.9, "GS건설 현장 적용 제안\nTBM · 다국어 소통 · 전자서명 · 법적 블랙박스", 14, False, (222, 230, 238))
textbox(s, 0.7, 6.55, 3.2, 0.28, "서원토건 / SAFE-LINK V2", 10, False, (210, 218, 226))

# 2 Why now
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, (248, 250, 252))
add_title(s, "왜 지금 SAFE-LINK인가", "GS건설의 안전·품질 중심 기조와 외국인 근로자 소통 이슈가 만나는 지점")
textbox(s, 0.8, 1.75, 5.25, 0.72, "현장 리스크는 “전달”보다\n“이해와 증빙”에서 발생합니다", 25, True, (20, 30, 40))
bullet_list(s, 0.83, 2.85, 5.35, 2.1, [
    "외국인 근로자에게 TBM이 전달되어도 실제 이해 여부는 별도 확인 필요",
    "안전지시·작업중지·질의응답은 사고 후 원문/번역문/시간 증빙이 중요",
    "관리자는 교육, 서명, 미확인자 추적, 보고서 작성까지 반복 행정 부담",
], 16)
rect(s, 6.75, 1.65, 5.8, 4.52, (17, 31, 48), transparency=0)
textbox(s, 7.15, 2.02, 4.9, 0.5, "GS건설에 필요한 다음 단계", 19, True, (255, 255, 255))
metric(s, 7.18, 2.9, "01", "다국어 안전 전달")
metric(s, 9.1, 2.9, "02", "확인·서명 이력")
metric(s, 11.0, 2.9, "03", "사후 증빙 패키지")
textbox(s, 7.18, 4.58, 4.85, 0.82, "번역 기능만이 아니라, 누가·언제·어떤 언어로·무엇을 확인했는지를 남기는 구조가 필요합니다.", 16, True, (255, 255, 255))
textbox(s, 0.78, 6.54, 11.2, 0.22, "참고: GS건설 2023 통합보고서, GS건설 Xi Voice 관련 공개 기사(2024.09)", 7.8, False, (110, 118, 126))
add_footer(s, 2)

# 3 Problem
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, (255, 255, 255))
add_title(s, "현장의 빈칸: 전달 후 증빙까지 이어지는가", "TBM 전달과 안전관리 기록 사이의 단절을 줄이는 것이 제안의 핵심")
labels = [
    ("TBM 작성", "관리자 한국어/음성 입력"),
    ("다국어 전달", "외국인 근로자 모국어 제공"),
    ("이해 확인", "퀴즈·응답·서명"),
    ("법적 증빙", "원문·번역문·시간·해시"),
]
for i, (a, b) in enumerate(labels):
    x = 0.82 + i * 3.08
    rect(s, x, 2.1, 2.45, 1.95, (245, 247, 250), line=(220, 226, 232))
    textbox(s, x + 0.22, 2.45, 2.0, 0.36, a, 18, True, (22, 33, 44), PP_ALIGN.CENTER)
    textbox(s, x + 0.22, 3.05, 2.0, 0.45, b, 11, False, (74, 84, 96), PP_ALIGN.CENTER)
    if i < 3:
        textbox(s, x + 2.52, 2.75, 0.5, 0.35, "→", 22, True, (245, 178, 35), PP_ALIGN.CENTER)
textbox(s, 1.05, 5.1, 11.2, 0.7, "기존 솔루션은 일부 기능을 제공합니다. SAFE-LINK는 이 흐름 전체를 하나의 안전관리 세션으로 묶습니다.", 22, True, (18, 28, 38), PP_ALIGN.CENTER)
add_footer(s, 3)

# 4 Solution
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, (15, 25, 36))
add_title(s, "SAFE-LINK 제안: 번역 앱이 아니라 현장 안전 커뮤니케이션 OS", "TBM, 확인, 서명, 대화, 작업중지를 하나의 증빙 세션으로 연결", dark=True)
for i, (title, body) in enumerate([
    ("검증 세션", "NFC·QR·HMAC URL로 근로자/현장/작업조를 특정"),
    ("다국어 TBM", "원문·정규화문·번역문을 근로자 언어별로 저장"),
    ("이해 확인", "확인·퀴즈·음성응답·전자서명으로 이수 판정"),
    ("블랙박스", "대화·서명·시간·해시를 법적 증빙 패키지화"),
]):
    x = 0.8 + (i % 2) * 6.05
    y = 1.85 + (i // 2) * 2.15
    rect(s, x, y, 5.25, 1.45, (255, 255, 255), transparency=88)
    textbox(s, x + 0.28, y + 0.22, 4.65, 0.34, title, 19, True, (245, 178, 35))
    textbox(s, x + 0.28, y + 0.78, 4.55, 0.34, body, 13, False, (235, 241, 247))
textbox(s, 0.85, 6.35, 11.7, 0.4, "제안 한 줄: 외국인 근로자 안전 커뮤니케이션을 ‘전달 완료’가 아니라 ‘증빙 완료’ 상태로 바꿉니다.", 18, True, (255, 255, 255), PP_ALIGN.CENTER)
add_footer(s, 4)

# 5 Workflow image
s = prs.slides.add_slide(blank)
s.shapes.add_picture(str(ASSET_DIR / "safe-link-workflow.png"), 0, 0, width=W, height=H)
rect(s, 0, 0, 13.333, 1.25, (8, 20, 34), transparency=8)
textbox(s, 0.65, 0.38, 8.7, 0.4, "현장 적용 흐름: 태그 한 번이 증빙 세션의 시작점", 25, True, (255, 255, 255))
textbox(s, 0.68, 6.72, 11.8, 0.32, "NFC/QR 검증 → 다국어 TBM → 확인·전자서명 → 해시 증거 데이터셋 → 관리자/HQ 리포트", 15, True, (255, 255, 255), PP_ALIGN.CENTER)

# 6 Evidence image
s = prs.slides.add_slide(blank)
s.shapes.add_picture(str(ASSET_DIR / "safe-link-evidence.png"), 0, 0, width=W, height=H)
rect(s, 0, 0, 5.1, 7.5, (7, 16, 26), transparency=4)
textbox(s, 0.65, 0.65, 4.1, 0.9, "사고 후 필요한 것은\n기억이 아니라 기록입니다", 27, True, (255, 255, 255))
bullet_list(s, 0.72, 2.25, 4.0, 2.4, [
    "TBM 원문 / 번역문",
    "근로자 확인·응답·서명",
    "1:1 안전 대화 로그",
    "현장·시간·단말·관리자",
    "개별 해시 / 세션 해시 / 감사 체인",
], 15, (235, 241, 247))
textbox(s, 0.72, 6.33, 4.0, 0.4, "법적 증빙 패키지 자동 생성", 18, True, (245, 178, 35))

# 7 Pilot proposal
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, (248, 250, 252))
add_title(s, "GS건설 적용 제안: 2주 파일럿", "작게 시작하되, 현장 적용성과 증빙 효과를 숫자로 확인")
steps = [
    ("1일차", "현장/근로자 등록", "근로자 언어·작업조·관리자 권한 세팅"),
    ("3일차", "TBM 전송 테스트", "원문 정규화·번역·서명 흐름 검증"),
    ("1주차", "실사용 운영", "미확인자 추적·1:1 대화·작업중지 요청"),
    ("2주차", "결과 리포트", "이수율·서명률·응답률·증빙 패키지 제출"),
]
for i, (a, b, c) in enumerate(steps):
    y = 1.75 + i * 1.15
    textbox(s, 0.95, y, 1.0, 0.32, a, 15, True, (245, 178, 35))
    rect(s, 2.0, y - 0.05, 0.03, 0.75, (245, 178, 35))
    textbox(s, 2.25, y, 2.65, 0.32, b, 17, True, (20, 30, 40))
    textbox(s, 5.15, y + 0.02, 6.55, 0.28, c, 13, False, (74, 84, 96))
rect(s, 0.9, 6.25, 11.6, 0.55, (17, 31, 48))
textbox(s, 1.16, 6.39, 11.05, 0.24, "파일럿 성공 기준: TBM 서명률, 미확인자 감소, 외국인 근로자 응답률, 사고 후 증빙자료 재현성", 13, True, (255, 255, 255), PP_ALIGN.CENTER)
add_footer(s, 7)

# 8 Close / ask
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, (12, 24, 38))
textbox(s, 0.75, 0.72, 8.4, 0.58, "요청 사항", 27, True, (245, 178, 35))
textbox(s, 0.75, 1.52, 8.5, 1.0, "GS건설 1개 현장 기준 SAFE-LINK 파일럿 협의", 32, True, (255, 255, 255))
bullet_list(s, 0.82, 3.03, 6.5, 1.8, [
    "외국인 근로자 포함 TBM 운영 현장 1곳 선정",
    "관리자 2~3명 / 근로자 20~50명 파일럿",
    "2주 운영 후 안전 커뮤니케이션 증빙 리포트 제출",
], 17, (224, 232, 240))
rect(s, 8.25, 1.5, 3.9, 3.8, (255, 255, 255), transparency=90)
textbox(s, 8.55, 1.9, 3.3, 0.42, "SAFE-LINK가 남기는 것", 18, True, (255, 255, 255))
textbox(s, 8.55, 2.65, 3.2, 1.65, "누가\n무엇을\n어떤 언어로\n언제 확인했는가", 24, True, (255, 255, 255), PP_ALIGN.CENTER)
textbox(s, 0.82, 6.55, 6.2, 0.3, "서원토건 / SAFE-LINK V2", 11, False, (190, 200, 210))
textbox(s, 8.55, 5.02, 3.2, 0.35, "전달 완료 → 증빙 완료", 18, True, (245, 178, 35), PP_ALIGN.CENTER)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
